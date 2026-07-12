"""Patch frozen groundtruth + safe eval course-literals to stay in sync with the
FULL-russified canvas seed (db/zzz_canvas_after_init.sql).

Single source of truth = scripts/canvas_relabel_map.py.

Two passes:
  1) GROUNDTRUTH FILES (xlsx/docx/json/csv/txt) under each canvas task's
     groundtruth_workspace/: whole-cell / field-aware replacement of English realia
     (course names incl. '&amp;', person full names w/ gender, regions, education,
     roles, departments, module/page/discussion/announcement phrases, comments).
     Conservative: header row skipped, formula cells skipped, person-name fires only
     when BOTH first+last atoms are canvas atoms, exact whole-value matches only.
  2) EVAL course-literals: in evaluation/main.py, replace the 22 full course-name
     strings + 7 bare subject strings (both '&' and '&amp;' forms), longest-first.
     These are long unique strings that only ever appear as data literals -> safe.
     SHORTER atoms in eval (person names, regions, departments, education) are NOT
     auto-patched here (risk of corrupting code); they are handled per-task by the
     review workflow which has full context.

Does NOT git commit. Reads task list from /tmp/canvas_tasks.txt.
Run:  uv run --with openpyxl --with python-docx python3 scripts/canvas_patch_groundtruth.py [--dry-run]
"""
import argparse, csv, glob, importlib.util, json, os

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TASKS_DIR = os.path.join(ROOT, "tasks", "finalpool")
MAP_PATH = os.path.join(ROOT, "scripts", "canvas_relabel_map.py")

_spec = importlib.util.spec_from_file_location("canvas_relabel_map", MAP_PATH)
M = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(M)

FLAT = M.FLAT_VALUE_MAP
FIRST = M.FIRST
LAST = M.LAST

issues = []
gt_files_changed = []
eval_files_changed = []


def log_issue(task, where, msg):
    issues.append(f"[{task}] {where}: {msg}")


# ---------------------------------------------------------------------------
# value mappers
# ---------------------------------------------------------------------------
def map_full_name_strict(s):
    """RU 'First Last' (gender-agreed) ONLY if BOTH atoms are canvas atoms; handles
    optional 'Dr.'/title prefix. Returns None otherwise (leave as-is)."""
    toks = s.split()
    if toks and toks[0] in M.TITLES:
        title_ok = True
    else:
        title_ok = False
    core = toks[1:] if title_ok else toks
    if len(core) != 2:
        return None
    f, l = core
    if f in FIRST and l in LAST:
        return M.map_full_name(s)
    return None


def map_cell_value(s):
    """Whole-cell map for xlsx/docx/txt data cells. Returns RU string or None."""
    if s in FLAT:
        return FLAT[s]
    r = M.map_course_name(s)            # course names (incl. &amp; / any season-year)
    if r:
        return r
    r = map_full_name_strict(s)         # person full names (both atoms)
    if r:
        return r
    r = M.map_sortable_name(s)          # 'Last, First'
    if r and ", " in s:
        return r
    return None


def map_json_field(key, value):
    if not isinstance(value, str):
        return None
    s = value.strip()
    if not s:
        return None
    kl = (key or "").lower()
    # person full-name fields
    if any(t in kl for t in ("student", "instructor", "teacher", "faculty", "user_name",
                             "full_name", "username_display", "author")) or kl == "name":
        r = map_full_name_strict(s)
        if r:
            return r
    if "first" in kl and "name" in kl:
        return FIRST.get(s)
    if "last" in kl and "name" in kl:
        return LAST.get(s)  # masculine; rare in canvas GT
    # course / subject
    if "course" in kl or kl in ("subject", "module_name"):
        r = M.map_course_name(s)
        if r:
            return r
        return FLAT.get(s)
    if "department" in kl or kl == "dept":
        return M.DEPARTMENTS.get(s)
    if "region" in kl:
        return M.REGIONS.get(s)
    if "education" in kl or kl == "edu":
        return M.EDU_LEVELS.get(s)
    # generic: whole-cell FLAT (course names, phrases, etc.)
    return FLAT.get(s)


# ---------------------------------------------------------------------------
# per-format GT patchers
# ---------------------------------------------------------------------------
def patch_xlsx(path, task):
    import openpyxl
    try:
        wb = openpyxl.load_workbook(path, data_only=False)
    except Exception as e:
        log_issue(task, os.path.basename(path), f"xlsx unreadable: {e}")
        return 0
    n = 0
    for ws in wb.worksheets:
        for row in ws.iter_rows(min_row=2):           # skip header row
            for c in row:
                v = c.value
                if not isinstance(v, str):
                    continue
                s = v.strip()
                if not s or s.startswith("="):
                    continue
                ru = map_cell_value(s)
                if ru and ru != v:
                    c.value = ru
                    n += 1
    if n:
        wb.save(path)
    return n


def patch_docx(path, task):
    from docx import Document
    try:
        doc = Document(path)
    except Exception as e:
        log_issue(task, os.path.basename(path), f"docx unreadable: {e}")
        return 0
    n = 0

    def patch_runs(paras):
        nonlocal n
        for p in paras:
            for r in p.runs:
                t = r.text
                if not t:
                    continue
                s = t.strip()
                ru = map_cell_value(s)
                if ru and ru != t:
                    r.text = t.replace(s, ru) if s in t else ru
                    n += 1

    patch_runs(doc.paragraphs)
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                patch_runs(cell.paragraphs)
    if n:
        doc.save(path)
    return n


def _walk_json(o, key=None):
    if isinstance(o, dict):
        for k, v in o.items():
            if isinstance(v, (dict, list)):
                yield from _walk_json(v, k)
            else:
                yield (o, k, k, v)
    elif isinstance(o, list):
        for i, v in enumerate(o):
            if isinstance(v, (dict, list)):
                yield from _walk_json(v, key)
            else:
                yield (o, i, key, v)


def patch_json(path, task):
    try:
        with open(path, encoding="utf-8") as f:
            data = json.load(f)
    except Exception as e:
        log_issue(task, os.path.basename(path), f"json unreadable: {e}")
        return 0
    n = 0
    for container, accessor, key, value in list(_walk_json(data)):
        if not isinstance(value, str):
            continue
        ru = map_json_field(key, value)
        if ru and ru != value:
            container[accessor] = ru
            n += 1
    if n:
        with open(path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    return n


def patch_csv(path, task):
    try:
        with open(path, newline="", encoding="utf-8") as f:
            rows = list(csv.reader(f))
    except Exception as e:
        log_issue(task, os.path.basename(path), f"csv unreadable: {e}")
        return 0
    if not rows:
        return 0
    header = rows[0]
    n = 0
    for r in rows[1:]:
        for ci, val in enumerate(r):
            col = header[ci] if ci < len(header) else ""
            ru = map_json_field(col, val)
            if ru and ru != val:
                r[ci] = ru
                n += 1
    if n:
        with open(path, "w", newline="", encoding="utf-8") as f:
            csv.writer(f).writerows(rows)
    return n


def _course_subst(text):
    """Apply the 22 course-name + 7 subject literals (both & forms), longest-first.
    Returns (new_text, changed)."""
    out = text
    for en, ru in _SORTED_COURSE:
        if en in out:
            out = out.replace(en, ru)
    return out, (out != text)


def patch_pptx(path, task):
    from pptx import Presentation
    try:
        prs = Presentation(path)
    except Exception as e:
        log_issue(task, os.path.basename(path), f"pptx unreadable: {e}")
        return 0
    n = 0

    def patch_runs(tf):
        nonlocal n
        for para in tf.paragraphs:
            for r in para.runs:
                t = r.text
                if not t:
                    continue
                s = t.strip()
                ru = map_cell_value(s)            # whole-cell realia / full name
                if ru and ru != t:
                    r.text = t.replace(s, ru) if s in t else ru
                    n += 1
                    continue
                nt, ch = _course_subst(t)          # embedded course-name substring
                if ch:
                    r.text = nt
                    n += 1

    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                patch_runs(sh.text_frame)
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        patch_runs(cell.text_frame)
    if n:
        prs.save(path)
    return n


def patch_txt(path, task):
    try:
        with open(path, encoding="utf-8") as f:
            txt = f.read()
    except Exception as e:
        log_issue(task, os.path.basename(path), f"txt unreadable: {e}")
        return 0
    new, changed = _course_subst(txt)
    if changed:
        with open(path, "w", encoding="utf-8") as f:
            f.write(new)
        return 1
    return 0


PATCHERS = {".xlsx": patch_xlsx, ".docx": patch_docx, ".json": patch_json,
            ".csv": patch_csv, ".pptx": patch_pptx, ".txt": patch_txt}


# ---------------------------------------------------------------------------
# eval course-literal auto-patch (long unique strings only)
# ---------------------------------------------------------------------------
def _build_course_literal_map():
    """{english_literal: ru} for the 22 full course names + 7 bare subjects, both
    '&' and '&amp;' forms. Sorted longest-first by caller."""
    m = {}
    for subj, se, yr in M.COURSE_VARIANTS:
        full = f"{subj} ({se} {yr})"
        ru = M.map_course_name(full)
        m[full] = ru
        m[full.replace("&", "&amp;")] = ru
    for en, ru in M.COURSE_SUBJECTS.items():
        m[en] = ru
        m[en.replace("&", "&amp;")] = ru
    return m


COURSE_LITERALS = _build_course_literal_map()
_SORTED_COURSE = sorted(COURSE_LITERALS.items(), key=lambda kv: -len(kv[0]))


def patch_eval_courses(task):
    py = os.path.join(TASKS_DIR, task, "evaluation", "main.py")
    if not os.path.isfile(py):
        return []
    with open(py, encoding="utf-8") as f:
        txt = f.read()
    orig = txt
    for en, ru in _SORTED_COURSE:
        if en in txt:
            txt = txt.replace(en, ru)
    if txt != orig:
        with open(py, "w", encoding="utf-8") as f:
            f.write(txt)
        return [py]
    return []


# ---------------------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--dry-run", action="store_true")
    ap.add_argument("--no-eval", action="store_true", help="skip eval course-literal pass")
    args = ap.parse_args()

    with open("/tmp/canvas_tasks.txt") as f:
        tasks = [l.strip() for l in f if l.strip()]

    per_task = {}
    for task in sorted(tasks):
        tdir = os.path.join(TASKS_DIR, task)
        gt = os.path.join(tdir, "groundtruth_workspace")
        changed_here = []
        if os.path.isdir(gt):
            for path in sorted(glob.glob(os.path.join(gt, "**", "*"), recursive=True)):
                if not os.path.isfile(path):
                    continue
                ext = os.path.splitext(path)[1].lower()
                fn = PATCHERS.get(ext)
                if not fn or args.dry_run:
                    continue
                k = fn(path, task)
                if k:
                    gt_files_changed.append(path)
                    changed_here.append(f"{os.path.relpath(path, tdir)} ({k} cells)")
        if not args.dry_run and not args.no_eval:
            for p in patch_eval_courses(task):
                eval_files_changed.append(p)
                changed_here.append(os.path.relpath(p, tdir) + " (course literals)")
        if changed_here:
            per_task[task] = changed_here

    print("=== PER-TASK CHANGES ===")
    for task in sorted(per_task):
        print(f"\n{task}:")
        for c in per_task[task]:
            print(f"   {c}")
    print("\n=== ISSUES ===")
    for i in issues:
        print(f"   {i}")
    print(f"\n=== COUNTS ===\n   tasks: {len(tasks)}  gt_files: {len(gt_files_changed)}  eval_files: {len(eval_files_changed)}")


if __name__ == "__main__":
    main()
