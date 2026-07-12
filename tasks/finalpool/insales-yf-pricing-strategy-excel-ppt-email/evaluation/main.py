"""Evaluation script for insales-yf-pricing-strategy-excel-ppt-email."""
import os
import argparse, json, os, sys
import openpyxl

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent", "password": "camel"
}

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []

def check(name, condition, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT, CRITICAL_FAILS
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS]{' [CRITICAL]' if critical else ''} {name}")
    else:
        FAIL_COUNT += 1
        detail_str = str(detail)[:200] if detail else ""
        print(f"  [FAIL]{' [CRITICAL]' if critical else ''} {name}: {detail_str}")
        if critical:
            CRITICAL_FAILS.append(name)

def safe_float(val, default=None):
    try:
        if val is None: return default
        return float(str(val).replace(",", "").replace("%", "").replace("$", "").strip())
    except (ValueError, TypeError):
        return default

def get_conn():
    import psycopg2
    return psycopg2.connect(**DB_CONFIG)

def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT
    PASS_COUNT = 0
    FAIL_COUNT = 0

    # Check Dynamic_Pricing_Strategy.xlsx
    excel_path = os.path.join(agent_workspace, "Dynamic_Pricing_Strategy.xlsx")
    check("Dynamic_Pricing_Strategy.xlsx exists", os.path.exists(excel_path), critical=True)
    if os.path.exists(excel_path):
        wb = openpyxl.load_workbook(excel_path)
        # CRITICAL: Market_Indicators must reflect the MOEX instruments (IMOEX.ME ~3000 scale, gold GLDRUB_TOM),
        # NOT the retired yahoo-finance instruments (Dow ~38000, Gold GC=F ~2050).
        if "Market_Indicators" in wb.sheetnames:
            mi = wb["Market_Indicators"]
            mi_headers = [str(c.value).strip().lower() if c.value else "" for c in mi[1]]
            mi_rows = list(mi.iter_rows(min_row=2, values_only=True))
            ind_idx = mi_headers.index("indicator") if "indicator" in mi_headers else 0
            val_idx = mi_headers.index("latest_value") if "latest_value" in mi_headers else 1
            blob = " ".join(str(r[ind_idx]) for r in mi_rows if r and r[ind_idx]).upper()
            check("Market_Indicators references IMOEX index (not Dow Jones)",
                  "IMOEX" in blob and "DOW" not in blob and "DJI" not in blob,
                  f"indicators: {blob[:120]}", critical=True)
            check("Market_Indicators references gold GLDRUB_TOM (not GC=F)",
                  ("GLDRUB" in blob or "GLD" in blob) and "GC=F" not in blob,
                  f"indicators: {blob[:120]}", critical=True)
            # Scale guard: an IMOEX-scale index row must be present (~3000), and no ~40000 Dow-scale leftover.
            vals = [safe_float(r[val_idx]) for r in mi_rows if r and len(r) > val_idx]
            vals = [v for v in vals if v is not None]
            check("Market_Indicators values on MOEX scale (no ~40000 Dow-scale leftover)",
                  all(v < 20000 for v in vals) if vals else False,
                  f"values: {vals}", critical=True)
        gt_path = os.path.join(groundtruth_workspace, "Dynamic_Pricing_Strategy.xlsx")
        gt_wb = openpyxl.load_workbook(gt_path) if os.path.exists(gt_path) else None

        if gt_wb:
            # NOTE: grading is intentionally ORDER-INSENSITIVE and matched by entity
            # (product name / indicator name), not by row index. Agent and GT may hold
            # the same data in a different row order; positional comparison is invalid.
            # Schema (headers) and row-count are still checked per sheet. Product_Analysis
            # and Market_Indicators are value-anchored to the live catalog / MOEX feed
            # (real, reproducible). Price_Recommendations and Impact_Forecast use
            # structural/range checks because the optimal-adjustment rule is unspecified
            # author judgment that no correct agent can reproduce cell-for-cell.

            def norm(s):
                return str(s).strip().lower() if s is not None else ""

            def find_col(headers, name):
                name = name.lower()
                for i, h in enumerate(headers):
                    if h == name:
                        return i
                return None

            # Per-sheet schema + row-count checks
            for sheet_name in gt_wb.sheetnames:
                check(f"{sheet_name} sheet exists", sheet_name in wb.sheetnames)
                if sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    gt_ws = gt_wb[sheet_name]
                    gt_headers = [norm(c.value) for c in gt_ws[1]]
                    headers = [norm(c.value) for c in ws[1]]
                    for h in gt_headers:
                        if h:
                            check(f"{sheet_name} has {h} column", h in headers, f"headers: {headers[:10]}")
                    gt_rows = list(gt_ws.iter_rows(min_row=2, values_only=True))
                    data_rows = list(ws.iter_rows(min_row=2, values_only=True))
                    min_rows = max(1, len(gt_rows) - 2)
                    check(f"{sheet_name} has >= {min_rows} data rows", len(data_rows) >= min_rows, f"got {len(data_rows)}")

            # --- Product_Analysis: order-insensitive, matched by Product_Name ---
            if "Product_Analysis" in wb.sheetnames and "Product_Analysis" in gt_wb.sheetnames:
                ws = wb["Product_Analysis"]; gt_ws = gt_wb["Product_Analysis"]
                a_h = [norm(c.value) for c in ws[1]]
                g_h = [norm(c.value) for c in gt_ws[1]]
                a_name = find_col(a_h, "product_name")
                # build agent index by name (substring-matchable)
                agent_rows = [r for r in ws.iter_rows(min_row=2, values_only=True) if r and a_name is not None and r[a_name]]
                num_cols = ["current_price", "total_sales", "stock_level", "revenue"]
                for gt_row in gt_ws.iter_rows(min_row=2, values_only=True):
                    if not gt_row:
                        continue
                    g_name_idx = find_col(g_h, "product_name")
                    gname = norm(gt_row[g_name_idx]) if g_name_idx is not None else ""
                    if not gname:
                        continue
                    # match by substring either direction (GT names mirror live catalog names)
                    match = None
                    for ar in agent_rows:
                        an = norm(ar[a_name])
                        if an and (gname in an or an in gname):
                            match = ar
                            break
                    short = gname[:30]
                    check(f"Product_Analysis has product '{short}'", match is not None,
                          f"agent products: {[norm(r[a_name])[:25] for r in agent_rows][:6]}")
                    if match is None:
                        continue
                    for col in num_cols:
                        gi = find_col(g_h, col); ai = find_col(a_h, col)
                        if gi is None or ai is None or gi >= len(gt_row) or ai >= len(match):
                            continue
                        gf = safe_float(gt_row[gi]); af = safe_float(match[ai])
                        if gf is None:
                            continue
                        tol = max(0.5, abs(gf) * 0.15)
                        check(f"Product_Analysis '{short}' {col} ~{gf:.2f}",
                              af is not None and abs(gf - af) <= tol, f"got {af}")

            # --- Market_Indicators: look up IMOEX/GLDRUB latest_value by name (live, real) ---
            # Trend / Change_Pct are derived from an unspecified window and the GT scale is
            # not authoritative, so they are NOT cell-checked here. The CRITICAL block above
            # already enforces the correct MOEX instruments and scale.
            if "Market_Indicators" in wb.sheetnames:
                ws = wb["Market_Indicators"]
                a_h = [norm(c.value) for c in ws[1]]
                ind_i = find_col(a_h, "indicator")
                val_i = find_col(a_h, "latest_value")
                rows = [r for r in ws.iter_rows(min_row=2, values_only=True) if r]
                def lookup_val(token):
                    for r in rows:
                        if ind_i is not None and r[ind_i] and token in str(r[ind_i]).upper():
                            return safe_float(r[val_i]) if val_i is not None else None
                    return None
                imoex_v = lookup_val("IMOEX")
                gld_v = lookup_val("GLDRUB") if lookup_val("GLDRUB") is not None else lookup_val("GLD")
                check("Market_Indicators IMOEX latest_value ~3052",
                      imoex_v is not None and abs(imoex_v - 3052.0) <= max(0.5, 3052.0 * 0.15),
                      f"got {imoex_v}")
                check("Market_Indicators Gold latest_value ~5093.3",
                      gld_v is not None and abs(gld_v - 5093.3) <= max(0.5, 5093.3 * 0.15),
                      f"got {gld_v}")

            # --- Price_Recommendations: structural / range checks (judgment, not reproducible) ---
            if "Price_Recommendations" in wb.sheetnames:
                ws = wb["Price_Recommendations"]
                a_h = [norm(c.value) for c in ws[1]]
                rows = [r for r in ws.iter_rows(min_row=2, values_only=True) if r and any(c is not None for c in r)]
                check("Price_Recommendations has >= 1 recommendation", len(rows) >= 1, f"got {len(rows)}")
                rec_i = find_col(a_h, "recommended_price")
                chg_i = find_col(a_h, "change_pct")
                # recommended prices must sit in a sane band relative to the live catalog
                # (catalog price range ~1.81..861.33; allow generous headroom up to 2000)
                if rec_i is not None:
                    recs = [safe_float(r[rec_i]) for r in rows if rec_i < len(r)]
                    recs = [v for v in recs if v is not None]
                    check("Price_Recommendations recommended prices within sane band",
                          bool(recs) and all(0 < v <= 2000 for v in recs),
                          f"prices: {recs[:8]}")
                if chg_i is not None:
                    chgs = [safe_float(r[chg_i]) for r in rows if chg_i < len(r)]
                    chgs = [v for v in chgs if v is not None]
                    check("Price_Recommendations change_pct values are reasonable",
                          bool(chgs) and all(-90 <= v <= 200 for v in chgs),
                          f"change_pct: {chgs[:8]}")

            # --- Impact_Forecast: structural / range checks (counts bounded by catalog size) ---
            if "Impact_Forecast" in wb.sheetnames:
                ws = wb["Impact_Forecast"]
                a_h = [norm(c.value) for c in ws[1]]
                m_i = find_col(a_h, "metric"); v_i = find_col(a_h, "value")
                metric_map = {}
                if m_i is not None and v_i is not None:
                    for r in ws.iter_rows(min_row=2, values_only=True):
                        if r and r[m_i] is not None:
                            metric_map[norm(r[m_i])] = safe_float(r[v_i])
                for key in ["products_price_increase", "products_price_decrease", "products_no_change"]:
                    present = any(key in k for k in metric_map)
                    check(f"Impact_Forecast has {key}", present, f"metrics: {list(metric_map)[:8]}")
                # the three product-count buckets must sum to within the live catalog size (82)
                def get_metric(key):
                    for k, v in metric_map.items():
                        if key in k:
                            return v
                    return None
                counts = [get_metric(k) for k in ["products_price_increase", "products_price_decrease", "products_no_change"]]
                if all(c is not None for c in counts):
                    total = sum(counts)
                    check("Impact_Forecast product counts sum within catalog size (<=82)",
                          0 < total <= 82, f"sum={total}")

    # Check Pricing_Strategy_Deck.pptx
    pptx_path = os.path.join(agent_workspace, "Pricing_Strategy_Deck.pptx")
    check("Pricing_Strategy_Deck.pptx exists", os.path.exists(pptx_path))
    if os.path.exists(pptx_path):
        from pptx import Presentation
        prs = Presentation(pptx_path)
        gt_pptx_path = os.path.join(groundtruth_workspace, "Pricing_Strategy_Deck.pptx")
        if os.path.exists(gt_pptx_path):
            gt_prs = Presentation(gt_pptx_path)
            gt_slide_count = len(gt_prs.slides)
            check(f"Pricing_Strategy_Deck.pptx has >= {gt_slide_count} slides", len(prs.slides) >= gt_slide_count, f"got {len(prs.slides)} slides")
            # Compare slide titles
            gt_titles = []
            for s in gt_prs.slides:
                if s.shapes.title:
                    gt_titles.append(s.shapes.title.text.strip().lower())
            agent_titles = []
            for s in prs.slides:
                if s.shapes.title:
                    agent_titles.append(s.shapes.title.text.strip().lower())
            for gt in gt_titles:
                found = any(gt in at or at in gt for at in agent_titles)
                check(f"Pricing_Strategy_Deck.pptx has slide \"{gt[:40]}\"", found, f"agent titles: {agent_titles[:5]}")
        else:
            check("Pricing_Strategy_Deck.pptx has >= 3 slides", len(prs.slides) >= 3, f"got {len(prs.slides)} slides")

    # Check Python script exists (terminal usage)
    py_files = [f for f in os.listdir(agent_workspace) if f.endswith(".py")]
    check("Python analysis script exists", len(py_files) >= 1, f"found: {py_files}")

    # Database checks
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute("SELECT subject, to_addr FROM email.messages WHERE folder_id = (SELECT id FROM email.folders WHERE name = 'Sent' LIMIT 1) AND subject ILIKE '%pricing%'")
        email_row = cur.fetchone()
        check("Email with correct subject sent", email_row is not None, "no matching email found", critical=True)
        if email_row:
            check("Email has recipient", email_row[1] is not None, f"to_addr: {email_row[1]}")
        # Reverse verification: noise emails should not be in Sent folder
        cur.execute("SELECT COUNT(*) FROM email.messages WHERE folder_id = (SELECT id FROM email.folders WHERE name = 'Sent' LIMIT 1) AND subject ILIKE '%newsletter%'")
        noise_sent = cur.fetchone()[0]
        check("No noise emails in Sent folder", noise_sent == 0, f"found {noise_sent} noise emails in Sent")
        conn.close()
    except Exception as e:
        check("DB checks", False, str(e))

    if CRITICAL_FAILS:
        msg = f"CRITICAL checks failed ({len(CRITICAL_FAILS)}): {CRITICAL_FAILS}. Passed {PASS_COUNT}/{PASS_COUNT + FAIL_COUNT} checks"
        print(f"  [GATE] {msg}")
        sys.exit(1)

    # Threshold-based pass: all CRITICAL checks must pass (enforced above) AND overall
    # accuracy must clear the bar. We no longer require FAIL_COUNT==0 because several
    # checks (price recommendations, impact forecast) are structural and a correct agent
    # may legitimately differ from one author's judgment on a few non-anchored cells.
    PASS_THRESHOLD = 0.85
    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total) if total else 0.0
    success = accuracy >= PASS_THRESHOLD
    return success, f"Passed {PASS_COUNT}/{total} checks (accuracy {accuracy:.2%}, threshold {PASS_THRESHOLD:.0%})"

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)
    sys.exit(0 if success else 1)

if __name__ == "__main__":
    main()