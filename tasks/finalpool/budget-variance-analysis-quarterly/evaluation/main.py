#!/usr/bin/env python3
"""Evaluation script for budget-variance-analysis-quarterly task validation.

Структура проверок:
  - НЕСКОЛЬКО CRITICAL (семантических) проверок: проверяют, что отклонения
    действительно ВЫЧИСЛЕНЫ (а не скопированы), что метки Status согласованы со
    знаком Variance $, что прогноз сложён арифметически верно, что .docx содержит
    раздел по каждому отделу, и что фазы 5-6 реально выполнены (письма владельцам
    бюджета + событие обзора бюджета в календаре). Любой провал CRITICAL =>
    sys.exit(1) ещё ДО порога точности.
  - Остальные структурные / GT-проверки идут к порогу accuracy >= 70.

Замечания по локали:
  - Метки Status принимаются и на русском, и на английском.
  - RU-ключевые слова в письмах/событиях ищутся в .lower() ОРИГИНАЛЬНОГО текста.
  - Ожидаемые значения variance/forecast ПЕРЕСЧИТЫВАЮТСЯ из вывода самого агента
    (Budget YTD / Actual YTD), а не берутся из возможно устаревшего groundtruth.
"""

from argparse import ArgumentParser
import json
import os
import sys

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILURES = []


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def record_critical(name, passed, detail=""):
    """Критическая проверка: фиксируется отдельно и валит весь evaluator."""
    record("[CRITICAL] " + name, passed, detail)
    if not passed:
        CRITICAL_FAILURES.append(name)


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def to_num(x):
    try:
        return float(x)
    except (TypeError, ValueError):
        return None


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


# Метки статуса: RU + EN.
FAVORABLE_WORDS = ("favorable", "благоприят")          # под бюджетом (variance < 0)
UNFAVORABLE_WORDS = ("unfavorable", "неблагоприят")    # сверх бюджета (variance > 0)
ONBUDGET_WORDS = ("on budget", "в бюджете", "по плану", "on-budget")

DEPARTMENTS = ["Operations", "Sales", "Marketing", "IT"]


def status_kind(label):
    """Возвращает 'fav' / 'unfav' / 'on' / None по тексту метки (RU+EN)."""
    if label is None:
        return None
    t = str(label).strip().lower()
    if any(w in t for w in ONBUDGET_WORDS):
        return "on"
    if any(w in t for w in UNFAVORABLE_WORDS):
        return "unfav"
    if any(w in t for w in FAVORABLE_WORDS):
        return "fav"
    return None


def find_header_row(rows, required_cols):
    """Находит индекс строки-заголовка, содержащей все required_cols (без регистра)."""
    req = [c.strip().lower() for c in required_cols]
    for i, r in enumerate(rows):
        cells = [str(c).strip().lower() if c is not None else "" for c in r]
        if all(any(rc == cell for cell in cells) for rc in req):
            return i, cells
    return None, None


def col_index(header_cells, name):
    name = name.strip().lower()
    for i, c in enumerate(header_cells):
        if c == name:
            return i
    return None


# ---------------------------------------------------------------------------
# CRITICAL: variance computed (not copied) + Status consistency
# ---------------------------------------------------------------------------

def check_variance_semantics(workspace):
    print("\n=== CRITICAL: семантика variance_analysis.xlsx ===")
    import openpyxl
    path = os.path.join(workspace, "variance_analysis.xlsx")
    if not os.path.isfile(path):
        record_critical("variance_analysis.xlsx присутствует", False, "Not found")
        record_critical("Variance $/% вычислены (1-я строка данных)", False, "no file")
        record_critical("Status согласован со знаком Variance $", False, "no file")
        return
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
    except Exception as e:
        record_critical("variance_analysis.xlsx читается", False, str(e))
        return

    ws = None
    for sn in wb.sheetnames:
        if sn.strip().lower() == "variance analysis":
            ws = wb[sn]
            break
    if ws is None:
        ws = wb[wb.sheetnames[0]]

    rows = list(ws.iter_rows(values_only=True))
    hdr_idx, hdr = find_header_row(
        rows, ["Budget YTD", "Actual YTD", "Variance $", "Variance %"]
    )
    if hdr_idx is None:
        record_critical("Variance $/% вычислены (1-я строка данных)", False,
                        "Заголовок со столбцами Budget YTD/Actual YTD/Variance $/Variance % не найден")
        record_critical("Status согласован со знаком Variance $", False, "no header")
        wb.close()
        return

    ci_b = col_index(hdr, "Budget YTD")
    ci_a = col_index(hdr, "Actual YTD")
    ci_v = col_index(hdr, "Variance $")
    ci_vp = col_index(hdr, "Variance %")
    ci_st = col_index(hdr, "Status")

    # Строки данных: после заголовка, где Budget YTD и Actual YTD числовые.
    data_rows = []
    for r in rows[hdr_idx + 1:]:
        b = to_num(r[ci_b]) if ci_b is not None and ci_b < len(r) else None
        a = to_num(r[ci_a]) if ci_a is not None and ci_a < len(r) else None
        if b is None or a is None:
            # пустая строка/начало блока сводки => конец таблицы данных
            if all((c is None or str(c).strip() == "") for c in r):
                if data_rows:
                    break
                continue
            # нечисловая строка (напр. 'Summary Metrics') => конец данных
            if data_rows:
                break
            continue
        # итоговые строки (TOTAL/Итого/Всего) не считаем строками данных
        label = " ".join(
            str(c).strip().lower() for c in r[:ci_b or 1] if c is not None
        )
        if "total" in label or "итог" in label or "всего" in label:
            continue
        data_rows.append(r)

    if not data_rows:
        record_critical("Variance $/% вычислены (1-я строка данных)", False,
                        "Не найдено ни одной строки данных с числовыми Budget/Actual YTD")
        record_critical("Status согласован со знаком Variance $", False, "no data rows")
        wb.close()
        return

    # --- Проверка 1: первая строка данных — variance вычислен верно. ---
    r0 = data_rows[0]
    b0 = to_num(r0[ci_b]); a0 = to_num(r0[ci_a]); v0 = to_num(r0[ci_v]) if ci_v is not None else None
    vp0 = to_num(r0[ci_vp]) if ci_vp is not None else None
    exp_v = a0 - b0
    var_ok = v0 is not None and num_close(v0, exp_v, 1.0)
    pct_ok = True
    if ci_vp is not None:
        exp_vp = round((exp_v / b0 * 100.0), 1) if b0 else 0.0
        # принимаем и процент (2.5), и долю (0.025 с percent-форматом ячейки)
        pct_ok = vp0 is not None and (
            num_close(vp0, exp_vp, 0.2) or num_close(vp0 * 100.0, exp_vp, 0.2)
        )
    record_critical(
        "Variance $/% вычислены (1-я строка данных = Actual YTD - Budget YTD)",
        var_ok and pct_ok,
        f"Budget={b0}, Actual={a0}, Variance$={v0} (ожид {exp_v}), Variance%={vp0}",
    )

    # --- Проверка 2: Status согласован со знаком Variance $ по всем строкам. ---
    if ci_st is None:
        record_critical("Status согласован со знаком Variance $", False,
                        "Столбец Status отсутствует")
    else:
        bad = []
        checked = 0
        for r in data_rows:
            b = to_num(r[ci_b]); a = to_num(r[ci_a])
            v = to_num(r[ci_v]) if ci_v is not None else (a - b)
            if v is None:
                v = a - b
            kind = status_kind(r[ci_st]) if ci_st < len(r) else None
            if kind is None:
                bad.append(f"{r[ci_st]!r}->нераспознан")
                continue
            checked += 1
            if abs(v) < 1.0:
                expected = "on"
            elif v < 0:
                expected = "fav"
            else:
                expected = "unfav"
            # 'On Budget' для околонулевых допустим как fav/unfav тоже не штрафуем,
            # но ненулевой variance с противоположной меткой — ошибка.
            if expected == "fav" and kind == "unfav":
                bad.append(f"V$={v}->{r[ci_st]!r}")
            elif expected == "unfav" and kind == "fav":
                bad.append(f"V$={v}->{r[ci_st]!r}")
        record_critical(
            "Status согласован со знаком Variance $ (Favorable/Unfavorable, RU+EN)",
            len(bad) == 0 and checked >= 1,
            f"расхождения: {bad}" if bad else f"проверено строк: {checked}",
        )
    wb.close()


# ---------------------------------------------------------------------------
# CRITICAL: forecast Total == sum of department columns
# ---------------------------------------------------------------------------

def check_forecast_semantics(workspace):
    print("\n=== CRITICAL: арифметика budget_forecast.xlsx ===")
    import openpyxl
    path = os.path.join(workspace, "budget_forecast.xlsx")
    if not os.path.isfile(path):
        record_critical("budget_forecast.xlsx: Total = сумма по отделам", False, "Not found")
        return
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
    except Exception as e:
        record_critical("budget_forecast.xlsx читается", False, str(e))
        return

    ws = None
    for sn in wb.sheetnames:
        if sn.strip().lower() == "budget forecast":
            ws = wb[sn]
            break
    if ws is None:
        ws = wb[wb.sheetnames[0]]

    rows = list(ws.iter_rows(values_only=True))
    hdr_idx, hdr = find_header_row(rows, ["Operations", "Sales", "Marketing", "IT", "Total"])
    if hdr_idx is None:
        record_critical("budget_forecast.xlsx: Total = сумма по отделам", False,
                        "Заголовок Operations/Sales/Marketing/IT/Total не найден")
        wb.close()
        return

    idx = {d: col_index(hdr, d) for d in DEPARTMENTS}
    idx_total = col_index(hdr, "Total")

    checked = 0
    bad = []
    for r in rows[hdr_idx + 1:]:
        dept_vals = [to_num(r[idx[d]]) if idx[d] is not None and idx[d] < len(r) else None
                     for d in DEPARTMENTS]
        total = to_num(r[idx_total]) if idx_total is not None and idx_total < len(r) else None
        if any(v is None for v in dept_vals) or total is None:
            continue
        checked += 1
        if not num_close(sum(dept_vals), total, max(abs(total) * 0.001, 1.0)):
            bad.append(f"строка {r[0]!r}: sum={sum(dept_vals)} != Total={total}")

    record_critical(
        "budget_forecast.xlsx: Total = Operations+Sales+Marketing+IT (>=1 строка сценария)",
        checked >= 1 and len(bad) == 0,
        f"расхождения: {bad}" if bad else f"проверено строк сценариев: {checked}",
    )
    wb.close()


# ---------------------------------------------------------------------------
# CRITICAL: docx has a per-department section with a YTD variance figure
# ---------------------------------------------------------------------------

def check_docx_departments(workspace):
    print("\n=== CRITICAL: dept_variance_reports.docx по отделам ===")
    from docx import Document
    path = os.path.join(workspace, "dept_variance_reports.docx")
    if not os.path.isfile(path):
        record_critical("dept_variance_reports.docx: раздел по каждому отделу", False, "Not found")
        return
    try:
        doc = Document(path)
    except Exception as e:
        record_critical("dept_variance_reports.docx читается", False, str(e))
        return

    full = " ".join(p.text for p in doc.paragraphs)
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                full += " " + cell.text
    low = full.lower()

    import re
    # RU/EN варианты названий отделов. Для IT используем границу слова, чтобы
    # не ловить "split"/"audit" и т.п.
    dept_aliases = {
        "Operations": [r"operations", r"операцион", r"эксплуатац"],
        "Sales": [r"sales", r"продаж", r"сбыт"],
        "Marketing": [r"marketing", r"маркетинг"],
        "IT": [r"\bit\b", r"\bит\b", r"\bит-", r"информацион"],
    }
    missing = [d for d, al in dept_aliases.items()
               if not any(re.search(a, low) for a in al)]
    record_critical(
        "dept_variance_reports.docx: упомянуты все 4 отдела (Operations/Sales/Marketing/IT)",
        len(missing) == 0,
        f"отсутствуют: {missing}",
    )

    # Должна присутствовать хотя бы одна цифра отклонения (variance / отклонение + число).
    has_variance_word = ("variance" in low) or ("отклонен" in low)
    has_number = re.search(r"\d{3,}", full) is not None
    record_critical(
        "dept_variance_reports.docx: есть показатель YTD-отклонения (слово+число)",
        has_variance_word and has_number,
        f"variance_word={has_variance_word}, number={has_number}",
    )


# ---------------------------------------------------------------------------
# CRITICAL: phases 5-6 deliverables (email to budget owner + calendar event)
# ---------------------------------------------------------------------------

def _load_budget_owner_emails(workspace, groundtruth_workspace):
    """Читает Manager_Email из cost_center_mapping.csv (workspace, иначе GT)."""
    import csv
    emails = []
    for base in (workspace, groundtruth_workspace, os.path.dirname(workspace)):
        if not base:
            continue
        p = os.path.join(base, "cost_center_mapping.csv")
        if os.path.isfile(p):
            try:
                with open(p, newline="", encoding="utf-8") as f:
                    for row in csv.DictReader(f):
                        e = (row.get("Manager_Email") or "").strip()
                        if e:
                            emails.append(e.lower())
                if emails:
                    return emails
            except Exception:
                pass
    # Запасной список (из исходного датасета).
    return [
        "robert.clark@company.com",
        "thomas.lee@company.com",
        "susan.kim@company.com",
        "james.johnson@company.com",
    ]


def check_phase56_email_calendar(workspace, groundtruth_workspace):
    print("\n=== CRITICAL: фазы 5-6 (email владельцам бюджета + событие обзора бюджета) ===")
    try:
        import psycopg2
    except Exception as e:
        record_critical("Письмо владельцу бюджета отправлено", False, f"psycopg2 import: {e}")
        record_critical("Событие обзора бюджета создано в календаре", False, "psycopg2 missing")
        return

    db = {
        "host": os.environ.get("PGHOST", "localhost"),
        "port": int(os.environ.get("PGPORT", "5432")),
        "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
        "user": os.environ.get("PGUSER", "eigent"),
        "password": os.environ.get("PGPASSWORD", "camel"),
    }

    owner_emails = _load_budget_owner_emails(workspace, groundtruth_workspace)

    # --- Email: хотя бы одно письмо на адрес владельца бюджета. ---
    try:
        conn = psycopg2.connect(**db)
        cur = conn.cursor()
        like_clauses = " OR ".join(["to_addr::text ILIKE %s"] * len(owner_emails))
        params = [f"%{e}%" for e in owner_emails]
        cur.execute(
            f"SELECT subject, to_addr, body_text FROM email.messages WHERE {like_clauses}",
            params,
        )
        rows = cur.fetchall()
        cur.close()
        conn.close()
        record_critical(
            "Письмо отправлено хотя бы одному владельцу бюджета (Manager_Email)",
            len(rows) >= 1,
            f"совпавших писем: {len(rows)}; адреса: {owner_emails}",
        )
    except Exception as e:
        record_critical("Письмо отправлено хотя бы одному владельцу бюджета", False, str(e))

    # --- Calendar: хотя бы одно событие обзора бюджета. ---
    review_kw = ["budget review", "budget-review", "обзор бюджет", "пересмотр бюджет",
                 "обзора бюджет", "review бюджет", "бюджетн", "variance", "отклонен", "budget"]
    try:
        conn = psycopg2.connect(**db)
        cur = conn.cursor()
        cur.execute("SELECT summary, description FROM gcal.events")
        events = cur.fetchall()
        cur.close()
        conn.close()
        matched = 0
        for summ, desc in events:
            text = ((summ or "") + " " + (desc or "")).lower()
            if any(kw in text for kw in review_kw):
                matched += 1
        record_critical(
            "Событие обзора бюджета создано в календаре (>=1)",
            matched >= 1,
            f"всего событий: {len(events)}, совпавших: {matched}",
        )
    except Exception as e:
        record_critical("Событие обзора бюджета создано в календаре", False, str(e))


# ---------------------------------------------------------------------------
# NON-critical: structural + loose groundtruth comparison
# ---------------------------------------------------------------------------

def check_xlsx_content(workspace, groundtruth_workspace="."):
    print("\n=== Check: XLSX files (структура + GT) ===")
    import openpyxl
    # Столбцы, чьи значения локализуемы/волатильны — не сравниваем их с GT строкой.
    VOLATILE_HEADERS = {"status", "trend", "notes", "scenario"}
    # Числовые столбцы, которые агент НЕ может воспроизвести из исходников
    # (фактические расходы отсутствуют в каком-либо источнике данных — финансовая
    # система/ClickHouse в этой задаче является фантомной зависимостью). Их
    # абсолютные значения с GT не сверяем: достаточно, что агент выдал число, а
    # внутренняя арифметика (Variance $ = Actual YTD - Budget YTD) проверяется в
    # CRITICAL-проверке. Budget YTD выводится из approved_budget.xlsx и сверяется.
    NONDERIVABLE_NUM_HEADERS = {"actual ytd", "variance $", "variance%", "variance %",
                                "january", "february", "march", "ytd variance"}
    for fname in ["variance_analysis.xlsx", "variance_tracking.xlsx", "budget_forecast.xlsx"]:
        xlsx_path = os.path.join(workspace, fname)
        if not os.path.isfile(xlsx_path):
            record(f"xlsx {fname} exists", False, "Not found")
            continue
        record(f"xlsx {fname} exists", True)
        try:
            wb = openpyxl.load_workbook(xlsx_path, data_only=True)
            for ws in wb.worksheets:
                rows = list(ws.iter_rows(values_only=True))
                record(f"xlsx {fname} '{ws.title}' has data", len(rows) >= 2, f"{len(rows)} rows")
        except Exception as e:
            record(f"xlsx {fname} readable", False, str(e))
            continue

        gt_path = os.path.join(groundtruth_workspace, fname)
        if not os.path.isfile(gt_path):
            wb.close()
            continue

        gt_wb = openpyxl.load_workbook(gt_path, data_only=True)
        for gt_sheet_name in gt_wb.sheetnames:
            gt_ws = gt_wb[gt_sheet_name]
            agent_ws = None
            for asn in wb.sheetnames:
                if asn.strip().lower() == gt_sheet_name.strip().lower():
                    agent_ws = wb[asn]
                    break
            if agent_ws is None:
                record(f"GT {fname} sheet '{gt_sheet_name}' exists in agent", False,
                       f"Available: {wb.sheetnames}")
                continue

            gt_all = list(gt_ws.iter_rows(values_only=True))
            agent_all = list(agent_ws.iter_rows(values_only=True))
            # Заголовок GT (для определения волатильных столбцов по индексу).
            gt_hdr_idx, gt_hdr = find_header_row(
                gt_all,
                ["Budget YTD", "Actual YTD"] if "analysis" in fname else
                (["Department", "January"] if "tracking" in fname else ["Operations", "Total"]),
            )
            volatile_cols = set()
            nonderivable_num_cols = set()
            if gt_hdr is not None:
                for i, h in enumerate(gt_hdr):
                    if h in VOLATILE_HEADERS:
                        volatile_cols.add(i)
                    if h in NONDERIVABLE_NUM_HEADERS:
                        nonderivable_num_cols.add(i)

            gt_rows = [r for r in gt_ws.iter_rows(min_row=2, values_only=True)
                       if any(c is not None for c in r)]
            agent_rows = [r for r in agent_ws.iter_rows(min_row=2, values_only=True)
                          if any(c is not None for c in r)]

            record(f"GT {fname} '{gt_sheet_name}' row count", len(agent_rows) == len(gt_rows),
                   f"Expected {len(gt_rows)}, got {len(agent_rows)}")

            check_indices = list(range(min(3, len(gt_rows))))
            if len(gt_rows) > 3:
                check_indices.append(len(gt_rows) - 1)
            for idx in check_indices:
                gt_row = gt_rows[idx]
                if idx < len(agent_rows):
                    a_row = agent_rows[idx]
                    row_ok = True
                    for col_idx in range(min(len(gt_row), len(a_row) if a_row else 0)):
                        gt_val = gt_row[col_idx]
                        a_val = a_row[col_idx]
                        if gt_val is None:
                            continue
                        if isinstance(gt_val, (int, float)):
                            if col_idx in nonderivable_num_cols:
                                # Не сверяем абсолют невоспроизводимых значений
                                # (фактические расходы/их производные) — достаточно
                                # наличия числа; арифметика проверена в CRITICAL.
                                ok = to_num(a_val) is not None
                            else:
                                ok = num_close(a_val, gt_val, max(abs(gt_val) * 0.1, 1.0))
                        else:
                            # Волатильные/локализуемые текстовые ячейки не сверяем по строке.
                            if col_idx in volatile_cols:
                                ok = True
                            else:
                                ok = str_match(a_val, gt_val) or a_val is not None
                        if not ok:
                            record(f"GT {fname} '{gt_sheet_name}' row {idx+1} col {col_idx+1}",
                                   False, f"Expected {gt_val}, got {a_val}")
                            row_ok = False
                            break
                    if row_ok:
                        record(f"GT {fname} '{gt_sheet_name}' row {idx+1} values match", True)
                else:
                    record(f"GT {fname} '{gt_sheet_name}' row {idx+1} exists", False,
                           "Row missing in agent")
        gt_wb.close()
        wb.close()


def check_docx_content(workspace):
    print("\n=== Check: DOCX files ===")
    from docx import Document
    path = os.path.join(workspace, "dept_variance_reports.docx")
    if not os.path.isfile(path):
        record("docx dept_variance_reports.docx exists", False, "Not found")
        return False
    record("docx dept_variance_reports.docx exists", True)
    try:
        doc = Document(path)
        record("docx has content", len(doc.paragraphs) > 0, f"{len(doc.paragraphs)} paragraphs")
    except Exception as e:
        record("docx readable", False, str(e))
    return True


def check_pptx_content(workspace):
    print("\n=== Check: PPTX executive_presentation.pptx ===")
    from pptx import Presentation
    path = os.path.join(workspace, "executive_presentation.pptx")
    if not os.path.isfile(path):
        record("pptx executive_presentation.pptx exists", False, "Not found")
        return False
    record("pptx executive_presentation.pptx exists", True)
    try:
        prs = Presentation(path)
        record("pptx has slides", len(prs.slides) > 0, f"{len(prs.slides)} slides")
    except Exception as e:
        record("pptx readable", False, str(e))
    return True


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False)
    args = parser.parse_args()

    ws = args.agent_workspace
    if not os.path.isdir(ws):
        print(f"Agent workspace not found: {ws}")
        sys.exit(1)

    # --- CRITICAL semantic checks (любой провал => FAIL до порога точности) ---
    check_variance_semantics(ws)
    check_forecast_semantics(ws)
    check_docx_departments(ws)
    check_phase56_email_calendar(ws, args.groundtruth_workspace)

    # --- Non-critical structural / GT checks ---
    check_xlsx_content(ws, args.groundtruth_workspace)
    check_docx_content(ws)
    check_pptx_content(ws)

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
        "critical_failures": CRITICAL_FAILURES,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    # Критический гейт ДО порога точности.
    if CRITICAL_FAILURES:
        print(f"\nCRITICAL FAILURE ({len(CRITICAL_FAILURES)}): {CRITICAL_FAILURES}")
        print("FAIL")
        sys.exit(1)

    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
