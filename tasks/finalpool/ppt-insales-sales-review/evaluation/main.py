"""
Скрипт оценки для задачи ppt-insales-sales-review.

Динамически запрашивает PostgreSQL, чтобы вычислить эталонные метрики
продаж (схема InSales wc.*) и рыночные данные MOEX (схема moex.*),
затем проверяет выходные файлы агента и письмо.

Структура чеков: структурные (NON-critical) + CRITICAL.
Любой провал CRITICAL => FAIL независимо от accuracy.
Иначе PASS при accuracy >= 70%.

Рыночный контекст переориентирован с ^DJI (Dow Jones) на инструмент MOEX
SBER.ME, так как форк moex-finance не содержит ^DJI.
"""

from argparse import ArgumentParser
import sys
import os
from pathlib import Path

import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "cowork_gym",
    "user": "eigent",
    "password": "camel",
}

# Инструмент MOEX, используемый как рыночный контекст вместо ^DJI.
MARKET_SYMBOL = "SBER.ME"

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []
CRITICAL_CHECKS = set()


def crit(name):
    """Пометить чек как CRITICAL и вернуть его имя."""
    CRITICAL_CHECKS.add(name)
    return name


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        marker = " [CRITICAL]" if name in CRITICAL_CHECKS else ""
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL]{marker} {name}{msg}")
        if name in CRITICAL_CHECKS:
            CRITICAL_FAILS.append(name)


def num_close(a, b, tol=1.0):
    """Compare two numeric values with tolerance."""
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def pct_close(a, b, tol=0.5):
    """Compare percentage values with tolerance."""
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def str_match(a, b):
    """Case-insensitive, whitespace-normalized comparison."""
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def get_expected_sales_data():
    """Query PostgreSQL for expected InSales sales metrics."""
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    # Total orders and revenue
    cur.execute("""
        SELECT COUNT(*) as total_orders,
               COUNT(*) FILTER (WHERE status='completed') as completed,
               SUM(total::float) as total_revenue,
               AVG(total::float) as avg_order_value
        FROM wc.orders
    """)
    row = cur.fetchone()
    summary = {
        "Total_Orders": int(row[0]),
        "Completed_Orders": int(row[1]),
        "Total_Revenue": round(float(row[2]), 2),
        "Avg_Order_Value": round(float(row[3]), 2),
    }

    # Top 10 products by revenue
    cur.execute("""
        SELECT li->>'name' as product_name,
               SUM((li->>'quantity')::int) as units_sold,
               SUM((li->>'total')::float) as revenue
        FROM wc.orders, jsonb_array_elements(line_items) as li
        GROUP BY li->>'name'
        ORDER BY revenue DESC
        LIMIT 10
    """)
    top_products = []
    for r in cur.fetchall():
        top_products.append({
            "Product_Name": r[0].strip(),
            "Units_Sold": int(r[1]),
            "Revenue": round(float(r[2]), 2),
        })

    conn.close()
    return summary, top_products


def get_expected_market_data():
    """Запрос эталонных рыночных данных MOEX (схема moex.*)."""
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    # Самая свежая цена закрытия
    cur.execute("""
        SELECT date, close FROM moex.stock_prices
        WHERE symbol=%s ORDER BY date DESC LIMIT 1
    """, (MARKET_SYMBOL,))
    recent = cur.fetchone()
    recent_close = float(recent[1])

    # Цена закрытия примерно месячной давности
    cur.execute("""
        SELECT date, close FROM moex.stock_prices
        WHERE symbol=%s AND date <= (
            SELECT MAX(date) - INTERVAL '28 days' FROM moex.stock_prices WHERE symbol=%s
        )
        ORDER BY date DESC LIMIT 1
    """, (MARKET_SYMBOL, MARKET_SYMBOL))
    month_ago = cur.fetchone()
    month_ago_close = float(month_ago[1])

    monthly_change_pct = round((recent_close - month_ago_close) / month_ago_close * 100, 2)

    conn.close()
    return {
        "Index": MARKET_SYMBOL,
        "Recent_Close": round(recent_close, 2),
        "Monthly_Change_Pct": monthly_change_pct,
    }


# ============================================================================
# Check 1: Excel file
# ============================================================================

def check_excel(workspace, expected_summary, expected_products, expected_dji):
    """Check Monthly_Sales_Review.xlsx for correctness."""
    import openpyxl

    print("\n=== Check 1: Excel File ===")

    xlsx_path = Path(workspace) / "Monthly_Sales_Review.xlsx"
    if not xlsx_path.exists():
        record("Excel file exists", False, f"Not found: {xlsx_path}")
        return

    record("Excel file exists", True)

    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    except Exception as e:
        record("Excel file readable", False, str(e))
        return

    # --- Sales Summary sheet ---
    found_ss = None
    for sn in wb.sheetnames:
        if sn.strip().lower() == "sales summary":
            found_ss = sn
            break
    if not found_ss:
        record("Sales Summary sheet exists", False, f"Sheets found: {wb.sheetnames}")
    else:
        record("Sales Summary sheet exists", True)
        ws = wb[found_ss]
        rows = list(ws.iter_rows(values_only=True))
        if len(rows) < 2:
            record("Sales Summary has data", False, "No data rows")
        else:
            # Find Metric and Value columns
            header = [str(h).strip().lower() if h else "" for h in rows[0]]
            metric_idx = None
            value_idx = None
            for i, h in enumerate(header):
                if "metric" in h:
                    metric_idx = i
                elif "value" in h:
                    value_idx = i

            if metric_idx is None or value_idx is None:
                record("Sales Summary columns", False,
                       f"Expected Metric/Value columns, found: {[str(h) for h in rows[0]]}")
            else:
                record("Sales Summary columns", True)
                # Build a dict from the sheet
                sheet_data = {}
                for row in rows[1:]:
                    if row[metric_idx]:
                        key = str(row[metric_idx]).strip()
                        val = row[value_idx]
                        sheet_data[key.lower()] = val

                for metric, expected_val in expected_summary.items():
                    # Total_Orders и Total_Revenue — суть deliverable из магазина => CRITICAL.
                    name = f"Sales Summary '{metric}'"
                    if metric in ("Total_Orders", "Total_Revenue"):
                        name = crit(name)
                    sheet_val = sheet_data.get(metric.lower())
                    if sheet_val is None:
                        record(name, False, "Metric not found")
                    elif metric in ("Total_Orders", "Completed_Orders"):
                        record(name, num_close(sheet_val, expected_val, tol=0.5),
                               f"expected {expected_val}, got {sheet_val}")
                    else:
                        record(name, num_close(sheet_val, expected_val, tol=5.0),
                               f"expected {expected_val}, got {sheet_val}")

    # --- Top Products sheet ---
    found_tp = None
    for sn in wb.sheetnames:
        if sn.strip().lower() == "top products":
            found_tp = sn
            break
    if not found_tp:
        record("Top Products sheet exists", False, f"Sheets found: {wb.sheetnames}")
    else:
        record("Top Products sheet exists", True)
        ws = wb[found_tp]
        rows = list(ws.iter_rows(values_only=True))
        if len(rows) < 2:
            record("Top Products has data", False, "No data rows")
        else:
            header = [str(h).strip().lower() if h else "" for h in rows[0]]
            name_idx = rev_idx = units_idx = None
            for i, h in enumerate(header):
                if "product" in h and "name" in h:
                    name_idx = i
                elif "revenue" in h or h == "revenue":
                    rev_idx = i
                elif "unit" in h:
                    units_idx = i

            if name_idx is None or rev_idx is None:
                record("Top Products columns", False,
                       f"Expected Product_Name/Revenue columns, found: {[str(h) for h in rows[0]]}")
            else:
                record("Top Products columns", True)
                data_rows = rows[1:]
                record("Top Products row count", len(data_rows) >= 10,
                       f"expected >= 10, got {len(data_rows)}")

                # Check top 3 products match by revenue (allow name truncation)
                for i in range(min(3, len(data_rows), len(expected_products))):
                    exp = expected_products[i]
                    row = data_rows[i]
                    agent_name = str(row[name_idx]).strip().lower() if row[name_idx] else ""
                    exp_name = exp["Product_Name"].strip().lower()
                    # Allow partial match (product names can be truncated)
                    name_ok = exp_name[:30] in agent_name or agent_name[:30] in exp_name
                    # Топ-1 товар по выручке — суть анализа топ-10 => CRITICAL (имя + выручка).
                    nm_check = f"Top Products #{i+1} name"
                    rev_check = f"Top Products #{i+1} revenue"
                    if i == 0:
                        nm_check = crit(nm_check)
                        rev_check = crit(rev_check)
                    record(nm_check, name_ok,
                           f"expected '{exp_name[:60]}...', got '{agent_name[:60]}...'")
                    if rev_idx is not None:
                        record(rev_check, num_close(row[rev_idx], exp["Revenue"], tol=5.0),
                               f"expected {exp['Revenue']}, got {row[rev_idx]}")

    # --- Market Context sheet ---
    found_mc = None
    for sn in wb.sheetnames:
        if sn.strip().lower() == "market context":
            found_mc = sn
            break
    if not found_mc:
        record("Market Context sheet exists", False, f"Sheets found: {wb.sheetnames}")
    else:
        record("Market Context sheet exists", True)
        ws = wb[found_mc]
        rows = list(ws.iter_rows(values_only=True))
        if len(rows) < 2:
            record("Market Context has data", False, "No data rows")
        else:
            header = [str(h).strip().lower() if h else "" for h in rows[0]]
            idx_col = close_col = change_col = None
            for i, h in enumerate(header):
                if h == "index" or "index" in h:
                    idx_col = i
                elif "close" in h:
                    close_col = i
                elif "change" in h or "pct" in h:
                    change_col = i

            if idx_col is None:
                record("Market Context columns", False,
                       f"Expected Index column, found: {[str(h) for h in rows[0]]}")
            else:
                record("Market Context columns", True)
                data_row = rows[1]
                idx_val = str(data_row[idx_col]).strip() if data_row[idx_col] else ""
                # Проверяем, что Index содержит выбранный тикер MOEX (SBER).
                idx_lower = idx_val.lower()
                record("Market Context index",
                       "sber" in idx_lower or MARKET_SYMBOL.lower() in idx_lower,
                       f"expected {MARKET_SYMBOL}-related, got '{idx_val}'")

                if close_col is not None:
                    # Recent_Close для инструмента MOEX — суть рыночного контекста => CRITICAL.
                    # Котировки MOEX порядка ~100, поэтому допуск ужесточён с 500 до 1.0.
                    record(crit("Market Context close"),
                           num_close(data_row[close_col], expected_dji["Recent_Close"], tol=1.0),
                           f"expected ~{expected_dji['Recent_Close']}, got {data_row[close_col]}")

                if change_col is not None:
                    record("Market Context change pct", pct_close(data_row[change_col], expected_dji["Monthly_Change_Pct"], tol=2.0),
                           f"expected ~{expected_dji['Monthly_Change_Pct']}%, got {data_row[change_col]}")

    wb.close()


# ============================================================================
# Check 2: PowerPoint file
# ============================================================================

def check_pptx(workspace, expected_summary):
    """Check Sales_Review_Presentation.pptx for correctness."""
    from pptx import Presentation

    print("\n=== Check 2: PowerPoint File ===")

    pptx_path = Path(workspace) / "Sales_Review_Presentation.pptx"
    if not pptx_path.exists():
        record("PPTX file exists", False, f"Not found: {pptx_path}")
        return

    record("PPTX file exists", True)

    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        record("PPTX file readable", False, str(e))
        return

    slides = list(prs.slides)
    record(crit("PPTX has >= 5 slides"), len(slides) >= 5,
           f"got {len(slides)} slides")

    # Collect all text from all slides
    all_slide_texts = []
    for slide in slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)
        all_slide_texts.append("\n".join(texts))

    full_text = "\n".join(all_slide_texts).lower()

    # Check title slide (RU+EN)
    if slides:
        title_text = all_slide_texts[0].lower()
        record("Title slide mentions sales review",
               any(t in title_text for t in
                   ("sales review", "monthly", "обзор продаж", "ежемесяч", "продаж")),
               f"Title text: {all_slide_texts[0][:100]}")

    # Check that key metrics appear somewhere in the presentation => CRITICAL
    total_orders_str = str(expected_summary["Total_Orders"])
    record(crit("PPTX contains total orders"),
           total_orders_str in full_text,
           f"Looking for '{total_orders_str}' in presentation text")

    # Check revenue appears (various formats)
    rev = expected_summary["Total_Revenue"]
    rev_strs = [
        f"{rev:,.2f}",
        f"{rev:,.0f}",
        f"{rev:.2f}",
        f"{rev:.0f}",
        str(int(rev)),
    ]
    found_rev = any(s in full_text for s in rev_strs)
    record(crit("PPTX contains total revenue"), found_rev,
           f"Looking for revenue ~{rev} in presentation")

    # Check for market content (RU+EN, repointed to MOEX/SBER)
    record("PPTX mentions market or MOEX",
           any(t in full_text for t in
               ("sber", "moex", "market", "рынок", "рыноч", "мосбирж",
                "московск", "сбер", "котировк", "индекс")),
           "No market/MOEX/SBER content found")

    # Check for recommendations slide (RU+EN)
    record("PPTX mentions recommendations/takeaways",
           any(t in full_text for t in
               ("recommend", "takeaway", "action", "insight", "conclusion",
                "рекоменд", "вывод", "итог", "заключ")),
           "No recommendations/takeaways content found")


# ============================================================================
# Check 3: Email
# ============================================================================

def check_email(expected_summary):
    """Check that email was sent to sales@company.com with correct subject."""
    print("\n=== Check 3: Email ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    # Look for the sales review email
    cur.execute("""
        SELECT id, subject, to_addr, body_text, from_addr
        FROM email.messages
        WHERE subject ILIKE '%monthly sales review%'
        ORDER BY date DESC
        LIMIT 5
    """)
    rows = cur.fetchall()

    if not rows:
        # Try broader search
        cur.execute("""
            SELECT id, subject, to_addr, body_text, from_addr
            FROM email.messages
            WHERE subject ILIKE '%sales%review%'
            ORDER BY date DESC
            LIMIT 5
        """)
        rows = cur.fetchall()

    if not rows:
        record(crit("Email sent"), False, "No email with 'sales review' in subject found")
        conn.close()
        return

    record(crit("Email sent"), True)

    # Check at least one email goes to sales@company.com
    found_correct_recipient = False
    matching_row = None
    for row in rows:
        msg_id, subject, to_addr, body_text, from_addr = row
        # to_addr is JSONB (list)
        if to_addr:
            recipients = to_addr if isinstance(to_addr, list) else [to_addr]
            for r in recipients:
                addr = str(r).strip().lower()
                if "sales@company.com" in addr:
                    found_correct_recipient = True
                    matching_row = row
                    break
        if found_correct_recipient:
            break

    record(crit("Email to sales@company.com"), found_correct_recipient,
           f"Recipients found: {[r[2] for r in rows]}")

    if matching_row:
        msg_id, subject, to_addr, body_text, from_addr = matching_row

        # Check subject (literal grep, остаётся английским)
        record(crit("Email subject correct"),
               "monthly sales review" in str(subject).lower(),
               f"Subject: {str(subject)[:100]}")

        # Check body contains key info => CRITICAL
        body = str(body_text).lower() if body_text else ""
        body_has_orders = str(expected_summary["Total_Orders"]) in body
        record(crit("Email body mentions total orders"), body_has_orders,
               f"Looking for {expected_summary['Total_Orders']} in body")

        # Check body mentions revenue (any format)
        rev = expected_summary["Total_Revenue"]
        rev_strs = [
            f"{rev:,.2f}", f"{rev:,.0f}", f"{rev:.2f}", f"{rev:.0f}", str(int(rev)),
        ]
        body_has_rev = any(s in body for s in rev_strs)
        record(crit("Email body mentions revenue"), body_has_rev,
               f"Looking for revenue ~{rev} in body")

    conn.close()


# ============================================================================
# Main
# ============================================================================

if __name__ == "__main__":
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False)
    args = parser.parse_args()

    workspace = args.agent_workspace
    if not workspace:
        print("Error: --agent_workspace is required")
        sys.exit(1)

    print("Fetching expected data from database...")
    try:
        expected_summary, expected_products = get_expected_sales_data()
        expected_dji = get_expected_market_data()
        print(f"  Sales: {expected_summary}")
        print(f"  Top product: {expected_products[0]['Product_Name'][:50]}...")
        print(f"  Market ({MARKET_SYMBOL}): {expected_dji}")
    except Exception as e:
        print(f"Error querying database: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

    # Run all checks
    check_excel(workspace, expected_summary, expected_products, expected_dji)
    check_pptx(workspace, expected_summary)
    check_email(expected_summary)

    # Summary
    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100) if total else 0.0
    print(f"\n{'='*60}")
    print(f"Results: {PASS_COUNT} passed, {FAIL_COUNT} failed  (accuracy {accuracy:.1f}%)")
    if CRITICAL_FAILS:
        print(f"CRITICAL failures: {CRITICAL_FAILS}")
    print(f"{'='*60}")

    if CRITICAL_FAILS:
        print("\nOverall: FAIL (провален CRITICAL-чек).")
        sys.exit(1)
    if accuracy >= 70:
        print("\nPass all tests!")
        sys.exit(0)
    print("\nOverall: FAIL (accuracy below 70%).")
    sys.exit(1)
