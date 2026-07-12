"""
Evaluation script for canvas-assignment-effectiveness-ppt-teamly-email task.

Checks:
1. Excel Assessment_Effectiveness.xlsx: 3 sheets, ~52 metric rows, DI values
   match groundtruth, Revision Needed populated.
2. PowerPoint Curriculum_Review.pptx with 6+ slides covering required sections.
3. Teamly (ex-Notion) tracker page 'Assignment Improvement Tracker' with a
   table of revision-needed assignments (Status=Pending, Target_DI=0.4).
4. Email to curriculum_committee@university.edu with effectiveness breakdown.

CRITICAL_CHECKS (semantic): any failure => overall FAIL regardless of accuracy.
Otherwise pass threshold: accuracy >= 70%.
"""
import argparse
import json
import os
import sys

import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "cowork_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0
FAILED_NAMES = []

# Critical checks reflect the task's analytical substance, not structure.
CRITICAL_CHECKS = {
    "DI values match groundtruth (>= 70%)",
    "Assignment Metrics has ~52 data rows",
    "Revision Needed has entries",
    "Teamly tracker page 'Assignment Improvement Tracker' exists",
    "Teamly tracker page lists revision assignments with Pending/Target_DI=0.4",
    "Email to curriculum_committee@university.edu found",
    "Email body reports effectiveness-category breakdown",
}


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        FAILED_NAMES.append(name)
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def di_close(a, b, tol=0.05):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


# ============================================================
# Check 1: Excel
# ============================================================
def check_excel(agent_workspace, gt_workspace):
    print("\n=== Checking Excel ===")
    xlsx_path = os.path.join(agent_workspace, "Assessment_Effectiveness.xlsx")
    if not os.path.isfile(xlsx_path):
        check("Assessment_Effectiveness.xlsx exists", False, f"Not found: {xlsx_path}")
        # critical checks depending on the file cannot pass
        check("Assignment Metrics has ~52 data rows", False, "xlsx missing")
        check("DI values match groundtruth (>= 70%)", False, "xlsx missing")
        check("Revision Needed has entries", False, "xlsx missing")
        return

    check("Assessment_Effectiveness.xlsx exists", True)

    try:
        import openpyxl
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        sheet_names = wb.sheetnames

        # Check sheets exist (structural)
        check("Sheet 'Assignment Metrics' exists",
              any("assignment" in s.lower() and "metric" in s.lower() for s in sheet_names),
              f"Sheets: {sheet_names}")
        check("Sheet 'Course Summary' exists",
              any("course" in s.lower() and "summary" in s.lower() for s in sheet_names),
              f"Sheets: {sheet_names}")
        check("Sheet 'Revision Needed' exists",
              any("revision" in s.lower() for s in sheet_names),
              f"Sheets: {sheet_names}")

        # Assignment Metrics sheet
        metrics_ws = None
        for s in sheet_names:
            if "assignment" in s.lower() and "metric" in s.lower():
                metrics_ws = wb[s]
                break
        if not metrics_ws:
            metrics_ws = wb[sheet_names[0]]

        rows = list(metrics_ws.iter_rows(values_only=True))
        # CRITICAL: ~52 data rows (one per Fall 2014 assignment + header).
        check("Assignment Metrics has ~52 data rows",
              len(rows) >= 40,
              f"Found {len(rows)} rows (expected ~52 incl. header)")

        if len(rows) > 1:
            all_text = " ".join(str(c) for row in rows for c in row if c).lower()
            check("Contains Fall 2014 course names",
                  "аналитик" in all_text or "биохими" in all_text,
                  f"Sample: {all_text[:200]}")
            check("Contains effectiveness labels",
                  "good" in all_text or "acceptable" in all_text or "poor" in all_text,
                  f"Sample: {all_text[:200]}")

        # CRITICAL: compare DI vs groundtruth. De-gated: missing GT => FAIL.
        gt_xlsx = os.path.join(gt_workspace, "Assessment_Effectiveness.xlsx")
        if not os.path.isfile(gt_xlsx):
            check("DI values match groundtruth (>= 70%)", False,
                  f"Groundtruth xlsx not found at {gt_xlsx} (required for DI spot-check)")
        else:
            gt_wb = openpyxl.load_workbook(gt_xlsx, data_only=True)
            gt_metrics = None
            for s in gt_wb.sheetnames:
                if "assignment" in s.lower() and "metric" in s.lower():
                    gt_metrics = gt_wb[s]
                    break
            if not gt_metrics:
                gt_metrics = gt_wb[gt_wb.sheetnames[0]]

            gt_rows = list(gt_metrics.iter_rows(values_only=True))
            agent_rows = rows

            check("Assignment count matches groundtruth",
                  abs(len(agent_rows) - len(gt_rows)) <= 3,
                  f"Agent: {len(agent_rows)}, GT: {len(gt_rows)}")

            gt_data = {}
            for row in gt_rows[1:]:
                if row[0] and row[1]:
                    key = (str(row[0]).strip(), str(row[1]).strip())
                    gt_data[key] = row

            matches = 0
            total_checked = 0
            for row in agent_rows[1:]:
                if row[0] and row[1]:
                    key = (str(row[0]).strip(), str(row[1]).strip())
                    if key in gt_data:
                        total_checked += 1
                        gt_row = gt_data[key]
                        if len(row) > 7 and len(gt_row) > 7:
                            if di_close(row[7], gt_row[7], tol=0.1):
                                matches += 1
            if total_checked > 0:
                match_rate = matches / total_checked
                check("DI values match groundtruth (>= 70%)",
                      match_rate >= 0.7,
                      f"{matches}/{total_checked} = {match_rate:.1%}")
            else:
                check("DI values match groundtruth (>= 70%)", False,
                      "No matching (Course, Assignment) rows found vs groundtruth")

        # Course Summary sheet (structural)
        summary_ws = None
        for s in sheet_names:
            if "course" in s.lower() and "summary" in s.lower():
                summary_ws = wb[s]
                break
        if summary_ws:
            summary_rows = list(summary_ws.iter_rows(values_only=True))
            check("Course Summary has 7 courses + header",
                  len(summary_rows) >= 7,
                  f"Found {len(summary_rows)} rows")

        # Revision Needed sheet — CRITICAL: must be populated.
        revision_ws = None
        for s in sheet_names:
            if "revision" in s.lower():
                revision_ws = wb[s]
                break
        if revision_ws:
            revision_rows = list(revision_ws.iter_rows(values_only=True))
            check("Revision Needed has entries",
                  len(revision_rows) >= 10,
                  f"Found {len(revision_rows)} rows (expected ~18 incl. header)")
        else:
            check("Revision Needed has entries", False, "Revision Needed sheet missing")

    except ImportError:
        check("openpyxl available", False, "Cannot parse Excel without openpyxl")
        check("Assignment Metrics has ~52 data rows", False, "no openpyxl")
        check("DI values match groundtruth (>= 70%)", False, "no openpyxl")
        check("Revision Needed has entries", False, "no openpyxl")
    except Exception as e:
        check("Excel parsing", False, str(e))
        check("Assignment Metrics has ~52 data rows", False, str(e))
        check("DI values match groundtruth (>= 70%)", False, str(e))
        check("Revision Needed has entries", False, str(e))


# ============================================================
# Check 2: PowerPoint
# ============================================================
def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint ===")
    pptx_path = os.path.join(agent_workspace, "Curriculum_Review.pptx")
    if not os.path.isfile(pptx_path):
        check("Curriculum_Review.pptx exists", False, f"Not found: {pptx_path}")
        return

    check("Curriculum_Review.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        check("PPT has at least 6 slides", slide_count >= 6,
              f"Found {slide_count} slides")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text.lower() + " "

        # Accept RU + EN wording.
        check("PPT mentions discrimination index or DI",
              any(t in all_text for t in ["discrimination", "дискриминац", " di ", "di=", "индекс"]),
              "Missing DI content")
        check("PPT mentions methodology or formula",
              any(t in all_text for t in ["methodol", "formula", "методик", "формул", "27%"]),
              "Missing methodology")
        check("PPT mentions revision or improvement",
              any(t in all_text for t in ["revision", "improv", "needs", "доработ", "улучш"]),
              "Missing revision content")
        check("PPT mentions recommendations",
              any(t in all_text for t in ["recommend", "suggest", "action", "рекоменд", "действ"]),
              "Missing recommendations")
        check("PPT mentions Fall 2014 or course names",
              "fall 2014" in all_text or "2014" in all_text,
              "Missing term reference")

    except ImportError:
        size = os.path.getsize(pptx_path)
        check("PPT file has content (>5KB)", size > 5000, f"Size: {size}")
    except Exception as e:
        check("PPT parsing", False, str(e))


# ============================================================
# Check 3: Teamly (ex-Notion) tracker page
# ============================================================
def check_teamly():
    print("\n=== Checking Teamly ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        # Tracker page: English title preserved per task.md, but accept RU too.
        cur.execute("""
            SELECT id, title, COALESCE(body, '')
            FROM teamly.pages
            WHERE title ILIKE '%%improvement%%tracker%%'
               OR title ILIKE '%%assignment%%tracker%%'
               OR title ILIKE '%%трекер%%доработ%%'
        """)
        pages = cur.fetchall()

        # CRITICAL: page exists.
        check("Teamly tracker page 'Assignment Improvement Tracker' exists",
              len(pages) >= 1,
              f"Found {len(pages)} matching pages")

        if not pages:
            check("Teamly tracker page lists revision assignments with Pending/Target_DI=0.4",
                  False, "No tracker page")
            cur.close()
            conn.close()
            return

        # Combined body of matching pages (original case -> lower for keyword search).
        combined = " ".join(str(b) for _, _, b in pages)
        combined_l = combined.lower()
        max_len = max(len(str(b)) for _, _, b in pages)

        check("Tracker page has non-trivial body", max_len >= 100,
              f"Longest matching page body is {max_len} chars")

        # Property/column names present (structural).
        check("Tracker mentions Assignment/Course columns",
              "assignment" in combined_l and "course" in combined_l,
              f"Body sample: {combined_l[:200]}")
        check("Tracker mentions Status/Target_DI columns",
              "status" in combined_l and ("target_di" in combined_l or "target di" in combined_l),
              f"Body sample: {combined_l[:200]}")

        # CRITICAL semantic: must list real revision-needed assignments and the
        # required follow-up values (Status=Pending, Target_DI=0.4) on multiple rows.
        pending_count = combined_l.count("pending")
        target_count = combined_l.count("0.4")
        # Assignment names come from live English Canvas: cma/tma/final etc.
        has_assignment_names = any(t in combined_l for t in ["cma", "tma", "final"])
        substantive = (pending_count >= 10 and target_count >= 10 and has_assignment_names)
        check("Teamly tracker page lists revision assignments with Pending/Target_DI=0.4",
              substantive,
              f"pending={pending_count} (need >=10), '0.4'={target_count} (need >=10), "
              f"assignment_names={has_assignment_names}")

        # Course references (structural).
        check("Tracker mentions course references",
              any(t in combined_l for t in ["финанс", "креативн", "вычислен",
                                            "аналитик", "биохими"]),
              f"Body sample: {combined_l[:200]}")

        cur.close()
        conn.close()
    except Exception as e:
        check("Teamly tracker page 'Assignment Improvement Tracker' exists", False, str(e))
        check("Teamly tracker page lists revision assignments with Pending/Target_DI=0.4",
              False, str(e))


# ============================================================
# Check 4: Email
# ============================================================
def check_email():
    print("\n=== Checking Email ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        cur.execute("""
            SELECT id, subject, from_addr, to_addr, body_text
            FROM email.messages
            WHERE to_addr::text ILIKE '%%curriculum_committee@university.edu%%'
               OR to_addr::text ILIKE '%%curriculum%%committee%%'
               OR subject ILIKE '%%assignment%%effectiveness%%'
               OR subject ILIKE '%%fall 2014%%assignment%%'
        """)
        emails = cur.fetchall()
        # CRITICAL: email exists.
        check("Email to curriculum_committee@university.edu found",
              len(emails) >= 1,
              f"Found {len(emails)} matching emails")

        if emails:
            email = emails[0]
            subject = str(email[1] or "").lower()
            body = str(email[4] or "").lower()

            check("Email subject mentions assignment or effectiveness",
                  "assignment" in subject or "effectiveness" in subject or "fall 2014" in subject,
                  f"Subject: {email[1]}")
            check("Email body has substantive content",
                  len(body) > 50,
                  f"Body length: {len(body)}")
            # CRITICAL semantic: body reports the effectiveness-category breakdown.
            # Accept EN labels (preserved identifiers Good/Acceptable/Poor) and RU prose.
            label_hits = sum(1 for t in ["good", "acceptable", "poor"] if t in body)
            mentions_revision = any(t in body for t in
                                    ["revision", "доработ", "completion", "выполнен",
                                     "discrimination", "дискриминац"])
            check("Email body reports effectiveness-category breakdown",
                  label_hits >= 2 and mentions_revision,
                  f"category-label hits={label_hits} (need >=2), revision/completion mention={mentions_revision}")
        else:
            check("Email body reports effectiveness-category breakdown", False, "No email")

        cur.close()
        conn.close()
    except Exception as e:
        check("Email to curriculum_committee@university.edu found", False, str(e))
        check("Email body reports effectiveness-category breakdown", False, str(e))


# ============================================================
# Main
# ============================================================
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=True)
    parser.add_argument("--groundtruth_workspace", required=False, default="")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    gt_ws = args.groundtruth_workspace or args.agent_workspace

    check_excel(args.agent_workspace, gt_ws)
    check_pptx(args.agent_workspace)
    check_teamly()
    check_email()

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100) if total else 0.0
    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}, Failed: {FAIL_COUNT} ({accuracy:.1f}%)")

    critical_failed = [n for n in FAILED_NAMES if n in CRITICAL_CHECKS]
    if critical_failed:
        print(f"  CRITICAL FAILURES: {len(critical_failed)}")
        for n in critical_failed:
            print(f"    - {n}")

    if args.res_log_file:
        result = {
            "total_passed": PASS_COUNT,
            "total_checks": total,
            "accuracy": accuracy,
            "critical_failed": critical_failed,
        }
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if total == 0:
        print("  Overall: FAIL (no checks performed)")
        sys.exit(1)
    if critical_failed:
        print("  Overall: FAIL (critical check failed)")
        sys.exit(1)
    if accuracy >= 70:
        print("  Overall: PASS")
        sys.exit(0)
    print("  Overall: FAIL")
    sys.exit(1)


if __name__ == "__main__":
    main()
