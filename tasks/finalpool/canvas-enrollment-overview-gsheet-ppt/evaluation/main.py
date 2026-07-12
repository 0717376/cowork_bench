"""
Evaluation script for canvas-enrollment-overview-gsheet-ppt task.

Critical checks (CRITICAL_CHECKS): any failure => overall FAIL regardless of
accuracy. Otherwise PASS requires accuracy >= 70%.

Checks:
1. Google Sheet "Fall 2013 Enrollment Overview" with correct per-course
   Student_Count/Teacher_Count, sorted by Student_Count descending.
2. PowerPoint Enrollment_Overview_F2013.pptx with title/course/summary slides.
3. Word document Enrollment_Report_F2013.docx with heading, summary, table.
4. Email to academic.office@university.edu with the correct subject.

Source-of-truth values (canvas, Fall 2013 courses), read here only to verify
the agent's honestly-derived numbers:
  AAA-2013J=383, BBB-2013J=2237, DDD-2013J=1938,
  EEE-2013J=1052, FFF-2013J=2283, GGG-2013J=952
  Total students = 8845, largest = FFF-2013J (Основы финансов).
"""

import argparse
import json
import os
import re
import sys

import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "cowork_gym",
    "user": "eigent",
    "password": "camel",
}

# Source-of-truth Fall 2013 enrollment (from canvas schema, deterministic).
EXPECTED = {
    "aaa-2013j": 383,
    "bbb-2013j": 2237,
    "ddd-2013j": 1938,
    "eee-2013j": 1052,
    "fff-2013j": 2283,
    "ggg-2013j": 952,
}
EXPECTED_TEACHERS = {
    "aaa-2013j": 1,
    "bbb-2013j": 3,
    "ddd-2013j": 1,
    "eee-2013j": 1,
    "fff-2013j": 1,
    "ggg-2013j": 2,
}
TOTAL_STUDENTS = 8845
LARGEST_CODE = "fff-2013j"
SORT_ORDER = ["fff-2013j", "bbb-2013j", "ddd-2013j",
              "eee-2013j", "ggg-2013j", "aaa-2013j"]

# Critical checks: any failure => overall FAIL regardless of accuracy.
CRITICAL_CHECKS = {
    "Sheet has all 6 Fall 2013 courses with correct Student_Count",
    "Sheet rows sorted by Student_Count descending",
    "Word doc reports total student enrollment 8845",
    "PPT reports total student enrollment 8845",
    "PPT summary identifies FFF-2013J as largest enrollment",
    "Email body reports total 8845 and 6 courses",
    "Email sent TO academic.office@university.edu with correct subject",
}

PASS_COUNT = 0
FAIL_COUNT = 0
FAILED_NAMES = []


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        FAILED_NAMES.append(name)
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def digits_in(text):
    """All integer values present in text, with thousands separators removed."""
    nums = set()
    # 8,845 or 8 845 style
    for m in re.findall(r"\d[\d.,\s]*\d|\d", text):
        cleaned = re.sub(r"[.,\s]", "", m)
        if cleaned.isdigit():
            nums.add(int(cleaned))
    return nums


# ============================================================================
# Check 1: Google Sheet
# ============================================================================

def check_gsheet():
    print("\n=== Checking Google Sheet ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("SELECT id, title FROM gsheet.spreadsheets")
    sheets = cur.fetchall()
    print(f"[check_gsheet] Found {len(sheets)} spreadsheets.")
    record("At least 1 spreadsheet created", len(sheets) >= 1)

    target = None
    for ss_id, title in sheets:
        if title and ("enrollment" in title.lower() or "fall 2013" in title.lower()):
            target = (ss_id, title)
            break

    if not target:
        record("Fall 2013 Enrollment spreadsheet found", False,
               f"Spreadsheets: {[(s[0], s[1]) for s in sheets]}")
        record("Sheet has all 6 Fall 2013 courses with correct Student_Count",
               False, "no spreadsheet")
        record("Sheet rows sorted by Student_Count descending", False, "no spreadsheet")
        cur.close()
        conn.close()
        return False

    ss_id, title = target
    record(f"Spreadsheet '{title}' found", True)

    cur.execute("""
        SELECT c.row_index, c.col_index, c.value
        FROM gsheet.cells c
        JOIN gsheet.sheets s ON c.sheet_id = s.id
        WHERE c.spreadsheet_id = %s
        ORDER BY c.row_index, c.col_index
    """, (ss_id,))
    cells = cur.fetchall()
    cur.close()
    conn.close()

    grid = {}
    for row_i, col_i, val in cells:
        grid.setdefault(row_i, {})[col_i] = val

    max_row = max(grid.keys()) if grid else 0
    record("Spreadsheet has at least 7 rows (header + 6 courses)",
           max_row >= 6, f"Max row index: {max_row}")

    # Locate the Course_Code and Student_Count columns from the header row.
    header_row = min(grid.keys()) if grid else 0
    header = {c: str(v).strip().lower() for c, v in grid.get(header_row, {}).items()}
    code_col = next((c for c, v in header.items() if "course_code" in v or v == "code"), None)
    student_col = next((c for c, v in header.items()
                        if "student_count" in v or v == "student" or v == "students"), None)
    teacher_col = next((c for c, v in header.items()
                        if "teacher_count" in v or v == "teacher" or v == "teachers"), None)

    record("Header has Course_Code, Student_Count, Teacher_Count columns",
           code_col is not None and student_col is not None and teacher_col is not None,
           f"header={header}")

    # Build per-course mapping from the data rows.
    found = {}  # code -> (student, teacher)
    ordered = []  # (code, student) in sheet order
    for r in sorted(grid.keys()):
        if r == header_row:
            continue
        row = grid[r]
        # Find a course code anywhere in the row (robust to column detection).
        row_code = None
        for v in row.values():
            sv = str(v).strip().lower()
            if sv in EXPECTED:
                row_code = sv
                break
        if not row_code:
            continue
        s_val = None
        t_val = None
        if student_col is not None and student_col in row:
            try:
                s_val = int(float(str(row[student_col]).replace(",", "").strip()))
            except (TypeError, ValueError):
                s_val = None
        if teacher_col is not None and teacher_col in row:
            try:
                t_val = int(float(str(row[teacher_col]).replace(",", "").strip()))
            except (TypeError, ValueError):
                t_val = None
        found[row_code] = (s_val, t_val)
        ordered.append((row_code, s_val))

    # CRITICAL: all 6 courses present with correct Student_Count.
    all_ok = True
    missing = []
    for code, exp in EXPECTED.items():
        s_val = found.get(code, (None, None))[0]
        if s_val != exp:
            all_ok = False
            missing.append(f"{code.upper()}: got {s_val}, want {exp}")
    record("Sheet has all 6 Fall 2013 courses with correct Student_Count",
           all_ok, "; ".join(missing) if missing else "")

    # Non-critical: teacher counts.
    teach_ok = all(found.get(code, (None, None))[1] == EXPECTED_TEACHERS[code]
                   for code in EXPECTED)
    record("Sheet has correct Teacher_Count per course", teach_ok,
           f"found={ {k: v[1] for k, v in found.items()} }")

    # CRITICAL: sort order by Student_Count descending.
    sheet_order = [c for c, _ in ordered]
    record("Sheet rows sorted by Student_Count descending",
           sheet_order == SORT_ORDER,
           f"got order {sheet_order}")

    return True


# ============================================================================
# Check 2: PowerPoint
# ============================================================================

def check_pptx(agent_workspace):
    print("\n=== Checking Enrollment_Overview_F2013.pptx ===")

    pptx_path = os.path.join(agent_workspace, "Enrollment_Overview_F2013.pptx")
    if not os.path.isfile(pptx_path):
        record("PPT file exists", False, f"Not found: {pptx_path}")
        record("PPT reports total student enrollment 8845", False, "no file")
        record("PPT summary identifies FFF-2013J as largest enrollment", False, "no file")
        return False
    record("PPT file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        record("PPT has at least 8 slides (title + 6 courses + summary)",
               slide_count >= 8, f"Found {slide_count} slides")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text.lower() + " "

        record("PPT mentions Fall 2013",
               "fall 2013" in all_text or "2013" in all_text, "Missing year 2013")
        record("PPT mentions enrollment",
               any(k in all_text for k in
                   ["enrollment", "enrolled", "students", "студент", "зачисл"]),
               "Missing enrollment content")

        for code in ["fff-2013j", "bbb-2013j"]:
            record(f"PPT mentions {code.upper()}", code in all_text,
                   f"Missing {code} in slides")

        nums = digits_in(all_text)
        record("PPT reports total student enrollment 8845",
               TOTAL_STUDENTS in nums,
               f"8845 not found as explicit number; nums sample={sorted(nums)[:15]}")

        # Largest enrollment correctly identified: FFF code or its course name.
        record("PPT summary identifies FFF-2013J as largest enrollment",
               "fff-2013j" in all_text or "основы финансов" in all_text,
               "FFF-2013J / Основы финансов not present in PPT")

        return slide_count >= 8

    except ImportError:
        size = os.path.getsize(pptx_path)
        record("PPT file has content (>5KB)", size > 5000, f"Size: {size} bytes")
        record("PPT reports total student enrollment 8845", False,
               "pptx lib unavailable")
        record("PPT summary identifies FFF-2013J as largest enrollment", False,
               "pptx lib unavailable")
        return size > 5000
    except Exception as e:
        record("PPT file readable", False, str(e))
        return False


# ============================================================================
# Check 3: Word document
# ============================================================================

def check_word(agent_workspace):
    print("\n=== Checking Enrollment_Report_F2013.docx ===")

    docx_path = os.path.join(agent_workspace, "Enrollment_Report_F2013.docx")
    if not os.path.isfile(docx_path):
        record("Word file exists", False, f"Not found: {docx_path}")
        record("Word doc reports total student enrollment 8845", False, "no file")
        return False
    record("Word file exists", True)

    try:
        from docx import Document
        doc = Document(docx_path)
        all_text = " ".join(p.text for p in doc.paragraphs)
        all_text_lc = all_text.lower()
        headings = " ".join(p.text for p in doc.paragraphs
                            if p.style.name.startswith("Heading")).lower()

        record("Word doc has substantial content",
               len(all_text_lc.strip()) >= 100,
               f"Content length: {len(all_text_lc.strip())}")
        record("Word doc mentions enrollment",
               any(k in all_text_lc for k in
                   ["enrollment", "enroll", "зачисл", "студент"]),
               "Missing 'enrollment' in document")
        # Accept English or Russian heading wording.
        record("Word doc has Fall 2013 Enrollment Report heading",
               ("fall 2013 enrollment report" in headings or "2013" in headings
                or "fall 2013 enrollment report" in all_text_lc),
               f"headings={headings[:120]}")

        # Include table cell text in the numeric scan (table may hold the figures).
        table_text = ""
        for t in doc.tables:
            for row in t.rows:
                for cell in row.cells:
                    table_text += " " + cell.text
        scan = all_text + " " + table_text

        record("Word doc has at least 1 table",
               len(doc.tables) >= 1, f"Found {len(doc.tables)} tables")

        nums = digits_in(scan)
        record("Word doc reports total student enrollment 8845",
               TOTAL_STUDENTS in nums,
               f"8845 not found as explicit number; nums sample={sorted(nums)[:15]}")
        # Non-critical: number of courses (6) mentioned.
        record("Word doc mentions 6 courses", 6 in nums,
               "course count 6 not found")

        return True

    except ImportError:
        size = os.path.getsize(docx_path)
        record("Word file has content (>3KB)", size > 3000, f"Size: {size} bytes")
        record("Word doc reports total student enrollment 8845", False,
               "docx lib unavailable")
        return size > 3000
    except Exception as e:
        record("Word file readable", False, str(e))
        return False


# ============================================================================
# Check 4: Email
# ============================================================================

def check_emails():
    print("\n=== Checking Emails ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
    all_emails = cur.fetchall()
    cur.close()
    conn.close()

    print(f"[check_emails] Found {len(all_emails)} total emails.")
    record("At least 1 email sent", len(all_emails) >= 1, f"Found {len(all_emails)}")

    # CRITICAL: AND of correct recipient and correct subject.
    target = None
    for subject, from_addr, to_addr, body_text in all_emails:
        to_str = str(to_addr or "").lower()
        subject_lc = (subject or "").lower()
        if ("academic.office@university.edu" in to_str
                and "fall 2013 semester enrollment summary" in subject_lc):
            target = (subject, from_addr, to_addr, body_text)
            break

    record("Email sent TO academic.office@university.edu with correct subject",
           target is not None,
           f"Emails: {[(e[0], str(e[2])[:60]) for e in all_emails[:3]]}")

    if not target:
        record("Email body reports total 8845 and 6 courses", False, "no matching email")
        return False

    subject, from_addr, to_addr, body_text = target
    body = body_text or ""
    body_lc = body.lower()
    record("Email body mentions enrollment content",
           any(k in body_lc for k in
               ["enrollment", "students", "courses", "студент", "зачисл", "курс"]),
           "Body missing enrollment content")

    nums = digits_in(body)
    # CRITICAL: total 8845 AND 6 courses present as explicit numbers.
    record("Email body reports total 8845 and 6 courses",
           TOTAL_STUDENTS in nums and 6 in nums,
           f"nums sample={sorted(nums)[:15]}")
    # Non-critical: largest course named.
    record("Email body names FFF-2013J / Основы финансов as largest",
           "fff-2013j" in body_lc or "основы финансов" in body_lc,
           "largest course not named")

    return True


# ============================================================================
# Main
# ============================================================================

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_gsheet()
    check_pptx(args.agent_workspace)
    check_word(args.agent_workspace)
    check_emails()

    total = PASS_COUNT + FAIL_COUNT
    accuracy = PASS_COUNT / total * 100 if total else 0
    critical_failed = [n for n in FAILED_NAMES if n in CRITICAL_CHECKS]

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}/{total} ({accuracy:.1f}%)")
    print(f"  Failed: {FAIL_COUNT}")
    if critical_failed:
        print(f"  CRITICAL FAILURES: {len(critical_failed)}")
        for n in critical_failed:
            print(f"    - {n}")

    success = (not critical_failed) and accuracy >= 70
    print(f"  Overall: {'PASS' if success else 'FAIL'}")

    if args.res_log_file:
        result = {
            "passed": PASS_COUNT,
            "failed": FAIL_COUNT,
            "total_checks": total,
            "accuracy": accuracy,
            "critical_failed": critical_failed,
            "success": success,
        }
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
