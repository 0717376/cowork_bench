"""
Evaluation for scholarly-arxiv-ppt-gsheet task.
Checks: PPT file, GSheet, Word document.

Критические проверки (CRITICAL_CHECKS): любой провал => общий FAIL независимо от
accuracy. В остальном PASS требует accuracy >= 70%.
"""
import argparse
import json
import os
import sys
import unicodedata

import psycopg2
from docx import Document
from pptx import Presentation

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "cowork_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0
FAILED_NAMES = []

# Три релевантные статьи (английские идентификаторы сохранены).
RELEVANT_ARXIV_IDS = ["2301.00234", "2302.11382", "2305.10601"]

# Шумовые статьи: их заголовки НЕ должны попасть в результаты (нижний регистр).
NOISE_TITLE_FRAGMENTS = [
    "segment anything",
    "scaling laws for neural machine translation",
    "image captioning",
    "super-resolution for satellite imagery",
]

# Критические проверки: любой провал => общий FAIL независимо от accuracy.
CRITICAL_CHECKS = {
    "Relevance filter: noise papers ABSENT from pptx",
    "Relevance filter: noise papers ABSENT from gsheet",
    "Relevance filter: noise papers ABSENT from docx",
    "All 3 relevant papers present in pptx AND gsheet AND docx",
    "GSheet header has required columns (Title/Authors/ArXiv_ID/Published_Date/Key_Method) and exactly 3 data rows",
    "GSheet contains correct arXiv IDs 2301.00234 / 2302.11382 / 2305.10601",
    "PPT has title slide (survey+reasoning) and final summary/comparison slide; total slides >= 5",
}


def normalize(s):
    """NFKD + cyr->lat транслитерация. ТОЛЬКО для сопоставления смешанных
    кириллица/латиница идентификаторов (например, ID). НЕ использовать для
    поиска русских ключевых слов."""
    if s is None:
        return ""
    s = unicodedata.normalize("NFKD", str(s))
    table = str.maketrans({
        "А": "A", "В": "B", "Е": "E", "К": "K", "М": "M", "Н": "H",
        "О": "O", "Р": "P", "С": "C", "Т": "T", "У": "Y", "Х": "X",
        "а": "a", "е": "e", "о": "o", "р": "p", "с": "c", "у": "y", "х": "x",
    })
    return s.translate(table)


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        FAILED_NAMES.append(name)
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


# Состояние, разделяемое между проверками для критических агрегатов.
STATE = {
    "ppt_text": "",
    "ppt_slides": 0,
    "gsheet_text": "",
    "gsheet_id_text": "",
    "docx_text": "",
}


def _has_all_relevant(text):
    has_chain = "chain-of-thought" in text or "chain of thought" in text
    has_self = "self-consistency" in text or "self consistency" in text
    has_tree = "tree of thoughts" in text or "tree of thought" in text
    return has_chain and has_self and has_tree


def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint ===")
    ppt_path = os.path.join(agent_workspace, "LLM_Reasoning_Survey.pptx")
    if not os.path.isfile(ppt_path):
        record("PPT file LLM_Reasoning_Survey.pptx exists", False, f"Not found at: {ppt_path}")
        return
    record("PPT file LLM_Reasoning_Survey.pptx exists", True)

    try:
        prs = Presentation(ppt_path)
    except Exception as e:
        record("PPT file readable", False, str(e))
        return
    record("PPT file readable", True)

    num_slides = len(prs.slides)
    STATE["ppt_slides"] = num_slides
    record("PPT has at least 5 slides", num_slides >= 5, f"Found {num_slides} slides")

    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text.append(shape.text.lower())
    full_text = " ".join(all_text)
    STATE["ppt_text"] = full_text

    has_title = "survey" in full_text and ("reasoning" in full_text or "llm" in full_text)
    has_summary = "summary" in full_text or "comparison" in full_text or "conclusion" in full_text \
        or "сравнен" in full_text or "итог" in full_text or "заключен" in full_text
    record("PPT contains survey/reasoning title content", has_title, "Looked for 'survey' and 'reasoning'")
    record("PPT has a summary/conclusion slide", has_summary)

    # CRITICAL: title + summary slide + >=5 slides
    record("PPT has title slide (survey+reasoning) and final summary/comparison slide; total slides >= 5",
           has_title and has_summary and num_slides >= 5)

    record("PPT mentions chain-of-thought paper",
           "chain-of-thought" in full_text or "chain of thought" in full_text)
    record("PPT mentions self-consistency paper",
           "self-consistency" in full_text or "self consistency" in full_text)
    record("PPT mentions tree of thoughts paper",
           "tree of thoughts" in full_text or "tree of thought" in full_text)


def check_gsheet():
    print("\n=== Checking Google Sheet ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        cur.execute("SELECT id, title FROM gsheet.spreadsheets")
        spreadsheets = cur.fetchall()

        target_ss = None
        for sid, title in spreadsheets:
            if title and ("llm" in title.lower() or "reasoning" in title.lower()) and "paper" in title.lower():
                target_ss = sid
                break

        record("GSheet 'LLM Reasoning Paper Tracker' exists",
               target_ss is not None,
               f"Found sheets: {[t for _, t in spreadsheets]}")

        if target_ss is None:
            conn.close()
            return

        cur.execute("SELECT id, title FROM gsheet.sheets WHERE spreadsheet_id = %s", (target_ss,))
        sheets = cur.fetchall()
        record("GSheet has at least one sheet", len(sheets) > 0, f"Found: {sheets}")
        if not sheets:
            conn.close()
            return
        sheet_id = sheets[0][0]

        # Header row (row_index = 0)
        cur.execute("""
            SELECT LOWER(value) FROM gsheet.cells
            WHERE spreadsheet_id = %s AND sheet_id = %s AND row_index = 0
        """, (target_ss, sheet_id))
        header_cells = [r[0] for r in cur.fetchall() if r[0]]
        header_text = " ".join(header_cells)
        required_cols = ["title", "authors", "arxiv_id", "published_date", "key_method"]
        has_cols = all(c in header_text for c in required_cols)

        # Data rows
        cur.execute("""
            SELECT COUNT(DISTINCT row_index) FROM gsheet.cells
            WHERE spreadsheet_id = %s AND sheet_id = %s AND row_index > 0
        """, (target_ss, sheet_id))
        data_rows = cur.fetchone()[0]
        record("GSheet has at least 3 data rows", data_rows >= 3, f"Found {data_rows} data rows")

        # CRITICAL: header columns + exactly 3 data rows
        record("GSheet header has required columns (Title/Authors/ArXiv_ID/Published_Date/Key_Method) and exactly 3 data rows",
               has_cols and data_rows == 3,
               f"header={header_cells} data_rows={data_rows}")

        # All cell values
        cur.execute("""
            SELECT LOWER(value) FROM gsheet.cells
            WHERE spreadsheet_id = %s AND sheet_id = %s
        """, (target_ss, sheet_id))
        cell_values = [row[0] for row in cur.fetchall() if row[0]]
        all_cells_text = " ".join(cell_values)
        STATE["gsheet_text"] = all_cells_text
        STATE["gsheet_id_text"] = normalize(all_cells_text)

        record("GSheet contains chain-of-thought paper entry",
               "chain-of-thought" in all_cells_text or "chain of thought" in all_cells_text)
        record("GSheet contains self-consistency paper entry",
               "self-consistency" in all_cells_text or "self consistency" in all_cells_text)
        record("GSheet contains tree of thoughts paper entry",
               "tree of thoughts" in all_cells_text or "tree of thought" in all_cells_text)

        # CRITICAL: correct arXiv IDs (normalize for id matching only)
        ids_ok = all(normalize(aid) in STATE["gsheet_id_text"] for aid in RELEVANT_ARXIV_IDS)
        record("GSheet contains correct arXiv IDs 2301.00234 / 2302.11382 / 2305.10601",
               ids_ok, f"id_text_present={STATE['gsheet_id_text'][:200]}")

        conn.close()
    except Exception as e:
        record("GSheet connection", False, str(e))


def check_word(agent_workspace):
    print("\n=== Checking Word Document ===")
    doc_path = os.path.join(agent_workspace, "LLM_Reasoning_Literature_Review.docx")
    if not os.path.isfile(doc_path):
        record("Word file LLM_Reasoning_Literature_Review.docx exists", False, f"Not found at: {doc_path}")
        return
    record("Word file LLM_Reasoning_Literature_Review.docx exists", True)

    try:
        doc = Document(doc_path)
    except Exception as e:
        record("Word file readable", False, str(e))
        return
    record("Word file readable", True)

    full_text = "\n".join(p.text for p in doc.paragraphs).lower()
    STATE["docx_text"] = full_text

    has_heading = "literature review" in full_text and ("llm" in full_text or "reasoning" in full_text)
    record("Word has 'Literature Review' heading with LLM/reasoning", has_heading)

    record("Word has substantial content (intro + sections)", len(full_text) > 300,
           f"Text length: {len(full_text)}")

    record("Word mentions chain-of-thought method",
           "chain-of-thought" in full_text or "chain of thought" in full_text)
    record("Word mentions self-consistency method",
           "self-consistency" in full_text or "self consistency" in full_text)
    record("Word mentions tree of thoughts method",
           "tree of thoughts" in full_text or "tree of thought" in full_text)


def check_critical_relevance():
    """CRITICAL: шумовые статьи отсутствуют + все три релевантные присутствуют везде."""
    print("\n=== Critical relevance checks ===")
    ppt = STATE["ppt_text"]
    gs = STATE["gsheet_text"]
    dx = STATE["docx_text"]

    ppt_clean = not any(frag in ppt for frag in NOISE_TITLE_FRAGMENTS)
    gs_clean = not any(frag in gs for frag in NOISE_TITLE_FRAGMENTS)
    dx_clean = not any(frag in dx for frag in NOISE_TITLE_FRAGMENTS)
    record("Relevance filter: noise papers ABSENT from pptx", ppt_clean,
           [f for f in NOISE_TITLE_FRAGMENTS if f in ppt])
    record("Relevance filter: noise papers ABSENT from gsheet", gs_clean,
           [f for f in NOISE_TITLE_FRAGMENTS if f in gs])
    record("Relevance filter: noise papers ABSENT from docx", dx_clean,
           [f for f in NOISE_TITLE_FRAGMENTS if f in dx])

    all_present = _has_all_relevant(ppt) and _has_all_relevant(gs) and _has_all_relevant(dx)
    record("All 3 relevant papers present in pptx AND gsheet AND docx", all_present)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=True)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False)
    args = parser.parse_args()

    check_pptx(args.agent_workspace)
    check_gsheet()
    check_word(args.agent_workspace)
    check_critical_relevance()

    total = PASS_COUNT + FAIL_COUNT
    accuracy = PASS_COUNT / total * 100 if total else 0
    critical_failed = [n for n in FAILED_NAMES if n in CRITICAL_CHECKS]

    print(f"\n=== Results: {PASS_COUNT}/{total} passed ({accuracy:.1f}%) ===")
    if critical_failed:
        print(f"CRITICAL FAILURES: {len(critical_failed)}")
        for n in critical_failed:
            print(f"  - {n}")

    if args.res_log_file:
        try:
            with open(args.res_log_file, "w") as f:
                json.dump({
                    "total_passed": PASS_COUNT, "total_checks": total,
                    "accuracy": accuracy, "critical_failed": critical_failed,
                }, f, indent=2)
        except Exception:
            pass

    success = (not critical_failed) and accuracy >= 70
    if success:
        print("All checks passed!")
        sys.exit(0)
    else:
        sys.exit(1)


if __name__ == "__main__":
    main()
