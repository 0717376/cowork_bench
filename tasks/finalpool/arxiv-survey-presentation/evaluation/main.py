"""
Evaluation for arxiv-survey-presentation task.
Checks Excel (Survey_Data.xlsx) and PowerPoint (NLP_Survey.pptx).

Two-tier gate:
  - CRITICAL_CHECKS: semantic correctness (paper selection, dataset values,
    Average_Sections, real per-paper content in pptx). Any critical failure =>
    sys.exit(1) before the accuracy gate.
  - Non-critical structural checks contribute to accuracy; PASS requires
    accuracy >= 70 AND no critical failure.
"""
import argparse
import os
import sys

import openpyxl

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILED = []

EXPECTED_PAPER_IDS = {"2404.00001", "2404.00002", "2404.00003", "2404.00004", "2404.00005"}
NOISE_PAPER_IDS = {"2404.00006", "2404.00007"}

# Primary dataset tokens per target paper (from injected LaTeX experiments sections).
PAPER_DATASETS = {
    "2404.00001": ["gsm8k", "math"],
    "2404.00002": ["mmlu", "bbh", "triviaqa"],
    "2404.00003": ["natural questions", "naturalquestions", "nq", "triviaqa"],
    "2404.00004": ["gsm8k", "svamp", "aqua"],
    "2404.00005": ["vqa", "gqa", "vizwiz"],
}

# All 5 target papers have exactly 4 LaTeX sections => mean = 4.0
EXPECTED_AVG_SECTIONS = 4.0

# Per-paper method keywords (RU+EN) for pptx body verification.
METHOD_KEYWORDS = [
    ["chain-of-thought", "chain of thought", "цепочк"],
    ["instruction tuning", "instruction-tuning", "инструкц"],
    ["retrieval", "retrieval-augmented", "поиск", "извлечен"],
    ["self-consistency", "self consistency", "самосогласован", "согласован"],
    ["multimodal", "multi-modal", "мультимодал"],
]


def record(name, passed, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT
    tag = "CRITICAL " if critical else ""
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {tag}{name}")
    else:
        FAIL_COUNT += 1
        if critical:
            CRITICAL_FAILED.append(name)
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {tag}{name}{msg}")


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower().replace(" ", "_") == sheet_name.strip().lower().replace(" ", "_"):
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
        if name.strip().lower().replace("_", " ") == sheet_name.strip().lower().replace("_", " "):
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def find_col(header, names):
    if not header:
        return None
    for i, cell in enumerate(header):
        if cell is None:
            continue
        cl = str(cell).strip().lower().replace(" ", "_")
        for n in names:
            if n.lower().replace(" ", "_") == cl:
                return i
    return None


def check_excel(workspace):
    print("\n=== Checking Excel ===")
    path = os.path.join(workspace, "Survey_Data.xlsx")
    if not os.path.isfile(path):
        record("Excel exists", False, f"Not found: {path}", critical=True)
        return False
    record("Excel exists", True)

    wb = openpyxl.load_workbook(path, data_only=True)

    # ---- Paper Comparison ----
    pc_rows = load_sheet_rows(wb, "Paper Comparison") or load_sheet_rows(wb, "Paper_Comparison")
    if pc_rows is None:
        record("Sheet 'Paper Comparison' exists", False, f"Sheets: {wb.sheetnames}", critical=True)
    else:
        record("Sheet 'Paper Comparison' exists", True)
        data = pc_rows[1:]
        record("Paper Comparison has 5 rows", len(data) == 5, f"Found {len(data)}")

        id_col = find_col(pc_rows[0], ["Paper_ID", "Paper ID", "ID"])
        found = set()
        if id_col is not None:
            found = {str(r[id_col]).strip() for r in data if id_col < len(r) and r[id_col]}

        # CRITICAL: exact paper selection (all 5 targets present, no noise papers).
        has_all_targets = EXPECTED_PAPER_IDS.issubset(found)
        no_noise = not (found & NOISE_PAPER_IDS)
        record("Correct paper selection (5 targets, no noise)",
               has_all_targets and no_noise,
               f"Found: {sorted(found)}", critical=True)

        method_col = find_col(pc_rows[0], ["Method", "method"])
        record("Method column exists", method_col is not None, f"Header: {pc_rows[0]}")

        dataset_col = find_col(pc_rows[0], ["Dataset", "dataset", "Dataset_Used"])
        record("Dataset column exists", dataset_col is not None, f"Header: {pc_rows[0]}")

        # CRITICAL: dataset values match the correct primary dataset per paper.
        if id_col is not None and dataset_col is not None:
            matches = 0
            for r in data:
                if id_col >= len(r) or dataset_col >= len(r):
                    continue
                pid = str(r[id_col]).strip()
                ds = str(r[dataset_col]).strip().lower() if r[dataset_col] else ""
                tokens = PAPER_DATASETS.get(pid)
                if tokens and any(t in ds for t in tokens):
                    matches += 1
            record("Dataset values correct (>=4/5 papers)", matches >= 4,
                   f"Matched {matches}/5", critical=True)
        else:
            record("Dataset values correct (>=4/5 papers)", False,
                   "Paper_ID or Dataset column missing", critical=True)

    # ---- Taxonomy ----
    tax_rows = load_sheet_rows(wb, "Taxonomy")
    if tax_rows is None:
        record("Sheet 'Taxonomy' exists", False, f"Sheets: {wb.sheetnames}")
    else:
        record("Sheet 'Taxonomy' exists", True)
        data = [r for r in tax_rows[1:] if r and r[0] and str(r[0]).strip()]
        record("Taxonomy has >= 2 non-empty categories", len(data) >= 2, f"Found {len(data)}")

        # Union of referenced paper IDs across the Papers column should cover the 5 targets.
        pap_col = find_col(tax_rows[0], ["Papers", "Paper_IDs", "Paper IDs"])
        if pap_col is None:
            pap_col = 1 if tax_rows and len(tax_rows[0]) > 1 else None
        referenced = ""
        if pap_col is not None:
            for r in data:
                if pap_col < len(r) and r[pap_col]:
                    referenced += " " + str(r[pap_col])
        covered = sum(1 for pid in EXPECTED_PAPER_IDS if pid in referenced)
        record("Taxonomy covers the 5 target papers", covered == 5,
               f"Covered {covered}/5 in: {referenced[:200]}")

    # ---- Summary Statistics ----
    ss_rows = load_sheet_rows(wb, "Summary Statistics") or load_sheet_rows(wb, "Summary_Statistics")
    if ss_rows is None:
        record("Sheet 'Summary Statistics' exists", False, f"Sheets: {wb.sheetnames}")
    else:
        record("Sheet 'Summary Statistics' exists", True)
        metrics = {}
        for row in ss_rows[1:]:
            if row and row[0]:
                metrics[str(row[0]).strip().lower().replace(" ", "_")] = row[1] if len(row) > 1 else None

        tp_key = next((k for k in metrics if "total" in k and "paper" in k), None)
        if tp_key:
            try:
                record("Total_Papers = 5", abs(float(metrics[tp_key]) - 5) < 1, f"Got {metrics[tp_key]}")
            except (TypeError, ValueError):
                record("Total_Papers is numeric", False, f"Got {metrics[tp_key]}")
        else:
            record("Total_Papers metric present", False, f"Metrics: {list(metrics)}")

        # CRITICAL: Average_Sections close to the true mean (4.0).
        avg_key = next((k for k in metrics if "average" in k and "section" in k), None)
        if avg_key is None:
            avg_key = next((k for k in metrics if "avg" in k and "section" in k), None)
        if avg_key is not None:
            try:
                val = float(str(metrics[avg_key]).replace(",", "."))
                record("Average_Sections ~= 4.0", abs(val - EXPECTED_AVG_SECTIONS) <= 0.5,
                       f"Got {metrics[avg_key]}", critical=True)
            except (TypeError, ValueError):
                record("Average_Sections ~= 4.0", False,
                       f"Non-numeric: {metrics[avg_key]}", critical=True)
        else:
            record("Average_Sections ~= 4.0", False,
                   f"Metric missing. Metrics: {list(metrics)}", critical=True)

    return True


def check_pptx(workspace):
    print("\n=== Checking PowerPoint ===")
    path = os.path.join(workspace, "NLP_Survey.pptx")
    if not os.path.isfile(path):
        record("PPTX exists", False, f"Not found: {path}", critical=True)
        return False
    record("PPTX exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = list(prs.slides)

        record("Has >= 6 slides", len(slides) >= 6, f"Found {len(slides)}", critical=True)

        all_text = []
        for slide in slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        texts.append(p.text)
            all_text.append("\n".join(texts))

        full = "\n".join(all_text).lower()

        # Title slide (RU+EN keywords).
        first = all_text[0].lower() if all_text else ""
        title_kws = ["survey", "nlp", "language model", "reasoning", "advances",
                     "обзор", "языковых моделей", "языковые модели", "рассужден"]
        record("Title slide has survey keywords", any(kw in first for kw in title_kws),
               f"First slide: {first[:200]}")

        # CRITICAL: real per-paper content -> at least 3 distinct method families present.
        distinct = sum(1 for group in METHOD_KEYWORDS if any(kw in full for kw in group))
        record("PPTX mentions >= 3 distinct paper methods", distinct >= 3,
               f"Distinct method families: {distinct}", critical=True)

        # Summary/conclusion slide has content.
        last = all_text[-1] if all_text else ""
        record("Last slide has content", len(last.strip()) > 10, f"Last slide: {last[:100]}")

        return True
    except Exception as e:
        record("PPTX readable", False, str(e), critical=True)
        return False


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100) if total else 0.0

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}, Failed: {FAIL_COUNT}, Accuracy: {accuracy:.1f}%")

    if CRITICAL_FAILED:
        print(f"  CRITICAL checks failed: {CRITICAL_FAILED}")
        sys.exit(1)

    sys.exit(0 if accuracy >= 70 else 1)


if __name__ == "__main__":
    main()
