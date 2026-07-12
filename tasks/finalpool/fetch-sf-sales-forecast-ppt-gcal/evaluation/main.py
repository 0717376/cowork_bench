"""
Evaluation script for fetch-sf-sales-forecast-ppt-gcal task (RU / ClickHouse).

Checks:
1. Sales_Forecast_Data.xlsx with Q1_Actuals, Q2_Forecast, Segment_Mix sheets
2. Q2_Sales_Forecast.pptx with forecast content
3. Calendar event for board presentation

The data warehouse (ClickHouse identity, PG schema sf_data, logical DB SALES_DW)
has its REGION / SEGMENT values russified centrally by db/zzz_clickhouse_after_init.sql.
The market projections JSON (files/mock_pages/api/projections.json) is russified to
the SAME region strings so the region join works. Eval RECOMPUTES the ground truth
from the live warehouse + projections rather than diffing a stale groundtruth file.

CRITICAL_CHECKS (semantic): any failure => sys.exit(1) before the accuracy gate.
"""

import argparse
import json
import os
import sys

import openpyxl
import psycopg2

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent",
    "password": "camel",
}

# Russified region labels (must match db/zzz_clickhouse_after_init.sql seed map).
REGION_RU = {
    "Asia Pacific": "Азиатско-Тихоокеанский регион",
    "Europe": "Европа",
    "Latin America": "Латинская Америка",
    "Middle East": "Ближний Восток",
    "North America": "Северная Америка",
}

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILED = []


def record(name, passed, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT
    tag = "CRITICAL " if critical else ""
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {tag}{name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {tag}{name}{msg}")
        if critical:
            CRITICAL_FAILED.append(name)


def num_close(a, b, tol=5000.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def find_num(value):
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


# ---------------------------------------------------------------------------
# Ground-truth recompute from live warehouse + projections JSON
# ---------------------------------------------------------------------------

def load_projections():
    """Load russified growth_rate per Russian region label.

    Prefer the extracted mock file; fall back to the live HTTP server.
    Returns {ru_region: growth_rate_pct}.
    """
    task_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    candidate = os.path.join(task_root, "tmp", "mock_pages", "api", "projections.json")
    data = None
    if os.path.isfile(candidate):
        try:
            with open(candidate, "r", encoding="utf-8") as f:
                data = json.load(f)
        except Exception:
            data = None
    if data is None:
        try:
            import urllib.request
            with urllib.request.urlopen("http://localhost:30209/api/projections.json", timeout=5) as r:
                data = json.loads(r.read().decode("utf-8"))
        except Exception:
            data = None
    gr = {}
    if data:
        for p in data.get("regional_projections", []):
            gr[str(p.get("region", "")).strip()] = float(p.get("growth_rate_pct"))
    return gr


def get_expected_data():
    """Recompute Q1 actuals + Q2 forecast from sf_data warehouse and projections.

    Returns dict with:
      q1_region_month: {(ru_region, month): (orders, revenue)}
      q1_region_total: {ru_region: (orders, revenue)}
      q2_forecast: {ru_region: (q1_rev, growth, q2_rev)}
      total_q2: float
      top_region: ru_region (highest growth)
      growth: {ru_region: growth_rate}
    """
    growth = load_projections()
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute('''
        SELECT c."REGION", EXTRACT(MONTH FROM o."ORDER_DATE")::int AS m,
               COUNT(*) AS orders,
               ROUND(SUM(o."TOTAL_AMOUNT"::float)::numeric, 2) AS revenue
        FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
        JOIN sf_data."SALES_DW__PUBLIC__CUSTOMERS" c
            ON o."CUSTOMER_ID" = c."CUSTOMER_ID"
        WHERE o."ORDER_DATE" >= '2026-01-01' AND o."ORDER_DATE" < '2026-04-01'
        GROUP BY c."REGION", m
        ORDER BY c."REGION", m
    ''')
    q1_region_month = {}
    q1_region_total = {}
    for region, m, orders, revenue in cur.fetchall():
        q1_region_month[(region, int(m))] = (int(orders), float(revenue))
        o, r = q1_region_total.get(region, (0, 0.0))
        q1_region_total[region] = (o + int(orders), r + float(revenue))

    conn.close()

    q2_forecast = {}
    total_q2 = 0.0
    for region, (orders, rev) in q1_region_total.items():
        g = growth.get(region)
        if g is None:
            continue
        q2_rev = round(rev * (1 + g / 100.0), 2)
        q2_forecast[region] = (round(rev, 2), g, q2_rev)
        total_q2 += q2_rev

    top_region = None
    if growth:
        top_region = max(growth, key=growth.get)

    return {
        "q1_region_month": q1_region_month,
        "q1_region_total": q1_region_total,
        "q2_forecast": q2_forecast,
        "total_q2": round(total_q2, 2),
        "top_region": top_region,
        "growth": growth,
    }


# ---------------------------------------------------------------------------
# Excel checks
# ---------------------------------------------------------------------------

def check_excel(agent_workspace, gt):
    print("\n=== Checking Excel Output ===")

    agent_file = os.path.join(agent_workspace, "Sales_Forecast_Data.xlsx")
    if not os.path.isfile(agent_file):
        record("Excel file exists", False, f"Not found: {agent_file}")
        return False

    record("Excel file exists", True)

    try:
        wb = openpyxl.load_workbook(agent_file, data_only=True)
    except Exception as e:
        record("Excel file readable", False, str(e))
        return False

    all_ok = True

    # ---- Q1_Actuals sheet ----
    q1_sheet = None
    for name in wb.sheetnames:
        if "q1" in name.lower() or "actual" in name.lower():
            q1_sheet = name
            break

    if not q1_sheet:
        record("Q1_Actuals sheet exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Q1_Actuals sheet exists", True)
        ws = wb[q1_sheet]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = [r for r in rows[1:] if r and r[0]] if len(rows) > 1 else []
        record(
            "Q1_Actuals has 15 rows (5 regions x 3 months)",
            len(data_rows) >= 15,
            f"Found {len(data_rows)} data rows",
        )
        if len(data_rows) < 15:
            all_ok = False

        # CRITICAL spot-value: Asia Pacific (RU label) January revenue must match
        # the recomputed warehouse value. Replaces the old silently-skipped EN check.
        ap_ru = REGION_RU["Asia Pacific"]
        gt_ap_jan = gt["q1_region_month"].get((ap_ru, 1))
        found_ap = False
        for r in data_rows:
            reg = str(r[0]) if r[0] is not None else ""
            month = r[1] if len(r) > 1 else None
            is_jan = month in (1, "1", "January", "Январь", "январь", "01")
            if ("азиат" in reg.lower() or "asia" in reg.lower()) and is_jan:
                found_ap = True
                got = r[3] if len(r) > 3 else None
                if gt_ap_jan is not None:
                    ok = num_close(got, gt_ap_jan[1], tol=200.0)
                    record(
                        "Q1_Actuals Asia-Pacific Jan revenue matches warehouse",
                        ok,
                        f"Got {got}, expected ~{gt_ap_jan[1]}",
                        critical=True,
                    )
                else:
                    record("Warehouse Asia-Pacific Jan value available", False,
                           "Could not recompute from sf_data", critical=True)
                break
        if not found_ap:
            record("Q1_Actuals Asia-Pacific January row present", False,
                   f"No Russian Asia-Pacific January row found", critical=True)

    # ---- Q2_Forecast sheet ----
    q2_sheet = None
    for name in wb.sheetnames:
        if "q2" in name.lower() or "forecast" in name.lower():
            q2_sheet = name
            break

    if not q2_sheet:
        record("Q2_Forecast sheet exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Q2_Forecast sheet exists", True)
        ws2 = wb[q2_sheet]
        rows2 = list(ws2.iter_rows(values_only=True))
        data_rows2 = [r for r in rows2[1:] if r and r[0]] if len(rows2) > 1 else []
        record(
            "Q2_Forecast has 5 region rows",
            len(data_rows2) >= 5,
            f"Found {len(data_rows2)} data rows",
        )
        if len(data_rows2) < 5:
            all_ok = False

        has_growth = False
        for r in data_rows2:
            if r and len(r) >= 3 and find_num(r[2]) is not None and 2.0 <= find_num(r[2]) <= 10.0:
                has_growth = True
                break
        record("Growth rates present", has_growth)

        # CRITICAL: per region (1) growth matches projections, (2) Q2 = Q1*(1+gr/100)
        # Match rows by Russian region label.
        gr_map = {}
        recompute_ok = 0
        recompute_total = 0
        for r in data_rows2:
            reg = str(r[0]).strip() if r[0] is not None else ""
            # locate the matching GT region by substring on the russified label
            match = None
            for ru in gt["q2_forecast"]:
                if ru.lower() in reg.lower() or reg.lower() in ru.lower():
                    match = ru
                    break
            if not match:
                continue
            q1_rev = find_num(r[1])
            gr = find_num(r[2])
            q2_rev = find_num(r[3])
            exp_q1, exp_gr, exp_q2 = gt["q2_forecast"][match]
            gr_map[match] = gr
            recompute_total += 1
            row_ok = True
            if gr is None or abs(gr - exp_gr) > 0.05:
                row_ok = False
            # recompute against the row's OWN q1 value (rounding tolerant)
            if q1_rev is not None and gr is not None and q2_rev is not None:
                expected_self = round(q1_rev * (1 + gr / 100.0), 2)
                if abs(q2_rev - expected_self) > 1.0:
                    row_ok = False
            else:
                row_ok = False
            if row_ok:
                recompute_ok += 1

        record(
            "Q2_Forecast growth rates map to correct regions (projections)",
            recompute_total >= 5 and all(
                gr_map.get(ru) is not None and abs(gr_map[ru] - gt["growth"].get(ru, -99)) <= 0.05
                for ru in gt["q2_forecast"]
            ),
            f"gr_map={gr_map}",
            critical=True,
        )
        record(
            "Q2_Forecast revenue == Q1*(1+gr/100) per region",
            recompute_total >= 5 and recompute_ok >= 5,
            f"{recompute_ok}/{recompute_total} region rows consistent",
            critical=True,
        )

    # ---- Segment_Mix sheet ----
    seg_sheet = None
    for name in wb.sheetnames:
        if "segment" in name.lower() or "mix" in name.lower():
            seg_sheet = name
            break

    if not seg_sheet:
        record("Segment_Mix sheet exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Segment_Mix sheet exists", True)
        ws3 = wb[seg_sheet]
        rows3 = list(ws3.iter_rows(values_only=True))
        data_rows3 = [r for r in rows3[1:] if r and r[0]] if len(rows3) > 1 else []
        record(
            "Segment_Mix has >= 20 rows (5 regions x 4 segments)",
            len(data_rows3) >= 20,
            f"Found {len(data_rows3)} data rows",
        )
        if len(data_rows3) < 20:
            all_ok = False

        # CRITICAL: per region the Revenue_Share_Pct across segments sums to ~100.
        # Share column is the LAST numeric column (index 3). Group by region (col 0).
        per_region = {}
        for r in data_rows3:
            reg = str(r[0]).strip() if r[0] is not None else ""
            share = find_num(r[3]) if len(r) > 3 else None
            if reg and share is not None:
                per_region.setdefault(reg, 0.0)
                per_region[reg] += share
        sums_ok = sum(1 for v in per_region.values() if abs(v - 100.0) <= 0.5)
        record(
            "Segment_Mix Revenue_Share_Pct sums to ~100 per region",
            len(per_region) >= 5 and sums_ok >= 5,
            f"region share sums: { {k: round(v,2) for k,v in per_region.items()} }",
            critical=True,
        )

    wb.close()
    return all_ok


# ---------------------------------------------------------------------------
# PPTX checks
# ---------------------------------------------------------------------------

def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint ===")

    pptx_file = os.path.join(agent_workspace, "Q2_Sales_Forecast.pptx")
    if not os.path.isfile(pptx_file):
        record("PowerPoint file exists", False, f"Not found: {pptx_file}")
        return False

    record("PowerPoint file exists", True)

    if Presentation is None:
        record("python-pptx available", False, "Cannot import pptx")
        return True

    try:
        prs = Presentation(pptx_file)
        slides = prs.slides

        record("PPT has >= 4 slides", len(slides) >= 4, f"Found {len(slides)} slides")

        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text.lower() + " "

        record(
            "PPT mentions Q2 forecast",
            "q2" in all_text and ("forecast" in all_text or "прогноз" in all_text),
        )
        record(
            "PPT mentions regions",
            any(r in all_text for r in [
                "asia", "europe", "north america",
                "азиат", "европ", "америк", "восток",
            ]),
        )
        record(
            "PPT mentions growth",
            "growth" in all_text or "рост" in all_text or "%" in all_text or "projection" in all_text or "прогноз" in all_text,
        )

        return True
    except Exception as e:
        record("PPT readable", False, str(e))
        return False


# ---------------------------------------------------------------------------
# Calendar checks
# ---------------------------------------------------------------------------

def check_calendar(gt):
    print("\n=== Checking Google Calendar ===")

    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute(
            "SELECT summary, description, start_datetime, end_datetime FROM gcal.events"
        )
        events = cur.fetchall()
        cur.close()
        conn.close()
    except Exception as e:
        record("Calendar DB accessible", False, str(e))
        return False

    found = False
    for summary, description, start_dt, end_dt in events:
        summary_lower = (summary or "").lower()
        if ("board" in summary_lower or "forecast" in summary_lower or "sales" in summary_lower
                or "совет" in summary_lower) and (
            "presentation" in summary_lower or "meeting" in summary_lower or "q2" in summary_lower
                or "презентац" in summary_lower):
            found = True
            record("Board presentation event exists", True)

            start_str = str(start_dt)
            end_str = str(end_dt)
            record("Event on March 28, 2026", "2026-03-28" in start_str, f"Start: {start_str}")
            # CRITICAL: exact start 10:00 and end 11:30
            record(
                "Event start 2026-03-28 10:00 and end 11:30",
                "2026-03-28" in start_str and "10:00" in start_str and "11:30" in end_str,
                f"start={start_str} end={end_str}",
                critical=True,
            )

            desc = description or ""
            desc_lower = desc.lower()
            has_info = any(kw in desc_lower for kw in [
                "revenue", "forecast", "growth", "region", "выручк", "прогноз", "рост", "регион",
            ])
            record("Event description has forecast info", has_info)

            # CRITICAL: description contains correct company-wide Q2 total + top region.
            import re
            nums = []
            for tok in re.findall(r"[\d][\d\s.,]*[\d]|\d", desc):
                no_space = tok.replace(" ", "")
                # Locale-aware: build both US (',' thousands) and RU (',' decimal) candidates.
                candidates = {no_space.replace(",", "")}
                if re.search(r",\d{2}$", no_space):
                    candidates.add(no_space.replace(".", "").replace(",", "."))
                for cand in candidates:
                    try:
                        nums.append(float(cand))
                    except ValueError:
                        pass
            total_ok = any(num_close(n, gt["total_q2"], tol=max(1000.0, gt["total_q2"] * 0.01)) for n in nums)
            record(
                "Calendar description has correct company-wide Q2 total",
                total_ok,
                f"expected ~{gt['total_q2']}, found nums {nums}",
                critical=True,
            )

            top = gt["top_region"] or ""
            top_ok = top.lower() in desc_lower or "азиат" in desc_lower
            record(
                "Calendar description names top-growth region (Азиатско-Тихоокеанский)",
                top_ok,
                f"expected '{top}'",
                critical=True,
            )
            break

    if not found:
        record(
            "Board presentation event exists",
            False,
            f"Found {len(events)} events but none for board/forecast/sales presentation",
        )

    return found


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    try:
        gt = get_expected_data()
        print(f"[gt] total_q2={gt['total_q2']} top_region={gt['top_region']} "
              f"growth={gt['growth']}")
    except Exception as e:
        print(f"[FATAL] Could not recompute ground truth from warehouse: {e}")
        sys.exit(1)

    excel_ok = check_excel(args.agent_workspace, gt)
    pptx_ok = check_pptx(args.agent_workspace)
    cal_ok = check_calendar(gt)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100.0) if total else 0.0

    print(f"\n=== SUMMARY ===")
    print(f"  Excel:    {'PASS' if excel_ok else 'FAIL'}")
    print(f"  PPT:      {'PASS' if pptx_ok else 'FAIL'}")
    print(f"  Calendar: {'PASS' if cal_ok else 'FAIL'}")
    print(f"  Passed: {PASS_COUNT}, Failed: {FAIL_COUNT}, Accuracy: {accuracy:.1f}%")

    if CRITICAL_FAILED:
        print(f"  CRITICAL checks FAILED: {CRITICAL_FAILED}")
        print(f"  Overall:  FAIL (critical)")
        sys.exit(1)

    overall = accuracy >= 70.0
    print(f"  Overall:  {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
