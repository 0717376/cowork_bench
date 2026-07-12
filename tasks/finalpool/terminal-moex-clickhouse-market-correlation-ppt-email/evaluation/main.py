"""Evaluation for terminal-moex-clickhouse-market-correlation-ppt-email (RU fork).

Structural (soft) checks + CRITICAL checks.

CRITICAL (any fail => overall FAIL regardless of accuracy):
  C1: PPTX stock-overview содержит реальную текущую цену ИЛИ 52-нед. максимум
      хотя бы для 3 из 5 тикеров MOEX, совпадающие с живыми moex.stock_info
      (с допуском). Цены читаются ЖИВЫМИ из БД, не хардкодятся.
  C2: PPTX называет КОРРЕКТНЫЙ крупнейший по выручке клиентский сегмент,
      вычисленный вживую из sf_data SALES_DW (ORDERS x CUSTOMERS), И содержит
      правдоподобную итоговую цифру выручки (масштаб совпадает с суммой
      ежемесячной выручки SALES_DW, с допуском).
  C3: Письмо на cfo@company.com содержит цифру выручки правдоподобного масштаба
      И называет крупнейший сегмент ИЛИ лучший тикер.
  C4: correlation_analysis.py существует И PPTX содержит >= 6 непустых слайдов,
      покрывающих обзор рынка, выручку, сегменты, корреляцию и рекомендации.

Иначе PASS, если accuracy >= 70%.

MOEX/sf_data значения читаются ЖИВЫМИ (не хардкодятся), т.к. волатильны/сидируются.
"""
import argparse
import json
import os
import re
import sys
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent", "password": "camel",
}

STOCKS = ["SBER.ME", "GAZP.ME", "LKOH.ME", "TCSG.ME", "MGNT.ME"]
BENCHMARK = "MTSS.ME"

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []
CRITICAL_CHECKS = {
    "CRITICAL: PPTX содержит реальные цены/52-нед.максимум для >=3 тикеров MOEX",
    "CRITICAL: PPTX называет корректный крупнейший сегмент + правдоподобную выручку",
    "CRITICAL: письмо CFO содержит выручку + крупнейший сегмент/лучший тикер",
    "CRITICAL: correlation_analysis.py существует и PPTX >=6 непустых слайдов с разделами",
}


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        marker = " [CRITICAL]" if name in CRITICAL_CHECKS else ""
        print(f"  [FAIL]{marker} {name}: {str(detail)[:200]}")
        if name in CRITICAL_CHECKS:
            CRITICAL_FAILS.append(name)


# ---------- live data readers ----------

def get_live_stock_metrics():
    """Return {symbol: {'price': float, 'high': float}} from moex.stock_info."""
    out = {}
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("SELECT symbol, data FROM moex.stock_info")
        for sym, data in cur.fetchall():
            d = data if isinstance(data, dict) else json.loads(data)
            price = d.get("currentPrice") or d.get("regularMarketPrice") or d.get("previousClose")
            high = d.get("fiftyTwoWeekHigh")
            out[sym] = {
                "price": float(price) if price is not None else None,
                "high": float(high) if high is not None else None,
            }
        cur.close()
        conn.close()
    except Exception as e:
        print(f"  [warn] could not read moex.stock_info: {e}")
    return out


def get_live_revenue_totals():
    """Return plausible total-revenue figures from SALES_DW MONTHLY_REVENUE."""
    totals = []
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute(
            'SELECT "MONTH_KEY", "TOTAL_REVENUE" '
            'FROM sf_data."SALES_DW__ANALYTICS__MONTHLY_REVENUE" ORDER BY "MONTH_KEY"'
        )
        rows = [(r[0], float(r[1])) for r in cur.fetchall()]
        cur.close()
        conn.close()
        if rows:
            all_total = sum(r[1] for r in rows)
            last12 = sum(r[1] for r in rows[-12:])
            last12_excl = sum(r[1] for r in rows[-13:-1])  # last 12 excl. partial newest
            totals = [all_total, last12, last12_excl]
    except Exception as e:
        print(f"  [warn] could not read MONTHLY_REVENUE: {e}")
    return totals


def get_live_top_segment():
    """Largest-revenue customer segment via ORDERS x CUSTOMERS (russified value)."""
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute(
            'SELECT c."SEGMENT", SUM(o."TOTAL_AMOUNT") AS rev '
            'FROM sf_data."SALES_DW__PUBLIC__ORDERS" o '
            'JOIN sf_data."SALES_DW__PUBLIC__CUSTOMERS" c '
            'ON o."CUSTOMER_ID" = c."CUSTOMER_ID" '
            'GROUP BY c."SEGMENT" ORDER BY rev DESC'
        )
        rows = cur.fetchall()
        cur.close()
        conn.close()
        if rows:
            return rows[0][0]
    except Exception as e:
        print(f"  [warn] could not read segment revenue: {e}")
    return None


def extract_numbers(text):
    """Numbers in text, tolerant of thousands separators and currency."""
    cleaned = re.sub(r"(?<=\d)[\s ,](?=\d{3}\b)", "", text)
    nums = []
    for m in re.findall(r"\d+(?:\.\d+)?", cleaned):
        try:
            nums.append(float(m))
        except Exception:
            pass
    return nums


def text_has_value(numbers, target, rel_tol=0.02, abs_tol=2.0):
    for n in numbers:
        if abs(n - target) <= max(abs_tol, abs(target) * rel_tol):
            return True
    return False


def text_has_any_total(numbers, totals, rel_tol=0.05):
    for t in totals:
        if text_has_value(numbers, t, rel_tol=rel_tol, abs_tol=1.0):
            return True
    return False


# ---------- PPTX ----------

def read_pptx_slides(path):
    from pptx import Presentation
    prs = Presentation(path)
    slide_texts = []
    for slide in prs.slides:
        chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        chunks.append(r.text or "")
            elif hasattr(shape, "text"):
                chunks.append(shape.text or "")
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        chunks.append(cell.text or "")
        slide_texts.append(" ".join(chunks))
    return slide_texts


def check_pptx(workspace, stock_metrics, rev_totals, top_segment):
    print("\n=== Check 1: Market_Correlation_Report.pptx ===")
    path = os.path.join(workspace, "Market_Correlation_Report.pptx")

    # explicit critical fails if missing/unreadable
    def fail_criticals(reason):
        check("CRITICAL: PPTX содержит реальные цены/52-нед.максимум для >=3 тикеров MOEX", False, reason)
        check("CRITICAL: PPTX называет корректный крупнейший сегмент + правдоподобную выручку", False, reason)

    if not os.path.exists(path):
        check("PPTX file exists", False, f"Not found at {path}")
        fail_criticals("PPTX missing")
        return None
    check("PPTX file exists", True)

    try:
        slide_texts = read_pptx_slides(path)
    except Exception as e:
        check("PPTX readable", False, str(e))
        fail_criticals("PPTX unreadable")
        return None
    check("PPTX readable", True)

    all_text = " ".join(slide_texts)
    all_lower = all_text.lower()
    n_slides = len(slide_texts)
    nums = extract_numbers(all_text)

    check("Has at least 6 slides", n_slides >= 6, f"Found {n_slides} slides")
    check("Mentions market or stock", any(k in all_lower for k in ["market", "stock", "рынок", "акци", "тикер"]),
          all_lower[:120])
    check("Mentions correlation or analysis",
          any(k in all_lower for k in ["correlation", "analysis", "корреляц", "анализ"]), all_lower[:120])
    check("Mentions MOEX tickers (SBER/GAZP/...)",
          any(t.split(".")[0].lower() in all_lower for t in STOCKS), all_lower[:120])
    check("Mentions revenue (выручка)",
          "revenue" in all_lower or "выручк" in all_lower, all_lower[:120])
    check("Mentions segment or consumer (сегмент)",
          any(k in all_lower for k in ["segment", "consumer", "сегмент", "частные клиенты",
                                       "корпоратив", "малый и средний"]), all_lower[:120])
    check("Has recommendations (рекомендации)",
          "recommend" in all_lower or "рекоменд" in all_lower, all_lower[:200])

    # ---- CRITICAL C1: real prices for >=3 tickers ----
    if stock_metrics:
        hit = 0
        for sym in STOCKS:
            m = stock_metrics.get(sym, {})
            targets = [v for v in (m.get("price"), m.get("high")) if v is not None]
            if any(text_has_value(nums, t, rel_tol=0.02, abs_tol=2.0) for t in targets):
                hit += 1
        check("CRITICAL: PPTX содержит реальные цены/52-нед.максимум для >=3 тикеров MOEX",
              hit >= 3, f"matched {hit}/5 tickers against live moex values")
    else:
        check("CRITICAL: PPTX содержит реальные цены/52-нед.максимум для >=3 тикеров MOEX",
              False, "no live moex metrics available")

    # ---- CRITICAL C2: correct top segment + plausible total revenue ----
    seg_ok = False
    if top_segment:
        seg_ok = top_segment.lower() in all_lower
    rev_ok = text_has_any_total(nums, rev_totals) if rev_totals else False
    check("CRITICAL: PPTX называет корректный крупнейший сегмент + правдоподобную выручку",
          bool(seg_ok and rev_ok),
          f"top_segment='{top_segment}' present={seg_ok}, plausible_total={rev_ok}")

    return slide_texts


def check_email(rev_totals, top_segment, stock_metrics):
    print("\n=== Check 2: Email to cfo@company.com ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
    messages = cur.fetchall()
    cur.close()
    conn.close()

    def to_str(to_addr):
        if isinstance(to_addr, list):
            return " ".join(str(r).lower() for r in to_addr)
        if isinstance(to_addr, str):
            try:
                parsed = json.loads(to_addr)
                return " ".join(str(r).lower() for r in parsed) if isinstance(parsed, list) else to_addr.lower()
            except Exception:
                return to_addr.lower()
        return str(to_addr).lower()

    matching = None
    for subject, from_addr, to_addr, body_text in messages:
        if "cfo@company.com" in to_str(to_addr):
            matching = (subject, from_addr, to_addr, body_text)
            break

    check("Email sent to cfo@company.com", matching is not None, f"Messages found: {len(messages)}")

    if not matching:
        check("CRITICAL: письмо CFO содержит выручку + крупнейший сегмент/лучший тикер",
              False, "no email to cfo@company.com")
        return

    subject, _, _, body_text = matching
    full = ((subject or "") + " " + (body_text or ""))
    low = full.lower()

    check("Email mentions correlation or market",
          any(k in low for k in ["correlation", "market", "корреляц", "рынок"]), f"Subject: {subject}")
    check("Email mentions revenue (выручка)", "revenue" in low or "выручк" in low, low[:120])
    check("Email mentions presentation or report",
          any(k in low for k in ["presentation", "report", "powerpoint", "презентац", "отчёт", "отчет"]),
          low[:120])

    # ---- CRITICAL C3 ----
    nums = extract_numbers(full)
    rev_ok = text_has_any_total(nums, rev_totals) if rev_totals else False
    seg_ok = bool(top_segment) and top_segment.lower() in low
    stock_ok = any(t.split(".")[0].lower() in low for t in STOCKS)
    check("CRITICAL: письмо CFO содержит выручку + крупнейший сегмент/лучший тикер",
          bool(rev_ok and (seg_ok or stock_ok)),
          f"plausible_total={rev_ok}, segment={seg_ok}, stock={stock_ok}")


def check_script_and_structure(workspace, slide_texts):
    print("\n=== Check 3: correlation_analysis.py + структура PPTX ===")
    script_path = os.path.join(workspace, "correlation_analysis.py")
    script_ok = os.path.exists(script_path)
    check("correlation_analysis.py exists", script_ok)

    structure_ok = False
    detail = "no slides"
    if slide_texts:
        non_empty = [t for t in slide_texts if t.strip()]
        joined = " ".join(slide_texts).lower()
        has_market = any(k in joined for k in ["market", "stock", "рынок", "акци", "тикер"])
        has_revenue = "revenue" in joined or "выручк" in joined
        has_segment = any(k in joined for k in ["segment", "сегмент", "частные клиенты", "корпоратив"])
        has_corr = any(k in joined for k in ["correlation", "корреляц"])
        has_reco = "recommend" in joined or "рекоменд" in joined
        structure_ok = (len(non_empty) >= 6 and has_market and has_revenue
                        and has_segment and has_corr and has_reco)
        detail = (f"non_empty={len(non_empty)}, market={has_market}, revenue={has_revenue}, "
                  f"segment={has_segment}, corr={has_corr}, reco={has_reco}")

    check("CRITICAL: correlation_analysis.py существует и PPTX >=6 непустых слайдов с разделами",
          bool(script_ok and structure_ok), detail)


def check_reverse_validation(workspace):
    print("\n=== Reverse Validation ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT COUNT(*) FROM email.messages
            WHERE (lower(subject) LIKE '%%correlation%%' OR lower(subject) LIKE '%%market%%'
                   OR lower(subject) LIKE '%%корреляц%%' OR lower(subject) LIKE '%%рынок%%')
              AND to_addr::text NOT ILIKE '%%cfo%%'
              AND to_addr::text NOT ILIKE '%%company.com%%'
        """)
        wrong = cur.fetchone()[0]
        check("No correlation emails to wrong recipients", wrong == 0, f"Found {wrong} misrouted emails")
        cur.close()
        conn.close()
    except Exception:
        pass

    path = os.path.join(workspace, "Market_Correlation_Report.pptx")
    if os.path.exists(path):
        try:
            slide_texts = read_pptx_slides(path)
            empty = sum(1 for t in slide_texts if not t.strip())
            check("No empty slides in PPTX", empty == 0, f"Found {empty} empty slides")
        except Exception:
            pass


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    stock_metrics = get_live_stock_metrics()
    rev_totals = get_live_revenue_totals()
    top_segment = get_live_top_segment()
    print(f"[eval] live top_segment={top_segment}, rev_totals={[round(t,2) for t in rev_totals]}")

    slide_texts = check_pptx(args.agent_workspace, stock_metrics, rev_totals, top_segment)
    check_email(rev_totals, top_segment, stock_metrics)
    check_script_and_structure(args.agent_workspace, slide_texts)
    check_reverse_validation(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": PASS_COUNT, "total_checks": total, "accuracy": accuracy,
        "critical_fails": CRITICAL_FAILS,
    }
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if CRITICAL_FAILS:
        print(f"FAIL: критичные чеки провалены ({len(CRITICAL_FAILS)}): {CRITICAL_FAILS}")
        sys.exit(1)
    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
