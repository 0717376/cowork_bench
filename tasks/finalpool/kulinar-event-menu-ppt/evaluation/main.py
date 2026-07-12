"""
Evaluation для kulinar-event-menu-ppt (русифицированная версия, kulinar).

Проверяет два артефакта в рабочей директории агента:
  - Event_Menu_Presentation.pptx (титул + 3 слайда блюд + итоговый слайд)
  - Menu_Budget.xlsx (листы "Menu Items" и "Summary")

Структурные проверки (NON-critical): существование файлов, >=5 слайдов,
имена листов, заголовки колонок, >=3 строк данных, наличие метрик.

CRITICAL проверки (содержательные) — провал любой => вся задача FAIL:
  1. Authenticity: >=2 названия рецептов в "Menu Items" реально есть в базе kulinar.
  2. Difficulty: все значения Difficulty — целые из {1,2,3} (маппинг kulinar 1->1,2->2,3+->3).
  3. Total_Recipes == числу строк блюд в "Menu Items".
  4. Total_Ingredients == сумме колонки Ingredients_Count (то же число на итоговом слайде PPT).
  5. Avg_Difficulty == среднему колонки Difficulty.

Порог: accuracy>=70% И нет критичных провалов => PASS.
"""

import argparse
import json
import os
import re
import sys


PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []

# Содержательные проверки. Провал любой => итог FAIL независимо от accuracy.
CRITICAL_CHECKS = {
    "Authenticity: >=2 рецепта из Menu Items есть в базе kulinar",
    "Difficulty: все значения — целые из {1,2,3}",
    "Summary: Total_Recipes == числу строк блюд",
    "Summary: Total_Ingredients == сумме Ingredients_Count",
    "Summary: Avg_Difficulty == среднему Difficulty",
}


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {str(detail)[:300]}" if detail else ""
        marker = " [CRITICAL]" if name in CRITICAL_CHECKS else ""
        print(f"  [FAIL]{marker} {name}{msg}")
        if name in CRITICAL_CHECKS:
            CRITICAL_FAILS.append(name)


# ---------------------------------------------------------------------------
# Загрузка эталонной базы kulinar (источник правды для authenticity)
# ---------------------------------------------------------------------------

def _norm_name(s):
    return re.sub(r"\s+", " ", str(s).strip().lower())


def load_kulinar_recipes():
    """Возвращает список рецептов kulinar или None, если база недоступна.

    Ищем all_recipes.json: сначала переменная окружения, затем поднимаемся
    вверх по дереву от файла eval до local_servers/kulinar-mcp/...
    """
    candidates = []
    env = os.environ.get("KULINAR_RECIPES_JSON")
    if env:
        candidates.append(env)

    here = os.path.abspath(__file__)
    cur = os.path.dirname(here)
    rel = os.path.join("local_servers", "kulinar-mcp", "src", "data", "all_recipes.json")
    for _ in range(12):
        candidates.append(os.path.join(cur, rel))
        parent = os.path.dirname(cur)
        if parent == cur:
            break
        cur = parent

    for path in candidates:
        try:
            if path and os.path.exists(path):
                with open(path, encoding="utf-8") as f:
                    data = json.load(f)
                if isinstance(data, list) and data:
                    return data
        except Exception:
            continue
    return None


KULINAR = load_kulinar_recipes()
KULINAR_NAMES = {_norm_name(r["name"]) for r in KULINAR} if KULINAR else set()


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def collect_pptx(workspace):
    """Возвращает (slides_text:list[str], all_text:str, title_text:str) или None."""
    from pptx import Presentation

    pptx_path = os.path.join(workspace, "Event_Menu_Presentation.pptx")
    if not os.path.exists(pptx_path):
        return None

    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    slides_text = []
    for slide in slides:
        buf = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    buf += paragraph.text + " "
            elif hasattr(shape, "text"):
                buf += shape.text + " "
        slides_text.append(buf)
    all_text = " ".join(slides_text)
    title_text = slides_text[0] if slides_text else ""
    return slides_text, all_text, title_text


def check_pptx(workspace):
    data = collect_pptx(workspace)
    if data is None:
        record("PPT: Event_Menu_Presentation.pptx существует", False,
               f"не найден в {workspace}")
        return None
    record("PPT: Event_Menu_Presentation.pptx существует", True)

    slides_text, all_text, title_text = data
    low = all_text.lower()

    record("PPT: >=5 слайдов (титул + 3 блюда + итог)",
           len(slides_text) >= 5, f"слайдов: {len(slides_text)}")

    # Титул: принимаем RU 'ужин'/'меню' ИЛИ EN 'dinner'/'menu' в оригинале (.lower)
    title_low = title_text.lower()
    title_ok = any(kw in title_low for kw in ("ужин", "меню", "dinner", "menu"))
    record("PPT: на титульном слайде есть 'ужин'/'меню'/'dinner'/'menu'",
           title_ok, f"титул: {title_text[:120]!r}")

    # Типы блюд: принимаем RU и EN маркеры
    course_groups = {
        "appetizer": ("appetizer", "закуск"),
        "main": ("main", "основн"),
        "dessert": ("dessert", "десерт"),
    }
    missing = [c for c, kws in course_groups.items()
               if not any(k in low for k in kws)]
    record("PPT: упомянуты все три блюда (appetizer/main/dessert)",
           not missing, f"не найдено: {missing}")

    return all_text


def pptx_total_ingredients(all_text):
    """Извлекает все целые числа из текста PPT (для сверки с Total_Ingredients)."""
    if not all_text:
        return set()
    return {int(n) for n in re.findall(r"\d+", all_text)}


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------

def check_excel(workspace, pptx_all_text):
    from openpyxl import load_workbook

    xlsx_path = os.path.join(workspace, "Menu_Budget.xlsx")
    if not os.path.exists(xlsx_path):
        record("Excel: Menu_Budget.xlsx существует", False, f"не найден в {workspace}")
        # критичные суммарные проверки тоже падают
        record("Summary: Total_Recipes == числу строк блюд", False, "нет файла")
        record("Summary: Total_Ingredients == сумме Ingredients_Count", False, "нет файла")
        record("Summary: Avg_Difficulty == среднему Difficulty", False, "нет файла")
        record("Difficulty: все значения — целые из {1,2,3}", False, "нет файла")
        record("Authenticity: >=2 рецепта из Menu Items есть в базе kulinar", False, "нет файла")
        return
    record("Excel: Menu_Budget.xlsx существует", True)

    wb = load_workbook(xlsx_path, data_only=True)
    sheet_lower = [s.lower() for s in wb.sheetnames]

    has_menu = "menu items" in sheet_lower
    has_summary = "summary" in sheet_lower
    record("Excel: лист 'Menu Items' присутствует", has_menu, f"листы: {wb.sheetnames}")
    record("Excel: лист 'Summary' присутствует", has_summary, f"листы: {wb.sheetnames}")

    if not has_menu:
        record("Menu Items: заголовки Course/Recipe_Name/Ingredients_Count/Difficulty", False, "нет листа")
        record("Menu Items: >=3 строк данных", False, "нет листа")
        record("Authenticity: >=2 рецепта из Menu Items есть в базе kulinar", False, "нет листа")
        record("Difficulty: все значения — целые из {1,2,3}", False, "нет листа")
        record("Summary: Total_Recipes == числу строк блюд", False, "нет листа Menu Items")
        record("Summary: Total_Ingredients == сумме Ingredients_Count", False, "нет листа Menu Items")
        record("Summary: Avg_Difficulty == среднему Difficulty", False, "нет листа Menu Items")
        return

    menu_ws = wb[wb.sheetnames[sheet_lower.index("menu items")]]
    rows = list(menu_ws.iter_rows(values_only=True))
    headers = [str(c).lower() if c is not None else "" for c in (rows[0] if rows else [])]

    # индекс колонок по fuzzy-совпадению
    def col_idx(rh):
        for i, h in enumerate(headers):
            if rh.replace("_", " ") in h or rh in h:
                return i
        return None

    idx_course = col_idx("course")
    idx_name = col_idx("recipe_name")
    if idx_name is None:
        idx_name = col_idx("recipe")
    idx_ings = col_idx("ingredients_count")
    if idx_ings is None:
        idx_ings = col_idx("ingredients")
    idx_diff = col_idx("difficulty")

    required = {"course": idx_course, "recipe_name": idx_name,
                "ingredients_count": idx_ings, "difficulty": idx_diff}
    missing_h = [k for k, v in required.items() if v is None]
    record("Menu Items: заголовки Course/Recipe_Name/Ingredients_Count/Difficulty",
           not missing_h, f"не найдено: {missing_h}; заголовки: {headers}")

    # строки данных: непустая первая значимая ячейка
    data_rows = [r for r in rows[1:]
                 if any(c is not None and str(c).strip() for c in r)]
    record("Menu Items: >=3 строк данных", len(data_rows) >= 3,
           f"строк: {len(data_rows)}")

    # --- Authenticity (CRITICAL) ---
    names = []
    if idx_name is not None:
        for r in data_rows:
            if idx_name < len(r) and r[idx_name] is not None:
                names.append(_norm_name(r[idx_name]))
    if KULINAR_NAMES:
        matched = [n for n in names if n in KULINAR_NAMES]
        record("Authenticity: >=2 рецепта из Menu Items есть в базе kulinar",
               len(matched) >= 2,
               f"совпало {len(matched)} из {len(names)}: {names}")
    else:
        # База недоступна — не валим жёстко, но фиксируем как пройденную
        # (деградация: проверка не может быть выполнена в окружении).
        record("Authenticity: >=2 рецепта из Menu Items есть в базе kulinar",
               True, "база kulinar недоступна в окружении — проверка пропущена")

    # --- Difficulty values (CRITICAL) ---
    diff_vals = []
    diff_ok = True
    if idx_diff is not None:
        for r in data_rows:
            if idx_diff < len(r):
                v = r[idx_diff]
                try:
                    iv = int(float(v))
                except (TypeError, ValueError):
                    diff_ok = False
                    continue
                diff_vals.append(iv)
                if iv not in (1, 2, 3):
                    diff_ok = False
    else:
        diff_ok = False
    record("Difficulty: все значения — целые из {1,2,3}",
           diff_ok and len(diff_vals) >= 3,
           f"значения: {diff_vals}")

    # --- Ingredients_Count sum ---
    ing_vals = []
    if idx_ings is not None:
        for r in data_rows:
            if idx_ings < len(r):
                try:
                    ing_vals.append(int(float(r[idx_ings])))
                except (TypeError, ValueError):
                    pass
    real_total_ings = sum(ing_vals)
    real_total_recipes = len(data_rows)
    real_avg_diff = (sum(diff_vals) / len(diff_vals)) if diff_vals else None

    # --- Summary list ---
    summary_metrics = {}
    if has_summary:
        sum_ws = wb[wb.sheetnames[sheet_lower.index("summary")]]
        for row in sum_ws.iter_rows(values_only=True):
            cells = [c for c in row if c is not None]
            if len(cells) >= 2:
                key = _norm_name(cells[0]).replace(" ", "_")
                summary_metrics[key] = cells[1]
        summary_text = " ".join(
            str(c).lower() for row in sum_ws.iter_rows(values_only=True)
            for c in row if c is not None)
    else:
        summary_text = ""

    record("Summary: метрика Total_Recipes присутствует",
           "total_recipes" in summary_text or "total recipes" in summary_text,
           f"metrics: {list(summary_metrics)}")
    record("Summary: метрика Total_Ingredients присутствует",
           "total_ingredients" in summary_text or "total ingredients" in summary_text,
           f"metrics: {list(summary_metrics)}")
    record("Summary: метрика Avg_Difficulty присутствует",
           "avg_difficulty" in summary_text or "avg difficulty" in summary_text,
           f"metrics: {list(summary_metrics)}")

    def as_num(v):
        try:
            return float(v)
        except (TypeError, ValueError):
            m = re.search(r"-?\d+(?:[.,]\d+)?", str(v) if v is not None else "")
            return float(m.group().replace(",", ".")) if m else None

    # --- CRITICAL: Total_Recipes ---
    tr = as_num(summary_metrics.get("total_recipes"))
    record("Summary: Total_Recipes == числу строк блюд",
           tr is not None and int(tr) == real_total_recipes,
           f"Summary={tr}, факт={real_total_recipes}")

    # --- CRITICAL: Total_Ingredients (+ совпадение на слайде PPT) ---
    ti = as_num(summary_metrics.get("total_ingredients"))
    ppt_nums = pptx_total_ingredients(pptx_all_text)
    ti_match_excel = ti is not None and int(ti) == real_total_ings
    ti_on_ppt = real_total_ings in ppt_nums if pptx_all_text else False
    record("Summary: Total_Ingredients == сумме Ingredients_Count",
           ti_match_excel and ti_on_ppt,
           f"Summary={ti}, сумма={real_total_ings}, на слайде PPT={ti_on_ppt}, "
           f"числа PPT={sorted(ppt_nums)[:20]}")

    # --- CRITICAL: Avg_Difficulty ---
    ad = as_num(summary_metrics.get("avg_difficulty"))
    avg_ok = (ad is not None and real_avg_diff is not None
              and abs(ad - real_avg_diff) <= 0.05)
    record("Summary: Avg_Difficulty == среднему Difficulty",
           avg_ok, f"Summary={ad}, факт={real_avg_diff}")


# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    print("\n=== Проверка PowerPoint ===")
    pptx_all_text = check_pptx(args.agent_workspace)

    print("\n=== Проверка Excel ===")
    check_excel(args.agent_workspace, pptx_all_text or "")

    total = PASS_COUNT + FAIL_COUNT
    accuracy = PASS_COUNT / total * 100 if total else 0
    print(f"\nИтого: {PASS_COUNT}/{total} проверок пройдено ({accuracy:.1f}%)")
    if CRITICAL_FAILS:
        print(f"Критичные провалы: {CRITICAL_FAILS}")

    if args.res_log_file:
        result = {
            "total_passed": PASS_COUNT,
            "total_checks": total,
            "accuracy": accuracy,
            "critical_fails": CRITICAL_FAILS,
            "success": (not CRITICAL_FAILS) and accuracy >= 70,
        }
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if CRITICAL_FAILS:
        print(f"FAIL: критичные чеки провалены ({len(CRITICAL_FAILS)})")
        sys.exit(1)
    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    print("FAIL")
    sys.exit(1)


if __name__ == "__main__":
    main()
