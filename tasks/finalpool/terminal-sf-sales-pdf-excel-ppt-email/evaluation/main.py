"""Evaluation for terminal-sf-sales-pdf-excel-ppt-email (ClickHouse / sf_data RU fork).

Checks:
1. Quarterly_Sales_Review.xlsx with 4 sheets (RU Region/Segment/Category values from sf_data)
2. Sales_Review_Presentation.pptx with >= 5 RU slides
3. Email sent to regional_managers@company.com

PASS = accuracy >= 70 AND no CRITICAL check failed.
CRITICAL checks reflect SUBSTANCE (correct values from the russified DB + PDF targets).
"""
import argparse
import os
import sys

import openpyxl
import psycopg2

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=5432,
          dbname=os.environ.get("PGDATABASE", "cowork_gym"),
          user="eigent", password="camel")

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILED = []


def check(name, condition, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS]{' [CRITICAL]' if critical else ''} {name}")
    else:
        FAIL_COUNT += 1
        if critical:
            CRITICAL_FAILED.append(name)
        print(f"  [FAIL]{' [CRITICAL]' if critical else ''} {name}: {str(detail)[:300]}")


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def get_sheet(wb, name):
    for s in wb.sheetnames:
        if s.strip().lower() == name.strip().lower():
            return wb[s]
    return None


def check_excel(agent_workspace, groundtruth_workspace):
    print("\n=== Checking Quarterly_Sales_Review.xlsx ===")
    agent_file = os.path.join(agent_workspace, "Quarterly_Sales_Review.xlsx")
    gt_file = os.path.join(groundtruth_workspace, "Quarterly_Sales_Review.xlsx")

    check("Excel file exists", os.path.isfile(agent_file), agent_file)
    if not os.path.isfile(agent_file):
        return

    try:
        agent_wb = openpyxl.load_workbook(agent_file, data_only=True)
        gt_wb = openpyxl.load_workbook(gt_file, data_only=True)
    except Exception as e:
        check("Excel readable", False, str(e))
        return

    # ---- Regional_Performance ----
    print("  Checking Regional_Performance...")
    a_sheet = get_sheet(agent_wb, "Regional_Performance")
    g_sheet = get_sheet(gt_wb, "Regional_Performance")
    check("Sheet 'Regional_Performance' exists", a_sheet is not None,
          f"Sheets: {agent_wb.sheetnames}")
    if a_sheet and g_sheet:
        a_rows = list(a_sheet.iter_rows(min_row=2, values_only=True))
        g_rows = list(g_sheet.iter_rows(min_row=2, values_only=True))
        check("Regional_Performance has 5 rows", len(a_rows) == 5, f"Got {len(a_rows)}")

        a_lookup = {str(r[0]).strip().lower(): r for r in a_rows if r and r[0]}
        # CRITICAL: every russified region present with correct Actual + Status
        region_substance_ok = True
        region_detail = ""
        for g_row in g_rows:
            if not g_row or not g_row[0]:
                continue
            key = str(g_row[0]).strip().lower()
            a_row = a_lookup.get(key)
            if a_row is None:
                check(f"Region '{g_row[0]}' present", False, "Missing")
                region_substance_ok = False
                region_detail = f"missing region {g_row[0]}"
                continue
            # Target_Revenue (non-critical structural)
            if len(a_row) > 1 and len(g_row) > 1:
                check(f"'{key}' Target",
                      num_close(a_row[1], g_row[1], 1000),
                      f"Expected {g_row[1]}, got {a_row[1]}")
            # Actual_Revenue
            if len(a_row) > 2 and len(g_row) > 2:
                ok = num_close(a_row[2], g_row[2], 500)
                check(f"'{key}' Actual", ok, f"Expected {g_row[2]}, got {a_row[2]}")
                if not ok:
                    region_substance_ok = False
                    region_detail = f"{key} actual {a_row[2]} vs {g_row[2]}"
            # Variance_Pct
            if len(a_row) > 4 and len(g_row) > 4:
                check(f"'{key}' Variance_Pct",
                      num_close(a_row[4], g_row[4], 1.0),
                      f"Expected {g_row[4]}, got {a_row[4]}")
            # Status (Met/Missed literal kept English)
            if len(a_row) > 5 and len(g_row) > 5:
                a_status = str(a_row[5]).strip().lower() if a_row[5] else ""
                g_status = str(g_row[5]).strip().lower() if g_row[5] else ""
                ok = a_status == g_status
                check(f"'{key}' Status", ok, f"Expected {g_status}, got {a_status}")
                if not ok:
                    region_substance_ok = False
                    region_detail = f"{key} status {a_status} vs {g_status}"
        check("CRITICAL Regional_Performance: all RU regions correct Actual+Status",
              region_substance_ok, region_detail, critical=True)

    # ---- Segment_Breakdown ----
    print("  Checking Segment_Breakdown...")
    a_sheet = get_sheet(agent_wb, "Segment_Breakdown")
    g_sheet = get_sheet(gt_wb, "Segment_Breakdown")
    check("Sheet 'Segment_Breakdown' exists", a_sheet is not None,
          f"Sheets: {agent_wb.sheetnames}")
    if a_sheet and g_sheet:
        a_rows = list(a_sheet.iter_rows(min_row=2, values_only=True))
        check("Segment_Breakdown has 4 rows", len(a_rows) == 4, f"Got {len(a_rows)}")

        a_lookup = {str(r[0]).strip().lower(): r for r in a_rows if r and r[0]}
        seg_substance_ok = True
        seg_detail = ""
        for g_row in g_sheet.iter_rows(min_row=2, values_only=True):
            if not g_row or not g_row[0]:
                continue
            key = str(g_row[0]).strip().lower()
            a_row = a_lookup.get(key)
            if a_row is None:
                check(f"Segment '{g_row[0]}' present", False, "Missing")
                seg_substance_ok = False
                seg_detail = f"missing segment {g_row[0]}"
                continue
            if len(a_row) > 2 and len(g_row) > 2:
                ok = num_close(a_row[2], g_row[2], 50)
                check(f"'{key}' Order_Count", ok, f"Expected {g_row[2]}, got {a_row[2]}")
                if not ok:
                    seg_substance_ok = False
                    seg_detail = f"{key} order_count {a_row[2]} vs {g_row[2]}"
            if len(a_row) > 3 and len(g_row) > 3:
                ok = num_close(a_row[3], g_row[3], 5000)
                check(f"'{key}' Total_Revenue", ok, f"Expected {g_row[3]}, got {a_row[3]}")
                if not ok:
                    seg_substance_ok = False
                    seg_detail = f"{key} total_revenue {a_row[3]} vs {g_row[3]}"
        check("CRITICAL Segment_Breakdown: all 4 RU segments with correct Order_Count+Total_Revenue",
              seg_substance_ok, seg_detail, critical=True)

    # ---- Top_Products ----
    print("  Checking Top_Products...")
    a_sheet = get_sheet(agent_wb, "Top_Products")
    check("Sheet 'Top_Products' exists", a_sheet is not None,
          f"Sheets: {agent_wb.sheetnames}")
    if a_sheet:
        a_rows = list(a_sheet.iter_rows(min_row=2, values_only=True))
        check("Top_Products has 5 rows", len(a_rows) == 5, f"Got {len(a_rows)}")
        if a_rows:
            revenues = [float(r[3]) for r in a_rows if r and len(r) > 3 and r[3]]
            if revenues:
                check("Top product revenue > 10000",
                      max(revenues) > 10000, f"Got max {max(revenues)}")

    # ---- Summary ----
    print("  Checking Summary...")
    a_sheet = get_sheet(agent_wb, "Summary")
    g_sheet = get_sheet(gt_wb, "Summary")
    check("Sheet 'Summary' exists", a_sheet is not None,
          f"Sheets: {agent_wb.sheetnames}")
    if a_sheet and g_sheet:
        a_data = {}
        for row in a_sheet.iter_rows(min_row=2, values_only=True):
            if row and row[0]:
                a_data[str(row[0]).strip().lower()] = row[1]
        g_data = {}
        for row in g_sheet.iter_rows(min_row=2, values_only=True):
            if row and row[0]:
                g_data[str(row[0]).strip().lower()] = row[1]

        check("Total_Target",
              num_close(a_data.get("total_target"), g_data.get("total_target"), 1000),
              f"Expected {g_data.get('total_target')}, got {a_data.get('total_target')}")

        # CRITICAL: core aggregates (Total_Actual within tol AND Regions_Met_Target exact)
        total_actual_ok = num_close(a_data.get("total_actual"), g_data.get("total_actual"), 5000)
        met_ok = num_close(a_data.get("regions_met_target"), g_data.get("regions_met_target"), 0)
        check("Total_Actual",
              total_actual_ok,
              f"Expected {g_data.get('total_actual')}, got {a_data.get('total_actual')}")
        check("Regions_Met_Target",
              met_ok,
              f"Expected {g_data.get('regions_met_target')}, got {a_data.get('regions_met_target')}")
        check("CRITICAL Summary core aggregates (Total_Actual + Regions_Met_Target)",
              total_actual_ok and met_ok,
              f"total_actual_ok={total_actual_ok}, met_ok={met_ok}", critical=True)

        # CRITICAL: Best_Region matches the russified groundtruth best region.
        # Robust: compare against groundtruth Summary (the broken hardcoded DB query
        # 'sf_data.orders'/'sales' never existed; REGION lives in CUSTOMERS, revenue in
        # ORDERS.TOTAL_AMOUNT). We also attempt a fixed DB confirmation (non-fatal).
        gt_best = str(g_data.get("best_region") or "").strip().lower()
        a_best = str(a_data.get("best_region") or "").strip().lower()
        check("CRITICAL Best_Region matches russified groundtruth",
              a_best != "" and gt_best != "" and (a_best == gt_best or gt_best in a_best),
              f"Expected '{g_data.get('best_region')}', got {a_data.get('best_region')}",
              critical=True)
        # Non-critical fixed DB confirmation that GT best region is the true argmax.
        try:
            conn_db = psycopg2.connect(**DB)
            cur_db = conn_db.cursor()
            cur_db.execute('''
                SELECT c."REGION"
                FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
                JOIN sf_data."SALES_DW__PUBLIC__CUSTOMERS" c
                  ON o."CUSTOMER_ID" = c."CUSTOMER_ID"
                WHERE o."ORDER_DATE" BETWEEN '2025-01-01' AND '2025-03-31'
                GROUP BY c."REGION"
                ORDER BY SUM(o."TOTAL_AMOUNT") DESC
                LIMIT 1
            ''')
            result = cur_db.fetchone()
            cur_db.close()
            conn_db.close()
            if result and result[0]:
                db_best = str(result[0]).strip().lower()
                check("DB argmax(SUM TOTAL_AMOUNT over CUSTOMERS.REGION) matches Best_Region",
                      a_best == db_best or db_best in a_best,
                      f"DB best={result[0]}, agent={a_data.get('best_region')}")
        except Exception:
            pass

        # Worst_Region (non-critical)
        gt_worst = str(g_data.get("worst_region") or "").strip().lower()
        a_worst = str(a_data.get("worst_region") or "").strip().lower()
        check("Worst_Region matches russified groundtruth",
              a_worst != "" and gt_worst != "" and (a_worst == gt_worst or gt_worst in a_worst),
              f"Expected '{g_data.get('worst_region')}', got {a_data.get('worst_region')}")


def check_pptx(agent_workspace):
    print("\n=== Checking Sales_Review_Presentation.pptx ===")
    pptx_path = os.path.join(agent_workspace, "Sales_Review_Presentation.pptx")
    check("PPTX file exists", os.path.isfile(pptx_path))
    if not os.path.isfile(pptx_path):
        return
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        check("Presentation has >= 5 slides", slide_count >= 5,
              f"Got {slide_count} slides")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text.lower() + " "

        def any_in(words):
            return any(w in all_text for w in words)

        check("Contains regional/territory reference (RU/EN)",
              any_in(["region", "territory", "регион", "территори"]))
        check("Contains sales/revenue reference (RU/EN)",
              any_in(["sales", "revenue", "продаж", "выручк"]))
        check("Contains recommendation (RU/EN)",
              any_in(["recommend", "next step", "action",
                      "рекоменд", "следующий шаг", "дальнейш", "меропри", "действи"]))
    except ImportError:
        check("python-pptx available", False)
    except Exception as e:
        check("PPTX readable", False, str(e))


def check_email():
    print("\n=== Checking Email ===")
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        # Subject preserved English per rule 3.
        cur.execute("""
            SELECT m.subject, m.to_addr, m.body_text
            FROM email.sent_log sl
            JOIN email.messages m ON sl.message_id = m.id
            WHERE lower(m.subject) LIKE '%%sales%%territory%%review%%'
               OR lower(m.subject) LIKE '%%q1%%sales%%'
        """)
        rows = cur.fetchall()
        if not rows:
            cur.execute("""
                SELECT subject, to_addr, body_text FROM email.messages
                WHERE lower(subject) LIKE '%%sales%%territory%%review%%'
                   OR lower(subject) LIKE '%%q1%%sales%%'
            """)
            rows = cur.fetchall()
        check("Sales review email sent", len(rows) > 0, f"Found {len(rows)}")
        to_ok = False
        if rows:
            to_str = str(rows[0][1]).lower() if rows[0][1] else ""
            to_ok = "regional_manager" in to_str
            check("Email to regional_managers", to_ok, f"To: {rows[0][1]}")
        # CRITICAL: review email exists AND addressed to a regional_manager address.
        check("CRITICAL Sales review email to regional_manager exists",
              len(rows) > 0 and to_ok,
              f"rows={len(rows)}, to_ok={to_ok}", critical=True)
        cur.close()
        conn.close()
    except Exception as e:
        check("Email check", False, str(e))
        check("CRITICAL Sales review email to regional_manager exists", False,
              str(e), critical=True)


def check_reverse_validation(workspace):
    """Verify things that should NOT exist in output."""
    print("\n=== Reverse Validation ===")

    path = os.path.join(workspace, "Quarterly_Sales_Review.xlsx")
    if os.path.isfile(path):
        wb = openpyxl.load_workbook(path, data_only=True)
        has_negative = False
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            for row in ws.iter_rows(min_row=2, values_only=True):
                for cell in row:
                    if isinstance(cell, (int, float)) and cell < -100000:
                        has_negative = True
                        break
                if has_negative:
                    break
            if has_negative:
                break
        check("No large negative values in sales Excel", not has_negative,
              "Found unexpectedly large negative revenue value")

    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT COUNT(*) FROM email.messages
            WHERE (lower(subject) LIKE '%%sales%%' OR lower(subject) LIKE '%%revenue%%')
              AND to_addr::text ILIKE '%%competitor%%'
        """)
        bad_count = cur.fetchone()[0]
        check("No sales emails to competitor addresses", bad_count == 0,
              f"Found {bad_count}")
        cur.close()
        conn.close()
    except Exception:
        pass


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    check_excel(args.agent_workspace, gt_dir)
    check_pptx(args.agent_workspace)
    check_email()
    check_reverse_validation(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100.0) if total else 0.0
    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Accuracy: {accuracy:.1f}%")

    if CRITICAL_FAILED:
        print(f"  CRITICAL checks failed: {CRITICAL_FAILED}")
        print("  Overall: FAIL (critical)")
        sys.exit(1)

    overall = accuracy >= 70.0
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
