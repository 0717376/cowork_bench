"""
Evaluation for kulinar-nutrition-ppt task (kulinar recipe DB).

The agent freely selects 5 recipes from the Russian kulinar database (>=4
categories), looks up ingredient nutrition in the workspace nutrition_guide.md,
and produces Nutrition_Comparison.xlsx + Healthy_Eating_Guide.pptx.

Because the chosen recipes are agent-selected, the SEMANTIC critical checks are
SELF-CONSISTENT / SOURCE-RECOMPUTED rather than hardcoded to specific recipes:
  - Health_Score is recomputed from each row's own Protein/Fiber/Calories via the
    guide formula (High=3, Medium=2, Low=1) — the formula must actually be applied.
  - The Recommendation list must be the true top-K of the 5 overview rows by the
    recomputed Health_Score, in strict descending order.

Scoring: accuracy >= 70% AND no CRITICAL check failed => PASS. Any critical
failure => FAIL regardless of accuracy.
"""

from argparse import ArgumentParser
import re
import sys
import os
from pathlib import Path

PASS_COUNT = 0
FAIL_COUNT = 0
FAILED_NAMES = []

# Semantic (critical) checks: any failure => overall FAIL regardless of accuracy.
CRITICAL_CHECKS = {
    "Excel file exists with required sheets/columns",
    "Recipe Overview: >=5 recipes spanning >=4 distinct categories",
    "Health_Score actually computed from guide formula",
    "Recommendation is the true top-K by Health_Score (descending)",
    "PPTX exists with required slides and recommended recipes",
}

POINTS = {"high": 3, "medium": 2, "low": 1}


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        FAILED_NAMES.append(name)
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    try:
        return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)
    except (TypeError, ValueError):
        return False


def parse_guide_calories(workspace):
    """Parse nutrition_guide.md -> {ingredient_lower: calories} from the WORKSPACE
    guide (the one the agent reads), not the groundtruth copy."""
    path = Path(workspace) / "nutrition_guide.md"
    cal = {}
    if not path.exists():
        return cal
    text = path.read_text(encoding="utf-8")
    # Lines like "Картофель: 77 калорий на 100 г, ..."
    for line in text.splitlines():
        m = re.match(r"^([^:]+):\s*(\d+(?:\.\d+)?)\s*кал", line.strip())
        if m:
            cal[m.group(1).strip().lower()] = float(m.group(2))
    return cal


def calories_feasible(value, guide_cal, kmin=1, kmax=5):
    """Is `value` reachable by summing the per-100g calories of some 1..kmax guide
    ingredients (within +/-15%)? Tolerant subset-sum feasibility — absorbs the
    agent's freedom in choosing main ingredients."""
    vals = sorted(guide_cal.values())
    if not vals:
        return True  # cannot validate -> do not penalize
    target = float(value)
    # DFS over combinations of size up to kmax (prune on exceeding upper tol).
    upper = target * 1.15
    n = len(vals)

    def dfs(start, count, total):
        if count >= kmin and num_close(total, target):
            return True
        if count >= kmax:
            return False
        for i in range(start, n):
            nt = total + vals[i]
            if nt > upper and count + 1 >= kmin:
                # vals sorted asc: further picks only grow the sum
                if total >= target * 0.85:
                    break
            if dfs(i + 1, count + 1, nt):
                return True
        return False

    return dfs(0, 0, 0.0)


def get_col_idx(header_lower, names):
    for n in names:
        if n in header_lower:
            return header_lower.index(n)
    return None


def check_excel(workspace):
    """Returns (overview_rows, rec_rows) dicts or None on hard failure."""
    import openpyxl

    xlsx_path = Path(workspace) / "Nutrition_Comparison.xlsx"
    if not xlsx_path.exists():
        record("Excel file exists with required sheets/columns", False, "Nutrition_Comparison.xlsx not found")
        return None, None

    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    sheet_names_lower = {s.lower(): s for s in wb.sheetnames}
    if "recipe overview" not in sheet_names_lower or "recommendation" not in sheet_names_lower:
        record("Excel file exists with required sheets/columns", False, f"Sheets: {wb.sheetnames}")
        return None, None

    ws1 = wb[sheet_names_lower["recipe overview"]]
    rows1 = list(ws1.iter_rows(values_only=True))
    if len(rows1) < 2:
        record("Excel file exists with required sheets/columns", False, "Recipe Overview empty")
        return None, None

    header1_lower = [(str(h).strip().lower().replace(" ", "_") if h else "") for h in rows1[0]]
    expected_overview = ["recipe_name", "category", "difficulty", "estimated_calories", "protein_level", "fiber_level"]
    missing = [c for c in expected_overview if c not in header1_lower]
    if missing:
        record("Excel file exists with required sheets/columns", False, f"Overview missing {missing}")
        return None, None
    idx = {c: header1_lower.index(c) for c in expected_overview}

    ws2 = wb[sheet_names_lower["recommendation"]]
    rows2 = list(ws2.iter_rows(values_only=True))
    if len(rows2) < 2:
        record("Excel file exists with required sheets/columns", False, "Recommendation empty")
        return None, None
    header2_lower = [(str(h).strip().lower().replace(" ", "_") if h else "") for h in rows2[0]]
    expected_rec = ["rank", "recipe_name", "health_score", "reason"]
    missing2 = [c for c in expected_rec if c not in header2_lower]
    if missing2:
        record("Excel file exists with required sheets/columns", False, f"Recommendation missing {missing2}")
        return None, None
    idx2 = {c: header2_lower.index(c) for c in expected_rec}

    record("Excel file exists with required sheets/columns", True)

    # --- Parse overview rows ---
    overview = []
    for r in rows1[1:]:
        name = str(r[idx["recipe_name"]]).strip() if r[idx["recipe_name"]] else ""
        if not name:
            continue
        overview.append({
            "name": name,
            "category": (str(r[idx["category"]]).strip() if r[idx["category"]] else ""),
            "difficulty": (str(r[idx["difficulty"]]).strip().lower() if r[idx["difficulty"]] else ""),
            "calories": r[idx["estimated_calories"]],
            "protein": (str(r[idx["protein_level"]]).strip().lower() if r[idx["protein_level"]] else ""),
            "fiber": (str(r[idx["fiber_level"]]).strip().lower() if r[idx["fiber_level"]] else ""),
        })

    # Structural: >=5 rows, valid enums, numeric calories in range
    record("Recipe Overview has >=5 data rows", len(overview) >= 5, f"got {len(overview)}")

    enums_ok = True
    cal_ok = True
    for row in overview:
        if row["difficulty"] not in ("easy", "medium", "hard"):
            enums_ok = False
        if row["protein"] not in ("high", "medium", "low"):
            enums_ok = False
        if row["fiber"] not in ("high", "medium", "low"):
            enums_ok = False
        try:
            cv = float(row["calories"])
            if cv < 50 or cv > 5000:
                cal_ok = False
        except (TypeError, ValueError):
            cal_ok = False
    record("Overview enum values valid (Easy/Medium/Hard, High/Medium/Low)", enums_ok)
    record("Overview Estimated_Calories numeric in 50-5000", cal_ok)

    categories = {row["category"].lower() for row in overview if row["category"]}
    record("Recipe Overview: >=5 recipes spanning >=4 distinct categories",
           len(overview) >= 5 and len(categories) >= 4,
           f"{len(overview)} recipes, {len(categories)} categories: {sorted(categories)}")

    # --- Parse recommendation rows ---
    rec = []
    for r in rows2[1:]:
        name = str(r[idx2["recipe_name"]]).strip() if r[idx2["recipe_name"]] else ""
        if not name:
            continue
        rec.append({
            "rank": r[idx2["rank"]],
            "name": name,
            "score": r[idx2["health_score"]],
            "reason": (str(r[idx2["reason"]]).strip() if r[idx2["reason"]] else ""),
        })

    record("Recommendation has >=5 data rows", len(rec) >= 5, f"got {len(rec)}")

    reasons_ok = all(r["reason"] for r in rec)
    record("All recommendation Reason cells non-empty", reasons_ok)

    # Recommended names must appear in overview
    overview_by_name = {row["name"].lower(): row for row in overview}
    cross_ok = all(r["name"].lower() in overview_by_name for r in rec)
    record("Recommended recipes appear in Recipe Overview", cross_ok)

    wb.close()
    return overview, rec


def check_calories_vs_guide(overview, workspace):
    """NON-critical: at least 4/5 overview calorie values are reachable as a sum of
    1..5 guide ingredient values (+/-15%), confirming the agent used the guide."""
    guide_cal = parse_guide_calories(workspace)
    if not guide_cal:
        record("Estimated_Calories reconcile with nutrition_guide.md (>=4/5)", True,
               "guide unreadable -> skipped")
        return
    feasible = 0
    for row in overview:
        try:
            if calories_feasible(float(row["calories"]), guide_cal):
                feasible += 1
        except (TypeError, ValueError):
            pass
    record("Estimated_Calories reconcile with nutrition_guide.md (>=4/5)",
           feasible >= min(4, len(overview)),
           f"{feasible}/{len(overview)} rows reconcile with guide")


def check_health_score(overview, rec):
    """CRITICAL: Health_Score must be the guide formula applied to each row's own
    Protein/Fiber/Calories; Recommendation must be the true top-K descending."""
    overview_by_name = {row["name"].lower(): row for row in overview}

    def expected_score(row):
        pp = POINTS.get(row["protein"])
        fp = POINTS.get(row["fiber"])
        try:
            cal = float(row["calories"])
        except (TypeError, ValueError):
            return None
        if pp is None or fp is None:
            return None
        return (pp + fp) * 10 - cal / 100.0

    # 1) Formula actually applied to each recommendation row.
    formula_ok = len(rec) > 0
    for r in rec:
        ov = overview_by_name.get(r["name"].lower())
        if ov is None:
            formula_ok = False
            break
        exp = expected_score(ov)
        try:
            got = float(r["score"])
        except (TypeError, ValueError):
            formula_ok = False
            break
        if exp is None or abs(got - exp) > 0.5:
            formula_ok = False
            break
    record("Health_Score actually computed from guide formula", formula_ok)

    # 2) Recommendation = true top-K by recomputed score, strictly descending.
    scored = []
    all_scorable = True
    for row in overview:
        s = expected_score(row)
        if s is None:
            all_scorable = False
        scored.append((row["name"].lower(), s))
    if not all_scorable:
        record("Recommendation is the true top-K by Health_Score (descending)", False,
               "some overview rows not scorable")
        return

    ranked = sorted(scored, key=lambda t: t[1], reverse=True)
    k = len(rec)
    true_topk = {n for n, _ in ranked[:k]}
    rec_names = [r["name"].lower() for r in rec]
    set_ok = set(rec_names) == true_topk

    # Strict descending order of the agent's listed scores.
    desc_ok = True
    prev = None
    for r in rec:
        try:
            sv = float(r["score"])
        except (TypeError, ValueError):
            desc_ok = False
            break
        if prev is not None and sv > prev + 0.01:
            desc_ok = False
            break
        prev = sv

    record("Recommendation is the true top-K by Health_Score (descending)",
           set_ok and desc_ok,
           f"top-{k} expected={sorted(true_topk)}, got={sorted(set(rec_names))}, desc={desc_ok}")


def check_pptx(workspace, rec):
    from pptx import Presentation

    pptx_path = Path(workspace) / "Healthy_Eating_Guide.pptx"
    if not pptx_path.exists():
        record("PPTX exists with required slides and recommended recipes", False, "not found")
        return

    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)

    all_text = []
    for slide in slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    texts.append(p.text)
        all_text.append("\n".join(texts))
    full_text = "\n".join(all_text).lower()
    first_slide = all_text[0].lower() if all_text else ""

    # Title keyword check: English mandated title 'Healthy Eating Guide', plus RU fallbacks.
    has_keyword = any(kw in first_slide for kw in [
        "healthy", "nutrition", "eating", "health", "wellness",
        "здоров", "питани",
    ])

    rec_names = [r["name"] for r in rec] if rec else []
    found = sum(1 for rn in rec_names if rn.lower() in full_text)

    last_slide = all_text[-1] if all_text else ""

    # Bundle the structural pieces; the critical line requires all to hold.
    record("PPTX slide count >= 7", len(slides) >= 7, f"got {len(slides)}")
    record("PPTX title slide has healthy-eating keyword", has_keyword, first_slide[:120])
    record("PPTX summary slide has content", len(last_slide.strip()) >= 10)

    recipes_ok = (not rec_names) or found >= min(3, len(rec_names))
    record("PPTX mentions recommended recipe names (>=3)", recipes_ok,
           f"{found}/{len(rec_names)} found")

    critical_ok = (len(slides) >= 7) and has_keyword and (len(last_slide.strip()) >= 10) and recipes_ok
    record("PPTX exists with required slides and recommended recipes", critical_ok)


if __name__ == "__main__":
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False)
    args = parser.parse_args()

    workspace = args.agent_workspace
    if not workspace:
        print("Error: --agent_workspace is required")
        sys.exit(1)

    print("\n--- Excel ---")
    overview, rec = None, None
    try:
        overview, rec = check_excel(workspace)
    except Exception as e:
        record("Excel file exists with required sheets/columns", False, str(e))
        import traceback
        traceback.print_exc()

    if overview is not None and rec is not None:
        print("\n--- Calorie/guide reconciliation ---")
        try:
            check_calories_vs_guide(overview, workspace)
        except Exception as e:
            record("Estimated_Calories reconcile with nutrition_guide.md (>=4/5)", False, str(e))

        print("\n--- Health Score formula & ranking ---")
        try:
            check_health_score(overview, rec)
        except Exception as e:
            record("Health_Score actually computed from guide formula", False, str(e))
            record("Recommendation is the true top-K by Health_Score (descending)", False, str(e))

    print("\n--- PPTX ---")
    try:
        check_pptx(workspace, rec)
    except Exception as e:
        record("PPTX exists with required slides and recommended recipes", False, str(e))
        import traceback
        traceback.print_exc()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: no checks performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\n=== Results: {PASS_COUNT}/{total} passed ({accuracy:.1f}%) ===")

    critical_failed = [n for n in FAILED_NAMES if n in CRITICAL_CHECKS]
    if critical_failed:
        print(f"CRITICAL FAILURES ({len(critical_failed)}):")
        for n in critical_failed:
            print(f"  - {n}")

    if args.res_log_file:
        import json
        with open(args.res_log_file, "w") as f:
            json.dump({
                "total_passed": PASS_COUNT,
                "total_checks": total,
                "accuracy": accuracy,
                "critical_failed": critical_failed,
            }, f, indent=2)

    if critical_failed:
        print("FAIL (critical check failed)")
        sys.exit(1)
    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    print("FAIL")
    sys.exit(1)
