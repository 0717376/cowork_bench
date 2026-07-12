"""
Evaluation for rzd-sber-gazp-investor-roadshow-spb-nvg-excel-ppt-email.

Checks split into structural (~30) and CRITICAL (3):
  CRITICAL #1: PPTX has >= 5 slides AND first 5 cover expected sections
                (title, travel, financials, comparison, consensus).
  CRITICAL #2: PPTX text mentions BOTH tickers (SBER.ME and GAZP.ME).
  CRITICAL #3: PPTX travel slide mentions BOTH train numbers (752, 818)
               AND BOTH cities (СПб, В.Новгород).

Any CRITICAL fail → overall FAIL regardless of accuracy %.
Otherwise PASS if accuracy >= 70%.

MOEX/financial values are checked softly (presence and sanity ranges only).
"""
import json
import os
import re
import sys
from argparse import ArgumentParser

import psycopg2
import openpyxl

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []
CRITICAL_CHECKS = {
    "PPTX: 5 слайдов с корректной структурой (title/travel/financials/comparison/consensus)",
    "PPTX: упомянуты оба тикера SBER.ME и GAZP.ME",
    "PPTX (travel slide): оба номера поездов (752/818) и оба города (СПб + В.Новгород)",
}


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    msg = f": {detail[:300]}" if detail else ""
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        marker = " [CRITICAL]" if name in CRITICAL_CHECKS else ""
        print(f"  [FAIL]{marker} {name}{msg}")
        if name in CRITICAL_CHECKS:
            CRITICAL_FAILS.append(name)


ISO_DATE_RE = re.compile(r"^\d{4}-\d{2}-\d{2}$")


def is_iso_date(s):
    return bool(s) and ISO_DATE_RE.match(str(s).strip())


def normalize_train(s):
    """752А / 752A → 752 (latin/cyrillic letter stripped)."""
    if s is None:
        return ""
    return re.sub(r"[^0-9]", "", str(s))


def check_excel(agent_workspace):
    print("\n=== Check 1: Excel Roadshow_Analysis.xlsx ===")

    xlsx_path = os.path.join(agent_workspace, "Roadshow_Analysis.xlsx")
    if not os.path.exists(xlsx_path):
        record("Roadshow_Analysis.xlsx exists", False, f"Not found at {xlsx_path}")
        return
    record("Roadshow_Analysis.xlsx exists", True)

    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    except Exception as e:
        record("Excel readable", False, str(e))
        return
    record("Excel readable", True)

    sn = [s.lower() for s in wb.sheetnames]
    has_travel = any("travel" in s for s in sn)
    has_stock = any("stock" in s or "summary" in s for s in sn)
    has_fin = any("financial" in s or "highlight" in s for s in sn)
    record("Excel has Travel_Plan sheet", has_travel, f"Sheets: {wb.sheetnames}")
    record("Excel has Stock_Summary sheet", has_stock, f"Sheets: {wb.sheetnames}")
    record("Excel has Financial_Highlights sheet", has_fin, f"Sheets: {wb.sheetnames}")

    # ----- Travel_Plan -----
    if has_travel:
        ws_name = wb.sheetnames[next(i for i, s in enumerate(sn) if "travel" in s)]
        ws = wb[ws_name]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = [r for r in rows[1:] if any(c is not None and str(c).strip() for c in r)]
        record("Travel_Plan: >= 2 строк", len(data_rows) >= 2, f"Found {len(data_rows)}")

        all_text = " ".join(str(c) for row in rows for c in row if c is not None)
        all_lower = all_text.lower()
        nums = " ".join(normalize_train(c) for row in rows for c in row if c is not None)

        record("Travel_Plan: содержит 752 (СПб)", "752" in nums, f"Sample: {all_text[:200]}")
        record("Travel_Plan: содержит 818 (Новгород)", "818" in nums, f"Sample: {all_text[:200]}")
        record("Travel_Plan: упомянут Санкт-Петербург", "петербург" in all_lower or "спб" in all_lower)
        record("Travel_Plan: упомянут Великий Новгород", "новгород" in all_lower)
        record("Travel_Plan: класс Бизнес", "бизнес" in all_lower)

        # Meeting_Date ISO
        date_cells = []
        for r in data_rows:
            for c in r:
                if c and re.match(r"^\d{4}-\d{2}-\d{2}$", str(c).strip()):
                    date_cells.append(str(c).strip())
        record(
            "Travel_Plan: даты в ISO YYYY-MM-DD",
            len(date_cells) >= 2,
            f"Found ISO dates: {date_cells}",
        )
        allowed = {"2026-03-10", "2026-03-12"}
        record(
            "Travel_Plan: даты ∈ {2026-03-10, 2026-03-12}",
            all(d in allowed for d in date_cells) and len(date_cells) >= 2,
            f"Got: {date_cells}",
        )

    # ----- Stock_Summary -----
    if has_stock:
        ws_name = wb.sheetnames[next(i for i, s in enumerate(sn) if "stock" in s or "summary" in s)]
        ws = wb[ws_name]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = [r for r in rows[1:] if any(c is not None and str(c).strip() for c in r)]
        record("Stock_Summary: >= 2 строк", len(data_rows) >= 2, f"Found {len(data_rows)}")

        all_text_upper = " ".join(str(c) for row in rows for c in row if c is not None).upper()
        record(
            "Stock_Summary: тикер SBER (.ME)",
            "SBER" in all_text_upper,
            f"Sample: {all_text_upper[:200]}",
        )
        record(
            "Stock_Summary: тикер GAZP (.ME)",
            "GAZP" in all_text_upper,
            f"Sample: {all_text_upper[:200]}",
        )

        # Soft check: numeric values present and > 0 in each row
        numeric_cells_total = 0
        for r in data_rows:
            for c in r:
                if isinstance(c, (int, float)) and c > 0:
                    numeric_cells_total += 1
        record(
            "Stock_Summary: достаточно положительных числовых значений",
            numeric_cells_total >= 6,
            f"positive numeric cells: {numeric_cells_total}",
        )

    # ----- Financial_Highlights -----
    if has_fin:
        ws_name = wb.sheetnames[next(i for i, s in enumerate(sn) if "financial" in s or "highlight" in s)]
        ws = wb[ws_name]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = [r for r in rows[1:] if any(c is not None and str(c).strip() for c in r)]
        record(
            "Financial_Highlights: >= 2 строк (Revenue/Net Income × 2 года ≥ 4 идеально)",
            len(data_rows) >= 2,
            f"Found {len(data_rows)}",
        )
        text_lower = " ".join(str(c) for row in rows for c in row if c is not None).lower()
        record("Financial_Highlights: упомянут revenue", "revenue" in text_lower or "выручка" in text_lower)
        record(
            "Financial_Highlights: упомянут net income",
            "net income" in text_lower or "чистая прибыль" in text_lower or "прибыль" in text_lower,
        )


def check_pptx(agent_workspace):
    print("\n=== Check 2: PPTX Investor_Roadshow.pptx ===")

    pptx_path = os.path.join(agent_workspace, "Investor_Roadshow.pptx")
    if not os.path.exists(pptx_path):
        record("Investor_Roadshow.pptx exists", False, f"Not found at {pptx_path}")
        # critical checks обязаны быть оценены: явно проваливаем оба
        record(
            "PPTX: 5 слайдов с корректной структурой (title/travel/financials/comparison/consensus)",
            False, "PPTX missing",
        )
        record("PPTX: упомянуты оба тикера SBER.ME и GAZP.ME", False, "PPTX missing")
        record(
            "PPTX (travel slide): оба номера поездов (752/818) и оба города (СПб + В.Новгород)",
            False, "PPTX missing",
        )
        return
    record("Investor_Roadshow.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        record("PPTX readable", False, str(e))
        return
    record("PPTX readable", True)

    slide_texts = []
    for slide in prs.slides:
        chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        chunks.append(r.text or "")
            elif hasattr(shape, "text"):
                chunks.append(shape.text or "")
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        chunks.append(cell.text or "")
        slide_texts.append(" ".join(chunks))

    all_text = " ".join(slide_texts)
    all_lower = all_text.lower()
    n_slides = len(slide_texts)

    record("PPTX: >= 5 слайдов", n_slides >= 5, f"Found {n_slides}")

    # ---- CRITICAL #1: 5 слайдов с корректной структурой ----
    # Гайд: 1=title (roadshow), 2=travel, 3=financials, 4=stock comparison, 5=consensus.
    # Допускаем порядок слайдов; ищем 5 явных секций в первых 5 слайдах суммарно.
    first5 = " ".join(slide_texts[:5]).lower() if n_slides >= 5 else ""

    has_title_section = "roadshow" in first5 or "инвест" in first5
    has_travel_section = (("travel" in first5 or "поезд" in first5 or "752" in first5 or "818" in first5
                          or "новгород" in first5 or "петербург" in first5)
                          and ("санкт" in first5 or "новгород" in first5))
    has_financials_section = (
        ("revenue" in first5 or "выручка" in first5)
        and ("eps" in first5 or "прибыль" in first5 or "net income" in first5)
    )
    has_comparison_section = ("sber" in first5 and "gazp" in first5)
    has_consensus_section = (
        "consensus" in first5
        or "консенсус" in first5
        or "buy" in first5
        or "hold" in first5
        or "рекомендац" in first5
    )

    structure_ok = (n_slides >= 5 and has_title_section and has_travel_section
                    and has_financials_section and has_comparison_section
                    and has_consensus_section)

    detail = (
        f"slides={n_slides}, title={has_title_section}, travel={has_travel_section}, "
        f"fin={has_financials_section}, cmp={has_comparison_section}, cons={has_consensus_section}"
    )
    record(
        "PPTX: 5 слайдов с корректной структурой (title/travel/financials/comparison/consensus)",
        structure_ok, detail,
    )

    # ---- CRITICAL #2: оба тикера ----
    has_sber = "sber" in all_lower
    has_gazp = "gazp" in all_lower
    record(
        "PPTX: упомянуты оба тикера SBER.ME и GAZP.ME",
        has_sber and has_gazp,
        f"sber={has_sber}, gazp={has_gazp}",
    )

    # ---- CRITICAL #3: travel slide содержит оба поезда и оба города ----
    # Определим travel slide: тот в первых 5, где есть "752" или "818" или "travel/поезд".
    travel_slide_text = ""
    for i, st in enumerate(slide_texts[:5]):
        sl = st.lower()
        if "752" in sl or "818" in sl or "travel" in sl or "поезд" in sl or "маршрут" in sl:
            travel_slide_text = sl
            break
    if not travel_slide_text:
        travel_slide_text = first5  # fallback: ищем в первых 5

    has_752 = "752" in travel_slide_text
    has_818 = "818" in travel_slide_text
    has_spb = "санкт" in travel_slide_text or "петербург" in travel_slide_text or "спб" in travel_slide_text
    has_nvg = "новгород" in travel_slide_text
    record(
        "PPTX (travel slide): оба номера поездов (752/818) и оба города (СПб + В.Новгород)",
        has_752 and has_818 and has_spb and has_nvg,
        f"752={has_752}, 818={has_818}, spb={has_spb}, nvg={has_nvg}",
    )

    # ---- Soft checks ----
    record("PPTX: упомянуто 'roadshow' или 'инвест'", has_title_section)
    record("PPTX: упомянуты финансовые метрики (revenue/EPS)", has_financials_section)
    record("PPTX: упомянут консенсус аналитиков", has_consensus_section)


def check_emails():
    print("\n=== Check 3: Emails sent ===")

    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
        messages = cur.fetchall()
        cur.close()
        conn.close()
    except Exception as e:
        record("Email DB readable", False, str(e))
        return

    def to_addresses(to_addr):
        if isinstance(to_addr, list):
            return " ".join(str(r).lower() for r in to_addr)
        elif to_addr:
            try:
                parsed = json.loads(str(to_addr))
                if isinstance(parsed, list):
                    return " ".join(str(r).lower() for r in parsed)
                return str(to_addr).lower()
            except Exception:
                return str(to_addr).lower()
        return ""

    to_investors = [m for m in messages if "investors@fundmanager.ru" in to_addresses(m[2])]
    to_spb = [m for m in messages if "spb_partners@finance.ru" in to_addresses(m[2])]

    record("Email отправлено на investors@fundmanager.ru", len(to_investors) >= 1,
           f"messages total: {len(messages)}")
    record("Email отправлено на spb_partners@finance.ru", len(to_spb) >= 1,
           f"messages total: {len(messages)}")

    if to_investors:
        subj, _, _, body = to_investors[0]
        content = ((subj or "") + " " + (body or "")).lower()
        has_topic = any(k in content for k in ["roadshow", "роадшоу", "расписан", "поездк", "презентац"])
        record("Investors email упоминает roadshow/расписание", has_topic, f"subject: {subj}")

    if to_spb:
        subj, _, _, body = to_spb[0]
        content = ((subj or "") + " " + (body or "")).lower()
        has_meeting = any(k in content for k in ["встреч", "10 март", "2026-03-10", "петербург", "спб"])
        record("SPB partner email упоминает встречу/СПб", has_meeting, f"subject: {subj}")


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)
    check_emails()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
        "critical_fails": CRITICAL_FAILS,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if CRITICAL_FAILS:
        print(f"FAIL: критичные чеки провалены ({len(CRITICAL_FAILS)}): {CRITICAL_FAILS}")
        sys.exit(1)
    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
