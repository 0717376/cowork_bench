"""Evaluation script for fetch-insales-market-analysis-excel-ppt-email."""
import os
import argparse, json, os, sys
import openpyxl

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)



DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent", "password": "camel"
}

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []


def check(name, condition, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT
    tag = "CRITICAL " if critical else ""
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {tag}{name}")
    else:
        FAIL_COUNT += 1
        detail_str = str(detail)[:200] if detail else ""
        print(f"  [FAIL] {tag}{name}: {detail_str}")
        if critical:
            CRITICAL_FAILS.append(name)


def safe_float(val, default=None):
    try:
        if val is None:
            return default
        return float(str(val).replace(',', '').replace('%', '').replace('$', '').strip())
    except (ValueError, TypeError):
        return default


def get_conn():
    import psycopg2
    return psycopg2.connect(**DB_CONFIG)


def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT, CRITICAL_FAILS
    PASS_COUNT = 0
    FAIL_COUNT = 0
    CRITICAL_FAILS = []

    # --- Excel checks ---
    excel_path = os.path.join(agent_workspace, "Competitive_Analysis.xlsx")
    check("Competitive_Analysis.xlsx exists", os.path.exists(excel_path))

    if os.path.exists(excel_path):
        wb = openpyxl.load_workbook(excel_path)
        gt_path = os.path.join(groundtruth_workspace, "Competitive_Analysis.xlsx")
        gt_wb = openpyxl.load_workbook(gt_path) if os.path.exists(gt_path) else None

        # Sheet 1: Category_Comparison
        check("Category_Comparison sheet exists", "Category_Comparison" in wb.sheetnames)
        if "Category_Comparison" in wb.sheetnames:
            ws = wb["Category_Comparison"]
            headers = [str(c.value).strip() if c.value else "" for c in ws[1]]
            data_rows = list(ws.iter_rows(min_row=2, values_only=True))
            check("Category_Comparison has 8 data rows", len(data_rows) == 8, f"got {len(data_rows)}")

            for expected_col in ["Category", "Own_Products", "Own_Avg_Price", "Market_Avg_Price",
                                 "Price_Position", "Own_Revenue", "Market_Revenue",
                                 "Market_Share_Pct", "Market_Growth_Rate"]:
                check(f"Category_Comparison has {expected_col}",
                      expected_col in headers, f"headers: {headers}")

            # Check row order matches groundtruth order (categories sorted by their
            # ENGLISH API name; the Category column carries the Russian store label,
            # so this is NOT Russian-alphabetical -- compare against GT order directly).
            if data_rows and gt_wb and "Category_Comparison" in gt_wb.sheetnames:
                cats = [str(r[0]) for r in data_rows if r[0]]
                gt_order = [str(r[0]) for r in
                            gt_wb["Category_Comparison"].iter_rows(min_row=2, values_only=True)
                            if r[0]]
                check("Category_Comparison row order matches expected (EN-name alphabetical)",
                      cats == gt_order, f"order: {cats}")

            # Compare values with groundtruth
            if gt_wb and "Category_Comparison" in gt_wb.sheetnames:
                gt_ws = gt_wb["Category_Comparison"]
                gt_rows = {str(r[0]): r for r in gt_ws.iter_rows(min_row=2, values_only=True)}
                agent_rows = {str(r[0]): r for r in data_rows}

                match_count = 0
                total_checks = 0
                num_match = 0
                num_total = 0
                pp_match = 0
                for cat in gt_rows:
                    if cat in agent_rows:
                        gt_r = gt_rows[cat]
                        ag_r = agent_rows[cat]
                        # Check numeric columns (indices 1-8)
                        for idx in [1, 2, 3, 5, 6, 7, 8]:
                            gt_val = safe_float(gt_r[idx])
                            ag_val = safe_float(ag_r[idx])
                            total_checks += 1
                            num_total += 1
                            if gt_val is not None and ag_val is not None:
                                if abs(gt_val - ag_val) <= max(abs(gt_val) * 0.05, 1.0):
                                    match_count += 1
                                    num_match += 1
                        # Check Price_Position (index 4)
                        total_checks += 1
                        if str(gt_r[4]).strip().lower() == str(ag_r[4]).strip().lower():
                            match_count += 1
                            pp_match += 1

                accuracy = match_count / total_checks if total_checks > 0 else 0
                check(f"Category_Comparison data accuracy >= 75%",
                      accuracy >= 0.75, f"{match_count}/{total_checks} = {accuracy:.1%}")

                # CRITICAL: numeric accuracy across all 8 RU-labeled categories.
                # Exercises store<->API cross-reference + EN->RU category join (core of task).
                num_acc = num_match / num_total if num_total > 0 else 0
                check("Category_Comparison numeric accuracy >= 85% (8 RU categories)",
                      num_total == 56 and num_acc >= 0.85,
                      f"{num_match}/{num_total} = {num_acc:.1%} (need all 8 categories joined)",
                      critical=True)
                # CRITICAL: Price_Position classification correct for >= 7/8 categories.
                check("Price_Position correct for >= 7/8 categories",
                      pp_match >= 7, f"{pp_match}/8",
                      critical=True)

        # Sheet 2: Strategic_Matrix
        check("Strategic_Matrix sheet exists", "Strategic_Matrix" in wb.sheetnames)
        if "Strategic_Matrix" in wb.sheetnames:
            ws = wb["Strategic_Matrix"]
            headers = [str(c.value).strip() if c.value else "" for c in ws[1]]
            data_rows = list(ws.iter_rows(min_row=2, values_only=True))
            check("Strategic_Matrix has 8 data rows", len(data_rows) == 8, f"got {len(data_rows)}")

            for expected_col in ["Category", "Price_Position", "Market_Share_Pct",
                                 "Market_Growth_Rate", "Growth_Opportunity",
                                 "Strategic_Priority", "Recommended_Action"]:
                check(f"Strategic_Matrix has {expected_col}",
                      expected_col in headers, f"headers: {headers}")

            # Verify Growth_Opportunity and Strategic_Priority logic
            if gt_wb and "Strategic_Matrix" in gt_wb.sheetnames:
                gt_ws = gt_wb["Strategic_Matrix"]
                gt_rows = {str(r[0]): r for r in gt_ws.iter_rows(min_row=2, values_only=True)}
                agent_rows = {str(r[0]): r for r in data_rows}

                go_match = 0
                sp_match = 0
                count = 0
                for cat in gt_rows:
                    if cat in agent_rows:
                        count += 1
                        gt_go = str(gt_rows[cat][4]).strip().lower()
                        ag_go = str(agent_rows[cat][4]).strip().lower()
                        if gt_go == ag_go:
                            go_match += 1
                        gt_sp = str(gt_rows[cat][5]).strip().lower()
                        ag_sp = str(agent_rows[cat][5]).strip().lower()
                        if gt_sp == ag_sp:
                            sp_match += 1

                if count > 0:
                    check(f"Growth_Opportunity accuracy >= 75%",
                          go_match / count >= 0.75, f"{go_match}/{count}")
                    check(f"Strategic_Priority accuracy >= 75%",
                          sp_match / count >= 0.75, f"{sp_match}/{count}")
                # CRITICAL: encodes the >10% growth & <5% share rule and the
                # High/Medium/Low priority rule; >= 7/8 each, all 8 categories joined.
                check("Strategic_Matrix Growth_Opportunity correct for >= 7/8 categories",
                      count == 8 and go_match >= 7, f"{go_match}/{count}",
                      critical=True)
                check("Strategic_Matrix Strategic_Priority correct for >= 7/8 categories",
                      count == 8 and sp_match >= 7, f"{sp_match}/{count}",
                      critical=True)

        # Sheet 3: Executive_Summary
        check("Executive_Summary sheet exists", "Executive_Summary" in wb.sheetnames)
        if "Executive_Summary" in wb.sheetnames:
            ws = wb["Executive_Summary"]
            data_rows = list(ws.iter_rows(min_row=2, values_only=True))
            check("Executive_Summary has >= 8 rows", len(data_rows) >= 8, f"got {len(data_rows)}")

            labels = {str(r[0]).strip(): r[1] for r in data_rows if r[0]}

            if gt_wb and "Executive_Summary" in gt_wb.sheetnames:
                gt_ws = gt_wb["Executive_Summary"]
                gt_labels = {str(r[0]).strip(): r[1]
                             for r in gt_ws.iter_rows(min_row=2, values_only=True) if r[0]}

                for key in ["Total_Own_Products", "Total_Own_Revenue", "Total_Market_Size",
                            "Overall_Market_Share_Pct"]:
                    gt_val = safe_float(gt_labels.get(key))
                    ag_val = safe_float(labels.get(key))
                    if gt_val is not None and ag_val is not None:
                        tol = max(abs(gt_val) * 0.05, 1.0)
                        check(f"Executive_Summary {key} matches",
                              abs(gt_val - ag_val) <= tol,
                              f"expected ~{gt_val}, got {ag_val}")
                    else:
                        check(f"Executive_Summary {key} present",
                              ag_val is not None, f"missing or non-numeric")

                # CRITICAL: key roll-up deliverables (overall share, growth-opp count,
                # high-priority count, total own revenue).
                share = safe_float(labels.get("Overall_Market_Share_Pct"))
                rev = safe_float(labels.get("Total_Own_Revenue"))
                goc = safe_float(labels.get("Growth_Opportunities_Count"))
                hpc = safe_float(labels.get("High_Priority_Categories"))
                rollup_ok = (
                    share is not None and abs(share - 1.4) <= 0.1 and
                    rev is not None and abs(rev - 69833.21) <= max(69833.21 * 0.05, 1.0) and
                    goc is not None and int(round(goc)) == 4 and
                    hpc is not None and int(round(hpc)) == 2
                )
                check("Executive_Summary roll-up (share~1.4, GO=4, HighPri=2, Revenue~69833)",
                      rollup_ok,
                      f"share={share}, rev={rev}, GO={goc}, HighPri={hpc}",
                      critical=True)

    # --- PPT checks ---
    ppt_path = os.path.join(agent_workspace, "Strategy_Presentation.pptx")
    check("Strategy_Presentation.pptx exists", os.path.exists(ppt_path))
    if os.path.exists(ppt_path):
        from pptx import Presentation
        prs = Presentation(ppt_path)
        slide_count = len(prs.slides)
        check("PPT has >= 6 slides", slide_count >= 6, f"got {slide_count}")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text.lower() + " "

        # Accept Russian or English deck wording.
        for variants in [["market", "рынок", "рыноч"], ["competitive", "конкурент"],
                         ["growth", "рост"], ["strategy", "стратег"],
                         ["priority", "приоритет"]]:
            check(f"PPT mentions '{variants[0]}' (RU/EN)",
                  any(v in all_text for v in variants), f"none of {variants}")

    # --- Email checks ---
    try:
        conn = get_conn()
        cur = conn.cursor()

        # Check CEO email
        cur.execute(
            "SELECT subject, to_addr, body_text FROM email.messages WHERE subject ILIKE %s",
            ('%executive summary%',)
        )
        ceo_emails = cur.fetchall()
        check("CEO email sent", len(ceo_emails) >= 1)
        if ceo_emails:
            to_str = str(ceo_emails[0][1]).lower()
            check("CEO email to ceo@company.com", "ceo@company.com" in to_str)

        # Check product team email
        cur.execute(
            "SELECT subject, to_addr, body_text FROM email.messages WHERE subject ILIKE %s",
            ('%detailed findings%',)
        )
        pt_emails = cur.fetchall()
        check("Product team email sent", len(pt_emails) >= 1)
        if pt_emails:
            to_str = str(pt_emails[0][1]).lower()
            check("Product team email to product_team@company.com",
                  "product_team@company.com" in to_str)

        # CRITICAL: email bodies must carry the substantive deliverables, not just
        # be sent to the right address with the right subject. Keywords checked on
        # the ORIGINAL (non-normalized) lowercased body, RU+EN alternatives.
        ceo_body = str(ceo_emails[0][2]).lower() if ceo_emails else ""
        pt_body = str(pt_emails[0][2]).lower() if pt_emails else ""

        ceo_body_ok = (
            any(v in ceo_body for v in ["overall", "общая доля", "доля рынка", "market share"]) and
            any(v in ceo_body for v in ["high priority", "высок", "приоритет"]) and
            any(v in ceo_body for v in ["growth", "рост", "возможност"])
        )
        check("CEO email body covers overall share + high-priority + growth opportunities (RU/EN)",
              ceo_body_ok, f"body: {ceo_body[:160]}", critical=True)

        pt_body_ok = (
            any(v in pt_body for v in ["price leader", "competitive", "premium",
                                       "позиционир", "ценов"]) and
            any(v in pt_body for v in ["доля", "share"]) and
            any(v in pt_body for v in ["priority", "приоритет"])
        )
        check("Product team email body enumerates price position + share + priority (RU/EN)",
              pt_body_ok, f"body: {pt_body[:160]}", critical=True)

        conn.close()
    except Exception as e:
        check("Email check", False, str(e))

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100) if total > 0 else 0

    # Any CRITICAL failure => hard FAIL regardless of overall accuracy.
    if CRITICAL_FAILS:
        return False, (f"CRITICAL FAIL ({len(CRITICAL_FAILS)}): "
                       f"{'; '.join(CRITICAL_FAILS)} | "
                       f"accuracy {PASS_COUNT}/{total} = {accuracy:.1f}%")

    # No critical failure: pass when overall accuracy >= 70%.
    success = accuracy >= 70
    return success, f"Passed {PASS_COUNT}/{total} checks = {accuracy:.1f}% (>=70 required)"


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
