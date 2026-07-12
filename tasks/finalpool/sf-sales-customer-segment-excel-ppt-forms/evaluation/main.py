"""Evaluation for sf-sales-customer-segment-excel-ppt-gform (RU / ClickHouse).

The data warehouse (ClickHouse identity, PG schema sf_data, logical DB SALES_DW)
has its SEGMENT / REGION values russified centrally by db/zzz_clickhouse_after_init.sql.
The agent therefore reads Russian segment/region labels from the warehouse and writes
them into the Excel/PPTX deliverables. Groundtruth artifacts in groundtruth_workspace/
were regenerated with the same Russian labels. To stay robust, segment/region row keys
are normalized through a RU<->EN map so that either spelling matches.

Forms MCP uses the gform.* PG schema (same as before the google_forms->forms swap).

CRITICAL_CHECKS (semantic): any failure => sys.exit(1) before the accuracy gate.
Structural checks (sheet exists, column present, slide count, question count) are NON-critical.
"""
import argparse
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent",
    "password": "camel",
}

# --- RU<->EN normalization (must match db/zzz_clickhouse_after_init.sql seed map) ---
SEG_RU = {
    "consumer": "частные клиенты",
    "enterprise": "корпоративный",
    "government": "государственный",
    "smb": "малый и средний бизнес",
}
SEG_RU2EN = {v: k for k, v in SEG_RU.items()}

REGION_RU = {
    "asia pacific": "азиатско-тихоокеанский регион",
    "europe": "европа",
    "latin america": "латинская америка",
    "middle east": "ближний восток",
    "north america": "северная америка",
}
REGION_RU2EN = {v: k for k, v in REGION_RU.items()}

# Growth_Indicator may be written EN or RU.
GROWTH_RU = {"high": "высокий", "low": "низкий"}
GROWTH_RU2EN = {v: k for k, v in GROWTH_RU.items()}


def seg_key(v):
    """Canonical EN-lowercase key for a segment label (accepts RU or EN)."""
    if v is None:
        return None
    s = str(v).strip().lower()
    return SEG_RU2EN.get(s, s)


def region_key(v):
    if v is None:
        return None
    s = str(v).strip().lower()
    return REGION_RU2EN.get(s, s)


def region_seg_key(c0, c1):
    """Order-independent key for the (Region, Segment) identifier pair.

    task.md fixes the SORT (region then segment) but not the COLUMN ORDER of the
    two identifier columns, so accept Region-first OR Segment-first. The label
    sets are disjoint: a value is either a known region or a known segment.
    """
    n0, n1 = str(c0).strip().lower(), str(c1).strip().lower()
    if n0 in REGION_RU2EN or n0 in REGION_RU:  # c0 is the region
        return f"{region_key(c0)}|{seg_key(c1)}"
    if n1 in REGION_RU2EN or n1 in REGION_RU:  # c1 is the region
        return f"{region_key(c1)}|{seg_key(c0)}"
    # Fallback: assume positional Region-first (matches GT layout).
    return f"{region_key(c0)}|{seg_key(c1)}"


def growth_norm(v):
    if v is None:
        return None
    s = str(v).strip().lower()
    s = GROWTH_RU2EN.get(s, s)
    # tolerate "high growth"/"low growth"
    if s.startswith("high"):
        return "high"
    if s.startswith("low"):
        return "low"
    return s


# Per-segment expected revenue figure (Consumer / Частные клиенты) — kept for the PPT check.
CONSUMER_REVENUE_TOKENS = ["839609", "839,609", "839 609"]


PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILED = []


def record(name, passed, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT
    tag = "CRITICAL " if critical else ""
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {tag}{name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {tag}{name}{msg}")
        if critical:
            CRITICAL_FAILED.append(name)


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def category_match(a, b):
    """Strategic_Category — kept English by design (Star/Cash Cow/Question Mark/Dog)."""
    return str_match(a, b)


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def check_excel(agent_ws, gt_dir):
    print("Checking Excel workbook...")
    agent_file = os.path.join(agent_ws, "Customer_Segment_Analysis.xlsx")
    gt_file = os.path.join(gt_dir, "Customer_Segment_Analysis.xlsx")

    if not os.path.exists(agent_file):
        record("Customer_Segment_Analysis.xlsx exists", False,
               "not found in agent workspace", critical=True)
        return
    record("Customer_Segment_Analysis.xlsx exists", True)
    if not os.path.exists(gt_file):
        record("groundtruth xlsx exists", False, "missing groundtruth", critical=True)
        return

    agent_wb = openpyxl.load_workbook(agent_file, data_only=True)
    gt_wb = openpyxl.load_workbook(gt_file, data_only=True)

    # ---- Sheet 1: Segment Performance ----
    a_rows = load_sheet_rows(agent_wb, "Segment Performance")
    g_rows = load_sheet_rows(gt_wb, "Segment Performance")
    if a_rows is None:
        record("Sheet 'Segment Performance' present", False, "not found", critical=False)
    elif g_rows is None:
        record("Sheet 'Segment Performance' groundtruth", False, "missing", critical=True)
    else:
        record("Sheet 'Segment Performance' present", True)
        a_data = a_rows[1:] if len(a_rows) > 1 else []
        g_data = g_rows[1:] if len(g_rows) > 1 else []
        a_lookup = {seg_key(r[0]): r for r in a_data if r and r[0] is not None}

        rev_ok = True
        share_ok = True
        rev_detail = []
        for g_row in g_data:
            if not g_row or g_row[0] is None:
                continue
            key = seg_key(g_row[0])
            a_row = a_lookup.get(key)
            if a_row is None:
                rev_ok = False
                rev_detail.append(f"missing segment {g_row[0]}")
                continue
            # Total_Revenue (col 3, tol 500) — CRITICAL core deliverable
            if len(a_row) > 3 and len(g_row) > 3 and not num_close(a_row[3], g_row[3], 500):
                rev_ok = False
                rev_detail.append(f"{key}.Total_Revenue {a_row[3]} vs {g_row[3]}")
            # Revenue_Share_Pct (col 7, tol 2) — CRITICAL
            if len(a_row) > 7 and len(g_row) > 7 and not num_close(a_row[7], g_row[7], 2):
                share_ok = False
                rev_detail.append(f"{key}.Revenue_Share_Pct {a_row[7]} vs {g_row[7]}")
        record("Segment Performance Total_Revenue per segment", rev_ok,
               "; ".join(rev_detail), critical=True)
        record("Segment Performance Revenue_Share_Pct per segment", share_ok,
               "; ".join(rev_detail), critical=True)

        # Non-critical secondary numbers
        sec_ok = True
        sec_detail = []
        for g_row in g_data:
            if not g_row or g_row[0] is None:
                continue
            a_row = a_lookup.get(seg_key(g_row[0]))
            if a_row is None:
                continue
            if len(a_row) > 1 and len(g_row) > 1 and not num_close(a_row[1], g_row[1], 5):
                sec_ok = False
                sec_detail.append(f"{seg_key(g_row[0])}.Customer_Count {a_row[1]} vs {g_row[1]}")
            if len(a_row) > 2 and len(g_row) > 2 and not num_close(a_row[2], g_row[2], 10):
                sec_ok = False
                sec_detail.append(f"{seg_key(g_row[0])}.Total_Orders {a_row[2]} vs {g_row[2]}")
            if len(a_row) > 4 and len(g_row) > 4 and not num_close(a_row[4], g_row[4], 5):
                sec_ok = False
                sec_detail.append(f"{seg_key(g_row[0])}.Avg_Order_Value {a_row[4]} vs {g_row[4]}")
        record("Segment Performance secondary metrics", sec_ok, "; ".join(sec_detail))

    # ---- Sheet 2: Segment by Region ----
    a_rows = load_sheet_rows(agent_wb, "Segment by Region")
    g_rows = load_sheet_rows(gt_wb, "Segment by Region")
    if a_rows is None:
        record("Sheet 'Segment by Region' present", False, "not found", critical=False)
    elif g_rows is None:
        record("Sheet 'Segment by Region' groundtruth", False, "missing", critical=True)
    else:
        record("Sheet 'Segment by Region' present", True)
        a_data = a_rows[1:] if len(a_rows) > 1 else []
        g_data = g_rows[1:] if len(g_rows) > 1 else []
        a_lookup = {}
        for r in a_data:
            if r and r[0] is not None and r[1] is not None:
                a_lookup[region_seg_key(r[0], r[1])] = r

        rev_ok = True
        prof_ok = True
        detail = []
        for g_row in g_data:
            if not g_row or g_row[0] is None:
                continue
            key = region_seg_key(g_row[0], g_row[1])
            a_row = a_lookup.get(key)
            if a_row is None:
                rev_ok = False
                prof_ok = False
                detail.append(f"missing {key}")
                continue
            # Revenue (col 4, tol 500)
            if len(a_row) > 4 and len(g_row) > 4 and not num_close(a_row[4], g_row[4], 500):
                rev_ok = False
                detail.append(f"{key}.Revenue {a_row[4]} vs {g_row[4]}")
            # Profitability_Index (col 6, tol 2) — CRITICAL (products-join formula)
            if len(a_row) > 6 and len(g_row) > 6 and not num_close(a_row[6], g_row[6], 2):
                prof_ok = False
                detail.append(f"{key}.Profitability_Index {a_row[6]} vs {g_row[6]}")
        record("Segment by Region Revenue per region+segment", rev_ok, "; ".join(detail))
        record("Segment by Region Profitability_Index per region+segment", prof_ok,
               "; ".join(detail), critical=True)

    # ---- Sheet 3: Strategic Matrix ----
    a_rows = load_sheet_rows(agent_wb, "Strategic Matrix")
    g_rows = load_sheet_rows(gt_wb, "Strategic Matrix")
    if a_rows is None:
        record("Sheet 'Strategic Matrix' present", False, "not found", critical=False)
    elif g_rows is None:
        record("Sheet 'Strategic Matrix' groundtruth", False, "missing", critical=True)
    else:
        record("Sheet 'Strategic Matrix' present", True)
        a_data = a_rows[1:] if len(a_rows) > 1 else []
        g_data = g_rows[1:] if len(g_rows) > 1 else []
        a_lookup = {seg_key(r[0]): r for r in a_data if r and r[0] is not None}

        cat_ok = True
        growth_ok = True
        contrib_ok = True
        detail = []
        for g_row in g_data:
            if not g_row or g_row[0] is None:
                continue
            key = seg_key(g_row[0])
            a_row = a_lookup.get(key)
            if a_row is None:
                cat_ok = False
                growth_ok = False
                detail.append(f"missing segment {g_row[0]}")
                continue
            # Revenue_Contribution_Pct (col 1, tol 2)
            if len(a_row) > 1 and len(g_row) > 1 and not num_close(a_row[1], g_row[1], 2):
                contrib_ok = False
                detail.append(f"{key}.Revenue_Contribution_Pct {a_row[1]} vs {g_row[1]}")
            # Growth_Indicator (col 2) — accept RU or EN
            if len(a_row) > 2 and len(g_row) > 2 and growth_norm(a_row[2]) != growth_norm(g_row[2]):
                growth_ok = False
                detail.append(f"{key}.Growth_Indicator {a_row[2]} vs {g_row[2]}")
            # Strategic_Category (col 4) — kept English
            if len(a_row) > 4 and len(g_row) > 4 and not category_match(a_row[4], g_row[4]):
                cat_ok = False
                detail.append(f"{key}.Strategic_Category {a_row[4]} vs {g_row[4]}")
        record("Strategic Matrix Strategic_Category per segment", cat_ok,
               "; ".join(detail), critical=True)
        record("Strategic Matrix Growth_Indicator per segment", growth_ok,
               "; ".join(detail), critical=True)
        record("Strategic Matrix Revenue_Contribution_Pct per segment", contrib_ok,
               "; ".join(detail))


def check_pptx(agent_ws):
    print("Checking PowerPoint presentation...")
    pptx_path = os.path.join(agent_ws, "QBR_Presentation.pptx")
    if not os.path.exists(pptx_path):
        record("QBR_Presentation.pptx exists", False, "not found", critical=True)
        return
    record("QBR_Presentation.pptx exists", True)

    from pptx import Presentation
    prs = Presentation(pptx_path)

    # slide count — NON-critical structural
    record("PPT has >=6 slides", len(prs.slides) >= 6,
           f"{len(prs.slides)} slides")

    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += " " + shape.text
    low = all_text.lower()

    # All four segments present (accept RU or EN) — CRITICAL substance
    missing = []
    for en, ru in SEG_RU.items():
        if en not in all_text and ru not in low:
            missing.append(en)
    record("PPT contains all four segment names", not missing,
           f"missing {missing}", critical=True)

    # At least two strategic categories present (kept English)
    cats_present = [c for c in ["Star", "Cash Cow", "Question Mark", "Dog"] if c in all_text]
    record("PPT contains strategic categories (>=2)", len(cats_present) >= 2,
           f"found {cats_present}", critical=True)

    # Consumer revenue figure
    rev_present = any(t in all_text for t in CONSUMER_REVENUE_TOKENS)
    record("PPT contains Consumer revenue figure (~839609)", rev_present)


def check_form():
    print("Checking survey form (forms MCP, gform.* schema)...")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
    except Exception as e:
        record("forms DB reachable", False, str(e), critical=True)
        return

    try:
        cur.execute("SELECT id FROM gform.forms WHERE LOWER(title) LIKE '%customer experience%'")
        rows = cur.fetchall()
        if not rows:
            record("Survey form 'Customer Experience Survey' exists", False,
                   "not found in gform.forms", critical=True)
            return
        record("Survey form 'Customer Experience Survey' exists", True)
        form_id = rows[0][0]

        cur.execute(
            "SELECT title, question_type FROM gform.questions WHERE form_id = %s",
            (form_id,),
        )
        qs = cur.fetchall()
        q_count = len(qs)
        # question count — NON-critical structural
        record("Survey has >=6 questions", q_count >= 6, f"{q_count} questions")

        # At least one MULTIPLE_CHOICE question — substance check
        types = [str(t).upper() for (_, t) in qs]
        has_mc = any("MULTIPLE" in t or "CHOICE" in t or "RADIO" in t for t in types)
        record("Survey has a multiple-choice (product brand) question", has_mc,
               f"types {types}", critical=True)
    except Exception as e:
        record("Survey form query", False, str(e), critical=True)
    finally:
        cur.close()
        conn.close()


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    agent_ws = args.agent_workspace or task_root
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    check_excel(agent_ws, gt_dir)
    check_pptx(agent_ws)
    check_form()

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100.0) if total else 0.0

    print("\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}, Failed: {FAIL_COUNT}, Accuracy: {accuracy:.1f}%")

    if CRITICAL_FAILED:
        print(f"  CRITICAL checks FAILED: {CRITICAL_FAILED}")
        print("  Overall: FAIL (critical)")
        sys.exit(1)

    overall = accuracy >= 70.0
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
