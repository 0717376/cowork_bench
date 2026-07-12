"""Evaluation for canvas-quiz-analysis-ppt.

Преамбула по русификации:
- canvas остаётся англоязычным (live demo-данные, не RU-форк), данные читаются "честно":
  eval сравнивает вывод агента с groundtruth_workspace с допусками, без хардкода
  волатильных баллов.
- Имена файлов, листов, заголовки столбцов и названия курсов с суффиксом (Fall 2014)
  — английские идентификаторы, по ним идёт substring/ключевое сопоставление.
- Текст слайдов агент может писать по-русски, поэтому substring-проверки PPT
  расширены до RU+EN. RU-ключевые слова ищутся в .lower() ОРИГИНАЛЬНОГО текста
  (без normalize()).
- CRITICAL_CHECKS: семантические проверки сути (агрегация по тестам, корректные
  сводные значения, корректное определение самого высокого/низкого курса, структура
  колоды). Любой провал критической проверки => немедленный FAIL (sys.exit(1))
  до порога точности. Порог точности по остальным проверкам: accuracy >= 70.
"""
import argparse
import os
import sys
import openpyxl
from pptx import Presentation


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    # Накопители проверок.
    #   checks: список (passed: bool, name: str) для порога точности
    #   critical_failures: список строк — провал любой => немедленный FAIL
    checks = []
    critical_failures = []
    all_errors = []

    def record(passed, name, critical=False):
        checks.append((bool(passed), name))
        if not passed:
            all_errors.append(name)
            if critical:
                critical_failures.append(name)

    # ---- Check Excel ----
    agent_excel = os.path.join(args.agent_workspace, "Quiz_Performance.xlsx")
    gt_excel = os.path.join(gt_dir, "Quiz_Performance.xlsx")

    agent_wb = None
    gt_wb = None
    if not os.path.exists(agent_excel):
        record(False, "Agent output Quiz_Performance.xlsx not found", critical=True)
    elif not os.path.exists(gt_excel):
        record(False, "Groundtruth Quiz_Performance.xlsx not found", critical=True)
    else:
        record(True, "Quiz_Performance.xlsx exists")
        agent_wb = openpyxl.load_workbook(agent_excel, data_only=True)
        gt_wb = openpyxl.load_workbook(gt_excel, data_only=True)

    # Check Quiz Details sheet (CRITICAL: ядро агрегации тестов из live Canvas)
    if agent_wb is not None and gt_wb is not None:
        print("  Checking Quiz Details...")
        a_rows = load_sheet_rows(agent_wb, "Quiz Details")
        g_rows = load_sheet_rows(gt_wb, "Quiz Details")
        if a_rows is None:
            record(False, "Sheet 'Quiz Details' not found in agent output", critical=True)
        elif g_rows is None:
            record(False, "Sheet 'Quiz Details' not found in groundtruth", critical=True)
        else:
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []
            record(len(a_data) == len(g_data),
                   f"Quiz Details row count: agent={len(a_data)}, expected={len(g_data)}")
            a_lookup = {}
            for row in a_data:
                if row and row[0] is not None and row[1] is not None:
                    key = (str(row[0]).strip().lower(), str(row[1]).strip().lower())
                    a_lookup[key] = row
            # CRITICAL: каждая строка (Course, Quiz_Title) groundtruth должна
            # присутствовать в выводе агента с Avg_Score в пределах допуска.
            missing_or_wrong = []
            for g_row in g_data:
                if not g_row or g_row[0] is None:
                    continue
                key = (str(g_row[0]).strip().lower(), str(g_row[1]).strip().lower())
                a_row = a_lookup.get(key)
                if a_row is None:
                    missing_or_wrong.append(f"missing {g_row[0]} / {g_row[1]}")
                    all_errors.append(f"Missing quiz row: {g_row[0]} / {g_row[1]}")
                    continue
                if len(a_row) > 2 and len(g_row) > 2:
                    ok = num_close(a_row[2], g_row[2], 5)
                    record(ok, f"{key}.Submissions: {a_row[2]} vs {g_row[2]}")
                if len(a_row) > 3 and len(g_row) > 3:
                    ok = num_close(a_row[3], g_row[3], 1.0)
                    if not ok:
                        missing_or_wrong.append(f"avg {key}")
                    record(ok, f"{key}.Avg_Score: {a_row[3]} vs {g_row[3]}")
                if len(a_row) > 4 and len(g_row) > 4:
                    record(num_close(a_row[4], g_row[4], 1.0),
                           f"{key}.Min_Score: {a_row[4]} vs {g_row[4]}")
                if len(a_row) > 5 and len(g_row) > 5:
                    record(num_close(a_row[5], g_row[5], 1.0),
                           f"{key}.Max_Score: {a_row[5]} vs {g_row[5]}")
            if missing_or_wrong:
                critical_failures.append(
                    "CRITICAL Quiz Details aggregation wrong: " + "; ".join(missing_or_wrong[:5]))
            else:
                print("    Quiz Details substance OK")

        # Check Course Summary sheet (CRITICAL: ключевой вычисляемый результат)
        print("  Checking Course Summary...")
        a_rows = load_sheet_rows(agent_wb, "Course Summary")
        g_rows = load_sheet_rows(gt_wb, "Course Summary")
        if a_rows is None:
            record(False, "Sheet 'Course Summary' not found in agent output", critical=True)
        elif g_rows is None:
            record(False, "Sheet 'Course Summary' not found in groundtruth", critical=True)
        else:
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []
            record(len(a_data) == len(g_data),
                   f"Course Summary row count: agent={len(a_data)}, expected={len(g_data)}")
            a_lookup = {}
            for row in a_data:
                if row and row[0] is not None:
                    a_lookup[str(row[0]).strip().lower()] = row
            summary_bad = []
            for g_row in g_data:
                if not g_row or g_row[0] is None:
                    continue
                key = str(g_row[0]).strip().lower()
                a_row = a_lookup.get(key)
                if a_row is None:
                    summary_bad.append(f"missing course {g_row[0]}")
                    all_errors.append(f"Missing course summary: {g_row[0]}")
                    continue
                if len(a_row) > 1 and len(g_row) > 1:
                    ok = num_close(a_row[1], g_row[1], 0)  # Total_Quizzes точно
                    if not ok:
                        summary_bad.append(f"Total_Quizzes {key}")
                    record(ok, f"{key}.Total_Quizzes: {a_row[1]} vs {g_row[1]}")
                if len(a_row) > 2 and len(g_row) > 2:
                    record(num_close(a_row[2], g_row[2], 10),
                           f"{key}.Total_Submissions: {a_row[2]} vs {g_row[2]}")
                if len(a_row) > 3 and len(g_row) > 3:
                    ok = num_close(a_row[3], g_row[3], 1.0)
                    if not ok:
                        summary_bad.append(f"Overall_Avg_Score {key}")
                    record(ok, f"{key}.Overall_Avg_Score: {a_row[3]} vs {g_row[3]}")
            if summary_bad:
                critical_failures.append(
                    "CRITICAL Course Summary wrong: " + "; ".join(summary_bad[:5]))
            else:
                print("    Course Summary substance OK")

    # ---- Check PowerPoint ----
    agent_ppt = os.path.join(args.agent_workspace, "Quiz_Report.pptx")
    if not os.path.exists(agent_ppt):
        record(False, "Agent output Quiz_Report.pptx not found", critical=True)
    else:
        record(True, "Quiz_Report.pptx exists")
        print("  Checking Quiz_Report.pptx...")
        prs = Presentation(agent_ppt)
        slides = list(prs.slides)

        # CRITICAL: титульный слайд + по слайду на курс + итоговый => >= 5 слайдов
        record(len(slides) >= 5,
               f"PPT has {len(slides)} slides, expected at least 5", critical=True)

        if len(slides) >= 1:
            # Титульный слайд: принимаем RU+EN.
            #   ('quiz' OR 'тест'/'викторин') AND ('fall 2014' OR 'осен'+2014)
            title_text = ""
            for shape in slides[0].shapes:
                if shape.has_text_frame:
                    title_text += shape.text_frame.text.lower() + " "
            has_quiz = ("quiz" in title_text or "тест" in title_text
                        or "викторин" in title_text)
            has_fall = ("fall 2014" in title_text
                        or ("осен" in title_text and "2014" in title_text))
            record(has_quiz and has_fall,
                   f"Title slide missing expected text (RU+EN). Found: {title_text[:100]}")

        if len(slides) >= 1:
            # Последний слайд: ключевые выводы. Принимаем RU+EN.
            last_text = ""
            for shape in slides[-1].shapes:
                if shape.has_text_frame:
                    last_text += shape.text_frame.text.lower() + " "
            takeaway_kw = ["takeaway", "key", "summary", "conclusion",
                           "вывод", "итог", "ключев", "заключ"]
            record(any(k in last_text for k in takeaway_kw),
                   "Last slide missing takeaways content (RU+EN)")

            # CRITICAL: корректное определение крайних значений.
            # Самый высокий = Глобальное управление и геополитика (85.87),
            # самый низкий = Креативные вычисления и культура (75.63).
            # Названия курсов берутся из русифицированного Canvas, поэтому ищем
            # дискриминирующие RU-подстроки названий курсов.
            has_global = "геополит" in last_text or "управлени" in last_text
            has_creative = "креативн" in last_text or "вычислен" in last_text
            if not (has_global and has_creative):
                critical_failures.append(
                    "CRITICAL: last slide must identify HIGHEST=Глобальное управление и геополитика "
                    "and LOWEST=Креативные вычисления и культура as extremes "
                    f"(global={has_global}, creative={has_creative})")
            record(has_global,
                   "Last slide should mention Глобальное управление и геополитика (highest avg)")
            record(has_creative,
                   "Last slide should mention Креативные вычисления и культура (lowest avg)")

        # CRITICAL: все 3 названия курсов присутствуют где-то в колоде.
        all_ppt_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_ppt_text += shape.text_frame.text.lower() + " "
        missing_courses = []
        # Дискриминирующие RU-подстроки трёх курсов Fall 2014 с тестами
        # (названия приходят из русифицированного Canvas).
        for course in ["креативн", "финанс", "геополит"]:
            present = course in all_ppt_text
            record(present, f"PPT missing course: {course}")
            if not present:
                missing_courses.append(course)
        if missing_courses:
            critical_failures.append(
                "CRITICAL: PPT missing course names: " + ", ".join(missing_courses))

    # ---- Итог: сначала критические проверки, затем порог точности ----
    total = len(checks)
    passed = sum(1 for ok, _ in checks if ok)
    accuracy = (passed / total * 100.0) if total else 0.0

    print(f"\nChecks passed: {passed}/{total} ({accuracy:.1f}%)")

    if critical_failures:
        print(f"\n=== RESULT: FAIL (critical) ===")
        for e in critical_failures[:15]:
            print(f"  CRITICAL: {e}")
        sys.exit(1)

    if accuracy < 70:
        print(f"\n=== RESULT: FAIL (accuracy {accuracy:.1f}% < 70%) ===")
        for e in all_errors[:15]:
            print(f"  {e}")
        sys.exit(1)

    print("\n=== RESULT: PASS ===")
    sys.exit(0)


if __name__ == "__main__":
    main()
