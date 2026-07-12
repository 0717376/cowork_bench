"""Evaluation for yf-portfolio-stress-test."""
import argparse, os, sys


def num_close(a, b, abs_tol=1.0, rel_tol=0.05):
    try:
        a_f, b_f = float(a), float(b)
        return abs(a_f - b_f) <= max(abs_tol, abs(b_f) * rel_tol)
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def check_excel(agent_workspace):
    """Returns (errors, critical_errors). critical_errors are a subset that must
    never fail regardless of overall accuracy."""
    errors = []
    critical = []
    import openpyxl
    path = os.path.join(agent_workspace, "Stress_Test.xlsx")
    if not os.path.exists(path):
        return (["Stress_Test.xlsx not found"], ["Stress_Test.xlsx not found"])
    try:
        wb = openpyxl.load_workbook(path, data_only=True)

        # Scenario Analysis
        rows = load_sheet_rows(wb, "Scenario Analysis")
        if rows is None:
            errors.append("Sheet 'Scenario Analysis' not found")
        else:
            data_rows = [r for r in rows[1:] if r and r[0] is not None]
            if len(data_rows) < 4:
                errors.append(f"Scenario Analysis has {len(data_rows)} rows, expected 4")
            scenarios = {str(r[0]).strip().lower(): r for r in data_rows if r[0]}
            # CRITICAL: Market Crash core values (price-independent, weights 25/25/20/15/10/5)
            if "market crash" in scenarios:
                r = scenarios["market crash"]
                if len(r) > 3 and not num_close(r[3], -27.75):
                    critical.append(f"Market Crash loss={r[3]}, expected ~-27.75")
                if len(r) > 4 and str(r[4]).strip().lower() != "yes":
                    critical.append(f"Market Crash Exceeds_Limit={r[4]}, expected Yes")
            else:
                critical.append("Market Crash scenario row missing in Scenario Analysis")
            # CRITICAL: two more scenarios value-checked to catch wrong shock-math
            if "recession" in scenarios:
                r = scenarios["recession"]
                if len(r) > 3 and not num_close(r[3], -18.0):
                    critical.append(f"Recession loss={r[3]}, expected ~-18.0")
                if len(r) > 4 and str(r[4]).strip().lower() != "no":
                    errors.append(f"Recession Exceeds_Limit={r[4]}, expected No")
            else:
                errors.append("Recession scenario row missing")
            if "rate hike" in scenarios:
                r = scenarios["rate hike"]
                if len(r) > 3 and not num_close(r[3], -9.35):
                    critical.append(f"Rate Hike loss={r[3]}, expected ~-9.35")
                if len(r) > 4 and str(r[4]).strip().lower() != "no":
                    errors.append(f"Rate Hike Exceeds_Limit={r[4]}, expected No")
            else:
                errors.append("Rate Hike scenario row missing")

        # Position Impact
        rows2 = load_sheet_rows(wb, "Position Impact")
        if rows2 is None:
            errors.append("Sheet 'Position Impact' not found")
        else:
            data_rows2 = [r for r in rows2[1:] if r and r[0] is not None]
            if len(data_rows2) < 24:
                errors.append(f"Position Impact has {len(data_rows2)} rows, expected 24")
            # CRITICAL: spot-check one Market Crash equity row so the 24 rows aren't placeholders.
            # Any equity position (weight w) under Market Crash equity shock -30: w * -30 / 100.
            # SBER.ME (25%) -> -7.5; LKOH.ME (20%) -> -6.0; MGNT.ME (10%) -> -3.0.
            expected_pos = {
                ("market crash", "sber.me"): 25 * -30 / 100,
                ("market crash", "lkoh.me"): 20 * -30 / 100,
                ("market crash", "mgnt.me"): 10 * -30 / 100,
                ("market crash", "gc=f"): 5 * 15 / 100,
            }
            pos_lookup = {}
            for r in data_rows2:
                if len(r) > 4 and r[0] and r[1]:
                    pos_lookup[(str(r[0]).strip().lower(), str(r[1]).strip().lower())] = r[4]
            spot_checked = 0
            for key, exp in expected_pos.items():
                if key in pos_lookup:
                    spot_checked += 1
                    if not num_close(pos_lookup[key], exp, abs_tol=0.5):
                        critical.append(
                            f"Position Impact {key} Position_Loss_Pct={pos_lookup[key]}, expected ~{exp}")
            if spot_checked == 0:
                critical.append("Position Impact: no recognizable Market Crash rows to spot-check (placeholder rows?)")

        # Risk Summary
        rows3 = load_sheet_rows(wb, "Risk Summary")
        if rows3 is None:
            errors.append("Sheet 'Risk Summary' not found")
        else:
            data_rows3 = [r for r in rows3[1:] if r and r[0] is not None]
            lookup = {str(r[0]).strip().lower(): r[1] for r in data_rows3 if r[0]}
            # CRITICAL: core summary metrics (worst case, count, rating, expected loss)
            if "worst_case_loss" in lookup:
                if not num_close(lookup["worst_case_loss"], 27.75):
                    critical.append(f"Worst_Case_Loss={lookup['worst_case_loss']}, expected ~27.75")
            else:
                critical.append("Worst_Case_Loss not found")
            if "scenarios_exceeding_limit" in lookup:
                if not num_close(lookup["scenarios_exceeding_limit"], 1, abs_tol=0):
                    critical.append(f"Scenarios_Exceeding_Limit={lookup['scenarios_exceeding_limit']}, expected 1")
            else:
                errors.append("Scenarios_Exceeding_Limit not found")
            if "expected_loss" in lookup:
                if not num_close(lookup["expected_loss"], -5.63):
                    critical.append(f"Expected_Loss={lookup['expected_loss']}, expected ~-5.63")
            else:
                critical.append("Expected_Loss not found")
            if "risk_rating" in lookup:
                if str(lookup["risk_rating"]).strip().lower() != "elevated":
                    critical.append(f"Risk_Rating={lookup['risk_rating']}, expected Elevated")
            else:
                errors.append("Risk_Rating not found")

    except Exception as e:
        errors.append(f"Error reading Excel: {e}")
    # critical errors are also overall errors
    return (errors + critical, critical)


def check_pptx(agent_workspace):
    errors = []
    path = os.path.join(agent_workspace, "Risk_Presentation.pptx")
    if not os.path.exists(path):
        return ["Risk_Presentation.pptx not found"]
    try:
        from pptx import Presentation
        prs = Presentation(path)
        if len(prs.slides) < 5:
            errors.append(f"Presentation has {len(prs.slides)} slides, expected at least 5")
        # Check title slide
        first_slide = prs.slides[0]
        title_text = ""
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                title_text += shape.text_frame.text.lower()
        if "stress test" not in title_text and "portfolio" not in title_text:
            errors.append("Title slide does not mention stress test or portfolio")
    except Exception as e:
        errors.append(f"Error reading PPTX: {e}")
    return errors


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    agent_ws = args.agent_workspace or os.path.join(os.path.dirname(__file__), "..", "groundtruth_workspace")

    # Each entry: (name, passed_bool). Accuracy = passed / total. >= 70% to pass,
    # but any CRITICAL failure forces FAIL regardless of accuracy.
    checks = []
    critical_errors = []

    print("  Checking Excel file...")
    excel_errs, excel_crit = check_excel(agent_ws)
    critical_errors.extend(excel_crit)
    checks.append(("excel", len(excel_errs) == 0))
    if excel_errs:
        for e in excel_errs[:8]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    print("  Checking PowerPoint...")
    pptx_errs = check_pptx(agent_ws)
    checks.append(("pptx", len(pptx_errs) == 0))
    if pptx_errs:
        for e in pptx_errs[:3]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    total = len(checks)
    passed = sum(1 for _, ok in checks if ok)
    accuracy = (passed / total * 100) if total else 0.0

    # CRITICAL gate: any critical failure => hard FAIL before accuracy is considered.
    if critical_errors:
        print(f"\n=== CRITICAL FAILURE ({len(critical_errors)}) ===")
        for e in critical_errors:
            print(f"  CRITICAL: {e}")
        print("=== RESULT: FAIL ===")
        sys.exit(1)

    print(f"\nAccuracy: {accuracy:.1f}% ({passed}/{total})")
    if accuracy >= 70:
        print("=== RESULT: PASS ===")
        sys.exit(0)
    else:
        print("=== RESULT: FAIL ===")
        sys.exit(1)


if __name__ == "__main__":
    main()
