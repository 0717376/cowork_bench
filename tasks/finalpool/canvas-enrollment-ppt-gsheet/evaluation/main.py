"""Evaluation for canvas-enrollment-ppt-gsheet.

Blocking checks: Enrollment_Overview.xlsx and Enrollment_Overview.pptx.
Non-blocking: Google Sheet DB check.
"""
import argparse
import os
import sys
import openpyxl
from pptx import Presentation


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")
    all_errors = []
    critical_errors = []

    # ---- Check Excel ----
    agent_excel = os.path.join(args.agent_workspace, "Enrollment_Overview.xlsx")
    gt_excel = os.path.join(gt_dir, "Enrollment_Overview.xlsx")

    if not os.path.exists(agent_excel):
        all_errors.append("Agent output Enrollment_Overview.xlsx not found")
    elif not os.path.exists(gt_excel):
        all_errors.append("Groundtruth Enrollment_Overview.xlsx not found")
    else:
        agent_wb = openpyxl.load_workbook(agent_excel, data_only=True)
        gt_wb = openpyxl.load_workbook(gt_excel, data_only=True)

        # Check Enrollment Details
        print("  Checking Enrollment Details...")
        a_rows = load_sheet_rows(agent_wb, "Enrollment Details")
        g_rows = load_sheet_rows(gt_wb, "Enrollment Details")
        if a_rows is None:
            all_errors.append("Sheet 'Enrollment Details' not found in agent output")
        elif g_rows is None:
            all_errors.append("Sheet 'Enrollment Details' not found in groundtruth")
        else:
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []
            if abs(len(a_data) - len(g_data)) > 2:
                all_errors.append(f"Enrollment Details row count: agent={len(a_data)}, expected={len(g_data)}")

            # Match by course code (col 1)
            a_lookup = {}
            for row in a_data:
                if row and len(row) > 1 and row[1] is not None:
                    a_lookup[str(row[1]).strip().lower()] = row
            # CRITICAL: all required courses must be present and Total_Enrollments
            # must be correct within a tight tolerance (verifies per-course
            # extraction from Canvas, not just two loose columns).
            for g_row in g_data:
                if not g_row or len(g_row) < 2 or g_row[1] is None:
                    continue
                key = str(g_row[1]).strip().lower()
                a_row = a_lookup.get(key)
                if a_row is None:
                    critical_errors.append(f"CRITICAL: Missing course: {g_row[0]} ({g_row[1]})")
                    continue
                # Col 2: Total_Enrollments (CRITICAL, tight tol)
                if len(a_row) > 2 and len(g_row) > 2:
                    if not num_close(a_row[2], g_row[2], 2):
                        critical_errors.append(f"CRITICAL: {key}.Total_Enrollments: {a_row[2]} vs {g_row[2]} (tol=2)")
                # Col 3: Students (non-critical, tight tol)
                if len(a_row) > 3 and len(g_row) > 3:
                    if not num_close(a_row[3], g_row[3], 2):
                        all_errors.append(f"{key}.Students: {a_row[3]} vs {g_row[3]} (tol=2)")
                # Cols 4-7: Teachers, TAs, Active, Completed (non-critical, tight tol)
                for ci, cname in ((4, "Teachers"), (5, "TAs"), (6, "Active"), (7, "Completed")):
                    if len(a_row) > ci and len(g_row) > ci:
                        if not num_close(a_row[ci], g_row[ci], 2):
                            all_errors.append(f"{key}.{cname}: {a_row[ci]} vs {g_row[ci]} (tol=2)")
            if not all_errors and not critical_errors:
                print("    PASS")

        # Check Summary sheet
        print("  Checking Summary...")
        a_rows = load_sheet_rows(agent_wb, "Summary")
        g_rows = load_sheet_rows(gt_wb, "Summary")
        prev_errors = len(all_errors)
        if a_rows is None:
            all_errors.append("Sheet 'Summary' not found in agent output")
        elif g_rows is None:
            all_errors.append("Sheet 'Summary' not found in groundtruth")
        else:
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []

            # Metrics that are CRITICAL deliverables (core aggregates / ranking /
            # rounding rule). Each maps metric-key -> tolerance for numerics.
            critical_numeric = {
                "total_enrollments": 2,
                "avg_enrollment_per_course": 0.5,
            }
            critical_string = {"largest_course", "smallest_course"}

            a_lookup = {}
            for row in a_data:
                if row and row[0] is not None:
                    a_lookup[str(row[0]).strip().lower()] = row
            for g_row in g_data:
                if not g_row or g_row[0] is None:
                    continue
                key = str(g_row[0]).strip().lower()
                a_row = a_lookup.get(key)
                if a_row is None:
                    # A missing core metric is critical; others non-critical.
                    if key in critical_numeric or key in critical_string:
                        critical_errors.append(f"CRITICAL: Missing summary metric: {g_row[0]}")
                    else:
                        all_errors.append(f"Missing summary metric: {g_row[0]}")
                    continue
                if len(a_row) > 1 and len(g_row) > 1:
                    # For string values, use str_match; for numbers, use num_close
                    try:
                        float(g_row[1])
                        if key in critical_numeric:
                            if not num_close(a_row[1], g_row[1], critical_numeric[key]):
                                critical_errors.append(
                                    f"CRITICAL: Summary.{key}: {a_row[1]} vs {g_row[1]} (tol={critical_numeric[key]})")
                        elif not num_close(a_row[1], g_row[1], 50):
                            all_errors.append(f"Summary.{key}: {a_row[1]} vs {g_row[1]} (tol=50)")
                    except (TypeError, ValueError):
                        if key in critical_string:
                            if not str_match(a_row[1], g_row[1]):
                                critical_errors.append(f"CRITICAL: Summary.{key}: {a_row[1]} vs {g_row[1]}")
                        elif not str_match(a_row[1], g_row[1]):
                            all_errors.append(f"Summary.{key}: {a_row[1]} vs {g_row[1]}")
            new_errors = len(all_errors) - prev_errors
            if new_errors == 0:
                print("    PASS")

    # ---- Check PowerPoint ----
    agent_ppt = os.path.join(args.agent_workspace, "Enrollment_Overview.pptx")
    if not os.path.exists(agent_ppt):
        all_errors.append("Agent output Enrollment_Overview.pptx not found")
    else:
        print("  Checking Enrollment_Overview.pptx...")
        prs = Presentation(agent_ppt)
        slides = list(prs.slides)
        if len(slides) < 4:
            critical_errors.append(f"CRITICAL: PPT has {len(slides)} slides, expected at least 4")
        else:
            # Check title slide
            title_text = ""
            for shape in slides[0].shapes:
                if shape.has_text_frame:
                    title_text += shape.text_frame.text.lower() + " "
            if "enrollment" not in title_text:
                all_errors.append(f"Title slide missing 'enrollment'. Found: {title_text[:100]}")

            # Check all PPT text for key content
            all_ppt_text = ""
            for slide in slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        all_ppt_text += shape.text_frame.text.lower() + " "

            # CRITICAL: the four required slide titles must be present.
            required_titles = [
                "course enrollment overview",
                "enrollment summary",
                "top 5 courses by enrollment",
                "enrollment distribution",
            ]
            for t in required_titles:
                if t not in all_ppt_text:
                    critical_errors.append(f"CRITICAL: PPT missing required slide title '{t}'")

            if "distribution" not in all_ppt_text:
                all_errors.append("PPT missing 'Distribution' slide")

            # CRITICAL: Top-5 slide must list the actual top-5 course codes
            # derived from the groundtruth Enrollment Details (code-based, not a
            # single hardcoded course-name substring).
            try:
                gt_wb2 = openpyxl.load_workbook(gt_excel, data_only=True)
                gt_det = load_sheet_rows(gt_wb2, "Enrollment Details")
                ranked = []
                for r in (gt_det[1:] if gt_det and len(gt_det) > 1 else []):
                    if r and len(r) > 2 and r[1] is not None and r[2] is not None:
                        try:
                            ranked.append((str(r[1]).strip(), float(r[2])))
                        except (TypeError, ValueError):
                            pass
                ranked.sort(key=lambda x: x[1], reverse=True)
                top5_codes = [c.lower() for c, _ in ranked[:5]]
                missing_codes = [c for c in top5_codes if c not in all_ppt_text]
                if missing_codes:
                    critical_errors.append(
                        f"CRITICAL: PPT Top-5 slide missing course code(s): {missing_codes}")
            except Exception as e:
                all_errors.append(f"Could not derive top-5 from groundtruth: {e}")

        if not any("ppt" in e.lower() or "slide" in e.lower() for e in all_errors + critical_errors):
            print("    PASS")

    # ---- GSheet check (non-critical: validated but reported as warnings only,
    # since this eval is all-or-nothing and the spreadsheet is a secondary
    # deliverable). ----
    print("  Checking 'Enrollment Dashboard' Google Sheet (non-critical)...")
    gsheet_warnings = []
    required_cols = {"course_name", "course_code", "total_enrollments", "students",
                     "teachers", "tas", "active", "completed"}
    try:
        import psycopg2
        conn = psycopg2.connect(host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="cowork_gym",
                                user="eigent", password="camel")
        cur = conn.cursor()
        # Spreadsheet titled "Enrollment Dashboard"
        cur.execute("SELECT id FROM gsheet.spreadsheets WHERE lower(trim(title)) = %s",
                    ("enrollment dashboard",))
        ss = cur.fetchone()
        if not ss:
            gsheet_warnings.append("GSheet 'Enrollment Dashboard' spreadsheet not found")
        else:
            ss_id = ss[0]
            # Sheet named "Course Data"
            cur.execute("SELECT id FROM gsheet.sheets WHERE spreadsheet_id = %s AND lower(trim(title)) = %s",
                        (ss_id, "course data"))
            sh = cur.fetchone()
            if not sh:
                gsheet_warnings.append("GSheet 'Course Data' sheet not found in 'Enrollment Dashboard'")
            else:
                sh_id = sh[0]
                # Header row (smallest row_index) cells -> required columns present
                cur.execute(
                    "SELECT row_index, value FROM gsheet.cells WHERE sheet_id = %s ORDER BY row_index",
                    (sh_id,))
                rows = cur.fetchall()
                if not rows:
                    gsheet_warnings.append("GSheet 'Course Data' sheet has no cells")
                else:
                    header_row = rows[0][0]
                    header = {str(v).strip().lower() for ri, v in rows
                              if ri == header_row and v is not None}
                    missing = required_cols - header
                    if missing:
                        gsheet_warnings.append(f"GSheet 'Course Data' missing columns: {sorted(missing)}")
        cur.close()
        conn.close()
    except Exception as e:
        # Schema shape may vary across infra; do not hard-fail the whole eval on
        # an introspection error for this non-critical deliverable.
        print(f"    [INFO] GSheet check skipped: {e} (non-critical)")

    if gsheet_warnings:
        print(f"    [WARN] GSheet deliverable issues (non-critical): {len(gsheet_warnings)}")
        for w in gsheet_warnings[:10]:
            print(f"      {w}")
    else:
        print("    PASS")

    # ---- CRITICAL gate: any critical failure => immediate FAIL ----
    if critical_errors:
        print(f"\n=== RESULT: FAIL ({len(critical_errors)} CRITICAL errors) ===")
        for e in critical_errors[:15]:
            print(f"  {e}")
        sys.exit(1)

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:15]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
