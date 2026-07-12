"""Evaluation для portfolio-diversification-analysis-engine (RU / moex-finance).

Агент строит отчёт о диверсификации портфеля по данным MCP `moex-finance`
(схема moex.*) и количеству позиций из holdings.csv:
  - Word  `Portfolio_Diversification_Report.docx`
  - PPTX  `Portfolio_Diversification_Deck.pptx`
  - письмо на investment-committee@capital-invest.ru с темой 'Portfolio Diversification Report'
  - событие календаря 'Strategy Review' 15.06.2026 11:00-12:00

Эталонные значения (последние цены закрытия) читаются ЧЕСТНО на лету из moex.*,
ничего не захардкожено — поэтому проверки устойчивы к ресиду цен.

CRITICAL_CHECKS (семантика): любой провал => немедленный FAIL (sys.exit(1))
ДО гейта по accuracy. Структурные проверки (наличие файлов, разделов, слайдов) —
не критические, идут в общий процент.

PASS: нет ни одного CRITICAL-фейла И accuracy >= 70%.
"""
import argparse
import csv
import json
import os
import sys

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent", "password": "camel",
}

TASK_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

TRACKED_TICKERS = ["SBER.ME", "GAZP.ME", "LKOH.ME", "TCSG.ME", "MGNT.ME", "MTSS.ME"]

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []

CRITICAL_CHECKS = {
    "CRITICAL: отчёт покрывает все 6 тикеров MOEX (.ME)",
    "CRITICAL: стоимости позиций ~= цена*количество по каждому тикеру",
    "CRITICAL: общая стоимость портфеля Total_Portfolio_Value совпадает с расчётом",
    "CRITICAL: письмо с темой 'Portfolio Diversification Report' на investment-committee@capital-invest.ru",
    "CRITICAL: событие 'Strategy Review' 15.06.2026 11:00-12:00",
}


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        marker = " [CRITICAL]" if name in CRITICAL_CHECKS else ""
        print(f"  [FAIL]{marker} {name}: {str(detail)[:300]}")
        if name in CRITICAL_CHECKS:
            CRITICAL_FAILS.append(name)


def num_close(a, b, rel=0.05, abs_tol=1.0):
    """Числовое сравнение с допуском (5% относительный или abs_tol)."""
    try:
        a, b = float(a), float(b)
    except (TypeError, ValueError):
        return False
    return abs(a - b) <= max(abs_tol, abs(b) * rel)


def get_conn():
    import psycopg2
    return psycopg2.connect(**DB_CONFIG)


def get_latest_prices():
    conn = get_conn()
    cur = conn.cursor()
    prices = {}
    for sym in TRACKED_TICKERS:
        cur.execute(
            "SELECT close FROM moex.stock_prices WHERE symbol = %s "
            "ORDER BY date DESC LIMIT 1", (sym,))
        row = cur.fetchone()
        if row and row[0] is not None:
            prices[sym] = float(row[0])
    cur.close()
    conn.close()
    return prices


def load_holdings():
    """Количество позиций из исходного holdings.csv (источник, не ответ)."""
    path = os.path.join(TASK_ROOT, "initial_workspace", "holdings.csv")
    qty = {}
    with open(path, newline="") as f:
        for row in csv.DictReader(f):
            qty[row["symbol"].strip().upper()] = float(row["quantity"])
    return qty


def to_addresses(to_addr):
    if isinstance(to_addr, list):
        return " ".join(str(r).lower() for r in to_addr)
    if to_addr:
        try:
            parsed = json.loads(str(to_addr))
            if isinstance(parsed, list):
                return " ".join(str(r).lower() for r in parsed)
        except Exception:
            pass
        return str(to_addr).lower()
    return ""


def check_word(agent_ws, prices, qty, total_value):
    print("\n=== Проверка 1: Portfolio_Diversification_Report.docx ===")
    path = os.path.join(agent_ws, "Portfolio_Diversification_Report.docx")
    check("Файл Portfolio_Diversification_Report.docx существует", os.path.isfile(path))
    if not os.path.isfile(path):
        check("CRITICAL: отчёт покрывает все 6 тикеров MOEX (.ME)", False, "нет файла")
        check("CRITICAL: стоимости позиций ~= цена*количество по каждому тикеру", False, "нет файла")
        check("CRITICAL: общая стоимость портфеля Total_Portfolio_Value совпадает с расчётом",
              False, "нет файла")
        return

    try:
        from docx import Document
        doc = Document(path)
    except Exception as e:
        check("Word-файл читается", False, str(e))
        check("CRITICAL: отчёт покрывает все 6 тикеров MOEX (.ME)", False, "файл не читается")
        check("CRITICAL: стоимости позиций ~= цена*количество по каждому тикеру", False, "файл не читается")
        check("CRITICAL: общая стоимость портфеля Total_Portfolio_Value совпадает с расчётом",
              False, "файл не читается")
        return

    paras = [p.text for p in doc.paragraphs]
    table_text = []
    numbers = []
    for t in doc.tables:
        for r in t.rows:
            for c in r.cells:
                table_text.append(c.text)
    full_text = " ".join(paras + table_text)
    low = full_text.lower()

    # Собираем все числа из таблиц (для проверки стоимостей)
    import re
    for cell in table_text + paras:
        for m in re.findall(r"-?\d[\d\s.,]*", cell):
            cleaned = m.replace(" ", "").replace(",", "")
            try:
                numbers.append(float(cleaned))
            except ValueError:
                pass

    # Структурные: заголовок и разделы (RU+EN)
    check("Заголовок содержит 'портфел' и 'диверсификац'",
          "портфел" in low and "диверсификац" in low, full_text[:120])
    check("Есть раздел 'Состав портфеля'/'Holdings'",
          "состав портфел" in low or "holdings" in low)
    check("Есть раздел 'Анализ диверсификации'/'Diversification'",
          "диверсификац" in low or "diversification" in low)
    check("Есть раздел 'Рекомендации'/'Recommendations'",
          "рекомендац" in low or "recommendation" in low)

    # CRITICAL: все 6 тикеров присутствуют
    missing = [t for t in TRACKED_TICKERS if t not in full_text]
    check("CRITICAL: отчёт покрывает все 6 тикеров MOEX (.ME)",
          not missing, f"отсутствуют: {missing}")

    # CRITICAL: стоимость каждой позиции ~= цена*количество (для тикеров с ценой)
    bad = []
    checked = 0
    for sym in TRACKED_TICKERS:
        if sym not in prices or sym not in qty:
            continue
        expected = prices[sym] * qty[sym]
        if any(num_close(n, expected, rel=0.05, abs_tol=max(1.0, expected * 0.02))
               for n in numbers):
            checked += 1
        else:
            bad.append(f"{sym}: ожидалось ~{expected:.0f}")
    check("CRITICAL: стоимости позиций ~= цена*количество по каждому тикеру",
          checked == len([s for s in TRACKED_TICKERS if s in prices and s in qty]) and not bad,
          f"не найдены: {bad}")

    # CRITICAL: общая стоимость портфеля
    found_total = any(num_close(n, total_value, rel=0.03,
                                abs_tol=max(1.0, total_value * 0.01))
                      for n in numbers)
    check("CRITICAL: общая стоимость портфеля Total_Portfolio_Value совпадает с расчётом",
          found_total, f"ожидалось ~{total_value:.0f}")


def check_pptx(agent_ws):
    print("\n=== Проверка 2: Portfolio_Diversification_Deck.pptx ===")
    path = os.path.join(agent_ws, "Portfolio_Diversification_Deck.pptx")
    check("Файл Portfolio_Diversification_Deck.pptx существует", os.path.isfile(path))
    if not os.path.isfile(path):
        return
    try:
        from pptx import Presentation
        prs = Presentation(path)
    except Exception as e:
        check("PPTX-файл читается", False, str(e))
        return
    slides = list(prs.slides)
    check("Презентация содержит >= 3 слайдов", len(slides) >= 3, f"слайдов: {len(slides)}")
    texts = []

    def _collect(shp):
        # Текст-фреймы, ячейки таблиц и группы (рекурсивно).
        if shp.has_text_frame:
            texts.append(shp.text_frame.text)
        if getattr(shp, "has_table", False):
            texts.extend(c.text for row in shp.table.rows for c in row.cells)
        for sub in getattr(shp, "shapes", ()):
            _collect(sub)

    for s in slides:
        for shp in s.shapes:
            _collect(shp)
    full = " ".join(texts)
    n_tickers = sum(1 for t in TRACKED_TICKERS if t in full)
    check("Слайды упоминают тикеры портфеля (>=4 из 6)", n_tickers >= 4,
          f"найдено тикеров: {n_tickers}")


def check_email():
    print("\n=== Проверка 3: письмо ===")
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute("SELECT subject, to_addr FROM email.messages")
        emails = cur.fetchall()
        conn.close()
    except Exception as e:
        check("CRITICAL: письмо с темой 'Portfolio Diversification Report' на investment-committee@capital-invest.ru",
              False, str(e))
        return
    exact = [m for m in emails
             if m[0] and str(m[0]).strip() == "Portfolio Diversification Report"
             and "investment-committee@capital-invest.ru" in to_addresses(m[1])]
    check("CRITICAL: письмо с темой 'Portfolio Diversification Report' на investment-committee@capital-invest.ru",
          len(exact) >= 1, f"subjects={[m[0] for m in emails]}")


def check_calendar():
    print("\n=== Проверка 4: событие календаря ===")
    try:
        conn = get_conn()
        cur = conn.cursor()
        # «11:00» для московского комитета — это локальное (МСК) время; task.md
        # не указывает UTC. timestamptz рендерится в timezone сессии PG, поэтому
        # нормализуем дату/часы к Europe/Moscow явно (08:00 UTC == 11:00 МСК).
        cur.execute("""
            SELECT summary,
                   (start_datetime AT TIME ZONE 'Europe/Moscow')::date AS start_date,
                   EXTRACT(HOUR   FROM start_datetime AT TIME ZONE 'Europe/Moscow') AS start_hour,
                   EXTRACT(MINUTE FROM start_datetime AT TIME ZONE 'Europe/Moscow') AS start_minute,
                   EXTRACT(HOUR   FROM end_datetime   AT TIME ZONE 'Europe/Moscow') AS end_hour,
                   EXTRACT(MINUTE FROM end_datetime   AT TIME ZONE 'Europe/Moscow') AS end_minute,
                   start_datetime, end_datetime
            FROM gcal.events
            WHERE summary ILIKE %s
        """, ('%Strategy Review%',))
        events = cur.fetchall()
        conn.close()
    except Exception as e:
        check("CRITICAL: событие 'Strategy Review' 15.06.2026 11:00-12:00", False, str(e))
        return
    ev_ok = False
    for (summ, start_date, start_hour, start_minute,
         end_hour, end_minute, raw_start, raw_end) in events:
        if not (summ and str(summ).strip() == "Strategy Review"):
            continue
        if start_date is None or start_hour is None or end_hour is None:
            continue
        if (str(start_date) == "2026-06-15"
                and int(start_hour) == 11 and int(start_minute) == 0
                and int(end_hour) == 12 and int(end_minute) == 0):
            ev_ok = True
            break
    check("CRITICAL: событие 'Strategy Review' 15.06.2026 11:00-12:00",
          ev_ok, f"events={[(e[0], str(e[6]), str(e[7])) for e in events]}")


def run_evaluation(agent_workspace):
    prices = get_latest_prices()
    qty = load_holdings()
    total_value = sum(prices[s] * qty[s] for s in TRACKED_TICKERS if s in prices and s in qty)
    print(f"Последние цены из moex.*: {prices}")
    print(f"Расчётная общая стоимость портфеля: {total_value:.2f}")

    if not os.path.isdir(agent_workspace):
        print(f"FAIL: рабочая директория агента не найдена: {agent_workspace}")
        sys.exit(1)

    check_word(agent_workspace, prices, qty, total_value)
    check_pptx(agent_workspace)
    check_email()
    check_calendar()

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100) if total else 0.0
    print(f"\nИтого: {PASS_COUNT}/{total} проверок пройдено ({accuracy:.1f}%)")

    if CRITICAL_FAILS:
        print(f"=== RESULT: FAIL (провалены критические проверки: {CRITICAL_FAILS}) ===")
        sys.exit(1)

    if accuracy >= 70:
        print("=== RESULT: PASS ===")
        sys.exit(0)
    print("=== RESULT: FAIL (accuracy < 70%) ===")
    sys.exit(1)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    run_evaluation(args.agent_workspace)


if __name__ == "__main__":
    main()
