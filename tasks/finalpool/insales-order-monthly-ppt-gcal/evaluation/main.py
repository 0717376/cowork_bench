"""Evaluation for insales-order-monthly-ppt-gcal (InSales / wc.* schema)."""
import argparse
import json
import os
import sys

import psycopg2


def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "cowork_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []  # имена провалившихся обязательных проверок

# Пары названий месяцев EN/RU (агент может писать по-русски).
MONTH_NAME_PAIRS = [
    ("january", "январ"), ("february", "феврал"), ("march", "март"),
    ("april", "апрел"), ("may", "ма"), ("june", "июн"),
    ("july", "июл"), ("august", "август"), ("september", "сентябр"),
    ("october", "октябр"), ("november", "ноябр"), ("december", "декабр"),
]

# Фактические значения выручки из БД (status-agnostic агрегация по wc.*).
# Группировка по date_created в UTC — ровно так, как woo_orders_list
# сериализует timestamptz (node-pg отдаёт UTC ISO, '...23:02:37Z').
# date_trunc('month', date_created AT TIME ZONE 'UTC'). Сверено с db/init.sql.gz.
EXPECTED_REVENUES = {
    "january": 5938.09, "february": 5615.02, "march": 2089.69,
    "april": 4349.71, "may": 6023.29, "june": 9555.69,
    "july": 3921.19, "august": 3873.37, "september": 5842.43,
    "october": 4447.76, "november": 1212.05, "december": 8843.75,
}
TOTAL_REVENUE = round(sum(EXPECTED_REVENUES.values()), 2)  # ≈ 61712.04

# Чеки, провал которых означает содержательное невыполнение задачи.
# Любой такой провал => итог FAIL независимо от accuracy.
CRITICAL_CHECKS = {
    "PPT: пиковый месяц июнь/June ~9555 + ключевое слово пик/peak/max",
    "PPT/Email: суммарная годовая выручка ~61712",
    "PPT: минимальный месяц ноябрь/November ~1212",
    "Calendar: событие 'Monthly Sales Review'/'Обзор продаж' на 2026-04-15 14:00",
    "Email: доставлен на ecommerce_manager@store.com с показателем выручки в теле",
}


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {str(detail)[:300]}" if detail else ""
        marker = " [CRITICAL]" if name in CRITICAL_CHECKS else ""
        print(f"  [FAIL]{marker} {name}{msg}")
        if name in CRITICAL_CHECKS:
            CRITICAL_FAILS.append(name)


def has_revenue_number(text, value):
    """Принять число выручки в форматах '9555', '9,555', '9555.69', '9 555' и т.п."""
    digits = str(int(round(value)))  # напр. '9555'
    variants = {
        digits,
        f"{int(round(value)):,}",            # 9,555
        f"{int(round(value)):,}".replace(",", " "),  # 9 555
        f"{value:.2f}",
        f"{value:.2f}".replace(".", ","),
    }
    # также первые 4 значащие цифры без округления, на случай '9555'
    variants.add(str(int(value)))
    return any(v in text for v in variants if v)


def check_ppt(agent_workspace):
    print("\n=== Checking PowerPoint ===")
    pptx_path = os.path.join(agent_workspace, "Monthly_Sales_Review.pptx")
    if not os.path.isfile(pptx_path):
        check("Monthly_Sales_Review.pptx exists", False, f"Not found: {pptx_path}")
        return ""
    check("Monthly_Sales_Review.pptx exists", True)

    all_text = ""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        check("PPT has at least 4 slides", slide_count >= 4, f"Found {slide_count} slides")

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        all_text += para.text + " "
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            all_text += cell.text + " "
        low = all_text.lower()

        # Заголовок (EN или RU)
        check("PPT contains title (monthly sales / обзор продаж / ежемесячн)",
              ("monthly sales" in low or "обзор продаж" in low
               or "ежемесячн" in low or "продаж" in low),
              f"Sample: {all_text[:300]}")

        # Названия месяцев (EN или RU основа)
        months_found = [en for en, ru in MONTH_NAME_PAIRS
                        if en in low or ru in low]
        check("PPT contains at least 6 month names (EN/RU)",
              len(months_found) >= 6,
              f"Found months: {months_found}")

        # Ключевые выводы (EN или RU)
        check("PPT contains key insights (insight/peak/revenue / вывод/пик/выручк/доход)",
              any(k in low for k in
                  ["insight", "peak", "revenue", "вывод", "пик", "выручк", "доход"]),
              f"Sample: {all_text[:300]}")

    except ImportError:
        check("PPT file has content", os.path.getsize(pptx_path) > 5000,
              f"Size: {os.path.getsize(pptx_path)}")
        return ""
    except Exception as e:
        check("PPT readable", False, str(e))
        return ""

    low = all_text.lower()

    # CRITICAL: пиковый месяц — июнь ~9555 + ключевое слово пик/peak/max
    june_ok = ("june" in low or "июн" in low)
    june_val_ok = has_revenue_number(all_text, EXPECTED_REVENUES["june"])
    peak_kw = any(k in low for k in ["peak", "пик", "max", "макс", "наибольш",
                                     "highest", "наивысш"])
    check("PPT: пиковый месяц июнь/June ~9555 + ключевое слово пик/peak/max",
          june_ok and june_val_ok and peak_kw,
          f"june={june_ok} val={june_val_ok} peak_kw={peak_kw}")

    # CRITICAL: минимальный месяц — ноябрь ~1212
    nov_ok = ("november" in low or "ноябр" in low)
    nov_val_ok = has_revenue_number(all_text, EXPECTED_REVENUES["november"])
    check("PPT: минимальный месяц ноябрь/November ~1212",
          nov_ok and nov_val_ok,
          f"nov={nov_ok} val={nov_val_ok}")

    return all_text


def check_total_revenue(ppt_text, email_body):
    """CRITICAL: суммарная годовая выручка ~61712 в PPT ИЛИ теле письма."""
    combined = (ppt_text or "") + " " + (email_body or "")
    ok = has_revenue_number(combined, TOTAL_REVENUE) or "61712" in combined or "61,712" in combined
    check("PPT/Email: суммарная годовая выручка ~61712", ok,
          f"Expected ~{TOTAL_REVENUE}")


def check_calendar():
    print("\n=== Checking Google Calendar ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT summary, start_datetime, end_datetime
            FROM gcal.events
            WHERE LOWER(summary) LIKE '%monthly sales%'
               OR LOWER(summary) LIKE '%sales review%'
               OR LOWER(summary) LIKE '%обзор продаж%'
               OR LOWER(summary) LIKE '%продаж%'
        """)
        events = cur.fetchall()

        # Найти событие на 2026-04-15 со стартом 14:00
        matched = None
        for ev in events:
            start = str(ev[1]) if ev[1] else ""
            if "2026-04-15" in start:
                matched = ev
                break

        check("Calendar: событие 'Monthly Sales Review'/'Обзор продаж' на 2026-04-15 14:00",
              matched is not None and ("14:00" in str(matched[1])),
              f"Found {len(events)} events; matched={matched}")

        cur.close()
        conn.close()
    except Exception as e:
        check("Calendar: событие 'Monthly Sales Review'/'Обзор продаж' на 2026-04-15 14:00",
              False, str(e))


def check_email():
    print("\n=== Checking Email ===")
    body = ""
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT id, subject, to_addr, body_text
            FROM email.messages
            WHERE to_addr::text ILIKE '%ecommerce_manager@store.com%'
        """)
        emails = cur.fetchall()

        email = emails[0] if emails else None
        subject = str(email[1]).lower() if email and email[1] else ""
        body = str(email[3]) if email and email[3] else ""

        # NON-critical: тема содержит 'sales'/'продаж'/'обзор'
        check("Email subject contains sales/продаж/обзор",
              bool(email) and ("sales" in subject or "продаж" in subject or "обзор" in subject),
              f"Subject: {email[1] if email else None}")

        # NON-critical: тело нетривиальной длины
        check("Email body has content", len(body) > 30,
              f"Body length: {len(body)}")

        # CRITICAL: доставлен адресату И в теле есть показатель выручки (9555 или 61712)
        rev_in_body = (has_revenue_number(body, EXPECTED_REVENUES["june"])
                       or has_revenue_number(body, TOTAL_REVENUE)
                       or "9555" in body or "61712" in body
                       or "9,555" in body or "61,712" in body)
        check("Email: доставлен на ecommerce_manager@store.com с показателем выручки в теле",
              bool(email) and rev_in_body,
              f"delivered={bool(email)} rev_in_body={rev_in_body}")

        cur.close()
        conn.close()
    except Exception as e:
        check("Email: доставлен на ecommerce_manager@store.com с показателем выручки в теле",
              False, str(e))
    return body


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=True)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    ppt_text = check_ppt(args.agent_workspace)
    check_calendar()
    email_body = check_email()
    check_total_revenue(ppt_text, email_body)

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\n=== Results: {PASS_COUNT}/{total} passed ({accuracy:.1f}%) ===")
    if CRITICAL_FAILS:
        print(f"Critical fails: {CRITICAL_FAILS}")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
        "critical_fails": CRITICAL_FAILS,
    }
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if CRITICAL_FAILS:
        print(f"FAIL: критичные чеки провалены ({len(CRITICAL_FAILS)})")
        sys.exit(1)
    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
