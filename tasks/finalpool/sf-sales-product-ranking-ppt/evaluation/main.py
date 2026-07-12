"""
Evaluation script for sf-sales-product-ranking-ppt task.

Data source is the ClickHouse-backed PostgreSQL schema sf_data (read live, so
PRODUCT_NAME stays English and CATEGORY values russified centrally are picked up
automatically -- the eval never hand-translates them).

Deliverables checked:
1. Product_Rankings.xlsx -- 'Top Products' (top-20 by revenue) + 'Category Summary'.
2. Product_Performance_Review.pptx -- >=4 slides, title marker, top product names.
3. Google Sheet 'Product Rankings Dashboard' with a 'Top 20' tab carrying the data.

Two-tier gating:
  CRITICAL_CHECKS (semantic correctness, computed live from the DB) -- any failure
  causes sys.exit(1) before the accuracy gate.
  NON-CRITICAL structural checks (file/sheet/column existence) feed the
  accuracy>=70% gate.

NOTE: file names (Product_Rankings.xlsx, Product_Performance_Review.pptx), sheet
names (Top Products, Category Summary, Top 20), the Google-Sheet title token and
the title-slide marker 'Q1 2026' are preserved English on purpose -- the agent is
instructed to keep them English and the eval greps them literally.
"""

import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "cowork_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_RESULTS = []  # (passed, name)


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def critical(name, passed, detail=""):
    CRITICAL_RESULTS.append((passed, name))
    record("[CRITICAL] " + name, passed, detail)


def num_close(a, b, rel_tol=0.05, abs_tol=1.0):
    try:
        a = float(a)
        b = float(b)
    except (TypeError, ValueError):
        return False
    return abs(a - b) <= max(abs_tol, abs(b) * rel_tol)


def norm(s):
    return str(s).strip().lower() if s is not None else ""


# ---------------------------------------------------------------------------
# Live DB expectations
# ---------------------------------------------------------------------------

def get_expected():
    """Compute top-20 products and category aggregates live from sf_data."""
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute('''
        SELECT p."PRODUCT_NAME", p."CATEGORY",
               SUM(o."TOTAL_AMOUNT") AS revenue,
               SUM(o."QUANTITY")     AS units
        FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
        JOIN sf_data."SALES_DW__PUBLIC__PRODUCTS" p
            ON o."PRODUCT_ID" = p."PRODUCT_ID"
        GROUP BY p."PRODUCT_NAME", p."CATEGORY"
        ORDER BY revenue DESC
        LIMIT 20
    ''')
    top20 = [(r[0], r[1], float(r[2]), float(r[3])) for r in cur.fetchall()]

    cur.execute('''
        SELECT p."CATEGORY",
               COUNT(DISTINCT p."PRODUCT_ID") AS product_count,
               SUM(o."TOTAL_AMOUNT")          AS revenue
        FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
        JOIN sf_data."SALES_DW__PUBLIC__PRODUCTS" p
            ON o."PRODUCT_ID" = p."PRODUCT_ID"
        GROUP BY p."CATEGORY"
        ORDER BY revenue DESC
    ''')
    categories = [(r[0], int(r[1]), float(r[2])) for r in cur.fetchall()]

    cur.close()
    conn.close()
    return top20, categories


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------

def find_sheet(wb, *needles):
    for s in wb.sheetnames:
        low = s.lower()
        if all(n in low for n in needles):
            return wb[s]
    return None


def check_excel(agent_workspace, top20, categories):
    print("\n=== Checking Excel Output ===")
    excel_path = os.path.join(agent_workspace, "Product_Rankings.xlsx")
    if not os.path.isfile(excel_path):
        record("Excel file exists", False, f"Expected {excel_path}")
        return
    record("Excel file exists", True)

    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
    except Exception as e:
        record("Excel file readable", False, str(e))
        return

    # --- Top Products sheet ---
    ws = find_sheet(wb, "top", "product") or find_sheet(wb, "top")
    record("Sheet 'Top Products' exists", ws is not None, f"Sheets: {wb.sheetnames}")

    if ws is not None:
        rows = [r for r in ws.iter_rows(min_row=2, values_only=True) if r and any(c is not None for c in r)]
        record("Top Products has 20 rows", len(rows) == 20, f"Got {len(rows)}")

        # Columns: Rank(0) Product_Name(1) Category(2) Total_Revenue(3) Units_Sold(4) Avg_Order_Value(5)
        # CRITICAL: #1 product name + revenue match the live DB top product.
        if rows and top20:
            top_name_exp = norm(top20[0][0])
            top_rev_exp = top20[0][2]
            r0 = rows[0]
            name_ok = len(r0) > 1 and norm(r0[1]) == top_name_exp
            rev_ok = len(r0) > 3 and num_close(r0[3], top_rev_exp, rel_tol=0.05)
            critical("Top Products #1 product name matches DB top product",
                     name_ok, f"Expected '{top20[0][0]}', got '{r0[1] if len(r0) > 1 else None}'")
            critical("Top Products #1 Total_Revenue within 5% of DB",
                     rev_ok, f"Expected ~{top_rev_exp:.2f}, got {r0[3] if len(r0) > 3 else None}")

        # CRITICAL: Rank 1..20 ascending AND revenue monotonically non-increasing.
        if len(rows) == 20:
            ranks_ok = True
            rev_mono_ok = True
            prev_rev = None
            for i, r in enumerate(rows):
                try:
                    if int(r[0]) != i + 1:
                        ranks_ok = False
                except (TypeError, ValueError, IndexError):
                    ranks_ok = False
                try:
                    rev = float(r[3])
                    if prev_rev is not None and rev > prev_rev * 1.001:
                        rev_mono_ok = False
                    prev_rev = rev
                except (TypeError, ValueError, IndexError):
                    rev_mono_ok = False
            critical("Rank column is 1..20 ascending", ranks_ok, "Rank values not 1..20 in order")
            critical("Revenue non-increasing down the ranking", rev_mono_ok,
                     "Revenue not monotonically non-increasing by rank")

    # --- Category Summary sheet ---
    ws2 = find_sheet(wb, "category") or find_sheet(wb, "summary")
    record("Category Summary sheet exists", ws2 is not None, f"Sheets: {wb.sheetnames}")

    if ws2 is not None and categories:
        crows = [r for r in ws2.iter_rows(min_row=2, values_only=True) if r and any(c is not None for c in r)]
        # Columns: Category(0) Product_Count(1) Total_Revenue(2) Avg_Revenue_Per_Product(3)
        # CRITICAL: top category's Total_Revenue aggregate matches a live GROUP BY.
        top_cat_name, _, top_cat_rev = categories[0]
        agent_cat = next((r for r in crows if norm(r[0]) == norm(top_cat_name)), None)
        if agent_cat is None:
            critical("Category Summary top category present", False,
                     f"Expected category '{top_cat_name}' not found in {[r[0] for r in crows]}")
        else:
            rev_ok = len(agent_cat) > 2 and num_close(agent_cat[2], top_cat_rev, rel_tol=0.05)
            critical("Category Summary top category Total_Revenue matches DB aggregate",
                     rev_ok, f"Expected ~{top_cat_rev:.2f} for '{top_cat_name}', got {agent_cat[2] if len(agent_cat) > 2 else None}")

    wb.close()


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def check_pptx(agent_workspace, top20):
    print("\n=== Checking PowerPoint Output ===")
    pptx_path = os.path.join(agent_workspace, "Product_Performance_Review.pptx")
    if not os.path.isfile(pptx_path):
        record("PPTX file exists", False, f"Expected {pptx_path}")
        return
    record("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        record("PPTX file readable", False, str(e))
        return

    slide_count = len(prs.slides)
    record("PPTX has >= 4 slides", slide_count >= 4, f"Got {slide_count} slides")

    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text + "\n"
            # Tables are graphicFrame shapes (has_text_frame=False); harvest cells too.
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        all_text += cell.text + "\n"
    low = all_text.lower()

    # Title marker 'Q1 2026' is preserved English in the prompt.
    record("PPTX contains 'Q1 2026' title marker", "q1 2026" in low,
           f"Text sample: {all_text[:200]}")

    # CRITICAL: the top-3 DB product names actually appear in the slide text.
    # Long PRODUCT_NAMEs (~100 chars) are truncated with an ellipsis on real and
    # groundtruth slides, so match a truncated prefix of the name (or accept any
    # slide line that is itself a prefix of the DB name) instead of the full name.
    if top20:
        slide_lines = [norm(ln) for ln in all_text.splitlines() if ln.strip()]
        for name, *_ in top20[:3]:
            n = norm(name)
            prefix = n[:40]
            present = (prefix and prefix in low) or any(
                ln and (n.startswith(ln) or ln.startswith(prefix)) for ln in slide_lines
            )
            critical(f"PPTX mentions top product '{name[:40]}'", present,
                     f"Product name (prefix '{prefix}') not found in slide text")


# ---------------------------------------------------------------------------
# Google Sheet
# ---------------------------------------------------------------------------

def check_gsheet(top20):
    print("\n=== Checking Google Sheet ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    # Broaden title match to RU + EN tokens.
    cur.execute("SELECT id, title FROM gsheet.spreadsheets")
    sheets = cur.fetchall()
    tokens = ("product", "ranking", "товар", "продукт", "рейтинг")
    matching = [s for s in sheets if s[1] and any(t in s[1].lower() for t in tokens)]

    critical("Google Sheet 'Product Rankings Dashboard' exists", len(matching) > 0,
             f"Sheet titles: {[s[1] for s in sheets]}")
    if not matching:
        cur.close()
        conn.close()
        return

    sheet_id = matching[0][0]

    # Prefer a 'Top 20' tab if present, else any tab.
    cur.execute("SELECT id, title FROM gsheet.sheets WHERE spreadsheet_id = %s", (sheet_id,))
    tabs = cur.fetchall()
    record("Spreadsheet has at least 1 tab", len(tabs) >= 1, f"Found {len(tabs)} tabs")
    if not tabs:
        cur.close()
        conn.close()
        return

    tab_id = next((t[0] for t in tabs if t[1] and "20" in t[1].lower()), tabs[0][0])

    cur.execute("""
        SELECT row_index, col_index, value FROM gsheet.cells
        WHERE sheet_id = %s ORDER BY row_index, col_index
    """, (tab_id,))
    cells = cur.fetchall()
    cur.close()
    conn.close()

    grid = {}
    for ri, ci, val in cells:
        grid.setdefault(ri, {})[ci] = val

    if grid:
        min_row = min(grid.keys())
        data_rows = {k: v for k, v in grid.items() if k > min_row}
    else:
        data_rows = {}
    record("Google Sheet 'Top 20' has data rows", len(data_rows) >= 1,
           f"Found {len(data_rows)} data rows")

    # CRITICAL: first data row carries the DB top product name + revenue.
    if data_rows and top20:
        first_key = min(data_rows.keys())
        first_vals = [str(v) for v in data_rows[first_key].values()]
        joined = " ".join(first_vals).lower()
        top_name = norm(top20[0][0])
        top_rev = top20[0][2]
        name_ok = top_name in joined
        rev_ok = any(num_close(v, top_rev, rel_tol=0.05)
                     for v in data_rows[first_key].values()
                     if _is_num(v))
        critical("Google Sheet 'Top 20' first row product name matches DB top product",
                 name_ok, f"Expected '{top20[0][0]}' in {first_vals}")
        critical("Google Sheet 'Top 20' first row revenue matches DB top product",
                 rev_ok, f"Expected ~{top_rev:.2f} in {first_vals}")


def _is_num(v):
    try:
        float(str(v).replace(",", "").replace(" ", ""))
        return True
    except (TypeError, ValueError):
        return False


# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    try:
        top20, categories = get_expected()
    except Exception as e:
        print(f"Error querying database: {e}")
        sys.exit(1)
    print(f"Expected top product: {top20[0][0] if top20 else 'N/A'}")

    check_excel(args.agent_workspace, top20, categories)
    check_pptx(args.agent_workspace, top20)
    check_gsheet(top20)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100.0) if total else 0.0

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Accuracy: {accuracy:.1f}%")

    critical_failures = [name for passed, name in CRITICAL_RESULTS if not passed]

    if args.res_log_file:
        result = {
            "passed": PASS_COUNT,
            "failed": FAIL_COUNT,
            "accuracy": accuracy,
            "critical_failures": critical_failures,
            "success": (not critical_failures) and accuracy >= 70.0,
        }
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if critical_failures:
        print("  Overall: FAIL")
        print(f"  CRITICAL CHECK(S) FAILED ({len(critical_failures)}):")
        for name in critical_failures:
            print(f"    - {name}")
        sys.exit(1)

    overall = accuracy >= 70.0
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
