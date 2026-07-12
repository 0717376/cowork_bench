"""
Evaluation для задачи notion-gform-gcal-onboarding (RU: teamly + forms).

Проверки:
1. PPTX Onboarding_Presentation.pptx: >=4 слайдов; титульный упоминает Welcome + 2026;
   слайд "New Team Members" с именами/почтами/отделами трёх новичков.
2. Календарь: события "Orientation Session" (9-12) и "Team Lunch" (12-13) на 2026-03-16;
   у Orientation участники — все три почты новичков.
3. Письма: >=3 отправлено, по одному на каждую из трёх почт, тема —
   "Welcome to Our Company - Onboarding Information", в теле имя + отдел + дата/ориентация.
4. Teamly: страница "New Employee Onboarding Checklist" получила раздел
   "March 2026 New Hires" с фамилиями всех трёх новичков.

CRITICAL_CHECKS: любой их fail => FAIL независимо от общей accuracy (sys.exit(1)
до accuracy-gate). Порог общей accuracy: >= 70%.
"""

import argparse
import json
import os
import sys

import psycopg2
from pptx import Presentation

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent",
    "password": "camel",
}

# Новые сотрудники: (имя_ru, email, отдел_ru, фамилия_альтернативы[ru, en])
NEW_HIRES = [
    ("Анна Парк", "sarah.park@company.com", "Инженерия", ("парк", "park")),
    ("Михаил Чен", "mike.chen@company.com", "Продажи", ("чен", "chen")),
    ("Мария Родригес", "amy.rodriguez@company.com", "Маркетинг", ("родригес", "rodriguez")),
]
HIRE_EMAILS = [h[1] for h in NEW_HIRES]

EMAIL_SUBJECT = "Welcome to Our Company - Onboarding Information"

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILED = []

# Имена (строки name из record()), фейл которых = FAIL всей задачи.
CRITICAL_CHECKS = {
    "CRITICAL: оба события (Orientation 9-12, Team Lunch 12-13) на 2026-03-16",
    "CRITICAL: участники Orientation Session — все три почты новичков",
    "CRITICAL: 3 письма с правильной темой — по одному на каждую почту новичка",
    "CRITICAL: каждое письмо упоминает имя получателя и его отдел",
    "CRITICAL: Teamly-страница чек-листа содержит раздел 'March 2026 New Hires' со всеми тремя фамилиями",
    "CRITICAL: PPTX содержит слайд New Team Members с именами, почтами и отделами всех троих",
}


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT, CRITICAL_FAILED
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")
        if name in CRITICAL_CHECKS:
            CRITICAL_FAILED.append(name)


def str_contains(haystack, needle):
    if haystack is None or needle is None:
        return False
    return needle.lower() in str(haystack).lower()


# ============================================================================
# Check 1: PowerPoint
# ============================================================================

def check_pptx(agent_workspace):
    print("\n=== Проверка PowerPoint ===")
    all_ok = True

    pptx_path = os.path.join(agent_workspace, "Onboarding_Presentation.pptx")
    if not os.path.isfile(pptx_path):
        record("PPT file exists", False, f"Не найдено: {pptx_path}")
        record("CRITICAL: PPTX содержит слайд New Team Members с именами, почтами и отделами всех троих",
               False, "нет файла")
        return False
    record("PPT file exists", True)

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        record("PPT file readable", False, str(e))
        record("CRITICAL: PPTX содержит слайд New Team Members с именами, почтами и отделами всех троих",
               False, "файл не читается")
        return False
    record("PPT file readable", True)

    slide_count = len(prs.slides)
    ok = slide_count >= 4
    record("PPT has at least 4 slides", ok, f"Найдено {slide_count} слайдов")
    all_ok = all_ok and ok

    # Титульный слайд
    if slide_count > 0:
        first_slide = prs.slides[0]
        slide_text = ""
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                slide_text += " " + shape.text_frame.text

        ok = str_contains(slide_text, "Welcome")
        record("Title slide mentions 'Welcome'", ok, f"Текст: {slide_text[:200]}")
        all_ok = all_ok and ok

        ok = str_contains(slide_text, "2026")
        record("Title slide mentions '2026'", ok, f"Текст: {slide_text[:200]}")
        all_ok = all_ok and ok
    else:
        record("Title slide exists", False, "Нет слайдов")
        all_ok = False

    # Весь текст презентации
    all_slide_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_slide_text += " " + shape.text_frame.text
    all_slide_lower = all_slide_text.lower()

    # Контент чек-листа: RU + EN ключевые слова
    checklist_keywords = [
        # RU стемы
        "hr", "документ", "рабоч", "учётн", "учетн", "команд",
        "вводн", "семинар", "справочник", "комплаенс", "обучен", "отдел",
        # EN (на случай англоязычной презентации)
        "paperwork", "workstation", "account", "team",
        "orientation", "handbook", "compliance", "training",
    ]
    checklist_found = sum(1 for kw in checklist_keywords if kw in all_slide_lower)
    ok = checklist_found >= 3
    record("PPT contains checklist content (>=3 keywords)", ok,
           f"Найдено {checklist_found} ключевых слов")
    all_ok = all_ok and ok

    # Упоминание новых сотрудников (RU + EN)
    ok = ("team member" in all_slide_lower or "new hire" in all_slide_lower
          or "new team" in all_slide_lower or "member" in all_slide_lower
          or "сотрудник" in all_slide_lower or "новичк" in all_slide_lower
          or "команд" in all_slide_lower)
    record("PPT mentions new team members", ok, "Нет упоминания о новых сотрудниках")
    all_ok = all_ok and ok

    # CRITICAL: слайд New Team Members — имена + почты + отделы всех троих
    emails_in = sum(1 for e in HIRE_EMAILS if e.lower() in all_slide_lower)
    surnames_in = sum(1 for h in NEW_HIRES if any(a in all_slide_lower for a in h[3]))
    depts_in = sum(1 for h in NEW_HIRES if h[2].lower() in all_slide_lower)
    crit_ok = emails_in == 3 and surnames_in == 3 and depts_in >= 2
    record("CRITICAL: PPTX содержит слайд New Team Members с именами, почтами и отделами всех троих",
           crit_ok,
           f"emails={emails_in}/3, фамилии={surnames_in}/3, отделы={depts_in}/3")
    all_ok = all_ok and crit_ok

    return all_ok


# ============================================================================
# Check 2: Google Calendar
# ============================================================================

def check_gcal():
    print("\n=== Проверка календаря ===")
    all_ok = True

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT summary, description, start_datetime, end_datetime, attendees::text
        FROM gcal.events
        ORDER BY start_datetime
    """)
    events = cur.fetchall()
    cur.close()
    conn.close()

    print(f"[check_gcal] Найдено {len(events)} событий.")
    for ev in events:
        print(f"  Событие: {ev[0]} | {ev[2]} - {ev[3]}")

    ok = len(events) >= 2
    record("At least 2 calendar events created", ok, f"Найдено {len(events)}")
    all_ok = all_ok and ok

    def find(kw):
        for summary, description, start_dt, end_dt, attendees in events:
            if kw in (summary or "").lower():
                return summary, description, start_dt, end_dt, attendees
        return None

    orient = find("orientation")
    lunch = find("lunch")

    record("Orientation Session event exists", orient is not None,
           "Нет события с 'Orientation' в названии")
    record("Team Lunch event exists", lunch is not None,
           "Нет события с 'Lunch' в названии")

    # CRITICAL: оба события на 2026-03-16 с верным временем
    def check_event(ev, exp_start_h, exp_end_h):
        if ev is None:
            return False
        _s, _d, start_dt, end_dt, _a = ev
        if start_dt is None or end_dt is None:
            return False
        if start_dt.strftime("%Y-%m-%d") != "2026-03-16":
            return False
        return start_dt.hour == exp_start_h and end_dt.hour == exp_end_h

    orient_time_ok = check_event(orient, 9, 12)
    lunch_time_ok = check_event(lunch, 12, 13)

    # Некритические подробности (для accuracy)
    record("Orientation on 2026-03-16 09:00-12:00", orient_time_ok,
           f"{orient[2]}..{orient[3]}" if orient else "нет события")
    record("Team Lunch on 2026-03-16 12:00-13:00", lunch_time_ok,
           f"{lunch[2]}..{lunch[3]}" if lunch else "нет события")

    crit_both = orient_time_ok and lunch_time_ok
    record("CRITICAL: оба события (Orientation 9-12, Team Lunch 12-13) на 2026-03-16",
           crit_both, f"orient={orient_time_ok}, lunch={lunch_time_ok}")
    all_ok = all_ok and crit_both

    # CRITICAL: участники Orientation — все три почты
    attendee_blob = (orient[4] or "").lower() if orient else ""
    att_found = sum(1 for e in HIRE_EMAILS if e.lower() in attendee_blob)
    crit_att = att_found == 3
    record("CRITICAL: участники Orientation Session — все три почты новичков",
           crit_att, f"найдено {att_found}/3 в attendees")
    all_ok = all_ok and crit_att

    return all_ok


# ============================================================================
# Check 3: Emails
# ============================================================================

def check_emails():
    print("\n=== Проверка писем ===")
    all_ok = True

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    # to_addr — JSONB; матчим по подстроке почты (как в эталонных RU-тасках).
    cur.execute("""
        SELECT subject, from_addr, to_addr::text, body_text
        FROM email.messages
    """)
    all_msgs = cur.fetchall()
    cur.close()
    conn.close()

    # Письма агента: те, что адресованы на наши три почты.
    sent = {}  # email -> list of (subject, body)
    for subject, from_addr, to_text, body in all_msgs:
        for e in HIRE_EMAILS:
            if e.lower() in (to_text or "").lower():
                sent.setdefault(e, []).append((subject or "", body or ""))

    total_sent = sum(len(v) for v in sent.values())
    print(f"[check_emails] Писем новичкам: {total_sent}")

    ok = total_sent >= 3
    record("At least 3 emails sent to hires", ok, f"Найдено {total_sent}")
    all_ok = all_ok and ok

    # CRITICAL: по одному письму на каждую почту с правильной темой
    per_hire_subject_ok = 0
    for e in HIRE_EMAILS:
        msgs = sent.get(e, [])
        if any(EMAIL_SUBJECT.lower() in (s or "").lower() for s, _b in msgs):
            per_hire_subject_ok += 1
    crit_subj = per_hire_subject_ok == 3
    record("CRITICAL: 3 письма с правильной темой — по одному на каждую почту новичка",
           crit_subj, f"корректных адресатов={per_hire_subject_ok}/3")
    all_ok = all_ok and crit_subj

    # CRITICAL: тело каждого письма упоминает имя получателя и его отдел
    name_dept_ok = 0
    for name, email, dept, alts in NEW_HIRES:
        msgs = sent.get(email, [])
        blob = " ".join((s + " " + b) for s, b in msgs).lower()
        name_hit = (name.lower() in blob
                    or any(a in blob for a in alts)
                    or name.split()[0].lower() in blob)
        dept_hit = dept.lower() in blob
        if name_hit and dept_hit:
            name_dept_ok += 1
    crit_nd = name_dept_ok == 3
    record("CRITICAL: каждое письмо упоминает имя получателя и его отдел",
           crit_nd, f"имя+отдел совпали для {name_dept_ok}/3")
    all_ok = all_ok and crit_nd

    # Некритическое: дата/ориентация в теле
    date_ok = 0
    for e in HIRE_EMAILS:
        blob = " ".join((s + " " + b) for s, b in sent.get(e, [])).lower()
        if ("март" in blob or "march" in blob or "2026" in blob
                or "ориентац" in blob or "вводн" in blob or "orientation" in blob):
            date_ok += 1
    record("Emails: тело упоминает дату/ориентацию (>=3)", date_ok >= 3,
           f"найдено {date_ok}/3")

    return all_ok


# ============================================================================
# Check 4: Teamly
# ============================================================================

def check_teamly():
    print("\n=== Проверка Teamly ===")
    all_ok = True

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("SELECT id, title, body FROM teamly.pages")
    pages = cur.fetchall()
    cur.close()
    conn.close()

    print(f"[check_teamly] Найдено {len(pages)} страниц.")

    # Страница чек-листа онбординга (заголовок английский).
    page = None
    for pid, title, body in pages:
        tl = (title or "").lower()
        if "onboarding" in tl and "checklist" in tl:
            page = (pid, title, body)
            break

    ok = page is not None
    record("Onboarding checklist page exists", ok,
           f"Нет страницы с 'onboarding' и 'checklist' среди {len(pages)} страниц")
    if not page:
        record("CRITICAL: Teamly-страница чек-листа содержит раздел 'March 2026 New Hires' со всеми тремя фамилиями",
               False, "нет страницы")
        return False

    body = (page[2] or "")
    body_lower = body.lower()

    # Раздел "March 2026 New Hires" (или RU-эквивалент).
    section_ok = ("march 2026" in body_lower
                  or ("март" in body_lower and "2026" in body_lower
                      and ("нович" in body_lower or "сотрудник" in body_lower
                           or "new hire" in body_lower)))
    record("Teamly blocks contain 'March 2026 New Hires' section", section_ok,
           "Раздел 'March 2026 New Hires' не найден")
    all_ok = all_ok and section_ok

    # CRITICAL: все три фамилии в разделе
    surnames_in = sum(1 for h in NEW_HIRES if any(a in body_lower for a in h[3]))
    crit_sec = section_ok and surnames_in == 3
    record("CRITICAL: Teamly-страница чек-листа содержит раздел 'March 2026 New Hires' со всеми тремя фамилиями",
           crit_sec, f"section={section_ok}, фамилии={surnames_in}/3")
    all_ok = all_ok and crit_sec

    return all_ok


# ============================================================================
# Main
# ============================================================================

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    pptx_ok = check_pptx(args.agent_workspace)
    gcal_ok = check_gcal()
    email_ok = check_emails()
    teamly_ok = check_teamly()

    total = PASS_COUNT + FAIL_COUNT
    pct = 100.0 * PASS_COUNT / total if total else 0.0

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Accuracy: {pct:.1f}%")

    success = (not CRITICAL_FAILED) and pct >= 70.0

    if args.res_log_file:
        result = {
            "passed": PASS_COUNT,
            "failed": FAIL_COUNT,
            "accuracy": pct,
            "success": success,
            "critical_failed": CRITICAL_FAILED,
            "details": {
                "pptx": pptx_ok,
                "gcal": gcal_ok,
                "email": email_ok,
                "teamly": teamly_ok,
            },
        }
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, ensure_ascii=False, indent=2)

    if CRITICAL_FAILED:
        print(f"CRITICAL FAIL: {CRITICAL_FAILED}")
        print("Overall: FAIL")
        sys.exit(1)
    if pct < 70.0:
        print("Overall: FAIL (accuracy < 70%)")
        sys.exit(1)
    print("Overall: PASS")
    sys.exit(0)


if __name__ == "__main__":
    main()
