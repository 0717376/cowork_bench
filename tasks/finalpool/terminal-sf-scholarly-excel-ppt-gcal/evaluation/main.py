"""Оценка для terminal-sf-scholarly-excel-ppt-gcal (форк ClickHouse).

Значения REGION/SEGMENT в sf_data русифицированы центральным отображением
(scripts/clickhouse_relabel_map.REGIONS/SEGMENTS), поэтому seed <-> eval <->
groundtruth остаются согласованными. Идентификаторы колонок/листов/таблиц/файлов,
теги "Priority Focus"/"On Track", имена владельцев из ростера и названия статей
остаются английскими (eval ищет их подстрокой).

Проверки:
1. Sales_Strategy_Analysis.xlsx (4 листа с корректными данными)
2. Sales_Strategy_Presentation.pptx (6 слайдов)
3. Событие Google Calendar "Q1 Sales Strategy Review"
4. Python-скрипты (analyze_sales_gaps.py, match_recommendations.py, generate_summary.py)
5. Выходные файлы (sales_gaps.json, research_recommendations.json, executive_summary.txt)

CRITICAL_CHECKS: любой провал критической проверки => sys.exit(1) до порога accuracy>=70.
"""
import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent", "password": "camel",
}

# Русифицированные имена регионов/сегментов (из scripts/clickhouse_relabel_map).
RU_EUROPE = "европа"
RU_LATAM = "латинская америка"
RU_MIDDLE = "ближний восток"
RU_NORTH = "северная америка"
RU_ASIA = "азиатско-тихоокеанский регион"
RU_REGIONS = [RU_NORTH, RU_EUROPE, RU_ASIA, RU_LATAM, RU_MIDDLE]
RU_CONSUMER = "частные клиенты"
RU_ENTERPRISE = "корпоративный"

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILED = []


def critical(name, condition, detail=""):
    """Критическая проверка: учитывается в accuracy И блокирует PASS при провале."""
    global CRITICAL_FAILED
    check(name, condition, detail)
    if not condition:
        CRITICAL_FAILED.append(name)
    return bool(condition)


def get_expected_from_db():
    """Query sf_data schema dynamically for regional revenue values."""
    defaults = {
        "europe_rev": 648798.0,
        "latam_rev": 549129.0,
        "total_rev": 3048998.0,
    }
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT c."REGION", SUM(o."TOTAL_AMOUNT") as rev
            FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
            JOIN sf_data."SALES_DW__PUBLIC__CUSTOMERS" c ON o."CUSTOMER_ID" = c."CUSTOMER_ID"
            GROUP BY c."REGION"
            ORDER BY rev DESC
        """)
        rows = cur.fetchall()
        total = 0.0
        for region, rev in rows:
            rev_f = float(rev)
            total += rev_f
            r = region.strip().lower() if region else ""
            if r == RU_EUROPE or r == "europe":
                defaults["europe_rev"] = rev_f
            elif "латинская" in r or "latin" in r:
                defaults["latam_rev"] = rev_f
        defaults["total_rev"] = total
        cur.close()
        conn.close()
    except Exception as e:
        print(f"  [WARN] DB query for expected values failed, using defaults: {e}")
    return defaults


EXPECTED = get_expected_from_db()


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:200]}")


def num_close(a, b, tol=2.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except Exception:
        return False


def check_excel(workspace):
    print("\n=== Check 1: Sales_Strategy_Analysis.xlsx ===")
    path = os.path.join(workspace, "Sales_Strategy_Analysis.xlsx")
    if not os.path.exists(path):
        check("Excel file exists", False, f"Not found at {path}")
        return
    check("Excel file exists", True)

    wb = openpyxl.load_workbook(path)
    sheets = wb.sheetnames
    check("Has 4 sheets", len(sheets) >= 4, f"Found {len(sheets)}: {sheets}")

    sheets_lower = [s.lower().replace(" ", "_") for s in sheets]

    # Sheet 1: Regional_Performance
    rp_idx = next((i for i, s in enumerate(sheets_lower) if "regional" in s or "region" in s), 0)
    ws1 = wb[sheets[rp_idx]]
    rows1 = list(ws1.iter_rows(values_only=True))
    data1 = [r for r in rows1[1:] if any(c for c in r)]
    # CRITICAL: ровно 5 строк регионов.
    critical("Regional_Performance has exactly 5 region rows", len(data1) == 5, f"Found {len(data1)}")

    all_text1 = " ".join(str(c) for r in rows1 for c in r if c).lower()
    check("Contains Europe region", RU_EUROPE in all_text1 or "europe" in all_text1)
    check("Contains Latin America region", "латинская" in all_text1 or "latin" in all_text1)
    check("Contains Priority Focus tag", "priority" in all_text1 and "focus" in all_text1,
          f"Text: {all_text1[:150]}")
    check("Contains On Track tag", "on track" in all_text1 or "on_track" in all_text1 or "ontrack" in all_text1,
          f"Text: {all_text1[:150]}")

    def region_match(cell, ru, en):
        c = str(cell).lower().strip()
        return ru in c or en in c

    # Check revenue values are reasonable (within tolerance of actual DB values)
    europe_rev = None
    latam_rev = None
    for row in data1:
        if row and region_match(row[0], RU_EUROPE, "europe"):
            europe_rev = row[1]
        if row and region_match(row[0], "латинская", "latin"):
            latam_rev = row[1]
    # CRITICAL: выручка Европы и Латинской Америки взята из живого склада (в пределах допуска).
    critical("Europe revenue matches live warehouse",
             europe_rev is not None and num_close(europe_rev, EXPECTED["europe_rev"], tol=5000),
             f"Got {europe_rev}, expected ~{EXPECTED['europe_rev']:.0f}")
    critical("Latin America revenue matches live warehouse",
             latam_rev is not None and num_close(latam_rev, EXPECTED["latam_rev"], tol=5000),
             f"Got {latam_rev}, expected ~{EXPECTED['latam_rev']:.0f}")

    # Check priority tags: Middle East and Latin America should be Priority Focus, Europe On Track
    priority_regions = []
    europe_tag = ""
    for row in data1:
        if row and "priority" in str(row[-1]).lower():
            priority_regions.append(str(row[0]).lower())
        if row and region_match(row[0], RU_EUROPE, "europe"):
            europe_tag = str(row[-1]).lower()
    # CRITICAL: корректная логика приоритета по медиане (Ближний Восток + Латинская Америка ниже медианы; Европа выше).
    critical("Middle East is Priority Focus",
             any(("ближний" in r or "middle" in r) for r in priority_regions),
             f"Priority regions: {priority_regions}")
    critical("Latin America is Priority Focus",
             any(("латинская" in r or "latin" in r) for r in priority_regions),
             f"Priority regions: {priority_regions}")
    critical("Europe is On Track (above median)",
             "priority" not in europe_tag and europe_tag != "",
             f"Europe tag: {europe_tag!r}")

    # Sheet 2: Segment_Analysis
    sa_idx = next((i for i, s in enumerate(sheets_lower) if "segment" in s), 1)
    if sa_idx < len(sheets):
        ws2 = wb[sheets[sa_idx]]
        rows2 = list(ws2.iter_rows(values_only=True))
        data2 = [r for r in rows2[1:] if any(c for c in r)]
        check("Segment_Analysis has 4 segment rows", len(data2) >= 4, f"Found {len(data2)}")
        all_text2 = " ".join(str(c) for r in rows2 for c in r if c).lower()
        check("Contains Consumer segment", RU_CONSUMER in all_text2 or "consumer" in all_text2)
        check("Contains Enterprise segment", RU_ENTERPRISE in all_text2 or "enterprise" in all_text2)

    # Sheet 3: Research_Insights
    ri_idx = next((i for i, s in enumerate(sheets_lower) if "research" in s or "insight" in s), 2)
    if ri_idx < len(sheets):
        ws3 = wb[sheets[ri_idx]]
        rows3 = list(ws3.iter_rows(values_only=True))
        data3 = [r for r in rows3[1:] if any(c for c in r)]
        check("Research_Insights has at least 2 rows", len(data3) >= 2, f"Found {len(data3)}")
        all_text3 = " ".join(str(c) for r in rows3 for c in r if c).lower()
        # Названия релевантных статей — английские реалии (статьи на английском).
        relevant_titles = [
            "data-driven territory management",
            "customer segmentation strategies",
            "discount optimization",
            "sales team specialization",
        ]
        title_hits = sum(1 for t in relevant_titles if t in all_text3)
        # CRITICAL: цитируются >=2 из 4 релевантных статей.
        critical("Research_Insights cites >=2 relevant papers", title_hits >= 2,
                 f"Matched {title_hits} relevant titles; text: {all_text3[:150]}")
        check("Research references territory/segmentation/optimization",
              any(k in all_text3 for k in
                  ["territory", "segmentation", "optimization", "территор", "сегмент", "оптимизаци"]),
              f"Text: {all_text3[:150]}")

    # Sheet 4: Action_Items
    ai_idx = next((i for i, s in enumerate(sheets_lower) if "action" in s), 3)
    if ai_idx < len(sheets):
        ws4 = wb[sheets[ai_idx]]
        rows4 = list(ws4.iter_rows(values_only=True))
        data4 = [r for r in rows4[1:] if any(c for c in r)]
        check("Action_Items has at least 2 rows", len(data4) >= 2, f"Found {len(data4)}")
        all_text4 = " ".join(str(c) for r in rows4 for c in r if c).lower()
        check("Action items reference research papers",
              any(k in all_text4 for k in
                  ["territory", "segmentation", "specialization", "pricing",
                   "территор", "сегмент", "специализ", "ценообраз"]),
              f"Text: {all_text4[:150]}")
        # CRITICAL: владельцы действий взяты из ростера (имена английские) для Priority Focus регионов
        # Ближний Восток -> Ahmed Hassan, Латинская Америка -> Carlos Rivera.
        critical("Action items owners match roster (Hassan/Rivera)",
                 ("ahmed" in all_text4 or "hassan" in all_text4) and
                 ("carlos" in all_text4 or "rivera" in all_text4),
                 f"Text: {all_text4[:200]}")
        # CRITICAL: research_basis ссылается на реальную статью.
        critical("Action items cite a real seeded paper title",
                 any(t in all_text4 for t in
                     ["data-driven territory management", "customer segmentation strategies",
                      "discount optimization", "sales team specialization"]),
                 f"Text: {all_text4[:200]}")


def check_pptx(workspace):
    print("\n=== Check 2: Sales_Strategy_Presentation.pptx ===")
    path = os.path.join(workspace, "Sales_Strategy_Presentation.pptx")
    if not os.path.exists(path):
        check("PPTX file exists", False, f"Not found at {path}")
        return
    check("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = list(prs.slides)
        check("Has 6 slides", len(slides) >= 6, f"Found {len(slides)} slides")

        # Check slide content
        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text + " "

        all_lower = all_text.lower()
        check("Title slide mentions Q1 or strategy",
              any(k in all_lower for k in ["q1", "strategy", "стратег"]),
              f"Text: {all_lower[:200]}")
        check("Contains regional data",
              any(k in all_lower for k in ["europe", "region", RU_EUROPE, "регион"]),
              f"Text snippet: {all_lower[:200]}")
        check("Contains segment data",
              any(k in all_lower for k in ["consumer", "segment", RU_CONSUMER, "сегмент"]),
              f"Text snippet: {all_lower[:200]}")
        check("Contains research findings",
              any(k in all_lower for k in
                  ["research", "finding", "study", "исследован", "вывод"]),
              f"Text snippet: {all_lower[:200]}")
        check("Contains action plan",
              any(k in all_lower for k in
                  ["action", "plan", "initiative", "действ", "план", "инициатив"]),
              f"Text snippet: {all_lower[:200]}")
        check("Contains timeline",
              any(k in all_lower for k in
                  ["timeline", "milestone", "week", "график", "веха", "недел", "срок"]),
              f"Text snippet: {all_lower[:200]}")
    except ImportError:
        check("python-pptx available", False, "Cannot import pptx module")


def check_gcal():
    print("\n=== Check 3: Calendar Event ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("""
        SELECT summary, start_datetime, end_datetime, description
        FROM gcal.events
        WHERE lower(summary) LIKE '%%q1%%strategy%%'
           OR lower(summary) LIKE '%%sales%%strategy%%review%%'
        ORDER BY start_datetime
    """)
    events = cur.fetchall()
    # CRITICAL: событие создано.
    critical("Q1 Sales Strategy Review event exists", len(events) >= 1,
             f"Found {len(events)} matching events")

    if events:
        ev = events[0]
        summary, start, end, desc = ev
        # Check it's 90 minutes
        if start and end:
            duration = (end - start).total_seconds() / 60
            check("Event is 90 minutes", num_close(duration, 90, tol=15),
                  f"Duration: {duration} minutes")

        # CRITICAL: рабочее время (9-17), будний день, без конфликтов с 14 засеянными событиями.
        ok_hours = bool(start) and 9 <= start.hour <= 15
        ok_weekday = bool(start) and start.weekday() < 5
        critical("Event during business hours on a weekday", ok_hours and ok_weekday,
                 f"Start: {start}")

        conflicts = None
        if start and end:
            cur.execute("""
                SELECT COUNT(*) FROM gcal.events
                WHERE id != (SELECT id FROM gcal.events WHERE lower(summary) LIKE '%%q1%%strategy%%' OR lower(summary) LIKE '%%sales%%strategy%%review%%' LIMIT 1)
                  AND start_datetime < %s AND end_datetime > %s
            """, (end, start))
            conflicts = cur.fetchone()[0]
        critical("No calendar conflicts", conflicts == 0,
                 f"Found {conflicts} conflicting events")

        # Check description mentions review/action (RU+EN)
        if desc:
            dl = desc.lower()
            check("Description mentions review or action",
                  any(k in dl for k in
                      ["review", "action", "regional",
                       "обзор", "действ", "регион", "результат"]),
                  f"Description: {str(desc)[:100]}")

    cur.close()
    conn.close()


def check_reverse_validation(workspace):
    print("\n=== Reverse Validation ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        # Check no gcal events on weekends
        cur.execute("""
            SELECT summary, start_datetime
            FROM gcal.events
            WHERE (lower(summary) LIKE '%%q1%%strategy%%'
               OR lower(summary) LIKE '%%sales%%strategy%%review%%')
              AND EXTRACT(DOW FROM start_datetime) IN (0, 6)
        """)
        weekend_events = cur.fetchall()
        check("No strategy review events on weekends", len(weekend_events) == 0,
              f"Found {len(weekend_events)} weekend events: {weekend_events}")

        # Check no duplicate strategy review events
        cur.execute("""
            SELECT COUNT(*)
            FROM gcal.events
            WHERE lower(summary) LIKE '%%q1%%strategy%%'
               OR lower(summary) LIKE '%%sales%%strategy%%review%%'
        """)
        event_count = cur.fetchone()[0]
        check("No duplicate strategy review events", event_count <= 1,
              f"Found {event_count} strategy review events, expected 1")

        # Check Research_Insights sheet does not contain noise papers
        # (healthcare, machine learning theory papers should be excluded)
        xlsx_path = os.path.join(workspace, "Sales_Strategy_Analysis.xlsx")
        if os.path.exists(xlsx_path):
            wb = openpyxl.load_workbook(xlsx_path)
            sheets_lower = [s.lower().replace(" ", "_") for s in wb.sheetnames]
            ri_idx = next((i for i, s in enumerate(sheets_lower) if "research" in s or "insight" in s), None)
            if ri_idx is not None:
                ws = wb[wb.sheetnames[ri_idx]]
                rows = list(ws.iter_rows(min_row=2, values_only=True))
                all_text = " ".join(str(c) for r in rows for c in r if c).lower()
                noise_topics = ["healthcare", "medical", "clinical trial", "genomic"]
                for topic in noise_topics:
                    check(f"Research_Insights does not contain noise topic '{topic}'",
                          topic not in all_text,
                          f"Found '{topic}' in Research_Insights")
    except Exception as e:
        check("Reverse validation", False, str(e))
    finally:
        cur.close()
        conn.close()


def check_scripts(workspace):
    print("\n=== Check 4: Python Scripts ===")
    check("analyze_sales_gaps.py exists",
          os.path.exists(os.path.join(workspace, "analyze_sales_gaps.py")))
    check("match_recommendations.py exists",
          os.path.exists(os.path.join(workspace, "match_recommendations.py")))
    check("generate_summary.py exists",
          os.path.exists(os.path.join(workspace, "generate_summary.py")))


def check_outputs(workspace):
    print("\n=== Check 5: Output Files ===")

    # sales_gaps.json
    sg_path = os.path.join(workspace, "sales_gaps.json")
    if os.path.exists(sg_path):
        check("sales_gaps.json exists", True)
        try:
            with open(sg_path) as f:
                data = json.load(f)
            check("sales_gaps.json is valid JSON", True)
            text = json.dumps(data).lower()
            check("sales_gaps contains priority tags",
                  "priority" in text or "focus" in text,
                  f"Keys: {list(data.keys()) if isinstance(data, dict) else 'list'}")
        except Exception as e:
            check("sales_gaps.json is valid JSON", False, str(e))
    else:
        check("sales_gaps.json exists", False)

    # research_recommendations.json
    rr_path = os.path.join(workspace, "research_recommendations.json")
    if os.path.exists(rr_path):
        check("research_recommendations.json exists", True)
        try:
            with open(rr_path) as f:
                data = json.load(f)
            check("research_recommendations.json is valid JSON", True)
            text = json.dumps(data).lower()
            check("Recommendations reference papers",
                  "territory" in text or "segmentation" in text or "optimization" in text,
                  f"Content: {text[:150]}")
        except Exception as e:
            check("research_recommendations.json is valid JSON", False, str(e))
    else:
        check("research_recommendations.json exists", False)

    # executive_summary.txt
    es_path = os.path.join(workspace, "executive_summary.txt")
    if os.path.exists(es_path):
        check("executive_summary.txt exists", True)
        with open(es_path) as f:
            content = f.read().lower()
        total_rev_str = f"{EXPECTED['total_rev']:.0f}"
        total_rev_comma = f"{EXPECTED['total_rev']:,.0f}"
        # Check for revenue mention or total revenue value (truncated to first 4 digits for flexibility)
        check("Summary mentions total revenue",
              "revenue" in content or total_rev_str[:4] in content or total_rev_comma[:5] in content,
              f"Content: {content[:150]}")
        check("Summary mentions priority regions",
              "priority" in content or "focus" in content or "underperform" in content,
              f"Content: {content[:150]}")
    else:
        check("executive_summary.txt exists", False)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)
    check_gcal()
    check_scripts(args.agent_workspace)
    check_outputs(args.agent_workspace)
    check_reverse_validation(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {"total_passed": PASS_COUNT, "total_checks": total, "accuracy": accuracy,
              "critical_failed": CRITICAL_FAILED}
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    # CRITICAL gate: любой провал критической проверки => FAIL до порога accuracy.
    if CRITICAL_FAILED:
        print(f"\nCRITICAL FAIL ({len(CRITICAL_FAILED)}): {CRITICAL_FAILED}")
        print("FAIL")
        sys.exit(1)

    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
