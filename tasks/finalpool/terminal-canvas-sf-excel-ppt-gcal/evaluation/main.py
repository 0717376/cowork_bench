"""Evaluation for terminal-canvas-sf-excel-ppt-gcal.
Checks:
1. Skills_Gap_Analysis.xlsx with 4 sheets and correct data
2. Skills_Gap_Presentation.pptx with 5+ slides
3. Google Calendar advisory board events
4. gap_analyzer.py script exists
"""
import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent", "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILURES = []

# 7 russified department names (db/zzz_clickhouse_after_init.sql) + EN originals.
DEPT_RU = ["инженерия", "финансы", "кадры", "операции", "ниокр", "продажи", "поддержка"]
DEPT_EN = ["engineering", "finance", "hr", "operations", "r&d", "sales", "support"]

# Industry benchmarks from gap_analysis_config.json (used to validate Gap_Matrix demand scores).
BENCHMARKS = [7.5, 6.8, 8.0, 5.5]


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:200]}")


def critical(name, condition, detail=""):
    """A semantic check that gates PASS/FAIL regardless of accuracy."""
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS][CRITICAL] {name}")
    else:
        FAIL_COUNT += 1
        CRITICAL_FAILURES.append(name)
        print(f"  [FAIL][CRITICAL] {name}: {str(detail)[:200]}")


def check_excel(workspace):
    print("\n=== Check 1: Skills_Gap_Analysis.xlsx ===")
    path = os.path.join(workspace, "Skills_Gap_Analysis.xlsx")
    if not os.path.exists(path):
        check("Excel file exists", False, f"Not found at {path}")
        return
    check("Excel file exists", True)

    wb = openpyxl.load_workbook(path)
    sheets = wb.sheetnames
    check("Has at least 4 sheets", len(sheets) >= 4, f"Found {len(sheets)}: {sheets}")

    sheets_lower = [s.lower() for s in sheets]

    # Curriculum_Coverage sheet
    cc_idx = next((i for i, s in enumerate(sheets_lower) if "curriculum" in s or "coverage" in s), 0)
    ws = wb[sheets[cc_idx]]
    rows = list(ws.iter_rows(values_only=True))
    data_rows = [r for r in rows[1:] if any(c for c in r)]
    check("Curriculum_Coverage has 4 course rows", len(data_rows) >= 4, f"Found {len(data_rows)}")

    all_text = " ".join(str(c) for r in rows for c in r if c).lower()
    check("Contains Applied Analytics course (RU/EN)",
          "applied analytics" in all_text or "analytics" in all_text or "аналитик" in all_text,
          f"Text sample: {all_text[:120]}")
    check("Contains Biochemistry course (RU/EN)",
          "biochemistry" in all_text or "bioinformatics" in all_text
          or "биохими" in all_text or "биоинформат" in all_text,
          f"Text sample: {all_text[:120]}")

    # Workforce_Profile sheet
    wp_idx = next((i for i, s in enumerate(sheets_lower) if "workforce" in s or "profile" in s), 1)
    if wp_idx < len(sheets):
        ws2 = wb[sheets[wp_idx]]
        rows2 = list(ws2.iter_rows(values_only=True))
        data_rows2 = [r for r in rows2[1:] if any(c for c in r)]
        check("Workforce_Profile has 7 department rows", len(data_rows2) >= 7, f"Found {len(data_rows2)}")
        all_text2 = " ".join(str(c) for r in rows2 for c in r if c).lower()
        check("Contains Engineering department (RU/EN)",
              "инженерия" in all_text2 or "engineering" in all_text2, f"Text: {all_text2[:120]}")

    # Gap_Matrix sheet
    gm_idx = next((i for i, s in enumerate(sheets_lower) if "gap" in s or "matrix" in s), 2)
    if gm_idx < len(sheets):
        ws3 = wb[sheets[gm_idx]]
        rows3 = list(ws3.iter_rows(values_only=True))
        data_rows3 = [r for r in rows3[1:] if any(c for c in r)]
        check("Gap_Matrix has 4 skill area rows", len(data_rows3) >= 4, f"Found {len(data_rows3)}")
        all_text3 = " ".join(str(c) for r in rows3 for c in r if c).lower()
        check("Contains Quantitative Analysis (RU/EN)",
              "quantitative" in all_text3 or "количественн" in all_text3, f"Text: {all_text3[:120]}")
        check("Contains priority levels (RU/EN)",
              any(p in all_text3 for p in ("critical", "high", "moderate",
                                           "критическ", "высок", "умеренн")),
              f"Text: {all_text3[:120]}")

    # Recommendations sheet
    rec_idx = next((i for i, s in enumerate(sheets_lower) if "recommend" in s), 3)
    if rec_idx < len(sheets):
        ws4 = wb[sheets[rec_idx]]
        rows4 = list(ws4.iter_rows(values_only=True))
        data_rows4 = [r for r in rows4[1:] if any(c for c in r)]
        check("Recommendations has at least 4 rows", len(data_rows4) >= 4, f"Found {len(data_rows4)}")


def check_pptx(workspace):
    print("\n=== Check 2: Skills_Gap_Presentation.pptx ===")
    path = os.path.join(workspace, "Skills_Gap_Presentation.pptx")
    if not os.path.exists(path):
        check("PPTX file exists", False, f"Not found at {path}")
        return
    check("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = list(prs.slides)
        check("Has at least 5 slides", len(slides) >= 5, f"Found {len(slides)} slides")

        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += " " + shape.text_frame.text
        all_text_lower = all_text.lower()
        check("Contains gap analysis content", "gap" in all_text_lower, f"Text: {all_text_lower[:150]}")
        check("Contains curriculum content", "curriculum" in all_text_lower or "course" in all_text_lower,
              f"Text: {all_text_lower[:150]}")
        check("Contains workforce content", "workforce" in all_text_lower or "department" in all_text_lower,
              f"Text: {all_text_lower[:150]}")
    except ImportError:
        check("python-pptx available", False, "python-pptx not installed")


def check_gcal():
    print("\n=== Check 3: Advisory Board Calendar Events ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("""
        SELECT summary, start_datetime, end_datetime
        FROM gcal.events
        WHERE lower(summary) LIKE '%%advisory%%'
        ORDER BY start_datetime
    """)
    events = cur.fetchall()
    check("At least 3 advisory board events", len(events) >= 3, f"Found {len(events)} events")

    if events:
        summaries = " ".join(str(e[0]) for e in events).lower()
        check("Events mention curriculum/review (RU/EN)",
              any(w in summaries for w in ("curriculum", "review", "учебн", "программ", "пересмотр")),
              f"Summaries: {summaries[:150]}")
        check("Events mention workforce/data/findings (RU/EN)",
              any(w in summaries for w in ("workforce", "data", "finding",
                                           "персонал", "данны", "вывод")),
              f"Summaries: {summaries[:150]}")
        check("Events mention recommendation/gap (RU/EN)",
              any(w in summaries for w in ("recommend", "gap", "рекомендац", "пробел")),
              f"Summaries: {summaries[:150]}")

    cur.close()
    conn.close()


def check_script(workspace):
    print("\n=== Check 4: gap_analyzer.py ===")
    path = os.path.join(workspace, "gap_analyzer.py")
    check("gap_analyzer.py exists", os.path.exists(path))


def check_reverse_validation(workspace):
    """Verify things that should NOT exist in output."""
    print("\n=== Reverse Validation ===")

    # Excel: no unexpected sheets beyond the 4 required
    path = os.path.join(workspace, "Skills_Gap_Analysis.xlsx")
    if os.path.isfile(path):
        wb = openpyxl.load_workbook(path)
        check("Excel has no more than 6 sheets", len(wb.sheetnames) <= 6,
              f"Found {len(wb.sheetnames)} sheets: {wb.sheetnames}")

    # GCal: no advisory events on weekends
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT COUNT(*) FROM gcal.events
            WHERE lower(summary) LIKE '%%advisory%%'
              AND EXTRACT(DOW FROM start_datetime) IN (0, 6)
        """)
        weekend_count = cur.fetchone()[0]
        check("No advisory events on weekends", weekend_count == 0,
              f"Found {weekend_count} weekend events")
        cur.close()
        conn.close()
    except Exception:
        pass


def _nums_in_row(row):
    """All numeric (int/float) cells in a row, including numeric-looking strings."""
    out = []
    for c in row:
        if isinstance(c, bool):
            continue
        if isinstance(c, (int, float)):
            out.append(float(c))
        elif isinstance(c, str):
            s = c.strip().replace(",", ".")
            try:
                out.append(float(s))
            except ValueError:
                pass
    return out


def _expected_priority(gap):
    if gap > 3.0:
        return ("critical", "критическ")
    if gap >= 1.5:
        return ("high", "высок")
    return ("moderate", "умеренн")


def critical_checks(workspace):
    """SEMANTIC critical checks. Any failure => sys.exit(1) before accuracy gate."""
    print("\n=== CRITICAL CHECKS ===")
    xlsx = os.path.join(workspace, "Skills_Gap_Analysis.xlsx")

    # --- C2: Workforce_Profile has >=7 russified department rows ---
    wf_ok = False
    try:
        if os.path.isfile(xlsx):
            wb = openpyxl.load_workbook(xlsx)
            sl = [s.lower() for s in wb.sheetnames]
            idx = next((i for i, s in enumerate(sl) if "workforce" in s or "profile" in s), 1)
            ws = wb[wb.sheetnames[idx]]
            rows = [r for r in ws.iter_rows(values_only=True) if any(c for c in r)]
            data = rows[1:]
            txt = " ".join(str(c) for r in rows for c in r if c).lower()
            ru_hits = sum(1 for d in DEPT_RU if d in txt)
            en_hits = sum(1 for d in DEPT_EN if d in txt)
            wf_ok = len(data) >= 7 and (ru_hits >= 5 or en_hits >= 5)
            critical("C2 Workforce_Profile >=7 dept rows w/ real dept names (RU/EN)", wf_ok,
                     f"rows={len(data)} ru_hits={ru_hits} en_hits={en_hits}")
        else:
            critical("C2 Workforce_Profile >=7 dept rows", False, "xlsx missing")
    except Exception as e:
        critical("C2 Workforce_Profile >=7 dept rows", False, f"err: {e}")

    # --- C1: Gap_Matrix 4 skill areas w/ numeric scores, demand=benchmark, gap=|curr-demand|, priority correct ---
    try:
        if os.path.isfile(xlsx):
            wb = openpyxl.load_workbook(xlsx)
            sl = [s.lower() for s in wb.sheetnames]
            idx = next((i for i, s in enumerate(sl) if "gap" in s or "matrix" in s), 2)
            ws = wb[wb.sheetnames[idx]]
            rows = [r for r in ws.iter_rows(values_only=True) if any(c for c in r)]
            data = rows[1:]
            # find rows with >=3 numbers (curriculum, demand, gap)
            scored = []
            for r in data:
                ns = _nums_in_row(r)
                if len(ns) >= 3:
                    scored.append((r, ns))
            ge4 = len(scored) >= 4
            # validate demand matches one of the benchmarks and gap=|curr-demand| for each scored row
            demand_ok = 0
            gap_ok = 0
            prio_ok = 0
            for r, ns in scored:
                rowtxt = " ".join(str(c) for c in r if c).lower()
                # Among the numbers, find a triple (curr, demand, gap) consistent w/ a benchmark.
                # Compare by POSITION (index), never by `is` identity: equal interned floats
                # (e.g. a row like 5.5, 5.5, 0) would otherwise be wrongly excluded.
                matched = False
                for di, demand in enumerate(ns):
                    if any(abs(demand - b) <= 0.1 for b in BENCHMARKS):
                        # find curr & gap among the remaining positions
                        for ci, curr in enumerate(ns):
                            if ci == di:
                                continue
                            exp_gap = abs(curr - demand)
                            if any(abs(ns[gi] - exp_gap) <= 0.15
                                   for gi in range(len(ns)) if gi != ci and gi != di):
                                matched = True
                                gap_ok += 1
                                ep_en, ep_ru = _expected_priority(exp_gap)
                                if ep_en in rowtxt or ep_ru in rowtxt:
                                    prio_ok += 1
                                break
                    if matched:
                        break
                if matched:
                    demand_ok += 1
            ok = ge4 and demand_ok >= 4 and gap_ok >= 4 and prio_ok >= 3
            critical("C1 Gap_Matrix 4 areas: demand=benchmark, gap=|curr-demand|, priority correct", ok,
                     f"scored={len(scored)} demand_ok={demand_ok} gap_ok={gap_ok} prio_ok={prio_ok}")
        else:
            critical("C1 Gap_Matrix semantic correctness", False, "xlsx missing")
    except Exception as e:
        critical("C1 Gap_Matrix semantic correctness", False, f"err: {e}")

    # --- C3: Curriculum_Coverage 4 course rows, numeric enrollment/assign/quiz, total=assign+quiz ---
    try:
        if os.path.isfile(xlsx):
            wb = openpyxl.load_workbook(xlsx)
            sl = [s.lower() for s in wb.sheetnames]
            idx = next((i for i, s in enumerate(sl) if "curriculum" in s or "coverage" in s), 0)
            ws = wb[wb.sheetnames[idx]]
            rows = [r for r in ws.iter_rows(values_only=True) if any(c for c in r)]
            data = rows[1:]
            four = len(data) >= 4
            consistent = 0
            for r in data[:4]:
                ns = [n for n in _nums_in_row(r)]
                # need at least 4 numbers: id, enrollment, assign, quiz, total(>=4 of the 5)
                # check that some pair sums to another value (assign+quiz=total)
                found = False
                # total = assignments + quizzes. Exclude the two addend POSITIONS, not values:
                # a valid row may have quizzes=0 so total==assignments (value collision).
                for i in range(len(ns)):
                    for j in range(i + 1, len(ns)):
                        if any(abs((ns[i] + ns[j]) - ns[k]) <= 0.001
                               for k in range(len(ns)) if k != i and k != j):
                            found = True
                if len(ns) >= 4 and found:
                    consistent += 1
            ok = four and consistent >= 4
            critical("C3 Curriculum_Coverage 4 rows w/ total_assessments=assignments+quizzes", ok,
                     f"rows={len(data)} consistent={consistent}")
        else:
            critical("C3 Curriculum_Coverage correctness", False, "xlsx missing")
    except Exception as e:
        critical("C3 Curriculum_Coverage correctness", False, f"err: {e}")

    # --- C4: Exactly 3 advisory events at the scheduled times, none on weekends ---
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT start_datetime, end_datetime
            FROM gcal.events
            WHERE lower(summary) LIKE '%%advisory%%'
            ORDER BY start_datetime
        """)
        ev = cur.fetchall()
        # expected (start, duration_hours)
        expected = [
            ("2026-03-16 10:00", 2.0),
            ("2026-03-18 14:00", 1.5),
            ("2026-03-20 09:00", 2.0),
        ]
        def matches(slot):
            for st, et in ev:
                if st is None:
                    continue
                if st.strftime("%Y-%m-%d %H:%M") == slot[0]:
                    if et is not None:
                        dur = (et - st).total_seconds() / 3600.0
                        if abs(dur - slot[1]) <= 0.1:
                            return True
            return False
        all_slots = all(matches(s) for s in expected)
        cur.execute("""
            SELECT COUNT(*) FROM gcal.events
            WHERE lower(summary) LIKE '%%advisory%%'
              AND EXTRACT(DOW FROM start_datetime) IN (0, 6)
        """)
        weekend = cur.fetchone()[0]
        ok = len(ev) == 3 and all_slots and weekend == 0
        critical("C4 Exactly 3 advisory events at Mon10/Wed14/Fri09 w/ correct durations, none on weekend",
                 ok, f"count={len(ev)} all_slots={all_slots} weekend={weekend}")
        cur.close()
        conn.close()
    except Exception as e:
        critical("C4 advisory events scheduling", False, f"err: {e}")

    # --- C5: gap_analyzer.py exists AND pptx >=5 slides covering core topics ---
    try:
        script_ok = os.path.isfile(os.path.join(workspace, "gap_analyzer.py"))
        pptx_ok = False
        pptx_path = os.path.join(workspace, "Skills_Gap_Presentation.pptx")
        if os.path.isfile(pptx_path):
            from pptx import Presentation
            prs = Presentation(pptx_path)
            slides = list(prs.slides)
            txt = ""
            for s in slides:
                for sh in s.shapes:
                    if sh.has_text_frame:
                        txt += " " + sh.text_frame.text
            t = txt.lower()
            has_curr = any(w in t for w in ("curriculum", "course", "учебн", "курс"))
            has_wf = any(w in t for w in ("workforce", "department", "персонал", "подразделен"))
            has_gap = "gap" in t or "пробел" in t
            has_rec = any(w in t for w in ("recommend", "рекоменд"))
            pptx_ok = len(slides) >= 5 and has_curr and has_wf and has_gap and has_rec
        critical("C5 gap_analyzer.py exists AND pptx >=5 slides cover curriculum/workforce/gap/recommendations",
                 script_ok and pptx_ok, f"script={script_ok} pptx_ok={pptx_ok}")
    except Exception as e:
        critical("C5 gap_analyzer.py + pptx coverage", False, f"err: {e}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)
    check_gcal()
    check_script(args.agent_workspace)
    check_reverse_validation(args.agent_workspace)
    critical_checks(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {"total_passed": PASS_COUNT, "total_checks": total, "accuracy": accuracy,
              "critical_failures": CRITICAL_FAILURES}
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if CRITICAL_FAILURES:
        print(f"\nFAIL: {len(CRITICAL_FAILURES)} critical check(s) failed: {CRITICAL_FAILURES}")
        sys.exit(1)

    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
