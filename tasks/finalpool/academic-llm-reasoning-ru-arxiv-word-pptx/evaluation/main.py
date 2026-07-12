"""
Evaluation for academic-llm-reasoning-ru-arxiv-word-pptx task.

Checks:
1. LLM_Reasoning_Review.docx: headings (RU+EN alt), summary table, required keywords
2. LLM_Reasoning_Slides.pptx: slide count, title (RU/EN alt), key terms
3. word_count.txt: exists with two lines containing word counts

Critical checks (any fail -> overall FAIL regardless of accuracy):
  - all 6 required heading groups found in Word doc
  - summary table has 5 data rows
  - all required keywords (chain-of-thought, GSM8K, MATH, RAP, process supervision)
    present in Word doc
  - slide count >= 7
"""
import os
import sys
import json
import re
from argparse import ArgumentParser
from datetime import datetime

from docx import Document
from pptx import Presentation


CRITICAL_CHECKS = {
    "All 6 heading groups found",
    "Summary table has 5 data rows",
    "All required keywords present in Word doc",
    "PPTX has 7+ slides",
}

FAILED_NAMES = []


def _check(name, condition, fail_msg=""):
    """Print PASS/FAIL line, return (passed, total) increment and track failures."""
    if condition:
        print(f"  PASS: {name}")
        return 1, 1
    else:
        print(f"  FAIL: {name}{(': ' + fail_msg) if fail_msg else ''}")
        FAILED_NAMES.append(name)
        return 0, 1


def check_word_doc(agent_workspace, gt_data):
    passed = 0
    total = 0
    filename = gt_data["review_doc"]["filename"]
    doc_path = os.path.join(agent_workspace, filename)

    p, t = _check(f"{filename} exists", os.path.exists(doc_path),
                  f"not found at {doc_path}")
    passed += p; total += t
    if not os.path.exists(doc_path):
        return passed, total

    doc = Document(doc_path)

    all_text = []
    headings_found = []
    for para in doc.paragraphs:
        all_text.append(para.text)
        if para.style and para.style.name and "Heading" in para.style.name:
            headings_found.append(para.text.strip())

    full_text = " ".join(all_text).lower()

    # Required headings — each entry is a list of acceptable alternatives.
    # Any alternative matching as substring (case-insensitive) of a heading
    # in the doc counts.
    required_heading_groups = gt_data["review_doc"]["required_headings"]
    all_groups_found = True
    for group in required_heading_groups:
        alts = group if isinstance(group, list) else [group]
        label = " / ".join(alts)
        found = any(
            any(alt.lower() in h.lower() for h in headings_found)
            for alt in alts
        )
        p, t = _check(f"Heading present: {label}", found,
                      f"none of {alts} found in headings={headings_found}")
        passed += p; total += t
        if not found:
            all_groups_found = False

    p, t = _check("All 6 heading groups found", all_groups_found,
                  "some heading groups missing (see above)")
    passed += p; total += t

    # Summary table
    tables = doc.tables
    expected_data_rows = gt_data["review_doc"]["required_table_rows"]
    if len(tables) == 0:
        p, t = _check(f"Summary table has {expected_data_rows} data rows",
                      False, "no tables in document")
        passed += p; total += t
        p, t = _check("All required table columns present", False,
                      "no tables to check columns")
        passed += p; total += t
    else:
        table = tables[0]
        data_rows = sum(
            1 for row in table.rows[1:]
            if " ".join(c.text.strip() for c in row.cells).strip()
        )
        p, t = _check(f"Summary table has {expected_data_rows} data rows",
                      data_rows >= expected_data_rows,
                      f"got {data_rows} data rows")
        passed += p; total += t

        header_cells = [c.text.strip().lower() for c in table.rows[0].cells]
        expected_cols = [c.lower() for c in gt_data["review_doc"]["required_table_columns"]]
        cols_found = sum(
            1 for ec in expected_cols if any(ec in hc for hc in header_cells)
        )
        p, t = _check("All required table columns present",
                      cols_found >= len(expected_cols),
                      f"only {cols_found}/{len(expected_cols)} cols, headers={header_cells}")
        passed += p; total += t

    # Required keywords
    required_mentions = gt_data["review_doc"]["required_mentions"]
    all_keywords_found = True
    for mention in required_mentions:
        ok = mention.lower() in full_text
        p, t = _check(f"Keyword in Word doc: {mention}", ok,
                      f"keyword '{mention}' missing")
        passed += p; total += t
        if not ok:
            all_keywords_found = False

    p, t = _check("All required keywords present in Word doc",
                  all_keywords_found,
                  "some keywords missing (see above)")
    passed += p; total += t

    return passed, total


def check_pptx(agent_workspace, gt_data):
    passed = 0
    total = 0
    filename = gt_data["slides"]["filename"]
    pptx_path = os.path.join(agent_workspace, filename)

    p, t = _check(f"{filename} exists", os.path.exists(pptx_path),
                  f"not found at {pptx_path}")
    passed += p; total += t
    if not os.path.exists(pptx_path):
        return passed, total

    prs = Presentation(pptx_path)

    slide_count = len(prs.slides)
    required_count = gt_data["slides"]["required_slide_count"]
    p, t = _check(f"PPTX has {required_count}+ slides",
                  slide_count >= required_count,
                  f"got {slide_count} slides")
    passed += p; total += t

    all_slide_text = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
        all_slide_text.append(" ".join(texts))

    full_slides_text = " ".join(all_slide_text).lower()
    first_slide_text = all_slide_text[0].lower() if all_slide_text else ""

    # Title — accept any of the alternatives, partial-match fallback.
    title_alts = gt_data["slides"].get("required_title_alternatives") or [
        gt_data["slides"].get("required_title", "")
    ]
    title_ok = False
    matched_alt = None
    for alt in title_alts:
        alt_lower = alt.lower()
        if alt_lower in first_slide_text or alt_lower in full_slides_text:
            title_ok = True
            matched_alt = alt
            break
        words = alt_lower.split()
        match_count = sum(1 for w in words if w in full_slides_text)
        if words and match_count >= len(words) - 1:
            title_ok = True
            matched_alt = f"{alt} (partial)"
            break
    p, t = _check(f"Title slide matches one of {title_alts}", title_ok,
                  f"first slide text: {first_slide_text[:100]}")
    if title_ok:
        print(f"    matched: {matched_alt}")
    passed += p; total += t

    for mention in gt_data["slides"]["required_mentions"]:
        ok = mention.lower() in full_slides_text
        p, t = _check(f"Keyword in slides: {mention}", ok,
                      f"keyword '{mention}' missing")
        passed += p; total += t

    return passed, total


def check_word_count(agent_workspace):
    passed = 0
    total = 0
    wc_path = os.path.join(agent_workspace, "word_count.txt")

    p, t = _check("word_count.txt exists", os.path.exists(wc_path),
                  f"not found at {wc_path}")
    passed += p; total += t
    if not os.path.exists(wc_path):
        return passed, total

    with open(wc_path, "r") as f:
        content = f.read().strip()
    lines = [line.strip() for line in content.split("\n") if line.strip()]

    p, t = _check("word_count.txt has >=2 lines", len(lines) >= 2,
                  f"got {len(lines)} lines")
    passed += p; total += t

    for line in lines:
        numbers = re.findall(r'\d+', line)
        ok = bool(numbers) and any(int(n) > 0 for n in numbers)
        p, t = _check(f"Line has positive word count: {line[:60]}", ok,
                      "no positive integer found in line")
        passed += p; total += t

    p, t = _check("Both filenames referenced",
                  "LLM_Reasoning_Review" in content
                  and "LLM_Reasoning_Slides" in content,
                  "one or both filenames missing")
    passed += p; total += t

    return passed, total


def main(args):
    gt_path = os.path.join(args.groundtruth_workspace, "expected_results.json")
    if not os.path.exists(gt_path):
        print(f"FAIL: expected_results.json not found at {gt_path}")
        sys.exit(1)

    with open(gt_path, "r") as f:
        gt_data = json.load(f)

    total_passed = 0
    total_checks = 0

    print("--- Check 1: Word Document (LLM_Reasoning_Review.docx) ---")
    p, t = check_word_doc(args.agent_workspace, gt_data)
    print(f"  Word Doc: {p}/{t} checks passed")
    total_passed += p; total_checks += t

    print("\n--- Check 2: PowerPoint (LLM_Reasoning_Slides.pptx) ---")
    p, t = check_pptx(args.agent_workspace, gt_data)
    print(f"  PowerPoint: {p}/{t} checks passed")
    total_passed += p; total_checks += t

    print("\n--- Check 3: Word Count File (word_count.txt) ---")
    p, t = check_word_count(args.agent_workspace)
    print(f"  Word Count: {p}/{t} checks passed")
    total_passed += p; total_checks += t

    if total_checks == 0:
        print("\nFAIL: No checks were performed.")
        accuracy = 0.0
    else:
        accuracy = total_passed / total_checks * 100
        print(f"\nOverall: {total_passed}/{total_checks} checks passed ({accuracy:.1f}%)")

    critical_failed = [n for n in FAILED_NAMES if n in CRITICAL_CHECKS]
    if critical_failed:
        print(f"CRITICAL FAILURES: {len(critical_failed)}")
        for n in critical_failed:
            print(f"  - {n}")

    result = {
        "total_passed": total_passed,
        "total_checks": total_checks,
        "accuracy": accuracy,
        "critical_failed": critical_failed,
        "timestamp": datetime.now().isoformat(),
    }
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)
        print(f"Report saved to {args.res_log_file}")

    if critical_failed:
        print("FAIL (critical check failed)")
        sys.exit(1)
    if accuracy >= 80:
        print("PASS")
        sys.exit(0)
    print("FAIL")
    sys.exit(1)


if __name__ == "__main__":
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    main(args)
