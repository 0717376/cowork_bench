"""Evaluation for terminal-clickhouse-moex-excel-ppt-gcal.
Checks:
1. Investment_Committee_Briefing.xlsx (4 sheets, correct data)
2. Committee_Briefing.pptx (6 slides)
3. Google Calendar briefing event
4. compute_growth.py and market_comparison.py scripts exist
5. briefing_notes.txt exists
6. market_comparison.json exists

Data layer:
- sf_data (ClickHouse swap): brand names stay English (lg/microsoft...),
  customer names + segments are russified centrally by db/zzz_clickhouse_after_init.sql,
  so expected realia values are queried LIVE, never hardcoded.
- moex (MOEX Finance swap): tickers SBER.ME/GAZP.ME/LKOH.ME/TCSG.ME/MGNT.ME,
  prices in moex.stock_prices.

CRITICAL_CHECKS (semantic, any fail => sys.exit(1) before accuracy gate):
- top brand sorted first by Q4 revenue desc AND its Q4 matches the warehouse value
- QoQ growth math recomputed from the sheet's own Q3/Q4 cells for the top brand
- top customer (live) appears and is the #1 row by Total_Spend
- Market_Context lists all 5 MOEX tickers with non-empty prices
- exactly one briefing event in the target week, ~2h, not overlapping any seeded conflict
"""
import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent", "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []

# Expected MOEX tickers the agent must report (from briefing_template.json / task.md).
MOEX_TICKERS = ["sber.me", "gazp.me", "lkoh.me", "tcsg.me", "mgnt.me"]

# Fallback values (only used if the live DB query fails). Brands stay English;
# the top customer is russified centrally (Ava Garcia -> Ева Григорьев).
_FALLBACK_SF = {
    "top_brand": "lg",
    "top_brand_q4_revenue": 0.0,
    "top_brand_q3_revenue": 0.0,
    "top_customer": "ева григорьев",
}


def _get_sf_expected():
    """Query sf_data live to compute expected top brand / top customer dynamically."""
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        try:
            # Q4 2025 revenue by brand (Oct-Dec 2025)
            cur.execute("""
                SELECT p."BRAND", SUM(o."TOTAL_AMOUNT") as q4_revenue
                FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
                JOIN sf_data."SALES_DW__PUBLIC__PRODUCTS" p ON o."PRODUCT_ID" = p."PRODUCT_ID"
                WHERE o."ORDER_DATE" >= '2025-10-01' AND o."ORDER_DATE" < '2026-01-01'
                GROUP BY p."BRAND"
                ORDER BY q4_revenue DESC
            """)
            q4_rows = cur.fetchall()
            top_brand = q4_rows[0][0].lower() if q4_rows else "lg"
            top_brand_q4 = float(q4_rows[0][1]) if q4_rows else 0.0

            # Q3 2025 revenue for the top brand (Jul-Sep 2025)
            cur.execute("""
                SELECT SUM(o."TOTAL_AMOUNT")
                FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
                JOIN sf_data."SALES_DW__PUBLIC__PRODUCTS" p ON o."PRODUCT_ID" = p."PRODUCT_ID"
                WHERE o."ORDER_DATE" >= '2025-07-01' AND o."ORDER_DATE" < '2025-10-01'
                  AND lower(p."BRAND") = %s
            """, (top_brand,))
            r = cur.fetchone()
            top_brand_q3 = float(r[0]) if r and r[0] is not None else 0.0

            # Top customer by total spend in 2025. Group at the customer-entity
            # level (one row per name+segment, matching task.md "имя, сегмент" and
            # the groundtruth Top_Customers sheet); name-only grouping would wrongly
            # merge distinct customers that happen to share a name.
            cur.execute("""
                SELECT c."CUSTOMER_NAME", SUM(o."TOTAL_AMOUNT") as total_spend
                FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
                JOIN sf_data."SALES_DW__PUBLIC__CUSTOMERS" c ON o."CUSTOMER_ID" = c."CUSTOMER_ID"
                WHERE o."ORDER_DATE" >= '2025-01-01' AND o."ORDER_DATE" < '2026-01-01'
                GROUP BY c."CUSTOMER_ID", c."CUSTOMER_NAME", c."SEGMENT"
                ORDER BY total_spend DESC
                LIMIT 1
            """)
            tc = cur.fetchone()
            top_customer = tc[0].lower() if tc else _FALLBACK_SF["top_customer"]

            return {
                "top_brand": top_brand,
                "top_brand_q4_revenue": top_brand_q4,
                "top_brand_q3_revenue": top_brand_q3,
                "top_customer": top_customer,
            }
        finally:
            cur.close()
            conn.close()
    except Exception:
        return _FALLBACK_SF


_SF_EXPECTED = _get_sf_expected()


def check(name, condition, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS]{' [CRIT]' if critical else ''} {name}")
    else:
        FAIL_COUNT += 1
        tag = " [CRIT]" if critical else ""
        print(f"  [FAIL]{tag} {name}: {str(detail)[:200]}")
        if critical:
            CRITICAL_FAILS.append(name)


def num_close(a, b, tol=2.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except Exception:
        return False


def _to_float(v):
    try:
        return float(str(v).replace(",", ".").replace("%", "").strip())
    except Exception:
        return None


def check_excel(workspace):
    print("\n=== Check 1: Investment_Committee_Briefing.xlsx ===")
    path = os.path.join(workspace, "Investment_Committee_Briefing.xlsx")
    if not os.path.exists(path):
        check("Excel file exists", False, f"Not found at {path}")
        return
    check("Excel file exists", True)

    wb = openpyxl.load_workbook(path)
    sheets = wb.sheetnames
    check("Has at least 4 sheets", len(sheets) >= 4, f"Found {len(sheets)}: {sheets}")

    sheets_lower = [s.lower() for s in sheets]

    # --- Revenue_By_Brand sheet ---
    rb_idx = next((i for i, s in enumerate(sheets_lower) if "revenue" in s or "brand" in s), 0)
    ws1 = wb[sheets[rb_idx]]
    rows1 = list(ws1.iter_rows(values_only=True))
    data1 = [r for r in rows1[1:] if any(c for c in r)]
    check("Revenue_By_Brand has brand rows", len(data1) >= 10, f"Found {len(data1)}")

    all_text1 = " ".join(str(c) for r in rows1 for c in r if c).lower()
    check("Contains LG brand", "lg" in all_text1, f"Text: {all_text1[:120]}")
    check("Contains Microsoft brand", "microsoft" in all_text1)
    check("Contains growth column", "growth" in all_text1 or "qoq" in all_text1
          or "рост" in all_text1 or "прирост" in all_text1,
          f"Headers: {rows1[0] if rows1 else 'none'}")

    expected_top_brand = _SF_EXPECTED["top_brand"]

    # CRITICAL: first row is the top brand sorted by Q4 desc, and its Q4 matches the warehouse.
    first_brand = str(data1[0][0]).lower() if data1 else ""
    check(f"[CRIT] First brand is {expected_top_brand} (highest Q4)",
          first_brand == expected_top_brand,
          f"First brand: {first_brand}", critical=True)

    top_row = data1[0] if data1 else None
    if top_row is not None and len(top_row) >= 4:
        exp_q4 = _SF_EXPECTED["top_brand_q4_revenue"]
        got_q4 = _to_float(top_row[2])
        tol_q4 = max(500.0, exp_q4 * 0.01)
        check(f"[CRIT] Top brand Q4 revenue ~{exp_q4:.0f}",
              got_q4 is not None and num_close(got_q4, exp_q4, tol=tol_q4),
              f"Got {top_row[2]}", critical=True)

        # CRITICAL: QoQ growth recomputed from the sheet's own Q3/Q4 cells.
        got_q3 = _to_float(top_row[1])
        got_growth = _to_float(top_row[3])
        if got_q3 and got_q3 != 0 and got_q4 is not None and got_growth is not None:
            expected_growth = (got_q4 - got_q3) / got_q3 * 100.0
            check("[CRIT] Top brand QoQ growth = (Q4-Q3)/Q3*100",
                  num_close(got_growth, expected_growth, tol=1.0),
                  f"Sheet growth {got_growth} vs recomputed {expected_growth:.2f}",
                  critical=True)
        else:
            check("[CRIT] Top brand QoQ growth = (Q4-Q3)/Q3*100", False,
                  f"Missing Q3/Q4/growth cells: {top_row[1:4]}", critical=True)
    else:
        check("[CRIT] Top brand Q4 revenue", False, "Top row malformed", critical=True)
        check("[CRIT] Top brand QoQ growth = (Q4-Q3)/Q3*100", False,
              "Top row malformed", critical=True)

    # --- Top_Customers sheet ---
    tc_idx = next((i for i, s in enumerate(sheets_lower) if "customer" in s or "клиент" in s), 1)
    if tc_idx < len(sheets):
        ws2 = wb[sheets[tc_idx]]
        rows2 = list(ws2.iter_rows(values_only=True))
        data2 = [r for r in rows2[1:] if any(c for c in r)]
        check("Top_Customers has 10 rows", len(data2) >= 10, f"Found {len(data2)}")
        all_text2 = " ".join(str(c) for r in rows2 for c in r if c).lower()
        expected_top_cust = _SF_EXPECTED["top_customer"]

        # CRITICAL: top customer (live) appears AND is the #1 row by spend.
        first_cust = str(data2[0][0]).lower() if data2 else ""
        check(f"[CRIT] #1 customer is {expected_top_cust} (top by spend)",
              expected_top_cust in all_text2 and expected_top_cust in first_cust,
              f"First row name: {first_cust}", critical=True)

        # Segment values are russified centrally (Корпоративный / Частные клиенты),
        # but accept English too in case the agent labels them differently.
        check("Contains segment info",
              "корпоратив" in all_text2 or "частн" in all_text2
              or "государствен" in all_text2 or "малый" in all_text2
              or "enterprise" in all_text2 or "consumer" in all_text2,
              f"Text: {all_text2[:160]}")

    # --- Market_Context sheet ---
    mc_idx = next((i for i, s in enumerate(sheets_lower) if "market" in s or "рынок" in s or "рыноч" in s), 2)
    if mc_idx < len(sheets):
        ws3 = wb[sheets[mc_idx]]
        rows3 = list(ws3.iter_rows(values_only=True))
        data3 = [r for r in rows3[1:] if any(c for c in r)]
        check("Market_Context has 5 rows", len(data3) >= 5, f"Found {len(data3)}")
        all_text3 = " ".join(str(c) for r in rows3 for c in r if c).lower()

        # CRITICAL: all 5 MOEX tickers present.
        missing = [t for t in MOEX_TICKERS if t not in all_text3]
        check("[CRIT] Contains all 5 MOEX tickers",
              not missing, f"Missing: {missing}", critical=True)

        # CRITICAL: each ticker row has a non-empty recent price.
        priced_ok = True
        bad = []
        for r in data3:
            sym = str(r[0]).lower() if r and r[0] else ""
            if sym in MOEX_TICKERS:
                price = _to_float(r[1]) if len(r) > 1 else None
                if price is None or price <= 0:
                    priced_ok = False
                    bad.append(sym)
        check("[CRIT] Each MOEX ticker has a non-empty Recent_Price",
              priced_ok, f"Missing/invalid prices for: {bad}", critical=True)

        # Trend indicator present (RU or EN).
        check("Contains trend indicator",
              any(k in all_text3 for k in
                  ("up", "down", "flat", "вверх", "вниз", "без изменен")),
              f"Text: {all_text3[:160]}")

        # Trend consistency with sign of Five_Day_Change_Pct (non-critical, soft).
        consistent = True
        for r in data3:
            if not r or len(r) < 4:
                continue
            sym = str(r[0]).lower()
            if sym not in MOEX_TICKERS:
                continue
            chg = _to_float(r[2])
            trend = str(r[3]).lower() if r[3] is not None else ""
            if chg is None:
                continue
            if chg > 0.5 and not ("up" in trend or "вверх" in trend):
                consistent = False
            elif chg < -0.5 and not ("down" in trend or "вниз" in trend):
                consistent = False
        check("Trend matches sign of 5-day change", consistent, "Trend/sign mismatch")

    # --- Executive_Summary sheet ---
    es_idx = next((i for i, s in enumerate(sheets_lower)
                   if "summary" in s or "executive" in s or "сводк" in s or "резюме" in s), 3)
    if es_idx < len(sheets):
        ws4 = wb[sheets[es_idx]]
        rows4 = list(ws4.iter_rows(values_only=True))
        data4 = [r for r in rows4[1:] if any(c for c in r)]
        check("Executive_Summary has at least 5 rows", len(data4) >= 5, f"Found {len(data4)}")
        all_text4 = " ".join(str(c) for r in rows4 for c in r if c).lower()
        check("Contains total Q4 revenue metric",
              "q4" in all_text4 and ("revenue" in all_text4 or "выручк" in all_text4),
              f"Text: {all_text4[:160]}")
        check("Contains growth metric",
              "growth" in all_text4 or "qoq" in all_text4
              or "рост" in all_text4 or "прирост" in all_text4)


def check_pptx(workspace):
    print("\n=== Check 2: Committee_Briefing.pptx ===")
    path = os.path.join(workspace, "Committee_Briefing.pptx")
    if not os.path.exists(path):
        check("PPTX file exists", False, f"Not found at {path}")
        return
    check("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = prs.slides
        check("Has 6 slides", len(slides) >= 6, f"Found {len(slides)}")

        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text + " "

        al = all_text.lower()
        check("Title slide mentions Q4", "q4" in al)
        check("Contains revenue content", "revenue" in al or "выручк" in al or "продаж" in al)
        check("Contains customer content", "customer" in al or "клиент" in al)
        check("Contains market content", "market" in al or "рынок" in al or "рыноч" in al)
        check("Contains next steps",
              "next" in al or "action" in al or "след" in al or "шаг" in al or "действ" in al)
    except ImportError:
        check("python-pptx available", False, "python-pptx not installed")


def _overlaps(a_start, a_end, b_start, b_end):
    return a_start < b_end and b_start < a_end


def check_gcal():
    print("\n=== Check 3: Calendar Briefing Event ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    # Match RU title 'Брифинг инвестиционного комитета за Q4' plus EN fallbacks.
    cur.execute("""
        SELECT summary, start_datetime, end_datetime, description
        FROM gcal.events
        WHERE lower(summary) LIKE '%%инвестиц%%комитет%%'
           OR lower(summary) LIKE '%%брифинг%%'
           OR lower(summary) LIKE '%%investment%%committee%%'
           OR lower(summary) LIKE '%%q4%%briefing%%'
           OR lower(summary) LIKE '%%committee%%briefing%%'
        ORDER BY start_datetime
    """)
    events = cur.fetchall()
    check("Briefing event exists", len(events) >= 1, f"Found {len(events)} matching events")

    # Target week Mar 9-13 2026: exactly one briefing event there.
    in_week = [e for e in events if e[1] and
               "2026-03-09" <= e[1].strftime("%Y-%m-%d") <= "2026-03-13"]
    check("[CRIT] Exactly one briefing event in target week (Mar 9-13)",
          len(in_week) == 1, f"Found {len(in_week)} in-week events", critical=True)

    if in_week:
        summary, start, end, desc = in_week[0]

        if start and end:
            duration = (end - start).total_seconds() / 3600
            check("[CRIT] Event is ~2 hours", 1.5 <= duration <= 2.5,
                  f"Duration: {duration} hours", critical=True)

            # CRITICAL: must not overlap any seeded conflict event in that week.
            cur.execute("""
                SELECT summary, start_datetime, end_datetime FROM gcal.events
                WHERE id <> (SELECT id FROM gcal.events WHERE start_datetime=%s AND summary=%s LIMIT 1)
                  AND start_datetime < %s AND end_datetime > %s
            """, (start, summary, end, start))
            conflicts = cur.fetchall()
            overlapping = [c for c in conflicts if c[1] and c[2]
                           and _overlaps(start, end, c[1], c[2])]
            check("[CRIT] Briefing does not overlap any other event",
                  not overlapping,
                  f"Overlaps: {[c[0] for c in overlapping]}", critical=True)

            # Soft structural: within business hours 8:00-17:00.
            check("Event within business hours (8:00-17:00 start)",
                  8 <= start.hour <= 17, f"Start hour: {start.hour}")

        if desc:
            dl = str(desc).lower()
            check("Description mentions briefing topic",
                  any(k in dl for k in ("sales", "market", "committee", "quarterly",
                                        "продаж", "рынок", "комитет", "квартал", "брифинг")),
                  f"Desc: {str(desc)[:100]}")

    cur.close()
    conn.close()


def check_scripts(workspace):
    print("\n=== Check 4: Python Scripts ===")
    check("compute_growth.py exists",
          os.path.exists(os.path.join(workspace, "compute_growth.py")))
    check("market_comparison.py exists",
          os.path.exists(os.path.join(workspace, "market_comparison.py")))


def check_outputs(workspace):
    print("\n=== Check 5: Additional Output Files ===")
    notes_path = os.path.join(workspace, "briefing_notes.txt")
    if os.path.exists(notes_path):
        check("briefing_notes.txt exists", True)
        with open(notes_path) as f:
            content = f.read().lower()
        check("Notes mention revenue",
              "revenue" in content or "sales" in content
              or "выручк" in content or "продаж" in content,
              f"Content: {content[:100]}")
        check("Notes mention market",
              "market" in content or "stock" in content
              or "рынок" in content or "рыноч" in content or "акци" in content,
              f"Content: {content[:100]}")
    else:
        check("briefing_notes.txt exists", False)

    mc_path = os.path.join(workspace, "market_comparison.json")
    if os.path.exists(mc_path):
        check("market_comparison.json exists", True)
        try:
            with open(mc_path) as f:
                data = json.load(f)
            check("market_comparison.json is valid JSON", True)
            text = json.dumps(data).lower()
            present = [t for t in MOEX_TICKERS if t in text]
            check("Contains MOEX ticker data (>=4 of 5)", len(present) >= 4,
                  f"Present: {present}")
        except Exception as e:
            check("market_comparison.json is valid JSON", False, str(e))
    else:
        check("market_comparison.json exists", False)


def check_reverse(workspace):
    print("\n=== Reverse Validation ===")
    path = os.path.join(workspace, "Investment_Committee_Briefing.xlsx")
    if os.path.isfile(path):
        wb = openpyxl.load_workbook(path)
        has_negative_revenue = False
        ws = wb[wb.sheetnames[0]]
        for row in ws.iter_rows(min_row=2, min_col=2, max_col=3, values_only=True):
            for cell in row:
                if isinstance(cell, (int, float)) and cell < 0:
                    has_negative_revenue = True
        check("No negative revenue values", not has_negative_revenue,
              "Found negative revenue")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)
    check_gcal()
    check_scripts(args.agent_workspace)
    check_outputs(args.agent_workspace)
    check_reverse(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {"total_passed": PASS_COUNT, "total_checks": total, "accuracy": accuracy}
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    # CRITICAL gate: any critical failure => FAIL regardless of accuracy.
    if CRITICAL_FAILS:
        print(f"\nFAIL: {len(CRITICAL_FAILS)} CRITICAL check(s) failed: {CRITICAL_FAILS}")
        sys.exit(1)

    sys.exit(0 if accuracy >= 70 else 1)


if __name__ == "__main__":
    main()
