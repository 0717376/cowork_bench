"""Evaluation script for moex-sf-market-sales-excel-ppt-gcal."""
import os
import argparse, json, os, sys
import openpyxl

# MOEX tickers seeded deterministically in moex.* schema.
MOEX_SYMBOLS = {"SBER.ME", "GAZP.ME", "LKOH.ME", "TCSG.ME", "MGNT.ME", "MTSS.ME"}

# Deterministic latest prices from moex.stock_info seed (currentPrice / regularMarketPrice).
MOEX_PRICES = {
    "SBER.ME": 133.30,
    "GAZP.ME": 198.00,
    "LKOH.ME": 3911.00,
    "TCSG.ME": 2013.00,
    "MGNT.ME": 4439.00,
    "MTSS.ME": 275.05,
}


def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent", "password": "camel"
}

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []


def check(name, condition, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT, CRITICAL_FAILS
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS]{' [CRITICAL]' if critical else ''} {name}")
    else:
        FAIL_COUNT += 1
        detail_str = str(detail)[:200] if detail else ""
        print(f"  [FAIL]{' [CRITICAL]' if critical else ''} {name}: {detail_str}")
        if critical:
            CRITICAL_FAILS.append(name)


def safe_float(val, default=None):
    try:
        if val is None: return default
        return float(str(val).replace(",", "").replace("%", "").replace("$", "").strip())
    except (ValueError, TypeError):
        return default


def get_conn():
    import psycopg2
    return psycopg2.connect(**DB_CONFIG)


def norm_sym(v):
    """Normalize a ticker-ish cell to uppercase token (strip spaces)."""
    if v is None:
        return ""
    return str(v).strip().upper()


def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT, CRITICAL_FAILS
    PASS_COUNT = 0
    FAIL_COUNT = 0
    CRITICAL_FAILS = []

    # ---------------------------------------------------------------
    # Structural + value checks for Market_Sales_Correlation.xlsx
    # ---------------------------------------------------------------
    excel_path = os.path.join(agent_workspace, "Market_Sales_Correlation.xlsx")
    check("Market_Sales_Correlation.xlsx exists", os.path.exists(excel_path))

    stock_symbols_found = set()
    corr_ok = False
    sales_consistent = False

    if os.path.exists(excel_path):
        wb = openpyxl.load_workbook(excel_path)
        gt_path = os.path.join(groundtruth_workspace, "Market_Sales_Correlation.xlsx")
        gt_wb = openpyxl.load_workbook(gt_path) if os.path.exists(gt_path) else None

        if gt_wb:
            for sheet_name in gt_wb.sheetnames:
                check(f"{sheet_name} sheet exists", sheet_name in wb.sheetnames)
                if sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    gt_ws = gt_wb[sheet_name]
                    gt_headers = [str(c.value).strip().lower() if c.value else "" for c in gt_ws[1]]
                    headers = [str(c.value).strip().lower() if c.value else "" for c in ws[1]]
                    for h in gt_headers:
                        if h:
                            check(f"{sheet_name} has {h} column", h in headers, f"headers: {headers[:10]}")
                    gt_rows = list(gt_ws.iter_rows(min_row=2, values_only=True))
                    data_rows = list(ws.iter_rows(min_row=2, values_only=True))
                    min_rows = max(1, len(gt_rows) - 2)
                    check(f"{sheet_name} has >= {min_rows} data rows", len(data_rows) >= min_rows, f"got {len(data_rows)}")

        # ----- Collect data needed for CRITICAL checks (read honestly) -----
        if "Stock_Overview" in wb.sheetnames:
            sov = wb["Stock_Overview"]
            sov_headers = [str(c.value).strip() if c.value else "" for c in sov[1]]
            hmap = {h.lower(): i for i, h in enumerate(sov_headers)}
            sym_i = hmap.get("symbol")
            price_i = hmap.get("latest_price")
            for row in sov.iter_rows(min_row=2, values_only=True):
                if sym_i is None or sym_i >= len(row):
                    continue
                sym = norm_sym(row[sym_i])
                if sym:
                    stock_symbols_found.add(sym)

        # Correlation_Matrix validity: square, diagonal==1, off-diag in [-1,1],
        # headers == MOEX symbols.
        if "Correlation_Matrix" in wb.sheetnames:
            cm = wb["Correlation_Matrix"]
            cm_rows = list(cm.iter_rows(values_only=True))
            if cm_rows and len(cm_rows) >= 2:
                col_syms = [norm_sym(c) for c in cm_rows[0][1:]]
                row_syms = [norm_sym(r[0]) for r in cm_rows[1:]]
                header_syms = set(s for s in col_syms if s)
                diag_ok = True
                rng_ok = True
                for ri, r in enumerate(cm_rows[1:]):
                    for ci, c in enumerate(r[1:]):
                        v = safe_float(c)
                        if v is None:
                            continue
                        if not (-1.001 <= v <= 1.001):
                            rng_ok = False
                        if ri < len(col_syms) and ci < len(row_syms):
                            if row_syms[ri] and col_syms[ci] and row_syms[ri] == col_syms[ci]:
                                if abs(v - 1.0) > 0.05:
                                    diag_ok = False
                corr_ok = (
                    diag_ok and rng_ok
                    and header_syms.issubset(MOEX_SYMBOLS)
                    and len(header_syms) >= 3
                    and set(s for s in row_syms if s) == header_syms
                )

        # Sales_Trends internal consistency: Avg_Order_Value ~= Total_Revenue/Order_Count
        if "Sales_Trends" in wb.sheetnames:
            st = wb["Sales_Trends"]
            st_headers = [str(c.value).strip().lower() if c.value else "" for c in st[1]]
            smap = {h: i for i, h in enumerate(st_headers)}
            rev_i = smap.get("total_revenue")
            cnt_i = smap.get("order_count")
            avg_i = smap.get("avg_order_value")
            consistent_months = 0
            total_months = 0
            if None not in (rev_i, cnt_i, avg_i):
                for row in st.iter_rows(min_row=2, values_only=True):
                    rev = safe_float(row[rev_i]) if rev_i < len(row) else None
                    cnt = safe_float(row[cnt_i]) if cnt_i < len(row) else None
                    avg = safe_float(row[avg_i]) if avg_i < len(row) else None
                    if rev is None or cnt in (None, 0) or avg is None:
                        continue
                    total_months += 1
                    if num_close(avg, rev / cnt, rel_tol=0.05, abs_tol=1.0):
                        consistent_months += 1
            sales_consistent = total_months >= 3 and consistent_months >= 3

    # ---------------------------------------------------------------
    # Market_Sales_Strategy.pptx (slide titles kept English in GT)
    # ---------------------------------------------------------------
    pptx_path = os.path.join(agent_workspace, "Market_Sales_Strategy.pptx")
    check("Market_Sales_Strategy.pptx exists", os.path.exists(pptx_path))
    if os.path.exists(pptx_path):
        from pptx import Presentation
        prs = Presentation(pptx_path)
        gt_pptx_path = os.path.join(groundtruth_workspace, "Market_Sales_Strategy.pptx")
        if os.path.exists(gt_pptx_path):
            gt_prs = Presentation(gt_pptx_path)
            gt_slide_count = len(gt_prs.slides)
            check(f"Market_Sales_Strategy.pptx has >= {gt_slide_count} slides",
                  len(prs.slides) >= gt_slide_count, f"got {len(prs.slides)} slides")
            gt_titles = [s.shapes.title.text.strip().lower() for s in gt_prs.slides if s.shapes.title]
            agent_titles = [s.shapes.title.text.strip().lower() for s in prs.slides if s.shapes.title]
            for gt in gt_titles:
                found = any(gt in at or at in gt for at in agent_titles)
                check(f"Market_Sales_Strategy.pptx has slide \"{gt[:40]}\"", found, f"agent titles: {agent_titles[:5]}")
        else:
            check("Market_Sales_Strategy.pptx has >= 3 slides", len(prs.slides) >= 3, f"got {len(prs.slides)} slides")

    # ---------------------------------------------------------------
    # Terminal artifacts
    # ---------------------------------------------------------------
    py_files = [f for f in os.listdir(agent_workspace) if f.endswith(".py")]
    check("Python analysis script exists", len(py_files) >= 1, f"found: {py_files}")

    corr_json_path = os.path.join(agent_workspace, "correlation_analysis.json")
    corr_json_ok = False
    if os.path.exists(corr_json_path):
        try:
            with open(corr_json_path) as f:
                cj = json.load(f)
            text = json.dumps(cj).upper()
            corr_json_ok = any(s in text for s in MOEX_SYMBOLS)
        except Exception:
            corr_json_ok = False
    check("correlation_analysis.json exists", os.path.exists(corr_json_path))

    # ---------------------------------------------------------------
    # Calendar checks (matched by DATE+TIME, NOT by english keyword)
    # ---------------------------------------------------------------
    cal_event_ok = False
    try:
        conn = get_conn()
        cur = conn.cursor()
        # Target deliverable: 2026-03-16 11:00 - 12:30 UTC, non-empty description.
        # Compare hour/minute in UTC explicitly (timestamptz is rendered in the
        # PG session timezone otherwise, which can shift the hour, e.g. MSK).
        cur.execute("""
            SELECT summary,
                   EXTRACT(HOUR   FROM start_datetime AT TIME ZONE 'UTC') AS start_hour,
                   EXTRACT(MINUTE FROM start_datetime AT TIME ZONE 'UTC') AS start_minute,
                   EXTRACT(HOUR   FROM end_datetime   AT TIME ZONE 'UTC') AS end_hour,
                   EXTRACT(MINUTE FROM end_datetime   AT TIME ZONE 'UTC') AS end_minute,
                   description, start_datetime, end_datetime
            FROM gcal.events
            WHERE (start_datetime AT TIME ZONE 'UTC')::date = DATE '2026-03-16'
        """)
        rows = cur.fetchall()
        for (summary, start_hour, start_minute, end_hour, end_minute,
             desc, start_dt, end_dt) in rows:
            if start_hour is None:
                continue
            start_ok = int(start_hour) == 11 and int(start_minute) == 0
            end_ok = (end_hour is not None
                      and int(end_hour) == 12 and int(end_minute) == 30)
            desc_ok = desc is not None and len(str(desc).strip()) > 0
            if start_ok and end_ok and desc_ok:
                cal_event_ok = True
                break
        check("Strategy review event on 2026-03-16 11:00-12:30 with description",
              cal_event_ok, f"events on date: {[(r[0], str(r[6]), str(r[7])) for r in rows]}")

        # Noise events (russified) must survive.
        cur.execute("SELECT COUNT(*) FROM gcal.events WHERE summary ILIKE '%стендап%' OR summary ILIKE '%перерыв%'")
        noise_events = cur.fetchone()[0]
        check("Noise events exist (not deleted by agent)", noise_events >= 1, f"noise events: {noise_events}")
        conn.close()
    except Exception as e:
        check("DB checks", False, str(e))

    # ---------------------------------------------------------------
    # CRITICAL CHECKS (semantic substance; any fail => hard FAIL)
    # ---------------------------------------------------------------
    # 1. Stock_Overview symbols are MOEX tickers (confirms the swapped MCP was used).
    check("CRITICAL: Stock_Overview symbols are MOEX tickers",
          len(stock_symbols_found) >= 3 and stock_symbols_found.issubset(MOEX_SYMBOLS),
          f"found: {sorted(stock_symbols_found)}", critical=True)

    # 2. Correlation_Matrix is a valid symmetric matrix over MOEX symbols.
    check("CRITICAL: Correlation_Matrix valid (diag=1, [-1,1], MOEX headers)",
          corr_ok, "matrix invalid or wrong symbols", critical=True)

    # 3. Sales_Trends rows internally consistent (Avg ~= Revenue/Count, >=3 months).
    check("CRITICAL: Sales_Trends consistent (Avg=Revenue/Count, >=3 months)",
          sales_consistent, "rows not internally consistent", critical=True)

    # 4. Calendar deliverable survives russification (date+time match).
    check("CRITICAL: Calendar event 2026-03-16 11:00-12:30 with description",
          cal_event_ok, "core calendar deliverable missing", critical=True)

    # 5. correlation_analysis.json produced by the terminal script, keyed by MOEX symbols.
    check("CRITICAL: correlation_analysis.json with MOEX symbols",
          corr_json_ok, "missing or no MOEX symbols in JSON", critical=True)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100) if total else 0.0
    msg = f"Passed {PASS_COUNT}/{total} checks ({accuracy:.1f}%)"
    if CRITICAL_FAILS:
        msg += f" | CRITICAL FAILS: {CRITICAL_FAILS}"
    return accuracy, msg


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    accuracy, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)

    total = PASS_COUNT + FAIL_COUNT
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump({"total_passed": PASS_COUNT, "total_checks": total,
                       "accuracy": accuracy, "critical_fails": CRITICAL_FAILS}, f, indent=2)

    # Any critical failure => hard FAIL regardless of accuracy.
    if CRITICAL_FAILS:
        print(f"FAIL: critical check(s) failed: {CRITICAL_FAILS}")
        sys.exit(1)

    sys.exit(0 if accuracy >= 70 else 1)


if __name__ == "__main__":
    main()
