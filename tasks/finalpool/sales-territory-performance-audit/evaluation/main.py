#!/usr/bin/env python3
"""Evaluation script for sales-territory-performance-audit task validation.

CRITICAL checks: any failure => overall FAIL regardless of accuracy.
Pass condition: no critical failure AND accuracy >= 70.

Data realia in sf_data are russified centrally; the scorecard rep/territory
names come from the read-only seed files (rep_assignments.csv with its
Customer_Region -> sf_data REGION mapping, quota_targets.xlsx), and the
numeric sales values are deterministic aggregates of the sf_data seed
(SUM of TOTAL_AMOUNT per region), so groundtruth comparisons below match
the agent output directly.
"""

from argparse import ArgumentParser
import json
import os
import sys

PASS_COUNT = 0
FAIL_COUNT = 0
FAILED_NAMES = []

# Underperforming reps that MUST receive coaching plans (worst % to Quota:
# John Martinez 93.3, David Anderson 85.7). Sourced from seed files, English.
UNDERPERFORMERS = ["John Martinez", "David Anderson"]

# Any failure of these => overall FAIL regardless of accuracy.
CRITICAL_CHECKS = {
    "CRITICAL: scorecard rep rows match groundtruth values",
    "CRITICAL: % to Quota internally consistent for >=3 reps",
    "CRITICAL: coaching_plans flags exactly the underperforming reps",
    "CRITICAL: report email sent to management",
    "CRITICAL: strategic review meeting scheduled on calendar",
}


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        FAILED_NAMES.append(name)
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def _norm_name(v):
    return str(v or "").strip().lower()


# ----------------------------------------------------------------------------
# XLSX structural + groundtruth checks
# ----------------------------------------------------------------------------

def _load_keyed_rows(ws):
    """Return {rep_name_lower: {col_lower: value}} keyed by 'Rep Name'.

    Header-aware: locates the header row containing 'Rep Name', then keeps only
    rows that carry a real rep name (skips title/blank/Summary/detail blocks).
    """
    rows = list(ws.iter_rows(values_only=True))
    header = None; hidx = None
    for i, r in enumerate(rows):
        cells = [str(c).strip().lower() if c is not None else "" for c in r]
        if "rep name" in cells:
            header = cells; hidx = i; break
    out = {}
    if header is None:
        return out, header
    known = {"john martinez", "sarah chen", "michael brown",
             "jennifer lee", "david anderson"}
    for r in rows[hidx + 1:]:
        if not any(c is not None for c in r):
            continue
        rec = {}
        for ci, col in enumerate(header):
            if col and ci < len(r):
                rec[col] = r[ci]
        name = _norm_name(rec.get("rep name"))
        if name in known:
            out[name] = rec
    return out, header


def check_xlsx_content(workspace, groundtruth_workspace="."):
    print("\n=== Check: XLSX files ===")
    import openpyxl
    for fname in ["performance_scorecard.xlsx", "coaching_plans.xlsx"]:
        xlsx_path = os.path.join(workspace, fname)
        if not os.path.isfile(xlsx_path):
            record(f"xlsx {fname} exists", False, "Not found")
            continue
        record(f"xlsx {fname} exists", True)
        try:
            wb = openpyxl.load_workbook(xlsx_path, data_only=True)
            for ws in wb.worksheets:
                rows = list(ws.iter_rows(values_only=True))
                # skip empty default sheets (e.g. an auto-created 'Sheet1')
                if not any(any(c is not None for c in r) for r in rows):
                    continue
                record(f"xlsx {fname} '{ws.title}' has data", len(rows) >= 2, f"{len(rows)} rows")

            # --- Groundtruth XLSX comparison (header-aware, keyed by Rep Name) ---
            gt_path = os.path.join(groundtruth_workspace, fname)
            if os.path.isfile(gt_path):
                gt_wb = openpyxl.load_workbook(gt_path, data_only=True)
                for gt_sname in gt_wb.sheetnames:
                    gt_ws = gt_wb[gt_sname]
                    a_ws = None
                    for asn in wb.sheetnames:
                        if asn.strip().lower() == gt_sname.strip().lower():
                            a_ws = wb[asn]; break
                    if a_ws is None:
                        record(f"GT {fname} sheet '{gt_sname}' exists in agent xlsx", False, f"Available: {wb.sheetnames}")
                        continue
                    gt_recs, _ = _load_keyed_rows(gt_ws)
                    a_recs, _ = _load_keyed_rows(a_ws)
                    # Agent must include every rep the groundtruth has.
                    record(f"GT {fname} '{gt_sname}' rep coverage",
                           set(gt_recs).issubset(set(a_recs)),
                           f"Expected reps {sorted(gt_recs)}, got {sorted(a_recs)}")
                    # Compare values only on reps present in BOTH, on GT's columns.
                    shared = sorted(set(gt_recs) & set(a_recs))
                    for rn in shared[:3]:
                        g, a = gt_recs[rn], a_recs[rn]
                        ok = True; bad = None
                        for col, gv in g.items():
                            if gv is None:
                                continue
                            av = a.get(col)
                            if isinstance(gv, (int, float)):
                                if not num_close(av, gv, max(abs(gv) * 0.1, 1.0)):
                                    ok = False; bad = (col, gv, av); break
                            else:
                                if not str_match(av, gv):
                                    ok = False; bad = (col, gv, av); break
                        record(f"GT {fname} '{gt_sname}' {rn} values", ok,
                               f"mismatch {bad}" if bad else "matched")
                gt_wb.close()

            wb.close()
        except Exception as e:
            record(f"xlsx {fname} readable", False, str(e))


# ----------------------------------------------------------------------------
# CRITICAL: scorecard rep values match groundtruth
# ----------------------------------------------------------------------------

def _load_scorecard_rows(path):
    """Return {rep_name_lower: {col_lower: value}} from a scorecard xlsx."""
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    ws = None
    for sn in wb.sheetnames:
        if sn.strip().lower() == "performance scorecard":
            ws = wb[sn]; break
    if ws is None:
        ws = wb.worksheets[0]
    rows = list(ws.iter_rows(values_only=True))
    # find header row containing 'Rep Name'
    header = None; hidx = None
    for i, r in enumerate(rows):
        cells = [str(c).strip().lower() if c is not None else "" for c in r]
        if "rep name" in cells:
            header = cells; hidx = i; break
    out = {}
    if header is None:
        wb.close(); return out
    for r in rows[hidx + 1:]:
        if not any(c is not None for c in r):
            continue
        rec = {}
        for ci, col in enumerate(header):
            if col and ci < len(r):
                rec[col] = r[ci]
        name = _norm_name(rec.get("rep name"))
        if name and name not in {"summary metrics", ""}:
            # skip non-rep summary rows (no territory)
            if rec.get("territory") is not None:
                out[name] = rec
    wb.close()
    return out


def check_scorecard_critical(workspace, groundtruth_workspace="."):
    print("\n=== CRITICAL Check: Performance Scorecard values ===")
    a_path = os.path.join(workspace, "performance_scorecard.xlsx")
    if not os.path.isfile(a_path):
        record("CRITICAL: scorecard rep rows match groundtruth values", False,
               "performance_scorecard.xlsx missing")
        record("CRITICAL: % to Quota internally consistent for >=3 reps", False,
               "performance_scorecard.xlsx missing")
        return

    try:
        a_rows = _load_scorecard_rows(a_path)
    except Exception as e:
        record("CRITICAL: scorecard rep rows match groundtruth values", False, str(e))
        record("CRITICAL: % to Quota internally consistent for >=3 reps", False, str(e))
        return

    # --- Groundtruth value match (YTD Sales + Quota within 1%) ---
    gt_path = os.path.join(groundtruth_workspace, "performance_scorecard.xlsx")
    if os.path.isfile(gt_path):
        try:
            gt_rows = _load_scorecard_rows(gt_path)
        except Exception as e:
            gt_rows = {}
            record("CRITICAL: scorecard rep rows match groundtruth values", False,
                   f"GT unreadable: {e}")
            gt_rows = None
        if gt_rows is not None:
            check_reps = ["john martinez", "sarah chen", "michael brown"]
            ok = True
            details = []
            for rn in check_reps:
                g = gt_rows.get(rn); a = a_rows.get(rn)
                if not g or not a:
                    ok = False; details.append(f"{rn}: missing (gt={bool(g)}, agent={bool(a)})"); continue
                for col in ("ytd sales", "quota"):
                    gv, av = g.get(col), a.get(col)
                    if not num_close(av, gv, max(abs(float(gv)) * 0.01, 1.0) if gv is not None else 1.0):
                        ok = False; details.append(f"{rn}.{col}: gt={gv} agent={av}")
            record("CRITICAL: scorecard rep rows match groundtruth values", ok,
                   "; ".join(details) or "all matched")
    else:
        # Groundtruth not mounted -> cannot verify source-derived values.
        record("CRITICAL: scorecard rep rows match groundtruth values", False,
               f"groundtruth scorecard not mounted at {gt_path}")

    # --- % to Quota internal consistency (YTD/Quota*100 ~ stated %) ---
    consistent = 0
    details = []
    for name, rec in a_rows.items():
        ytd = rec.get("ytd sales"); quota = rec.get("quota"); pct = rec.get("% to quota")
        try:
            ytd = float(ytd); quota = float(quota); pct = float(pct)
        except (TypeError, ValueError):
            continue
        if quota == 0:
            continue
        expected = ytd / quota * 100.0
        # allow tolerance of 1.5 pct points (adjusted-target vs quota basis + rounding)
        if abs(expected - pct) <= 1.5 or num_close(pct, expected, max(expected * 0.02, 1.0)):
            consistent += 1
        else:
            details.append(f"{name}: stated={pct} ytd/quota={expected:.1f}")
    record("CRITICAL: % to Quota internally consistent for >=3 reps", consistent >= 3,
           f"{consistent} consistent; mism: {details}")


# ----------------------------------------------------------------------------
# CRITICAL: coaching plans flag the underperformers
# ----------------------------------------------------------------------------

def check_coaching_critical(workspace):
    print("\n=== CRITICAL Check: Coaching Plans target reps ===")
    import openpyxl
    path = os.path.join(workspace, "coaching_plans.xlsx")
    if not os.path.isfile(path):
        record("CRITICAL: coaching_plans flags exactly the underperforming reps", False,
               "coaching_plans.xlsx missing")
        return
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
    except Exception as e:
        record("CRITICAL: coaching_plans flags exactly the underperforming reps", False, str(e))
        return

    ws = None
    for sn in wb.sheetnames:
        if sn.strip().lower() == "coaching plans":
            ws = wb[sn]; break
    if ws is None:
        ws = wb.worksheets[0]
    rows = list(ws.iter_rows(values_only=True))
    wb.close()

    # locate header row with 'Rep Name'
    hidx = None
    for i, r in enumerate(rows):
        cells = [str(c).strip().lower() if c is not None else "" for c in r]
        if "rep name" in cells:
            hidx = i; break

    planned = set()
    if hidx is not None:
        cells = [str(c).strip().lower() if c is not None else "" for c in rows[hidx]]
        name_ci = cells.index("rep name")
        for r in rows[hidx + 1:]:
            if name_ci < len(r) and r[name_ci] is not None:
                nm = _norm_name(r[name_ci])
                # stop at detail/section blocks that lack a recognized rep
                planned.add(nm)

    expected = {_norm_name(n) for n in UNDERPERFORMERS}
    # plan rows for the two underperformers must be present; arbitrary extra
    # top-performer plan rows are not allowed.
    table_reps = {p for p in planned if p in {
        "john martinez", "sarah chen", "michael brown",
        "jennifer lee", "david anderson"}}
    ok = expected.issubset(table_reps) and not (table_reps - expected)
    record("CRITICAL: coaching_plans flags exactly the underperforming reps", ok,
           f"expected={sorted(expected)}, found_rep_rows={sorted(table_reps)}")


# ----------------------------------------------------------------------------
# DOCX / PPTX structural
# ----------------------------------------------------------------------------

def check_docx_content(workspace):
    print("\n=== Check: DOCX files ===")
    from docx import Document
    path = os.path.join(workspace, "coaching_analysis.docx")
    if not os.path.isfile(path):
        record("docx coaching_analysis.docx exists", False, "Not found")
        return False
    record("docx coaching_analysis.docx exists", True)
    try:
        doc = Document(path)
        record("docx has content", len(doc.paragraphs) > 0, f"{len(doc.paragraphs)} paragraphs")
    except Exception as e:
        record("docx readable", False, str(e))
    return True


def check_pptx_content(workspace):
    print("\n=== Check: PPTX leadership_presentation.pptx ===")
    from pptx import Presentation
    path = os.path.join(workspace, "leadership_presentation.pptx")
    if not os.path.isfile(path):
        record("pptx leadership_presentation.pptx exists", False, "Not found")
        return False
    record("pptx leadership_presentation.pptx exists", True)
    try:
        prs = Presentation(path)
        record("pptx has slides", len(prs.slides) > 0, f"{len(prs.slides)} slides")
    except Exception as e:
        record("pptx readable", False, str(e))
    return True


# ----------------------------------------------------------------------------
# CRITICAL: phase-6 email + calendar deliverables (honestly queried)
# ----------------------------------------------------------------------------

def _pg():
    import psycopg2
    cfg = {
        "host": os.environ.get("PGHOST", "localhost"),
        "port": int(os.environ.get("PGPORT", "5432")),
        "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
        "user": os.environ.get("PGUSER", "eigent"),
        "password": os.environ.get("PGPASSWORD", "camel"),
    }
    return psycopg2.connect(**cfg)


def check_email_critical():
    print("\n=== CRITICAL Check: report email to management ===")
    rows = []
    try:
        conn = _pg(); cur = conn.cursor()
        cur.execute("SELECT subject, to_addr, body_text FROM email.messages")
        rows = cur.fetchall()
        cur.close(); conn.close()
    except Exception as e:
        record("CRITICAL: report email sent to management", False, f"DB error: {e}")
        return
    print(f"  [INFO] Found {len(rows)} emails.")
    # The report email: subject/body about the territory/sales audit report.
    # Accept RU or EN keywords (original text, lowercased — never normalized).
    found = False
    for subject, to_addr, body in rows:
        text = ((subject or "") + " " + (body or "")).lower()
        topical = any(k in text for k in (
            "territor", "территор", "sales", "продаж", "audit", "аудит",
            "scorecard", "performance", "эффективн"))
        reportish = any(k in text for k in (
            "report", "отчёт", "отчет", "recommend", "рекоменд",
            "review", "пересмотр", "scorecard", "audit", "аудит"))
        has_recipient = bool((to_addr is not None) and str(to_addr).strip()
                             and str(to_addr).strip() not in ("[]", "{}", "null"))
        if topical and reportish and has_recipient:
            found = True
            break
    record("CRITICAL: report email sent to management", found,
           "No territory/sales audit report email with a recipient found")


def check_calendar_critical():
    print("\n=== CRITICAL Check: strategic review meeting ===")
    rows = []
    try:
        conn = _pg(); cur = conn.cursor()
        cur.execute("SELECT summary, description FROM gcal.events")
        rows = cur.fetchall()
        cur.close(); conn.close()
    except Exception as e:
        record("CRITICAL: strategic review meeting scheduled on calendar", False, f"DB error: {e}")
        return
    print(f"  [INFO] Found {len(rows)} calendar events.")
    found = False
    for summary, description in rows:
        text = ((summary or "") + " " + (description or "")).lower()
        reviewish = any(k in text for k in (
            "review", "пересмотр", "обсужден", "совещан", "встреч", "meeting", "strateg", "стратег"))
        topical = any(k in text for k in (
            "territor", "территор", "sales", "продаж", "performance",
            "эффективн", "audit", "аудит", "recommend", "рекоменд"))
        if reviewish and topical:
            found = True
            break
    record("CRITICAL: strategic review meeting scheduled on calendar", found,
           "No strategic review meeting about the sales/territory audit found")


# ----------------------------------------------------------------------------
# Main
# ----------------------------------------------------------------------------

def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False)
    args = parser.parse_args()

    ws = args.agent_workspace
    if not os.path.isdir(ws):
        print(f"Agent workspace not found: {ws}")
        sys.exit(1)

    check_xlsx_content(ws, args.groundtruth_workspace)
    check_scorecard_critical(ws, args.groundtruth_workspace)
    check_coaching_critical(ws)
    check_docx_content(ws)
    check_pptx_content(ws)
    check_email_critical()
    check_calendar_critical()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    critical_failed = [n for n in FAILED_NAMES if n in CRITICAL_CHECKS]
    all_passed = (not critical_failed) and accuracy >= 70

    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")
    if critical_failed:
        print(f"  CRITICAL FAILURES ({len(critical_failed)}):")
        for n in critical_failed:
            print(f"    - {n}")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
        "critical_failed": critical_failed,
        "success": all_passed,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if all_passed:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
