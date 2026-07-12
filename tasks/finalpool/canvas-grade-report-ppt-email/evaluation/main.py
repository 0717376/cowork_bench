"""
Evaluation for canvas-grade-report-ppt-email task.

Checks:
1. PowerPoint Grade_Report.pptx with correct, per-slide grade data:
   - Slide 2 'Grade Distribution': per-category count AND 2-decimal average.
   - Slide 3 'Course Summary': total enrolled, number graded, overall average.
   - Slide 4 'Key Metrics': pass rate %, Distinction count, Fail count,
     largest-category name + count.
2. Email to instructors@university.edu (DB check) with required figures and
   exact subject.

Canvas data is the read-only English base fixture (course id=1 =
'Прикладная аналитика и алгоритмы (Осень 2013)'). Expected values are computed
live from PostgreSQL, never hardcoded.

CRITICAL_CHECKS (semantic): any failure => overall FAIL regardless of accuracy.
Otherwise pass threshold: accuracy >= 70%.
"""
import argparse
import json
import sys
import os
from pathlib import Path

PASS_COUNT = 0
FAIL_COUNT = 0
FAILED_NAMES = []

COURSE_ID = 1

EXACT_SUBJECT = "Grade Analysis Report for Applied Analytics and Algorithms"

# Semantic checks reflecting the task's analytical substance and core
# deliverables. Any failure forces an overall FAIL.
CRITICAL_CHECKS = {
    "Slide 'Grade Distribution': all categories with correct count+average",
    "Slide 'Course Summary': enrolled, graded and overall average correct",
    "Slide 'Key Metrics': pass rate, Distinction, Fail and largest category correct",
    "Email to instructors@university.edu has overall avg, pass rate and total enrolled",
    "Email subject equals the required exact string",
}


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        FAILED_NAMES.append(name)
        d = (str(detail)[:300]) if detail else ""
        print(f"  [FAIL] {name}: {d}")


def fmt2(x):
    """2-decimal string as it would naturally render (e.g. 87.37, 68.69)."""
    return f"{x:.2f}"


def fmt1(x):
    return f"{x:.1f}"


def get_expected_data():
    """Query PostgreSQL for expected course grade data."""
    import psycopg2

    conn = psycopg2.connect(
        host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="cowork_gym",
        user="eigent", password="camel"
    )
    cur = conn.cursor()

    cur.execute("SELECT name FROM canvas.courses WHERE id = %s", (COURSE_ID,))
    course_name = cur.fetchone()[0]

    cur.execute("""
        SELECT
            grades->>'current_grade' as grade,
            COUNT(*) as cnt,
            ROUND(AVG((grades->>'current_score')::float)::numeric, 2) as avg_score
        FROM canvas.enrollments
        WHERE course_id = %s AND type = 'StudentEnrollment'
          AND grades->>'current_grade' IS NOT NULL
        GROUP BY grades->>'current_grade'
        ORDER BY grades->>'current_grade'
    """, (COURSE_ID,))
    grade_dist = [(r[0], int(r[1]), float(r[2])) for r in cur.fetchall()]

    cur.execute("""
        SELECT COUNT(*) as total_enrolled,
               SUM(CASE WHEN grades->>'current_grade' IS NOT NULL THEN 1 ELSE 0 END) as graded,
               ROUND(AVG(CASE WHEN grades->>'current_score' IS NOT NULL
                   THEN (grades->>'current_score')::float END)::numeric, 2) as overall_avg
        FROM canvas.enrollments
        WHERE course_id = %s AND type = 'StudentEnrollment'
    """, (COURSE_ID,))
    row = cur.fetchone()
    summary = (int(row[0]), int(row[1]), float(row[2]))

    conn.close()
    return course_name, grade_dist, summary


def _slide_texts(prs):
    """Return list of per-slide text blobs (joined paragraphs)."""
    out = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    texts.append(p.text)
        out.append("\n".join(texts))
    return out


def _find_slide(slide_texts, title_kw):
    """Return text of the slide whose text contains title_kw (lower), else None."""
    kw = title_kw.lower()
    for t in slide_texts:
        if kw in t.lower():
            return t
    return None


def check_pptx(workspace, course_name, grade_dist, summary):
    """Check Grade_Report.pptx for correctness, per-slide / per-category."""
    print("\n=== Checking PowerPoint ===")
    pptx_path = Path(workspace) / "Grade_Report.pptx"

    if not pptx_path.exists():
        check("PPTX file exists", False, f"Not found: {pptx_path}")
        # Critical checks depending on the deck cannot pass.
        check("Slide 'Grade Distribution': all categories with correct count+average",
              False, "pptx missing")
        check("Slide 'Course Summary': enrolled, graded and overall average correct",
              False, "pptx missing")
        check("Slide 'Key Metrics': pass rate, Distinction, Fail and largest category correct",
              False, "pptx missing")
        return

    check("PPTX file exists", True)

    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    slides = _slide_texts(prs)

    # Structural (non-critical)
    check("Has at least 4 slides", len(slides) >= 4, f"Got {len(slides)}")

    # --- Title slide (structural, non-critical) ---
    title_txt = slides[0] if slides else ""
    check("Title slide has 'Grade Analysis'",
          "grade analysis" in title_txt.lower(),
          f"Text: {title_txt[:120]}")

    # --- Slide 2: Grade Distribution (CRITICAL, per-category) ---
    dist = _find_slide(slides, "Grade Distribution")
    if dist is None:
        check("Slide 'Grade Distribution' exists", False, "no slide titled Grade Distribution")
        check("Slide 'Grade Distribution': all categories with correct count+average",
              False, "slide missing")
    else:
        check("Slide 'Grade Distribution' exists", True)
        dist_ok = True
        for grade, cnt, avg_s in grade_dist:
            cnt_ok = str(cnt) in dist
            avg_ok = fmt2(avg_s) in dist
            check(f"Grade '{grade}': count {cnt} on Distribution slide",
                  cnt_ok, f"'{cnt}' not on slide")
            check(f"Grade '{grade}': avg {fmt2(avg_s)} on Distribution slide",
                  avg_ok, f"'{fmt2(avg_s)}' not on slide")
            dist_ok = dist_ok and cnt_ok and avg_ok
        check("Slide 'Grade Distribution': all categories with correct count+average",
              dist_ok, "one or more category count/avg missing on the slide")

    total_enrolled, graded, overall_avg = summary

    # --- Slide 3: Course Summary (CRITICAL) ---
    csum = _find_slide(slides, "Course Summary")
    if csum is None:
        check("Slide 'Course Summary' exists", False, "no slide titled Course Summary")
        check("Slide 'Course Summary': enrolled, graded and overall average correct",
              False, "slide missing")
    else:
        check("Slide 'Course Summary' exists", True)
        enr_ok = str(total_enrolled) in csum
        grd_ok = str(graded) in csum
        avg_ok = fmt2(overall_avg) in csum
        check("Course Summary: total enrolled", enr_ok, f"Expected {total_enrolled}")
        check("Course Summary: number graded", grd_ok, f"Expected {graded}")
        check("Course Summary: overall average", avg_ok, f"Expected {fmt2(overall_avg)}")
        check("Slide 'Course Summary': enrolled, graded and overall average correct",
              enr_ok and grd_ok and avg_ok, "one or more summary value missing on the slide")

    # --- Slide 4: Key Metrics (CRITICAL) ---
    total_graded = sum(cnt for _, cnt, _ in grade_dist)
    pass_count = sum(cnt for g, cnt, _ in grade_dist if g != 'Fail')
    pass_rate = round(pass_count / total_graded * 100, 1) if total_graded > 0 else 0.0
    distinction_cnt = next((c for g, c, _ in grade_dist if g == 'Distinction'), 0)
    fail_cnt = next((c for g, c, _ in grade_dist if g == 'Fail'), 0)
    largest = max(grade_dist, key=lambda r: r[1]) if grade_dist else ("", 0, 0.0)
    largest_name, largest_cnt = largest[0], largest[1]

    km = _find_slide(slides, "Key Metrics")
    if km is None:
        check("Slide 'Key Metrics' exists", False, "no slide titled Key Metrics")
        check("Slide 'Key Metrics': pass rate, Distinction, Fail and largest category correct",
              False, "slide missing")
    else:
        check("Slide 'Key Metrics' exists", True)
        pr_ok = fmt1(pass_rate) in km
        dist_ok = str(distinction_cnt) in km
        fail_ok = str(fail_cnt) in km
        # largest category: name AND count both present on the slide
        large_ok = (largest_name in km) and (str(largest_cnt) in km)
        check("Key Metrics: pass rate %", pr_ok, f"Expected {fmt1(pass_rate)}")
        check("Key Metrics: Distinction count", dist_ok, f"Expected {distinction_cnt}")
        check("Key Metrics: Fail count", fail_ok, f"Expected {fail_cnt}")
        check("Key Metrics: largest category name+count",
              large_ok, f"Expected {largest_name} ({largest_cnt})")
        check("Slide 'Key Metrics': pass rate, Distinction, Fail and largest category correct",
              pr_ok and dist_ok and fail_ok and large_ok,
              "one or more key metric missing on the slide")


def check_email(course_name, grade_dist, summary):
    """Check email was sent (DB check)."""
    import psycopg2

    print("\n=== Checking Email ===")
    try:
        conn = psycopg2.connect(
            host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="cowork_gym",
            user="eigent", password="camel"
        )
        cur = conn.cursor()
    except Exception as e:
        check("DB connection", False, str(e))
        check("Email to instructors@university.edu has overall avg, pass rate and total enrolled",
              False, "db error")
        check("Email subject equals the required exact string", False, "db error")
        return

    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
    """)
    all_emails = cur.fetchall()
    cur.close()
    conn.close()

    target = "instructors@university.edu"
    found = None
    for subj, from_addr, to_addr, body in all_emails:
        to_str = str(to_addr or "").lower()
        if target in to_str:
            found = (subj, from_addr, to_addr, body)
            break

    # Structural (non-critical): an email to the target exists.
    check(f"Email sent to {target}", found is not None,
          f"Found {len(all_emails)} total emails")

    if not found:
        check("Email to instructors@university.edu has overall avg, pass rate and total enrolled",
              False, "no email to target")
        check("Email subject equals the required exact string", False, "no email to target")
        return

    subj, _, _, body = found
    body = body or ""
    subj = subj or ""

    # Exact subject (CRITICAL). Compare on stripped value.
    subj_ok = subj.strip() == EXACT_SUBJECT
    check("Email subject equals the required exact string",
          subj_ok, f"Subject: {subj[:120]!r}")

    # Required figures in body (CRITICAL): overall avg, pass rate %, total enrolled.
    total_enrolled, graded, overall_avg = summary
    total_graded = sum(cnt for _, cnt, _ in grade_dist)
    pass_count = sum(cnt for g, cnt, _ in grade_dist if g != 'Fail')
    pass_rate = round(pass_count / total_graded * 100, 1) if total_graded > 0 else 0.0

    avg_ok = fmt2(overall_avg) in body
    pr_ok = fmt1(pass_rate) in body
    enr_ok = str(total_enrolled) in body

    check("Email body: overall average score", avg_ok, f"Expected {fmt2(overall_avg)}")
    check("Email body: pass rate %", pr_ok, f"Expected {fmt1(pass_rate)}")
    check("Email body: total enrolled", enr_ok, f"Expected {total_enrolled}")
    check("Email to instructors@university.edu has overall avg, pass rate and total enrolled",
          avg_ok and pr_ok and enr_ok, "one or more required figure missing in body")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    print("Fetching expected data...")
    course_name, grade_dist, summary = get_expected_data()
    print(f"  Course: {course_name}")
    print(f"  Grade categories: {len(grade_dist)}")
    print(f"  Summary: enrolled={summary[0]}, graded={summary[1]}, avg={summary[2]}")

    check_pptx(args.agent_workspace, course_name, grade_dist, summary)
    check_email(course_name, grade_dist, summary)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100) if total else 0.0

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}, Failed: {FAIL_COUNT} ({accuracy:.1f}%)")

    critical_failed = [n for n in FAILED_NAMES if n in CRITICAL_CHECKS]
    if critical_failed:
        print(f"  CRITICAL FAILURES: {len(critical_failed)}")
        for n in critical_failed:
            print(f"    - {n}")

    if args.res_log_file:
        result = {
            "total_passed": PASS_COUNT,
            "total_checks": total,
            "accuracy": accuracy,
            "critical_failed": critical_failed,
        }
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if total == 0:
        print("  Overall: FAIL (no checks performed)")
        sys.exit(1)
    if critical_failed:
        print("  Overall: FAIL (critical check failed)")
        sys.exit(1)
    if accuracy >= 70:
        print("  Overall: PASS")
        sys.exit(0)
    print("  Overall: FAIL")
    sys.exit(1)


if __name__ == "__main__":
    main()
