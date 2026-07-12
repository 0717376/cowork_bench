"""Evaluation for fetch-sf-sales-territory-ppt-email-gcal (ClickHouse / russified DB)."""
import argparse
import os
import sys

import psycopg2

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=5432,
          dbname=os.environ.get("PGDATABASE", "cowork_gym"),
          user="eigent", password="camel")

# Region/segment names may be written in Russian (sourced from the russified
# ClickHouse DB + CRM JSON) OR in English. Each entry is a list of accepted
# alternatives; presence = at least one alternative appears in the text.
REGIONS = [
    ["азиатско-тихоокеанский", "asia pacific"],
    ["европа", "europe"],
    ["латинская америка", "latin america"],
    ["ближний восток", "middle east"],
    ["северная америка", "north america"],
]
SEGMENTS = [
    ["корпоративный", "enterprise"],
    ["малый и средний бизнес", "smb"],
    ["частные клиенты", "consumer"],
    ["государственный", "government"],
]

# Overall company quota attainment computed from russified DB revenue / total quota:
# total revenue 3,048,998.33 / company_total_quota 3,070,000 = 99.3%.
# Accept the 1-dp percent or its integer rounding.
OVERALL_ATTAINMENT = ["99.3", "99.2", "99.4", "99"]


def _has_any(text, alternatives):
    return any(a in text for a in alternatives)


def check_pptx(agent_workspace):
    """Returns (errors, critical_errors)."""
    errors = []
    critical = []
    path = os.path.join(agent_workspace, "Territory_Scorecard.pptx")
    if not os.path.exists(path):
        return (["Territory_Scorecard.pptx not found"],
                ["Territory_Scorecard.pptx not found"])
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = list(prs.slides)
        if len(slides) < 5:
            errors.append(f"Expected 5 slides, found {len(slides)}")

        def _shape_text(shape):
            # Text frames, table cells and grouped shapes (recursively).
            t = ""
            if shape.has_text_frame:
                t += shape.text_frame.text + "\n"
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    t += " | ".join(c.text for c in row.cells) + "\n"
            for sub in getattr(shape, "shapes", ()):
                t += _shape_text(sub)
            return t

        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                all_text += _shape_text(shape)
        all_lower = all_text.lower()

        # --- Structural (non-critical) ---
        # Title slide anchor: "Q1 2026" is a stable EN-or-RU anchor.
        if "territory performance" not in all_lower and "q1 2026" not in all_lower:
            errors.append("Title slide missing expected title text")
        if "pipeline" not in all_lower and "coverage" not in all_lower and "покрыт" not in all_lower and "воронк" not in all_lower:
            errors.append("Pipeline coverage not discussed")
        if "recommendation" not in all_lower and "рекоменда" not in all_lower:
            errors.append("Recommendations slide not found")

        # --- CRITICAL: all 5 regions present (RU or EN) ---
        for alts in REGIONS:
            if not _has_any(all_lower, alts):
                critical.append(f"Region missing from pptx: {alts}")

        # --- CRITICAL: all 4 segments present (RU or EN) ---
        for alts in SEGMENTS:
            if not _has_any(all_lower, alts):
                critical.append(f"Segment missing from pptx: {alts}")

        # --- CRITICAL: quota attainment anchors (correct revenue/quota join) ---
        if not any(v in all_text for v in ("103.7", "103.6", "103.8")):
            critical.append("Asia Pacific attainment ~103.7% not found")
        if not any(v in all_text for v in ("94.7", "94.6", "94.8")):
            critical.append("Latin America attainment ~94.7% not found")

        # --- CRITICAL: overall company quota attainment present ---
        if not _has_any(all_text, OVERALL_ATTAINMENT):
            critical.append("Overall company quota attainment (~99.3%) not found in pptx")

    except Exception as e:
        errors.append(f"Error reading PPTX: {e}")
        critical.append(f"Error reading PPTX: {e}")
    return errors, critical


def check_gcal():
    """Returns (errors, critical_errors)."""
    errors = []
    critical = []
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT summary, description, start_datetime, end_datetime FROM gcal.events
            WHERE start_datetime::date = '2026-03-28'
        """)
        rows = cur.fetchall()
        cur.close()
        conn.close()
        if not rows:
            errors.append("No GCal event found on 2026-03-28")
            critical.append("No GCal event found on 2026-03-28")
        else:
            summaries = [r[0].lower() if r[0] else "" for r in rows]
            if not any("territory" in s or "review" in s or "executive" in s
                       or "территор" in s or "обзор" in s for s in summaries):
                errors.append(f"No territory review event (found: {[r[0] for r in rows]})")
            # CRITICAL: description (ORIGINAL .lower()) contains the overall
            # company quota-attainment percentage figure.
            descs = [(r[1] or "").lower() for r in rows]
            if not any(_has_any(d, OVERALL_ATTAINMENT) for d in descs):
                critical.append("GCal event description missing overall quota attainment (~99.3%)")
    except Exception as e:
        errors.append(f"Error checking GCal: {e}")
        critical.append(f"Error checking GCal: {e}")
    return errors, critical


def check_email():
    """Returns (errors, critical_errors)."""
    errors = []
    critical = []
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT subject, body_text FROM email.messages
            WHERE to_addr::text ILIKE '%%executive_team@company.com%%'
            ORDER BY id DESC LIMIT 5
        """)
        rows = cur.fetchall()
        cur.close()
        conn.close()
        if not rows:
            errors.append("No email found to executive_team@company.com")
            critical.append("No email found to executive_team@company.com")
        else:
            subjects = [r[0].lower() if r[0] else "" for r in rows]
            if not any("territory" in s or "q1" in s or "performance" in s
                       or "территор" in s for s in subjects):
                errors.append(f"Email subject doesn't match (found: {[r[0] for r in rows]})")
            # CRITICAL: body (ORIGINAL .lower()) contains the overall quota
            # attainment figure AND top + bottom performing region names (RU-or-EN).
            bodies = [(r[1] or "").lower() for r in rows]
            if not any(_has_any(b, OVERALL_ATTAINMENT) for b in bodies):
                critical.append("Email body missing overall quota attainment (~99.3%)")
            # Top region by revenue: Europe; bottom: Latin America.
            top_alts = ["европа", "europe"]
            bottom_alts = ["латинская америка", "latin america"]
            if not any(_has_any(b, top_alts) for b in bodies):
                critical.append("Email body missing top-performing region (Europe)")
            if not any(_has_any(b, bottom_alts) for b in bodies):
                critical.append("Email body missing bottom-performing region (Latin America)")
    except Exception as e:
        errors.append(f"Error checking email: {e}")
        critical.append(f"Error checking email: {e}")
    return errors, critical


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    agent_ws = args.agent_workspace or os.path.join(os.path.dirname(__file__), "..", "groundtruth_workspace")

    all_errors = []
    all_critical = []

    print("  Checking PowerPoint...")
    errs, crit = check_pptx(agent_ws)
    all_errors.extend(errs)
    all_critical.extend(crit)
    if errs or crit:
        for e in (errs + crit)[:8]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    print("  Checking GCal event...")
    errs, crit = check_gcal()
    all_errors.extend(errs)
    all_critical.extend(crit)
    if errs or crit:
        for e in (errs + crit)[:4]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    print("  Checking email...")
    errs, crit = check_email()
    all_errors.extend(errs)
    all_critical.extend(crit)
    if errs or crit:
        for e in (errs + crit)[:4]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    # Total checks (for accuracy gate). Structural + critical are all counted.
    total_checks = 21  # 3 structural pptx + 5 regions + 4 segments + 3 numeric anchors
    #               + 2 gcal (event, attainment) + 4 email (subject, attainment, top, bottom)
    failed = len(all_errors) + len(all_critical)
    # De-dup: critical items already counted in all_critical; structural in all_errors.
    accuracy = max(0.0, 100.0 * (1 - failed / total_checks))

    # CRITICAL GATE: any critical failure => FAIL regardless of accuracy.
    if all_critical:
        print(f"\n=== RESULT: FAIL (critical checks failed: {len(all_critical)}) ===")
        for e in all_critical[:10]:
            print(f"  CRITICAL: {e}")
        sys.exit(1)

    print(f"\n  Accuracy: {accuracy:.1f}% ({failed} non-critical issues / {total_checks})")
    if accuracy >= 70:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)
    else:
        print(f"\n=== RESULT: FAIL (accuracy {accuracy:.1f}% < 70%) ===")
        for e in all_errors[:10]:
            print(f"  {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
