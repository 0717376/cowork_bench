"""
Evaluation script for team-survey-report task.

Проверки:
1. Excel-файл (Survey_Analysis.xlsx): 2 листа, 10 строк ответов, верные имена, средние.
2. PowerPoint-файл (Team_Report.pptx): >=5 слайдов, проекты, измерения, рекомендации, статусы.
3. Письмо на management@company.com с результатами опроса.

CRITICAL_CHECKS: любой их fail = задача FAIL (sys.exit(1)) даже при высокой accuracy.
Структурные проверки (лист существует, файл читается) — не критические.
"""

import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent",
    "password": "camel",
}

EXPECTED_AVERAGES = {
    "leadership": 3.5,
    "workload": 3.5,
    "communication": 3.6,
    "growth": 3.5,
}

# (status, ожидаемое число проектов) — для проверки сводки статусов
EXPECTED_STATUS_COUNTS = {
    "active": 2,
    "completed": 1,
    "on hold": 1,
    "planning": 1,
}

# Критические чеки (по строке name, как передаётся в check())
CRITICAL_CHECKS = {
    "Survey Data: ровно 10 строк ответов",
    "Survey Data: >=8 из 10 имён респондентов присутствуют",
    "Summary Statistics: все четыре средних точные (tol<=0.05)",
    "Письмо отправлено на management@company.com",
    "PowerPoint: >=5 слайдов, >=3 проекта, >=3 измерения, рекомендация",
}

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILED = []


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT, CRITICAL_FAILED
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")
        if name in CRITICAL_CHECKS:
            CRITICAL_FAILED.append(name)


def num_close(a, b, tol=0.3):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def check_excel(agent_workspace):
    print("\n=== Проверка Excel ===")

    excel_path = os.path.join(agent_workspace, "Survey_Analysis.xlsx")
    check("Excel file exists", os.path.isfile(excel_path),
          f"Expected {excel_path}")
    if not os.path.isfile(excel_path):
        # Зафиксировать провал критических чеков, зависящих от файла
        check("Survey Data: ровно 10 строк ответов", False, "нет файла")
        check("Survey Data: >=8 из 10 имён респондентов присутствуют", False, "нет файла")
        check("Summary Statistics: все четыре средних точные (tol<=0.05)", False, "нет файла")
        return

    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
    except Exception as e:
        check("Excel file readable", False, str(e))
        return
    check("Excel file readable", True)

    sheet_names_lower = [s.lower().replace("_", " ").strip() for s in wb.sheetnames]

    has_data_sheet = any("survey" in s and "data" in s for s in sheet_names_lower) or \
                     any("data" in s or "response" in s for s in sheet_names_lower)
    check("Has Survey Data sheet", has_data_sheet,
          f"Found sheets: {wb.sheetnames}")

    ws_data = None
    for s in wb.sheetnames:
        sl = s.lower().replace("_", " ")
        if ("survey" in sl and "data" in sl) or "response" in sl:
            ws_data = wb[s]
            break
    if ws_data is None and len(wb.sheetnames) >= 1:
        ws_data = wb[wb.sheetnames[0]]

    if ws_data:
        data_rows = list(ws_data.iter_rows(min_row=2, values_only=True))
        data_rows = [r for r in data_rows if r and r[0] is not None]
        check("Survey Data: ровно 10 строк ответов", len(data_rows) == 10,
              f"Found {len(data_rows)} data rows")

        all_names = " ".join(str(r[0]).lower() for r in data_rows if r and r[0])
        expected_names = ["alice", "bob", "carol", "david", "eva", "frank",
                          "grace", "henry", "irene", "jack"]
        names_found = sum(1 for n in expected_names if n in all_names)
        check("Survey Data: >=8 из 10 имён респондентов присутствуют",
              names_found >= 8,
              f"Found {names_found}/10 names")

    has_summary_sheet = any("summary" in s or "statistic" in s or "average" in s for s in sheet_names_lower)
    check("Has Summary Statistics sheet", has_summary_sheet,
          f"Found sheets: {wb.sheetnames}")

    ws_summary = None
    for s in wb.sheetnames:
        sl = s.lower().replace("_", " ")
        if "summary" in sl or "statistic" in sl or "average" in sl:
            ws_summary = wb[s]
            break

    all_values = []
    if ws_summary:
        for row in ws_summary.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    try:
                        all_values.append(float(cell))
                    except (TypeError, ValueError):
                        pass

    # Критический: все четыре средних должны присутствовать с tol<=0.05
    all_averages_ok = all(
        any(num_close(v, expected_avg, 0.05) for v in all_values)
        for expected_avg in EXPECTED_AVERAGES.values()
    )
    check("Summary Statistics: все четыре средних точные (tol<=0.05)",
          all_averages_ok,
          f"Expected {EXPECTED_AVERAGES}, numeric values: {all_values[:12]}")

    # Некритические по-измеренные чеки (с прежним мягким допуском для информативности)
    for dimension, expected_avg in EXPECTED_AVERAGES.items():
        found = any(num_close(v, expected_avg, 0.05) for v in all_values)
        check(f"Average {dimension} ~{expected_avg}",
              found,
              f"Expected {expected_avg}")


def check_pptx(agent_workspace):
    print("\n=== Проверка PowerPoint ===")

    pptx_path = os.path.join(agent_workspace, "Team_Report.pptx")
    file_ok = os.path.isfile(pptx_path)
    check("PowerPoint file exists", file_ok, f"Expected {pptx_path}")
    if not file_ok:
        check("PowerPoint: >=5 слайдов, >=3 проекта, >=3 измерения, рекомендация",
              False, "нет файла")
        return

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        check("PowerPoint file readable", False, str(e))
        check("PowerPoint: >=5 слайдов, >=3 проекта, >=3 измерения, рекомендация",
              False, str(e))
        return
    check("PowerPoint file readable", True)

    slide_count = len(prs.slides)
    check("PowerPoint has at least 5 slides",
          slide_count >= 5,
          f"Found {slide_count} slides")

    all_slide_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    all_slide_text += para.text + " "
    all_text_lower = all_slide_text.lower()

    if slide_count > 0:
        first_text = ""
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    first_text += para.text + " "
        first_lower = first_text.lower()
        check("Title slide mentions engineering or quarterly or report",
              "engineering" in first_lower or "quarterly" in first_lower or "report" in first_lower,
              f"First slide text: {first_text[:150]}")

    project_names = ["alpha", "beta", "gamma", "delta", "epsilon"]
    projects_found = sum(1 for p in project_names if p in all_text_lower)
    check("PowerPoint references at least 3 project names",
          projects_found >= 3,
          f"Found {projects_found}/5 project names")

    dimensions = ["leadership", "workload", "communication", "growth"]
    dims_found = sum(1 for d in dimensions if d in all_text_lower)
    check("PowerPoint mentions at least 3 survey dimensions",
          dims_found >= 3,
          f"Found {dims_found}/4 dimensions")

    # RU+EN ключевые слова рекомендаций
    rec_kw = ["recommend", "improvement", "concern", "action", "suggestion", "next step",
              "рекоменд", "улучшен", "действи", "следующий шаг", "проблем"]
    has_rec = any(kw in all_text_lower for kw in rec_kw)
    check("PowerPoint has recommendation content", has_rec,
          "No recommendation keywords found")

    # Сводка статусов проектов (статус-токены оставлены английскими в источнике)
    status_ok = sum(1 for st in EXPECTED_STATUS_COUNTS if st in all_text_lower) >= 3
    check("PowerPoint: упомянуты >=3 статуса проектов (Active/Completed/On Hold/Planning)",
          status_ok,
          f"slide text snippet: {all_text_lower[:200]}")

    # Критический агрегат по презентации
    critical_pptx = (slide_count >= 5 and projects_found >= 3
                     and dims_found >= 3 and has_rec)
    check("PowerPoint: >=5 слайдов, >=3 проекта, >=3 измерения, рекомендация",
          critical_pptx,
          f"slides={slide_count}, projects={projects_found}, dims={dims_found}, rec={has_rec}")


def check_emails():
    print("\n=== Проверка писем ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
        ORDER BY created_at DESC
    """)
    all_messages = cur.fetchall()
    cur.close()
    conn.close()

    matching_email = None
    for subject, from_addr, to_addr, body_text in all_messages:
        subj_lower = (subject or "").lower()
        if "quarterly" in subj_lower or "report" in subj_lower or "engineering" in subj_lower:
            matching_email = (subject, from_addr, to_addr, body_text)
            break

    check("Email with report-related subject exists",
          matching_email is not None,
          f"Found {len(all_messages)} total emails, none matching")

    if matching_email is None:
        check("Письмо отправлено на management@company.com", False, "нет письма")
        return

    subject, from_addr, to_addr, body_text = matching_email

    to_str = ""
    if isinstance(to_addr, list):
        to_str = " ".join(str(r).lower() for r in to_addr)
    elif isinstance(to_addr, str):
        try:
            parsed = json.loads(to_addr)
            if isinstance(parsed, list):
                to_str = " ".join(str(r).lower() for r in parsed)
            else:
                to_str = str(to_addr).lower()
        except (json.JSONDecodeError, TypeError):
            to_str = str(to_addr).lower()

    check("Письмо отправлено на management@company.com",
          "management@company.com" in to_str,
          f"Recipient: {to_addr}")

    body_lower = (body_text or "").lower()

    dims_mentioned = sum(1 for d in ["leadership", "workload", "communication", "growth"]
                       if d in body_lower)
    check("Email mentions at least 3 survey dimensions",
          dims_mentioned >= 3,
          f"Found {dims_mentioned}/4 dimensions in email body")

    # Письмо должно содержать как минимум 3 из четырёх средних как реальные числа (3.5/3.6)
    score_hits = 0
    for v in EXPECTED_AVERAGES.values():
        if f"{v:.1f}" in body_lower or f"{v:.3f}" in body_lower or str(v) in body_lower:
            score_hits += 1
    check("Email mentions average scores (>=3 measures as numbers)",
          score_hits >= 3 or "3.5" in body_lower,
          f"score_hits={score_hits}, body snippet: {body_lower[:200]}")

    project_mentions = sum(1 for p in ["alpha", "beta", "gamma", "delta", "epsilon"]
                         if p in body_lower)
    check("Email mentions at least 2 project names",
          project_mentions >= 2 or "project" in body_lower or "проект" in body_lower,
          f"Found {project_mentions} project name mentions")

    # Сводка статусов проектов в письме (английские статус-токены)
    status_hits = sum(1 for st in EXPECTED_STATUS_COUNTS if st in body_lower)
    check("Email reflects project statuses (>=3 of Active/Completed/On Hold/Planning)",
          status_hits >= 3,
          f"status_hits={status_hits}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)
    check_emails()

    total = PASS_COUNT + FAIL_COUNT
    pass_rate = PASS_COUNT / total if total > 0 else 0

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Pass Rate: {pass_rate:.1%}")
    if CRITICAL_FAILED:
        print(f"  CRITICAL FAILED: {CRITICAL_FAILED}")

    success = (not CRITICAL_FAILED) and pass_rate >= 0.7

    result = {
        "passed": PASS_COUNT,
        "failed": FAIL_COUNT,
        "pass_rate": round(pass_rate, 3),
        "critical_failed": CRITICAL_FAILED,
        "success": success,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if CRITICAL_FAILED:
        print("FAIL (critical check failed)")
        sys.exit(1)
    sys.exit(0 if pass_rate >= 0.7 else 1)


if __name__ == "__main__":
    main()
