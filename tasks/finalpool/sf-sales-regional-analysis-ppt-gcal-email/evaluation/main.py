"""Evaluation for sf-sales-regional-analysis-ppt-gcal-email (ClickHouse / RU).

Checks:
1. PPT file Regional_Sales_Review.pptx with at least 5 slides
   - Title slide with "Regional Sales Performance Review"
   - Names the top region (Europe / Европа) and bottom region (Latin America / Латинская Америка)
   - Contains correct top-region revenue figure (648/642/606 family) derived from ORDERS
   - Contains a recommendations / actions slide
2. GCal event "Regional Sales Review Meeting" ~14 days from launch_time (RU+EN title match)
3. Email to sales-leadership@company.example.com from reporting@company.example.com
   Subject: "Regional Sales Performance Review"; body highlights top/bottom region

CRITICAL checks (any failure => FAIL regardless of accuracy):
 - PPT exists
 - PPT names top region (Europe/Европа) AND bottom region (Latin America/Латинская Америка)
 - PPT contains correct top-region revenue figure (648/642/606)
 - Email sent to sales-leadership@company.example.com with correct subject
"""
import argparse
import json
import os
import sys
from datetime import datetime, timedelta, timezone

import psycopg2

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="cowork_gym", user="eigent", password="camel")

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_ERRORS = []


def check(name, condition, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        detail_str = f": {str(detail)[:200]}" if detail else ""
        tag = "CRITICAL" if critical else "FAIL"
        print(f"  [{tag}] {name}{detail_str}")
        if critical:
            CRITICAL_ERRORS.append(name)


def check_ppt(agent_workspace):
    print("\n=== Checking PPT File ===")
    ppt_path = os.path.join(agent_workspace, "Regional_Sales_Review.pptx")
    exists = os.path.isfile(ppt_path)
    check("Regional_Sales_Review.pptx exists", exists,
          f"Expected at {ppt_path}", critical=True)
    if not exists:
        return

    try:
        from pptx import Presentation
        prs = Presentation(ppt_path)
    except Exception as e:
        check("PPT file readable", False, str(e), critical=True)
        return

    check("PPT has at least 5 slides", len(prs.slides) >= 5,
          f"Found {len(prs.slides)} slides")

    all_text = " ".join(
        shape.text.lower()
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )

    check("PPT contains 'Regional Sales Performance Review'",
          "regional sales" in all_text and "review" in all_text,
          "Title not found")
    # Region names: agent reads Russian REGION values from the russified DB,
    # but may write either RU or EN -> accept both.
    top_ok = "europe" in all_text or "европ" in all_text
    bottom_ok = ("latin america" in all_text or "latin" in all_text
                 or "латинская америка" in all_text or "латин" in all_text)
    check("PPT names top region (Europe/Европа)", top_ok,
          "Top region not mentioned", critical=True)
    check("PPT names bottom region (Latin America/Латинская Америка)", bottom_ok,
          "Bottom region not mentioned", critical=True)
    check("PPT contains correct top-region revenue figure",
          "648" in all_text or "642" in all_text or "606" in all_text,
          "Revenue figures not found (proves DB was queried)", critical=True)
    check("PPT contains recommendations / actions slide",
          any(k in all_text for k in (
              "recommend", "action", "strategy",
              "рекомендац", "действ", "стратег")),
          "No recommendations found")


def check_gcal(launch_time_str=None):
    print("\n=== Checking Google Calendar ===")
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()

    # Title may be English ("Regional Sales Review Meeting") or Russian
    # ("Региональный обзор продаж") -> match RU+EN substrings.
    cur.execute("""
        SELECT summary, start_datetime, end_datetime, description
        FROM gcal.events
        WHERE (LOWER(summary) LIKE '%regional%' OR LOWER(summary) LIKE '%региональн%')
          AND (LOWER(summary) LIKE '%sales%' OR LOWER(summary) LIKE '%продаж%')
    """)
    events = cur.fetchall()
    check("Regional Sales Review Meeting event created", len(events) >= 1,
          f"Found {len(events)} matching events")

    if events and launch_time_str:
        try:
            launch_time = datetime.fromisoformat(launch_time_str)
            if launch_time.tzinfo is None:
                launch_time = launch_time.replace(tzinfo=timezone.utc)
            target_date = launch_time + timedelta(days=14)
            for event in events:
                event_start = event[1]
                if event_start.tzinfo is None:
                    event_start = event_start.replace(tzinfo=timezone.utc)
                diff_days = abs((event_start.date() - target_date.date()).days)
                if diff_days <= 3:
                    check("Review Meeting is ~14 days from launch", True)
                    break
            else:
                check("Review Meeting is ~14 days from launch", False,
                      f"Closest event at {events[0][1]}, expected ~{target_date.date()}")
        except Exception as e:
            check("Review Meeting date check", False, str(e))

    cur.close()
    conn.close()


def check_emails():
    print("\n=== Checking Emails ===")
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()

    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
    """)
    all_emails = cur.fetchall()
    conn.close()

    def parse_recipients(to_addr):
        if to_addr is None:
            return []
        if isinstance(to_addr, list):
            return [str(r).strip().lower() for r in to_addr]
        to_str = str(to_addr).strip()
        try:
            parsed = json.loads(to_str)
            if isinstance(parsed, list):
                return [str(r).strip().lower() for r in parsed]
            return [to_str.lower()]
        except (json.JSONDecodeError, TypeError):
            return [to_str.lower()]

    target = "sales-leadership@company.example.com"
    found = None
    for subj, from_addr, to_addr, body in all_emails:
        recipients = parse_recipients(to_addr)
        if target in recipients:
            found = (subj, from_addr, to_addr, body)
            break

    check("Email sent to sales-leadership@company.example.com", found is not None,
          critical=True)
    if found:
        subj, from_addr, to_addr, body = found
        check("Email from reporting@company.example.com",
              "reporting@company.example.com" in (from_addr or "").lower(),
              f"From: {from_addr}")
        check("Subject is 'Regional Sales Performance Review'",
              "regional sales" in (subj or "").lower() and "review" in (subj or "").lower(),
              f"Subject: {subj}", critical=True)
        # Body should highlight Europe (top) and Latin America (attention) — RU+EN.
        body_l = (body or "").lower()
        body_top = "europe" in body_l or "европ" in body_l
        body_bottom = ("latin" in body_l or "латин" in body_l)
        check("Email body highlights top region (Europe/Европа)", body_top,
              "Top region not in body")
        check("Email body highlights bottom region (Latin America/Латинская Америка)",
              body_bottom, "Bottom region not in body")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    print("=" * 70)
    print("SF SALES REGIONAL ANALYSIS PPT GCAL EMAIL - EVALUATION")
    print("=" * 70)

    check_ppt(args.agent_workspace)
    check_gcal(args.launch_time)
    check_emails()

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100.0) if total else 0.0

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Accuracy: {accuracy:.1f}%")

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump({
                "passed": PASS_COUNT,
                "failed": FAIL_COUNT,
                "accuracy": accuracy,
                "critical_errors": CRITICAL_ERRORS,
                "success": (not CRITICAL_ERRORS) and accuracy >= 70,
            }, f, indent=2)

    # CRITICAL gate: any critical failure => FAIL regardless of accuracy.
    if CRITICAL_ERRORS:
        print(f"\n=== CRITICAL FAILURES ({len(CRITICAL_ERRORS)}) ===")
        for e in CRITICAL_ERRORS:
            print(f"  [CRITICAL] {e}")
        print("\n=== RESULT: FAIL (critical check failed) ===")
        sys.exit(1)

    if accuracy >= 70:
        print("=== RESULT: PASS ===")
        sys.exit(0)
    else:
        print(f"=== RESULT: FAIL (accuracy {accuracy:.1f}% < 70) ===")
        sys.exit(1)


if __name__ == "__main__":
    main()
