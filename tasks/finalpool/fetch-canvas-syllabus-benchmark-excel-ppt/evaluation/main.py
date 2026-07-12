"""
Оценка задачи fetch-canvas-syllabus-benchmark-excel-ppt.

Проверки (с пересчётом эталона из источников, БЕЗ хардкода волатильных данных):
  1. Course_Benchmark_Analysis.xlsx: 3 листа, корректные значения.
  2. Academic_Benchmark_Presentation.pptx: титул + 7 слайдов по типам + summary.
  3. Письмо с бенчмарк-отчётом отправлено.

Источники эталона:
  - National Benchmarks: benchmarks.json (mock API на порту 30203; при недоступности
    читаем файл tmp/mock_pages/api/benchmarks.json напрямую).
  - Our Courses / Comparison: пересчёт из canvas.* вживую:
        Enrollment       = canvas.courses.total_students
        Assignment_Count = COUNT(canvas.assignments по course_id)
        Quiz_Count       = COUNT(canvas.quizzes по course_id)
        Course_Type      = название курса до '(' (split('(')[0].strip())
        Our_Avg_*        = среднее по курсам данного типа
        *_Diff           = Our_Avg - National_Avg

CRITICAL_CHECKS: любой провал критической проверки => sys.exit(1) до порога accuracy.
Порог: accuracy >= 70 И нет критических провалов => PASS.
"""

import argparse
import json
import os
import sys
import urllib.request

import openpyxl
import psycopg2

# Канвас-данные русифицированы; benchmarks.json приходит с АНГЛИЙСКИМИ названиями
# типов курсов, поэтому ключи бенчмарков приводим к тем же РУ-значениям, что и
# course_type, выводимый из русифицированных названий курсов canvas.
# .../cowork_gym/tasks/finalpool/<task>/evaluation/main.py -> .../cowork_gym/scripts
_COWORK_ROOT = os.path.abspath(__file__)
for _ in range(5):  # evaluation -> <task> -> finalpool -> tasks -> cowork_gym
    _COWORK_ROOT = os.path.dirname(_COWORK_ROOT)
_SCRIPTS_DIR = os.path.join(_COWORK_ROOT, "scripts")
if _SCRIPTS_DIR not in sys.path:
    sys.path.insert(0, _SCRIPTS_DIR)
try:
    from canvas_relabel_map import COURSE_SUBJECTS as _COURSE_SUBJECTS
except Exception:
    _COURSE_SUBJECTS = {}


_CT_RU = {v.strip().lower() for v in _COURSE_SUBJECTS.values()}
_CT_EN2RU = {k.replace("&amp;", "&").strip().lower(): v.strip().lower()
             for k, v in _COURSE_SUBJECTS.items()}


def norm_ct(name):
    """Course_Type -> канонический РУ-ключ в lower (EN и RU эквивалентны)."""
    if name is None:
        return None
    n = str(name).replace("&amp;", "&").strip().lower()
    return _CT_EN2RU.get(n, n)

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent",
    "password": "camel",
}

BENCHMARK_URL = "http://localhost:30203/api/benchmarks.json"

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAIL = False


def record(name, passed, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT, CRITICAL_FAIL
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS]{' [CRIT]' if critical else ''} {name}")
    else:
        FAIL_COUNT += 1
        if critical:
            CRITICAL_FAIL = True
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL]{' [CRIT]' if critical else ''} {name}{msg}")


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def str_contains(haystack, needle):
    if haystack is None or needle is None:
        return False
    return needle.strip().lower() in str(haystack).strip().lower()


def course_type_of(name):
    """Тип курса = часть названия до '(' (например 'Основы финансов')."""
    return str(name).split("(")[0].strip()


# ---------------------------------------------------------------------------
# Эталон из источников (НЕ хардкод)
# ---------------------------------------------------------------------------

def fetch_benchmarks():
    """Национальные бенчмарки из mock API; при недоступности — из файла."""
    data = None
    try:
        with urllib.request.urlopen(BENCHMARK_URL, timeout=5) as resp:
            data = json.loads(resp.read().decode("utf-8"))
    except Exception:
        task_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        fpath = os.path.join(task_root, "tmp", "mock_pages", "api", "benchmarks.json")
        with open(fpath, "r", encoding="utf-8") as f:
            data = json.load(f)
    out = {}
    for b in data["benchmarks"]:
        ct = norm_ct(b["course_type"])
        out[ct] = {
            "discipline": b["discipline"],
            "enrollment": float(b["national_avg_enrollment"]),
            "assignments": float(b["national_avg_assignments"]),
            "quizzes": float(b["national_avg_quizzes"]),
        }
    return out


def compute_our_from_canvas():
    """Пересчёт пер-курсовых и пер-типовых средних из canvas.* вживую."""
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT c.id, c.name, COALESCE(c.total_students, 0),
               (SELECT COUNT(*) FROM canvas.assignments a WHERE a.course_id = c.id),
               (SELECT COUNT(*) FROM canvas.quizzes q WHERE q.course_id = c.id)
        FROM canvas.courses c
        ORDER BY c.name
    """)
    rows = cur.fetchall()
    cur.close()
    conn.close()

    per_course = {}   # name(lower) -> dict
    groups = {}       # course_type(lower) -> {enr:[], asg:[], qz:[]}
    for cid, name, enr, asg, qz in rows:
        ct = course_type_of(name)
        # Курсы вне 7 типов COURSE_SUBJECTS (например id 9991) — вне эталона.
        if _CT_RU and norm_ct(ct) not in _CT_RU:
            continue
        per_course[str(name).strip().lower()] = {
            "course_type": ct,
            "enrollment": float(enr),
            "assignments": float(asg),
            "quizzes": float(qz),
        }
        g = groups.setdefault(norm_ct(ct), {"enr": [], "asg": [], "qz": []})
        g["enr"].append(float(enr))
        g["asg"].append(float(asg))
        g["qz"].append(float(qz))

    avgs = {}
    for ct, g in groups.items():
        n = len(g["enr"])
        avgs[ct] = {
            "enrollment": sum(g["enr"]) / n,
            "assignments": sum(g["asg"]) / n,
            "quizzes": sum(g["qz"]) / n,
            "count": n,
        }
    return per_course, avgs


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------

def find_sheet(wb, *needles, exclude=()):
    for name in wb.sheetnames:
        low = name.lower()
        if any(n in low for n in needles) and not any(e in low for e in exclude):
            return name
    return None


def sheet_grid(ws):
    rows = list(ws.iter_rows(values_only=True))
    if not rows:
        return [], []
    header = [str(h).strip().lower() if h is not None else "" for h in rows[0]]
    data = [r for r in rows[1:] if any(v not in (None, "") for v in r)]
    return header, data


def col_idx(header, *keys):
    for i, h in enumerate(header):
        if any(k == h for k in keys):
            return i
    for i, h in enumerate(header):
        if any(k in h for k in keys):
            return i
    return None


def check_excel(agent_workspace, benchmarks, per_course, avgs):
    print("\n=== Проверка Excel ===")
    fpath = os.path.join(agent_workspace, "Course_Benchmark_Analysis.xlsx")
    if not os.path.isfile(fpath):
        record("Excel-файл существует", False, f"Не найден: {fpath}", critical=True)
        return
    record("Excel-файл существует", True)

    try:
        wb = openpyxl.load_workbook(fpath, data_only=True)
    except Exception as e:
        record("Excel читается", False, str(e), critical=True)
        return
    record("Excel читается", True)

    # --- Sheet 1: National Benchmarks ---
    bench_sheet = find_sheet(wb, "benchmark", "national",
                             exclude=("comparison", "compare"))
    if not bench_sheet:
        record("Лист National Benchmarks существует", False, f"Листы: {wb.sheetnames}")
        record("National Benchmarks: значения == benchmarks.json", False,
               "лист отсутствует", critical=True)
    else:
        record("Лист National Benchmarks существует", True)
        header, data = sheet_grid(wb[bench_sheet])
        record("National Benchmarks: 7 строк данных", len(data) == 7,
               f"Найдено {len(data)}")
        ct_i = col_idx(header, "course_type")
        enr_i = col_idx(header, "national_avg_enrollment")
        asg_i = col_idx(header, "national_avg_assignments")
        qz_i = col_idx(header, "national_avg_quizzes")
        if None in (ct_i, enr_i, asg_i, qz_i):
            record("National Benchmarks: значения == benchmarks.json", False,
                   f"Не найдены столбцы; шапка: {header}", critical=True)
        else:
            matched = 0
            for row in data:
                ct = norm_ct(row[ct_i]) if row[ct_i] else None
                b = benchmarks.get(ct)
                if not b:
                    continue
                if (num_close(row[enr_i], b["enrollment"], 0.5) and
                        num_close(row[asg_i], b["assignments"], 0.5) and
                        num_close(row[qz_i], b["quizzes"], 0.5)):
                    matched += 1
            record("National Benchmarks: значения == benchmarks.json (все 7)",
                   matched == 7, f"Совпало {matched}/7", critical=True)

    # --- Sheet 2: Our Courses ---
    our_sheet = find_sheet(wb, "our", "course",
                           exclude=("benchmark", "comparison", "national"))
    if not our_sheet:
        record("Лист Our Courses существует", False, f"Листы: {wb.sheetnames}")
        record("Our Courses: 22 курса, группировка по типу верна", False,
               "лист отсутствует", critical=True)
    else:
        record("Лист Our Courses существует", True)
        header, data = sheet_grid(wb[our_sheet])
        record("Our Courses: 22 строки данных", len(data) == 22, f"Найдено {len(data)}")
        name_i = col_idx(header, "course_name")
        ct_i = col_idx(header, "course_type")
        enr_i = col_idx(header, "enrollment")
        asg_i = col_idx(header, "assignment_count", "assignment")
        qz_i = col_idx(header, "quiz_count", "quiz")
        if None in (name_i, ct_i, enr_i, asg_i, qz_i):
            record("Our Courses: 22 курса, значения и группировка верны", False,
                   f"Не найдены столбцы; шапка: {header}", critical=True)
        else:
            distinct_types = set()
            value_ok = 0
            value_total = 0
            for row in data:
                nm = str(row[name_i]).strip().lower() if row[name_i] else None
                if not nm:
                    continue
                distinct_types.add(norm_ct(course_type_of(row[name_i])))
                exp = per_course.get(nm)
                if not exp:
                    continue
                value_total += 1
                # Group correctness: Course_Type == name before '(' (EN/RU эквивалентны)
                ct_ok = (norm_ct(row[ct_i]) == norm_ct(exp["course_type"]))
                if (ct_ok and
                        num_close(row[enr_i], exp["enrollment"], 0.5) and
                        num_close(row[asg_i], exp["assignments"], 0.5) and
                        num_close(row[qz_i], exp["quizzes"], 0.5)):
                    value_ok += 1
            expected_types = set(avgs.keys())
            types_ok = distinct_types == expected_types
            record("Our Courses: 7 типов курсов, группировка по имени-до-скобки верна",
                   types_ok and len(data) == 22,
                   f"типы={sorted(distinct_types)} ожид.={sorted(expected_types)}",
                   critical=True)
            record("Our Courses: значения enrollment/assignment/quiz верны для всех курсов",
                   value_total == 22 and value_ok == 22,
                   f"верно {value_ok}/{value_total}", critical=True)

    # --- Sheet 3: Comparison ---
    comp_sheet = find_sheet(wb, "comparison", "compare")
    if not comp_sheet:
        record("Лист Comparison существует", False, f"Листы: {wb.sheetnames}")
        record("Comparison: Our_Avg_Enrollment и Enrollment_Diff верны (все 7)", False,
               "лист отсутствует", critical=True)
    else:
        record("Лист Comparison существует", True)
        header, data = sheet_grid(wb[comp_sheet])
        record("Comparison: 7 строк данных", len(data) == 7, f"Найдено {len(data)}")
        ct_i = col_idx(header, "course_type")
        oenr_i = col_idx(header, "our_avg_enrollment")
        ediff_i = col_idx(header, "enrollment_diff")
        oasg_i = col_idx(header, "our_avg_assignments")
        adiff_i = col_idx(header, "assignment_diff")
        oqz_i = col_idx(header, "our_avg_quizzes")
        qdiff_i = col_idx(header, "quiz_diff")
        if None in (ct_i, oenr_i, ediff_i):
            record("Comparison: Our_Avg_Enrollment и Enrollment_Diff верны (все 7)",
                   False, f"Не найдены столбцы; шапка: {header}", critical=True)
        else:
            enr_ok = 0
            diff_ok = 0
            extra_ok = 0   # вторичные средние/разницы (некритично)
            extra_total = 0
            for row in data:
                ct = norm_ct(row[ct_i]) if row[ct_i] else None
                a = avgs.get(ct)
                b = benchmarks.get(ct)
                if not a or not b:
                    continue
                if num_close(row[oenr_i], a["enrollment"], 1.0):
                    enr_ok += 1
                exp_diff = a["enrollment"] - b["enrollment"]
                if num_close(row[ediff_i], exp_diff, 1.0):
                    diff_ok += 1
                # вторичные (assignments/quizzes) — некритично
                for ci, exp_val in (
                        (oasg_i, a["assignments"]), (oqz_i, a["quizzes"]),
                        (adiff_i, a["assignments"] - b["assignments"]),
                        (qdiff_i, a["quizzes"] - b["quizzes"])):
                    if ci is not None:
                        extra_total += 1
                        if num_close(row[ci], exp_val, 0.2):
                            extra_ok += 1
            record("Comparison: Our_Avg_Enrollment верен (все 7 типов)",
                   enr_ok == 7, f"верно {enr_ok}/7", critical=True)
            record("Comparison: Enrollment_Diff == Our_Avg - National (все 7 типов)",
                   diff_ok == 7, f"верно {diff_ok}/7", critical=True)
            record("Comparison: вторичные средние/разницы (assign/quiz) корректны",
                   extra_total > 0 and extra_ok >= int(extra_total * 0.8),
                   f"верно {extra_ok}/{extra_total}")

    wb.close()


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def check_pptx(agent_workspace, avgs):
    print("\n=== Проверка PowerPoint ===")
    fpath = os.path.join(agent_workspace, "Academic_Benchmark_Presentation.pptx")
    if not os.path.isfile(fpath):
        record("PPT-файл существует", False, f"Не найден: {fpath}", critical=True)
        return
    record("PPT-файл существует", True)

    try:
        from pptx import Presentation
        prs = Presentation(fpath)
    except Exception as e:
        record("PPT читается", False, str(e), critical=True)
        return

    slide_count = len(prs.slides)
    record("PPT: >= 9 слайдов (титул + 7 типов + summary)", slide_count >= 9,
           f"Найдено {slide_count}")

    first_slide = prs.slides[0]
    title_text = ""
    for shape in first_slide.shapes:
        if shape.has_text_frame:
            title_text += shape.text_frame.text.lower()
    record("PPT: титульный слайд содержит benchmark/course",
           "benchmark" in title_text or "course" in title_text,
           f"Заголовок: {title_text[:100]}")

    # Названия типов курсов могут быть на РУ (русифицированный canvas) или на EN
    # (benchmarks.json). Для каждого из 7 типов принимаем любой из вариантов;
    # подстроки дискриминирующие — однозначно определяют свой тип.
    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += " " + shape.text_frame.text.lower()
    # (RU-подстрока, EN-подстрока) на каждый из 7 типов курсов
    type_keywords = [
        ("аналитик", "analytics"),
        ("биохими", "biochemistry"),
        ("креативн", "computing"),
        ("на основе данных", "data-driven"),
        ("эколог", "environmental"),
        ("финанс", "finance"),
        ("геополит", "governance"),
    ]
    found = [pair for pair in type_keywords
             if pair[0] in all_text or pair[1] in all_text]
    record("PPT: упомянуты все 7 типов курсов", len(found) == 7,
           f"Найдено {len(found)}/7: {[p[0] for p in found]}", critical=True)


# ---------------------------------------------------------------------------
# Email
# ---------------------------------------------------------------------------

def check_email(benchmarks, avgs):
    print("\n=== Проверка Email ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
        emails = cur.fetchall()
        cur.close()
        conn.close()
    except Exception as e:
        record("Email БД доступна", False, str(e), critical=True)
        return

    # Множество типов выше национального уровня по enrollment (пересчёт).
    above = set()
    for ct, a in avgs.items():
        b = benchmarks.get(ct)
        if b and a["enrollment"] > b["enrollment"]:
            above.add(ct)

    target = None
    for subject, from_addr, to_addr, body_text in emails:
        subj = (subject or "").lower()
        if "benchmark" in subj or "course" in subj:
            target = (subject, from_addr, to_addr, body_text)
            break

    if not target:
        record("Письмо с бенчмарк-отчётом отправлено", False,
               f"Найдено {len(emails)} писем, нет benchmark/course в теме", critical=True)
        return

    subject, from_addr, to_addr, body_text = target
    record("Письмо с бенчмарк-отчётом отправлено", True)

    from_ok = str_contains(from_addr, "academic") or str_contains(from_addr, "university")
    record("Письмо от академического адреса", from_ok, f"From: {from_addr}")

    to_ok = "dean" in str(to_addr).lower()
    record("Письмо адресовано dean", to_ok, f"To: {to_addr}", critical=True)

    body = (body_text or "").lower()
    # Тело на русском -> расширенный набор ключевых слов RU + EN.
    kw = ["enrollment", "benchmark", "бенчмарк", "above", "below",
          "превыша", "выше", "ниже", "численност", "набор"]
    body_kw_ok = any(k in body for k in kw)
    record("Письмо обсуждает бенчмарки (RU/EN ключевые слова)", body_kw_ok,
           f"Превью: {(body_text or '')[:200]}")

    # Семантика: письмо должно корректно называть типы, превышающие нац. бенчмарк.
    # above-множество (по данным) должно упоминаться; below пусто => не требуем.
    above_named = sum(1 for ct in above if ct in body)
    # Допускаем упоминание по краткой дискриминирующей подстроке. Ключи above —
    # русифицированные типы курсов; для каждого даём РУ-подстроку (на случай иной
    # формулировки в письме) с EN-вариантом как запасным.
    short_kw = {
        "прикладная аналитика и алгоритмы": ("аналитик", "analytics"),
        "биохимия и биоинформатика": ("биохими", "biochemistry"),
        "креативные вычисления и культура": ("креативн", "computing"),
        "проектирование на основе данных": ("на основе данных", "data-driven"),
        "экологическая экономика и этика": ("эколог", "environmental"),
        "основы финансов": ("финанс", "finance"),
        "глобальное управление и геополитика": ("геополит", "governance"),
    }
    above_named_kw = sum(
        1 for ct in above
        if any(sub in body for sub in short_kw.get(ct, ("\x00",)))
    )
    named = max(above_named, above_named_kw)
    # Все 7 типов выше нац. уровня => тело должно назвать хотя бы большинство.
    threshold = max(1, int(len(above) * 0.6)) if above else 0
    record("Письмо называет типы курсов выше нац. enrollment",
           len(above) == 0 or named >= threshold,
           f"above={len(above)}, названо={named}, порог={threshold}", critical=True)


# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    benchmarks = fetch_benchmarks()
    per_course, avgs = compute_our_from_canvas()

    check_excel(args.agent_workspace, benchmarks, per_course, avgs)
    check_pptx(args.agent_workspace, avgs)
    check_email(benchmarks, avgs)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = PASS_COUNT / total * 100 if total else 0.0
    print(f"\n=== ИТОГ ===")
    print(f"  Пройдено: {PASS_COUNT}/{total} ({accuracy:.1f}%)")
    print(f"  Критический провал: {CRITICAL_FAIL}")

    result = {"total_passed": PASS_COUNT, "total_checks": total,
              "accuracy": accuracy, "critical_fail": CRITICAL_FAIL}
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if CRITICAL_FAIL:
        print("FAIL (критическая проверка не пройдена)")
        sys.exit(1)

    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    print("FAIL")
    sys.exit(1)


if __name__ == "__main__":
    main()
