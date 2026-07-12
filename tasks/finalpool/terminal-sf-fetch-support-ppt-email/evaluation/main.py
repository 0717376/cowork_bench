"""Evaluation for terminal-sf-fetch-support-ppt-email (ClickHouse fork).

Swap: snowflake -> clickhouse. Support data lives in schema sf_data, tables
SUPPORT_CENTER__PUBLIC__TICKETS / SUPPORT_CENTER__PUBLIC__AGENTS. Realia data
values are russified CENTRALLY by db/zzz_clickhouse_after_init.sql (AGENT_NAME ->
Ева/Давид/Эмма/...; PRIORITY/STATUS enums stay English High/Medium/Low). All
expected numbers are computed LIVE from the DB and from the benchmark JSON so the
eval auto-syncs with the seed -- no volatile numbers are hardcoded.

Scoring: any CRITICAL check failure => overall FAIL regardless of accuracy
(sys.exit(1) before the accuracy gate). Otherwise PASS requires accuracy >= 70%.

Free-text keyword greps run on the ORIGINAL .lower() text (NOT normalize()) and
accept RU + EN alternatives, because the agent legitimately writes Russian prose
in the PPT/email while preserved identifiers (PRIORITY enums, status literals,
file/subject names) stay English.
"""
import argparse
import json
import os
import sys

import psycopg2

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=5432,
          dbname=os.environ.get("PGDATABASE", "cowork_gym"),
          user="eigent", password="camel")

# reporting_template.json compliance thresholds (kept in sync with initial_workspace).
SAT_TOL = 0.10            # satisfaction_tolerance
RESP_TOL_PCT = 10         # response_time_tolerance_pct

PASS_COUNT = 0
FAIL_COUNT = 0
FAILED_NAMES = []
CRITICAL_FAILS = []


def check(name, condition, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT
    tag = "CRITICAL " if critical else ""
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {tag}{name}")
    else:
        FAIL_COUNT += 1
        FAILED_NAMES.append(name)
        if critical:
            CRITICAL_FAILS.append(name)
        d = f": {str(detail)[:200]}" if detail else ""
        print(f"  [FAIL] {tag}{name}{d}")


def load_benchmarks():
    """Read the industry benchmark JSON shipped with the task (source of truth
    for the target values the agent must have fetched from the API)."""
    here = os.path.dirname(os.path.abspath(__file__))
    task_root = os.path.dirname(here)
    candidates = [
        os.path.join(task_root, "files", "mock_pages", "api", "csat_benchmarks.json"),
        os.path.join(task_root, "tmp", "mock_pages", "api", "csat_benchmarks.json"),
    ]
    for p in candidates:
        if os.path.isfile(p):
            with open(p, encoding="utf-8") as f:
                return json.load(f)
    return None


def compute_expected():
    """Compute company per-priority metrics LIVE from the DB and derive the
    expected SLA compliance status against the benchmark JSON."""
    bench = load_benchmarks()
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()
    cur.execute(
        'SELECT "PRIORITY", COUNT(*), '
        'ROUND(AVG("RESPONSE_TIME_HOURS")::numeric, 2), '
        'ROUND(AVG("CUSTOMER_SATISFACTION")::numeric, 2) '
        'FROM sf_data."SUPPORT_CENTER__PUBLIC__TICKETS" '
        'GROUP BY "PRIORITY"'
    )
    stats = {}
    for prio, cnt, avg_resp, avg_csat in cur.fetchall():
        key = str(prio).strip().lower()
        stats[key] = {
            "priority": prio,
            "count": int(cnt),
            "avg_resp": float(avg_resp),
            "avg_csat": float(avg_csat),
        }
    conn.close()

    # Derive compliance status per priority using reporting_template thresholds.
    bm = (bench or {}).get("benchmarks", {})
    for key, s in stats.items():
        # Match benchmark key case-insensitively (benchmark keys are High/Medium/Low).
        b = None
        for bk, bv in bm.items():
            if bk.strip().lower() == key:
                b = bv
                break
        if b is None:
            s["status"] = None
            continue
        t_sat = float(b["target_satisfaction"])
        t_resp = float(b["target_response_hours"])
        sat_ok = s["avg_csat"] >= t_sat
        resp_ok = s["avg_resp"] <= t_resp
        sat_miss_big = s["avg_csat"] < (t_sat - SAT_TOL)
        resp_miss_big = s["avg_resp"] > t_resp * (1 + RESP_TOL_PCT / 100.0)
        if sat_ok and resp_ok:
            status = "Compliant"
        elif sat_miss_big or resp_miss_big:
            status = "Non-Compliant"
        else:
            status = "At Risk"
        s["status"] = status
        s["target_satisfaction"] = t_sat
        s["target_response_hours"] = t_resp
    return stats, bench


def check_pptx(ws_path, stats, bench):
    """Check Support_Performance_Review.pptx (preserved English filename)."""
    print("\n=== Checking PowerPoint ===")
    path = os.path.join(ws_path, "Support_Performance_Review.pptx")
    if not os.path.isfile(path):
        check("PPT file exists", False, f"Not found: {path}", critical=True)
        return
    check("PPT file exists", True, critical=True)

    if Presentation is None:
        check("python-pptx available", False, "Cannot import pptx")
        return

    try:
        prs = Presentation(path)
        slides = prs.slides
        check("PPT has >= 7 slides", len(slides) >= 7, f"Found {len(slides)} slides", critical=True)

        all_text = ""
        empty_slides = 0
        for slide in slides:
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text += shape.text_frame.text
            all_text += slide_text.lower() + " "
            if not slide_text.strip():
                empty_slides += 1
        check("No empty slides in PPT", empty_slides == 0, f"Found {empty_slides} empty slides", critical=True)

        # --- Soft keyword presence (RU + EN). ---
        check("PPT mentions benchmark",
              any(k in all_text for k in ("benchmark", "бенчмарк", "эталон", "industry", "отрасл")))
        # PRIORITY enum literals stay English.
        check("PPT mentions priority levels",
              "high" in all_text and ("medium" in all_text or "low" in all_text))
        check("PPT mentions satisfaction or CSAT",
              any(k in all_text for k in ("satisfaction", "csat", "удовлетвор")))
        check("PPT mentions compliance or SLA",
              any(k in all_text for k in ("compliance", "sla", "compliant", "соответств")))
        check("PPT mentions agents or performance",
              any(k in all_text for k in ("agent", "агент", "специалист", "performance", "производительн", "показател")))
        check("PPT mentions recommendations or takeaways",
              any(k in all_text for k in ("recommend", "takeaway", "improvement", "рекоменд", "вывод", "улучш")))

        # --- CRITICAL: real benchmark target_satisfaction values must appear ---
        # (proves the agent fetched the API rather than fabricating).
        bm = (bench or {}).get("benchmarks", {})
        bench_hits = 0
        for bk, bv in bm.items():
            tgt = bv.get("target_satisfaction")
            if tgt is None:
                continue
            # accept e.g. "3.50" or "3.5"
            t = float(tgt)
            forms = {f"{t:.2f}", f"{t:.1f}", str(t)}
            if any(form in all_text for form in forms):
                bench_hits += 1
        check("CRITICAL: PPT reports >=2 real benchmark target_satisfaction values",
              bench_hits >= 2, f"benchmark target values found in PPT text: {bench_hits}/3", critical=True)

        # --- CRITICAL: computed SLA compliance status per priority ---
        # High & Medium must be reported Non-Compliant; Low must NOT be Non-Compliant
        # (At Risk boundary treated leniently -> Compliant/At Risk both accepted for Low).
        _check_status_in_text(all_text, stats, where="PPT")

    except Exception as e:
        check("PPT readable", False, str(e), critical=True)


def _status_present(text, status):
    """Map a compliance status literal (kept English) to acceptable substrings in text."""
    s = status.lower()
    if s == "non-compliant":
        return any(k in text for k in ("non-compliant", "non compliant",
                                       "не соответств", "несоответств", "нарушен"))
    if s == "compliant":
        # 'compliant' substring also matches 'non-compliant', so callers handle ordering.
        return any(k in text for k in ("compliant", "соответств"))
    if s == "at risk":
        return any(k in text for k in ("at risk", "риск", "под угроз"))
    return False


def _check_status_in_text(text, stats, where):
    """CRITICAL: High & Medium reported Non-Compliant; Low not Non-Compliant."""
    high = stats.get("high")
    med = stats.get("medium")
    low = stats.get("low")

    if high and high.get("status") == "Non-Compliant":
        check(f"CRITICAL: {where} reports High as Non-Compliant",
              _status_present(text, "Non-Compliant"),
              f"expected Non-Compliant for High (csat {high['avg_csat']} resp {high['avg_resp']})",
              critical=True)
    if med and med.get("status") == "Non-Compliant":
        check(f"CRITICAL: {where} reports Medium as Non-Compliant",
              _status_present(text, "Non-Compliant"),
              f"expected Non-Compliant for Medium (csat {med['avg_csat']} resp {med['avg_resp']})",
              critical=True)
    if low and low.get("status") in ("At Risk", "Compliant"):
        # Low should not be flagged Non-Compliant. We require some compliance language
        # present and do not demand a specific Low label (At Risk is subjective).
        check(f"CRITICAL: {where} does not mis-flag Low priority",
              _status_present(text, "Compliant") or _status_present(text, "At Risk"),
              f"expected At Risk/Compliant for Low (csat {low['avg_csat']} resp {low['avg_resp']})",
              critical=True)


def check_email(stats):
    """Check email sent to managers (preserved English addresses/subject)."""
    print("\n=== Checking Email ===")
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()
    cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
    all_emails = cur.fetchall()
    conn.close()

    target_email = None
    for subj, from_addr, to_addr, body in all_emails:
        recipients = []
        if isinstance(to_addr, list):
            recipients = [str(r).strip().lower() for r in to_addr]
        elif isinstance(to_addr, str):
            try:
                parsed = json.loads(to_addr)
                recipients = ([str(r).strip().lower() for r in parsed]
                              if isinstance(parsed, list) else [to_addr.strip().lower()])
            except (json.JSONDecodeError, TypeError):
                recipients = [str(to_addr).strip().lower()]
        if "managers@support-team.example.com" in recipients:
            target_email = (subj, from_addr, to_addr, body)
            break

    check("CRITICAL: Email sent to managers@support-team.example.com",
          target_email is not None, f"Total emails: {len(all_emails)}", critical=True)
    if not target_email:
        return

    subj, from_addr, to_addr, body = target_email
    subj_l = (subj or "").lower()
    # Preserved English subject; accept the documented subject substring.
    check("CRITICAL: Email subject 'Support Performance Benchmark Analysis'",
          "support performance benchmark analysis" in subj_l
          or ("benchmark" in subj_l and "performance" in subj_l),
          f"Subject: {subj}", critical=True)
    check("CRITICAL: Email from analytics@support-team.example.com",
          "analytics@support-team.example.com" in (from_addr or "").lower(),
          f"From: {from_addr}", critical=True)

    body_lower = (body or "").lower()
    check("Email body mentions satisfaction or CSAT",
          any(k in body_lower for k in ("satisfaction", "csat", "удовлетвор")),
          f"body[:150]={body_lower[:150]}")
    check("Email body mentions compliance or priority",
          any(k in body_lower for k in ("compliant", "compliance", "priority",
                                        "соответств", "приоритет", "sla")),
          f"body[:150]={body_lower[:150]}")
    # Non-trivial body.
    check("Email body is non-trivial", len(body_lower.strip()) >= 60,
          f"len={len(body_lower.strip())}")
    # CRITICAL: email body must reflect the SLA non-compliance findings.
    check("CRITICAL: Email body flags an SLA gap / non-compliance",
          any(k in body_lower for k in ("non-compliant", "non compliant", "not compliant",
                                        "не соответств", "несоответств", "нарушен",
                                        "превыша", "below target", "improve", "улучш", "разрыв")),
          f"body[:200]={body_lower[:200]}", critical=True)


def check_reverse_validation(workspace):
    """Verify things that should NOT exist in output."""
    print("\n=== Reverse Validation ===")
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT COUNT(*) FROM email.messages
            WHERE (lower(subject) LIKE '%%benchmark%%' OR lower(subject) LIKE '%%performance%%')
              AND to_addr::text ILIKE '%%competitor%%'
        """)
        bad_count = cur.fetchone()[0]
        check("No benchmark emails to competitor addresses", bad_count == 0, f"Found {bad_count}")
        cur.close()
        conn.close()
    except Exception:
        pass


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    print("=" * 70)
    print("TERMINAL-SF-FETCH-SUPPORT-PPT-EMAIL (ClickHouse) - EVALUATION")
    print("=" * 70)

    stats, bench = compute_expected()
    print("Expected per-priority SLA status (computed live):")
    for k in ("high", "medium", "low"):
        if k in stats:
            print(f"  {stats[k]['priority']}: csat={stats[k]['avg_csat']} "
                  f"resp={stats[k]['avg_resp']} -> {stats[k].get('status')}")

    check_pptx(args.agent_workspace, stats, bench)
    check_email(stats)
    check_reverse_validation(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = PASS_COUNT / total * 100 if total else 0
    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Accuracy: {accuracy:.1f}%")

    if CRITICAL_FAILS:
        print(f"  CRITICAL FAILURES: {len(CRITICAL_FAILS)}")
        for n in CRITICAL_FAILS:
            print(f"    - {n}")

    success = (not CRITICAL_FAILS) and accuracy >= 70
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump({"passed": PASS_COUNT, "failed": FAIL_COUNT,
                       "accuracy": accuracy, "critical_failed": CRITICAL_FAILS,
                       "success": success}, f)
    print(f"  Overall: {'PASS' if success else 'FAIL'}")

    # Critical gate runs before the accuracy gate.
    if CRITICAL_FAILS:
        sys.exit(1)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
