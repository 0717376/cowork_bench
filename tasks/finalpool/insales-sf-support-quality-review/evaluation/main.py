"""Evaluation for insales-sf-support-quality-review (RU fork: InSales + ClickHouse).

Check-based model: each check increments PASS_COUNT/total. CRITICAL_CHECKS are
semantic checks; any critical failure => overall FAIL (sys.exit(1)) regardless of
accuracy. Otherwise PASS requires accuracy >= 70%.

Data values are russified CENTRALLY (db/zzz_wc_after_init.sql, db/zzz_clickhouse_after_init.sql):
- wc product category 'Electronics' -> 'Электроника'
- sf_data ISSUE_TYPE 'Bug' -> 'Ошибка'
Priority slugs (High/Medium/Low) stay English as stored, so they are matched as-is.
We match against RU+EN literals (never hand-edit groundtruth).
"""
import argparse
import os
import sys

PASS_COUNT = 0
TOTAL = 0
FAILED_NAMES = []

# Semantic checks: any failure => overall FAIL regardless of accuracy.
CRITICAL_CHECKS = {
    "order_issues_electronics",
    "support_high_response_target",
    "support_three_priorities",
    "cross_ref_bug_type",
    "exec_summary_present",
}


def record(name, ok, msg=""):
    global PASS_COUNT, TOTAL
    TOTAL += 1
    if ok:
        PASS_COUNT += 1
    else:
        FAILED_NAMES.append(name)
        tag = "CRITICAL" if name in CRITICAL_CHECKS else "check"
        print(f"    FAIL [{tag}] {name}: {msg}")


def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


def nums_close(a, b, abs_tol=1.0, rel_tol=0.05):
    try:
        a, b = float(a), float(b)
    except (TypeError, ValueError):
        return False
    if abs(a - b) <= abs_tol:
        return True
    if b != 0 and abs(a - b) / abs(b) <= rel_tol:
        return True
    return False


def to_float(v):
    try:
        return float(v)
    except (TypeError, ValueError):
        return None


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


# Central-map aware matchers (RU label OR original EN slug).
def is_electronics(v):
    s = str(v).lower()
    return "электрон" in s or "electronics" in s


def is_bug(v):
    s = str(v).lower()
    return "ошибк" in s or "bug" in s


def is_high(v):
    s = str(v).lower()
    return "high" in s or "высок" in s


def check_excel(agent_workspace):
    import openpyxl

    path = os.path.join(agent_workspace, "CX_Quality_Review.xlsx")
    if not os.path.exists(path):
        record("excel_exists", False, "CX_Quality_Review.xlsx not found")
        # Mark dependent critical checks failed.
        for n in ("order_issues_electronics", "support_high_response_target",
                  "support_three_priorities", "cross_ref_bug_type",
                  "exec_summary_present"):
            record(n, False, "workbook missing")
        return
    record("excel_exists", True)

    try:
        wb = openpyxl.load_workbook(path, data_only=True)
    except Exception as e:
        record("excel_readable", False, f"Error reading Excel: {e}")
        for n in ("order_issues_electronics", "support_high_response_target",
                  "support_three_priorities", "cross_ref_bug_type",
                  "exec_summary_present"):
            record(n, False, "workbook unreadable")
        return
    record("excel_readable", True)

    # --- Sheet 1: Order Issues ---
    rows = load_sheet_rows(wb, "Order Issues")
    record("order_issues_sheet", rows is not None, "sheet not found")
    if rows is not None:
        data_rows = [r for r in rows[1:] if r and r[0] is not None]
        record("order_issues_rowcount", len(data_rows) >= 4,
               f"{len(data_rows)} rows, expected >= 4")

        elec = [r for r in data_rows if r[0] and is_electronics(r[0])]
        if not elec:
            record("order_issues_electronics", False,
                   "Electronics (Электроника) category not found")
        else:
            row = elec[0]
            total_ok = len(row) > 1 and row[1] is not None and \
                nums_close(row[1], 83, abs_tol=10)
            rate = to_float(row[3]) if len(row) > 3 else None
            rate_ok = rate is not None and 0 < rate < 100
            # Sanity: issue rate ~= issue/total*100 if both present.
            consistent = True
            tot = to_float(row[1]) if len(row) > 1 else None
            iss = to_float(row[2]) if len(row) > 2 else None
            if tot and iss is not None and rate is not None and tot > 0:
                consistent = num_close(rate, iss / tot * 100, rel_tol=0.2,
                                       abs_tol=5)
            record("order_issues_electronics", total_ok and rate_ok and consistent,
                   f"Electronics total={row[1] if len(row) > 1 else None} "
                   f"(~83), rate={rate} (0<rate<100, consistent={consistent})")

        # Non-critical: some valid issue-rate value exists.
        has_rate = any(
            len(r) > 3 and to_float(r[3]) is not None
            and 0 < to_float(r[3]) < 100
            for r in data_rows
        )
        record("order_issues_rate_col", has_rate,
               "no valid issue-rate values")
    else:
        record("order_issues_electronics", False, "Order Issues sheet missing")

    # --- Sheet 2: Support Metrics ---
    rows2 = load_sheet_rows(wb, "Support Metrics")
    record("support_metrics_sheet", rows2 is not None, "sheet not found")
    if rows2 is not None:
        data_rows2 = [r for r in rows2[1:] if r and r[0] is not None]

        # All three priority rows present.
        def has_prio(pred):
            return any(r[0] and pred(r[0]) for r in data_rows2)
        three_ok = (has_prio(is_high)
                    and has_prio(lambda v: "medium" in str(v).lower()
                                 or "средн" in str(v).lower())
                    and has_prio(lambda v: "low" in str(v).lower()
                                 or "низк" in str(v).lower()))
        # Resolution rate populated for all three rows.
        res_ok = sum(1 for r in data_rows2 if len(r) > 4
                     and to_float(r[4]) is not None) >= 3
        # Response target sourced from portal (4/12/24 across rows).
        targets = sorted({to_float(r[5]) for r in data_rows2
                          if len(r) > 5 and to_float(r[5]) is not None})
        target_ok = (len(data_rows2) >= 3
                     and any(t and nums_close(t, 4, abs_tol=0.6) for t in targets)
                     and any(t and nums_close(t, 12, abs_tol=1) for t in targets)
                     and any(t and nums_close(t, 24, abs_tol=1) for t in targets))
        record("support_three_priorities", three_ok and res_ok and target_ok,
               f"three_prios={three_ok}, resolution_populated={res_ok}, "
               f"portal_targets(4/12/24)={target_ok} (got {targets})")

        # High priority: avg response ~6.23 AND Meets_Response_Target == No
        high = [r for r in data_rows2 if r[0] and is_high(r[0])]
        if not high:
            record("support_high_response_target", False,
                   "High priority row not found")
        else:
            row = high[0]
            resp_ok = len(row) > 2 and row[2] is not None and \
                nums_close(row[2], 6.23, abs_tol=1.0)
            # Meets_Response_Target should be No (6.23h > 4h target).
            meets_val = str(row[6]).strip().lower() if len(row) > 6 \
                and row[6] is not None else ""
            meets_no = meets_val in ("no", "нет", "false", "0")
            record("support_high_response_target", resp_ok and meets_no,
                   f"High avg_response={row[2] if len(row) > 2 else None} "
                   f"(~6.23), meets_target='{meets_val}' (expected No/Нет)")
    else:
        record("support_three_priorities", False, "Support Metrics sheet missing")
        record("support_high_response_target", False,
               "Support Metrics sheet missing")

    # --- Sheet 3: Cross Reference ---
    rows3 = load_sheet_rows(wb, "Cross Reference")
    record("cross_ref_sheet", rows3 is not None, "sheet not found")
    if rows3 is not None:
        data_rows3 = [r for r in rows3[1:] if r and r[0] is not None]
        record("cross_ref_rowcount", len(data_rows3) >= 5,
               f"{len(data_rows3)} rows, expected >= 5")
        bug = [r for r in data_rows3 if r[0] and is_bug(r[0])]
        if not bug:
            record("cross_ref_bug_type", False,
                   "Bug (Ошибка) issue type not found")
        else:
            row = bug[0]
            tc = to_float(row[1]) if len(row) > 1 else None
            rt = to_float(row[2]) if len(row) > 2 else None
            sat = to_float(row[3]) if len(row) > 3 else None
            populated = tc is not None and rt is not None and sat is not None
            record("cross_ref_bug_type", populated,
                   f"Bug row Ticket_Count={tc}, Avg_Response_Time={rt}, "
                   f"Avg_Satisfaction={sat} (all must be populated)")
    else:
        record("cross_ref_bug_type", False, "Cross Reference sheet missing")

    # --- Sheet 4: Executive Summary ---
    rows4 = load_sheet_rows(wb, "Executive Summary")
    record("exec_summary_sheet", rows4 is not None, "sheet not found")
    if rows4 is not None:
        data_rows4 = [r for r in rows4[1:] if r and r[0] is not None]
        # At least 5 metric rows, each with a non-empty value.
        valued = [r for r in data_rows4 if len(r) > 1 and r[1] is not None
                  and str(r[1]).strip() != ""]
        record("exec_summary_present", len(valued) >= 5,
               f"{len(valued)} populated metric rows, expected >= 5")
    else:
        record("exec_summary_present", False, "Executive Summary sheet missing")


def check_pptx(agent_workspace):
    from pptx import Presentation

    path = os.path.join(agent_workspace, "Quality_Review.pptx")
    if not os.path.exists(path):
        record("pptx_exists", False, "Quality_Review.pptx not found")
        return
    record("pptx_exists", True)

    try:
        prs = Presentation(path)
    except Exception as e:
        record("pptx_readable", False, f"Error reading PowerPoint: {e}")
        return
    record("pptx_readable", True)

    record("pptx_slidecount", len(prs.slides) >= 6,
           f"{len(prs.slides)} slides, expected >= 6")

    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += shape.text.lower() + " "

    def any_in(*subs):
        return any(s in all_text for s in subs)

    # RU + EN alternatives (agent legitimately writes Russian).
    record("pptx_benchmarks",
           any_in("benchmark", "target", "бенчмарк", "целев", "целью"),
           "no benchmark/target mention")
    record("pptx_order_issues",
           any_in("issue", "order", "проблем", "заказ"),
           "no order-issue mention")
    record("pptx_support",
           any_in("support", "ticket", "поддержк", "тикет", "обращен"),
           "no support/ticket mention")
    record("pptx_recommendations",
           any_in("recommend", "рекоменд"),
           "no recommendations mention")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    agent_ws = args.agent_workspace or os.path.join(
        os.path.dirname(__file__), "..", "groundtruth_workspace"
    )

    print("  Checking Excel file...")
    check_excel(agent_ws)
    print("  Checking PowerPoint...")
    check_pptx(agent_ws)

    accuracy = PASS_COUNT / TOTAL * 100 if TOTAL else 0
    critical_failed = [n for n in FAILED_NAMES if n in CRITICAL_CHECKS]

    print(f"\n=== Results: {PASS_COUNT}/{TOTAL} passed ({accuracy:.1f}%) ===")
    if critical_failed:
        print(f"CRITICAL FAILURES: {len(critical_failed)} -> {critical_failed}")

    success = (not critical_failed) and accuracy >= 70
    if success:
        print("=== RESULT: PASS ===")
        sys.exit(0)
    else:
        print("=== RESULT: FAIL ===")
        sys.exit(1)


if __name__ == "__main__":
    main()
