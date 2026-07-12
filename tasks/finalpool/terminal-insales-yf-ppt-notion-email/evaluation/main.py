"""Evaluation for terminal-insales-yf-ppt-notion-email.

Checks:
1. Market_Strategy_Presentation.pptx (7 slides, correct per-slide roles)
2. Teamly "Market Strategy Tracker" page with 5 initiatives (correct markers)
3. Email to ceo@company.com
4. Email to marketing_team@company.com
5. market_correlation.py / category_analysis.py scripts + JSON outputs

Sources of truth (recomputed live):
  - wc.orders / wc.products  (e-commerce, InSales-backed, schema wc.*)
  - moex.stock_prices        (MOEX historical prices, symbol/close/date)
  - teamly.pages / teamly.spaces  (corporate knowledge base)
  - email.messages

CRITICAL_CHECKS (semantic): any failure => overall FAIL regardless of accuracy.
Pass threshold otherwise: accuracy >= 70%.
"""
import argparse
import json
import os
import sys
from collections import defaultdict

import psycopg2
from pptx import Presentation

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent", "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0
FAILED_NAMES = []

# Consumer/retail benchmark symbol after the YF->MOEX swap (AMZN -> Магнит).
CONSUMER_SYMBOL = "MGNT.ME"
# Financial-sector indicator (JPM -> Сбербанк).
FINANCIAL_SYMBOL = "SBER.ME"
# Equal-weight broad-market benchmark over all available .ME stocks (^DJI ->).
BENCH_SYMBOLS = ["SBER.ME", "GAZP.ME", "LKOH.ME", "MGNT.ME", "MTSS.ME", "TCSG.ME"]

# The three recommendation strings are eval-defined tokens (kept English).
REC_ALIGNED = "market-aligned pricing"
REC_COUNTER = "counter-cyclical promotions"
REC_INDEPENDENT = "independent pricing strategy"
ALLOWED_RECS = {REC_ALIGNED, REC_COUNTER, REC_INDEPENDENT}

# Critical checks: any failure => overall FAIL regardless of accuracy.
CRITICAL_CHECKS = {
    "market_correlation.json recommendation is one of the three allowed strings",
    "market_correlation.json recommendation matches recomputed correlation rule",
    "Teamly 'Market Strategy Tracker' page exists",
    "Teamly tracker lists exactly the 5 required initiatives with correct markers",
    "category_market_analysis.json top category == DB-computed top category",
    "CEO email exists with pricing recommendation matching market_correlation.json",
    "Marketing email references the DB-computed top category",
}


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        FAILED_NAMES.append(name)
        print(f"  [FAIL] {name}: {str(detail)[:300]}")


def num_close(a, b, tol=2.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except Exception:
        return False


# ---------------------------------------------------------------------------
# Live source-of-truth recomputation
# ---------------------------------------------------------------------------

def _monthly_revenue(cur):
    """Monthly e-commerce revenue {YYYY-MM: revenue} from wc.orders."""
    cur.execute("""
        SELECT to_char(o.date_created, 'YYYY-MM') AS month,
               SUM((li->>'total')::numeric) AS rev
        FROM wc.orders o, jsonb_array_elements(o.line_items) li
        WHERE o.status NOT IN ('cancelled','refunded','failed')
        GROUP BY month
    """)
    return {m: float(r) for m, r in cur.fetchall() if m}


def _monthly_close(cur, symbol):
    """Monthly average close {YYYY-MM: avg_close} for one symbol."""
    cur.execute("""
        SELECT to_char(date, 'YYYY-MM') AS month, AVG(close) AS c
        FROM moex.stock_prices WHERE symbol = %s GROUP BY month
    """, (symbol,))
    return {m: float(c) for m, c in cur.fetchall() if m and c is not None}


def _monthly_close_benchmark(cur):
    """Equal-weight benchmark: per-day mean close across BENCH_SYMBOLS, then
    averaged per month. Returns {YYYY-MM: avg}."""
    cur.execute("""
        SELECT date, AVG(close) AS dayavg
        FROM moex.stock_prices WHERE symbol = ANY(%s)
        GROUP BY date
    """, (BENCH_SYMBOLS,))
    by_month = defaultdict(list)
    for d, v in cur.fetchall():
        if v is None:
            continue
        by_month[d.strftime("%Y-%m")].append(float(v))
    return {m: sum(vs) / len(vs) for m, vs in by_month.items() if vs}


def _pearson(xs, ys):
    n = len(xs)
    if n < 2:
        return None
    mx = sum(xs) / n
    my = sum(ys) / n
    cov = sum((x - mx) * (y - my) for x, y in zip(xs, ys))
    vx = sum((x - mx) ** 2 for x in xs)
    vy = sum((y - my) ** 2 for y in ys)
    if vx <= 0 or vy <= 0:
        return None
    return cov / (vx ** 0.5 * vy ** 0.5)


def _rec_from_corr(corr):
    """Apply the task's recommendation rule to a MGNT correlation value.
    When the correlation is undefined (degenerate overlap / <2 shared months),
    the task's 'otherwise' branch applies: there is no significant market
    correlation, so the recommendation is the independent pricing strategy."""
    if corr is None:
        return REC_INDEPENDENT
    if abs(corr) > 0.5:
        return REC_ALIGNED
    if corr < -0.3:
        return REC_COUNTER
    return REC_INDEPENDENT


def get_expected_from_db():
    """Recompute top category, MGNT pct-change, and the expected recommendation
    dynamically. No silent hard-coded defaults: on DB error these are None and
    the dependent critical checks fail loudly."""
    exp = {"top_category": None, "consumer_pct_change": None,
           "expected_rec": None, "consumer_corr": None}
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        # Top category by revenue (completed-ish orders).
        cur.execute("""
            SELECT p.categories->0->>'name' AS cat, SUM((li->>'total')::numeric) AS rev
            FROM wc.orders o, jsonb_array_elements(o.line_items) li
            JOIN wc.products p ON (li->>'product_id')::int = p.id
            WHERE o.status NOT IN ('cancelled','refunded','failed')
            GROUP BY cat ORDER BY rev DESC LIMIT 1
        """)
        row = cur.fetchone()
        if row and row[0]:
            exp["top_category"] = row[0].lower()

        # MGNT.ME pct change (earliest -> latest close).
        cur.execute("""
            SELECT
              (SELECT close FROM moex.stock_prices WHERE symbol=%s ORDER BY date ASC LIMIT 1),
              (SELECT close FROM moex.stock_prices WHERE symbol=%s ORDER BY date DESC LIMIT 1)
        """, (CONSUMER_SYMBOL, CONSUMER_SYMBOL))
        row = cur.fetchone()
        if row and row[0] and row[1]:
            exp["consumer_pct_change"] = float((row[1] - row[0]) / row[0] * 100)

        # Monthly correlation MGNT.ME vs e-commerce revenue over overlapping months.
        rev = _monthly_revenue(cur)
        close = _monthly_close(cur, CONSUMER_SYMBOL)
        months = sorted(set(rev) & set(close))
        if len(months) >= 2:
            corr = _pearson([rev[m] for m in months], [close[m] for m in months])
            exp["consumer_corr"] = corr
            exp["expected_rec"] = _rec_from_corr(corr)

        cur.close()
        conn.close()
    except Exception as e:
        print(f"  [WARN] DB recomputation failed: {e}")
    return exp


EXPECTED = get_expected_from_db()


# ---------------------------------------------------------------------------
# Checks
# ---------------------------------------------------------------------------

def check_pptx(workspace):
    print("\n=== Check 1: Market_Strategy_Presentation.pptx ===")
    path = os.path.join(workspace, "Market_Strategy_Presentation.pptx")
    if not os.path.exists(path):
        check("PPTX file exists", False, f"Not found at {path}")
        return
    check("PPTX file exists", True)

    prs = Presentation(path)
    slides = list(prs.slides)
    check("Has 7 slides", len(slides) == 7, f"Found {len(slides)}")

    all_texts = []
    for slide in slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    slide_text.append(para.text)
        all_texts.append(" ".join(slide_text).lower())
    full_text = " ".join(all_texts)

    def has(idx, *subs):
        return len(all_texts) > idx and any(s in all_texts[idx] for s in subs)

    check("Slide 1 has positioning title",
          has(0, "market", "positioning", "позицион", "стратег"),
          all_texts[0][:120] if all_texts else "no slides")
    check("Slide 2 mentions MGNT/Магнит", has(1, "mgnt", "магнит"),
          all_texts[1][:120] if len(all_texts) > 1 else "")
    check("Slide 2 mentions SBER/Сбербанк", has(1, "sber", "сбер"),
          all_texts[1][:120] if len(all_texts) > 1 else "")
    check("Slide 2 mentions broad-market benchmark",
          has(1, "benchmark", "бенчмарк", "equal", "равновзвеш", "average", "средн"),
          all_texts[1][:120] if len(all_texts) > 1 else "")
    check("Slide 3 has revenue data",
          has(2, "revenue", "выручк", "2026", "2025"),
          all_texts[2][:120] if len(all_texts) > 2 else "")
    check("Slide 4 mentions categories",
          has(3, "categor", "категор", "electronics", "tv", "audio"),
          all_texts[3][:120] if len(all_texts) > 3 else "")
    check("Slide 5 mentions correlation",
          has(4, "correlation", "корреляц", "pearson", "пирсон"),
          all_texts[4][:120] if len(all_texts) > 4 else "")
    check("Slide 6 has recommendations",
          has(5, "recommend", "рекоменд", "pricing", "strategy", "стратег"),
          all_texts[5][:120] if len(all_texts) > 5 else "")
    check("Slide 7 has next steps",
          has(6, "next", "action", "step", "следующ", "шаг", "действ"),
          all_texts[6][:120] if len(all_texts) > 6 else "")

    # Recommendation slide must state the same recommendation as the JSON output.
    if EXPECTED.get("expected_rec"):
        check("Recommendation slide states the recomputed recommendation",
              EXPECTED["expected_rec"] in full_text,
              f"Expected '{EXPECTED['expected_rec']}' somewhere in deck")

    if EXPECTED.get("top_category"):
        check(f"Mentions top category '{EXPECTED['top_category']}'",
              EXPECTED["top_category"] in full_text,
              f"No '{EXPECTED['top_category']}' mention")


def _teamly_tracker_pages(cur):
    """Return list of (id, title, body) for the Market Strategy Tracker page(s)."""
    cur.execute("""
        SELECT id, title, COALESCE(body, '')
        FROM teamly.pages
        WHERE title ILIKE '%%market%%strategy%%tracker%%'
           OR title ILIKE '%%трекер%%страт%%'
    """)
    return cur.fetchall()


def check_teamly():
    print("\n=== Check 2: Teamly Market Strategy Tracker ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        pages = _teamly_tracker_pages(cur)
        check("Teamly 'Market Strategy Tracker' page exists", len(pages) >= 1,
              "No page titled 'Market Strategy Tracker'")
        if not pages:
            check("Teamly tracker lists exactly the 5 required initiatives with correct markers",
                  False, "No tracker page")
            return

        combined = " ".join((str(t) + " " + str(b)) for _, t, b in pages)
        low = combined.lower()

        # The 5 required initiatives (English markers preserved per task.md).
        required_initiatives = [
            "tv & home theater", "electronics", "audio", "camera", "home appliances",
        ]
        present = sum(1 for kw in required_initiatives if kw in low)
        check("Tracker mentions all 5 initiative subjects",
              present >= 5, f"Found {present}/5 initiative subjects")

        # Structured field markers (English tokens, exact case-insensitive).
        category_tokens = ["pricing", "marketing", "inventory", "expansion", "cost management"]
        condition_tokens = ["bull", "bear", "neutral"]
        cat_hits = sum(1 for c in category_tokens if c in low)
        cond_hits = sum(1 for c in condition_tokens if c in low)

        # Expected_Impact numbers from the 5 prescribed initiatives.
        impact_tokens = ["5000", "3000", "2500", "1500", "2000"]
        impact_hits = sum(1 for v in impact_tokens if v in combined)

        # Verify at least the two anchor initiatives precisely:
        #  TV & Home Theater => Pricing / Bull / 5000 ; Audio => Expansion / Bull / 2500
        anchor_ok = (
            "tv & home theater" in low and "pricing" in low and "bull" in low
            and "5000" in combined and "audio" in low and "expansion" in low
            and "2500" in combined
        )

        all_markers = (present >= 5 and cat_hits >= 5 and cond_hits >= 2
                       and impact_hits >= 5 and anchor_ok)
        check("Teamly tracker lists exactly the 5 required initiatives with correct markers",
              all_markers,
              f"initiatives={present}/5 categories={cat_hits}/5 "
              f"conditions={cond_hits}/3 impacts={impact_hits}/5 anchor={anchor_ok}")

        # Non-critical structural detail.
        check("Tracker has Category/Market_Condition field markers",
              ("category" in low or "категори" in low)
              and ("market_condition" in low or "market condition" in low
                   or "condition" in low or "услови" in low),
              "Field-name markers not found")
    except Exception as e:
        check("Teamly 'Market Strategy Tracker' page exists", False, str(e))
        check("Teamly tracker lists exactly the 5 required initiatives with correct markers",
              False, str(e))
    finally:
        cur.close()
        conn.close()


def check_emails():
    print("\n=== Check 3: Emails ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        cur.execute("""
            SELECT subject, body_text FROM email.messages
            WHERE to_addr::text LIKE '%%ceo@company.com%%'
            ORDER BY date DESC LIMIT 1
        """)
        ceo = cur.fetchone()
        ceo_subj = (ceo[0] or "").lower() if ceo else ""
        ceo_body = (ceo[1] or "").lower() if ceo else ""

        # CRITICAL: CEO email exists AND carries the recommendation matching the JSON.
        rec = EXPECTED.get("expected_rec")
        ceo_has_rec = bool(ceo) and (rec is not None) and (rec in ceo_body)
        check("CEO email exists with pricing recommendation matching market_correlation.json",
              ceo_has_rec,
              f"ceo_exists={bool(ceo)} expected_rec={rec}")

        if ceo:
            check("CEO email subject mentions strategy/summary",
                  any(s in ceo_subj for s in ["strategy", "summary", "executive",
                                              "стратег", "резюме"]),
                  f"Subject: {ceo[0]}")
            check("CEO email body mentions MGNT/Магнит",
                  "mgnt" in ceo_body or "магнит" in ceo_body, "No MGNT mention")
            if EXPECTED.get("top_category"):
                check("CEO email body mentions top category",
                      EXPECTED["top_category"] in ceo_body,
                      f"No '{EXPECTED['top_category']}' mention")

        cur.execute("""
            SELECT subject, body_text FROM email.messages
            WHERE to_addr::text LIKE '%%marketing_team@company.com%%'
            ORDER BY date DESC LIMIT 1
        """)
        mkt = cur.fetchone()
        mkt_subj = (mkt[0] or "").lower() if mkt else ""
        mkt_body = (mkt[1] or "").lower() if mkt else ""

        check("Marketing email exists", mkt is not None,
              "No email to marketing_team@company.com")

        # CRITICAL: marketing email references the DB-computed top category.
        top = EXPECTED.get("top_category")
        mkt_has_top = bool(mkt) and (top is not None) and (top in mkt_body)
        check("Marketing email references the DB-computed top category",
              mkt_has_top, f"mkt_exists={bool(mkt)} top={top}")

        if mkt:
            check("Marketing email subject mentions category/performance/Q2",
                  any(s in mkt_subj for s in ["category", "performance", "q2",
                                              "категор", "показател"]),
                  f"Subject: {mkt[0]}")
            check("Marketing email references underperforming categories",
                  ("camera" in mkt_body or "watch" in mkt_body
                   or "камер" in mkt_body or "час" in mkt_body),
                  "No underperforming-category detail")
    except Exception as e:
        check("CEO email exists with pricing recommendation matching market_correlation.json",
              False, str(e))
        check("Marketing email references the DB-computed top category", False, str(e))
    finally:
        cur.close()
        conn.close()


def check_scripts(workspace):
    print("\n=== Check 4: Scripts and JSON outputs ===")
    check("market_correlation.py exists",
          os.path.exists(os.path.join(workspace, "market_correlation.py")))
    check("category_analysis.py exists",
          os.path.exists(os.path.join(workspace, "category_analysis.py")))

    # --- market_correlation.json ---
    corr_path = os.path.join(workspace, "market_correlation.json")
    corr_data = None
    if os.path.exists(corr_path):
        check("market_correlation.json exists", True)
        try:
            with open(corr_path) as f:
                corr_data = json.load(f)
            check("market_correlation.json valid JSON", True)
        except Exception as e:
            check("market_correlation.json valid JSON", False, str(e))
    else:
        check("market_correlation.json exists", False, f"Not found at {corr_path}")

    rec_value = None
    if isinstance(corr_data, dict):
        data_str = json.dumps(corr_data).lower()
        check("market_correlation.json mentions MGNT correlation",
              "mgnt" in data_str or "магнит" in data_str,
              f"Content: {data_str[:200]}")
        # Locate the recommendation string anywhere in the JSON.
        for r in ALLOWED_RECS:
            if r in data_str:
                rec_value = r
                break

    # CRITICAL: recommendation is exactly one of the three allowed strings.
    check("market_correlation.json recommendation is one of the three allowed strings",
          rec_value is not None,
          f"Recommendation not found among {sorted(ALLOWED_RECS)}")

    # CRITICAL: recommendation matches the recomputed correlation rule.
    expected_rec = EXPECTED.get("expected_rec")
    check("market_correlation.json recommendation matches recomputed correlation rule",
          (rec_value is not None) and (expected_rec is not None)
          and (rec_value == expected_rec),
          f"Got '{rec_value}', expected '{expected_rec}' "
          f"(corr={EXPECTED.get('consumer_corr')})")

    # --- category_market_analysis.json ---
    cat_path = os.path.join(workspace, "category_market_analysis.json")
    cat_data = None
    if os.path.exists(cat_path):
        check("category_market_analysis.json exists", True)
        try:
            with open(cat_path) as f:
                cat_data = json.load(f)
            check("category_market_analysis.json valid JSON", True)
        except Exception as e:
            check("category_market_analysis.json valid JSON", False, str(e))
    else:
        check("category_market_analysis.json exists", False, f"Not found at {cat_path}")

    # CRITICAL: identified top category == DB-computed top category (dynamic).
    top = EXPECTED.get("top_category")
    cat_str = json.dumps(cat_data).lower() if cat_data is not None else ""
    check("category_market_analysis.json top category == DB-computed top category",
          (top is not None) and (top in cat_str),
          f"Expected top category '{top}' in JSON; got: {cat_str[:200]}")


def check_reverse_validation():
    print("\n=== Reverse Validation ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        # No emails to noise recipients.
        noise_recipients = [
            "all-staff@company.com", "hr@company.com",
            "newsletter@company.com", "sales_team@company.com",
        ]
        for addr in noise_recipients:
            cur.execute(
                "SELECT COUNT(*) FROM email.messages WHERE to_addr::text ILIKE %s",
                (f"%{addr}%",))
            cnt = cur.fetchone()[0]
            check(f"No email sent to noise recipient {addr}", cnt == 0,
                  f"Found {cnt} emails to {addr}")

        # Tracker page should not carry unrelated noise categories.
        pages = _teamly_tracker_pages(cur)
        if pages:
            txt = " ".join((str(t) + " " + str(b)).lower() for _, t, b in pages)
            for noise_cat in ["healthcare", "real estate", "energy", "недвижим"]:
                check(f"Tracker has no noise category '{noise_cat}'",
                      noise_cat not in txt, f"Found '{noise_cat}'")
    except Exception as e:
        check("Reverse validation", False, str(e))
    finally:
        cur.close()
        conn.close()


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    print(f"[eval] recomputed expectations: top_category={EXPECTED.get('top_category')} "
          f"consumer_corr={EXPECTED.get('consumer_corr')} "
          f"expected_rec={EXPECTED.get('expected_rec')}")

    check_pptx(args.agent_workspace)
    check_teamly()
    check_emails()
    check_scripts(args.agent_workspace)
    check_reverse_validation()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    critical_failed = [n for n in FAILED_NAMES if n in CRITICAL_CHECKS]
    if critical_failed:
        print(f"  CRITICAL FAILURES: {len(critical_failed)}")
        for n in critical_failed:
            print(f"    - {n}")

    result = {
        "total_passed": PASS_COUNT, "total_checks": total, "accuracy": accuracy,
        "critical_failed": critical_failed,
    }
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if critical_failed:
        print("FAIL (critical check failed)")
        sys.exit(1)
    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    print("FAIL")
    sys.exit(1)


if __name__ == "__main__":
    main()
