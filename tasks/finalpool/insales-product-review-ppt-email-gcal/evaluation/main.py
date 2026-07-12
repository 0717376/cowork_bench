"""Evaluation for insales-product-review-ppt-email-gcal (InSales store).

Checks:
1. PPT file Product_Review_Report.pptx with at least 5 slides
   - Contains "Product Quality Review" and "Q1" title
   - Mentions an in-scope category by name (accept RU or EN)
   - Contains a best/top-reviewed section AND a recommendations section
2. GCal event "Product Quality Review Meeting" ~7 days from launch_time
   with a non-empty description
3. Email to product-team@store.example.com from analytics@store.example.com
   Subject contains "Product Quality Review", body names a highest/lowest category

NOTE: store category names are russified centrally (db/zzz_wc_after_init.sql via
scripts/wc_relabel_map.py): Headphones->Наушники, Audio->Аудио, Cameras->Камеры,
Speakers->Колонки, ... The agent reads live RU names and writes them into the PPT/email,
so all category matching accepts RU OR EN tokens.

CRITICAL_CHECKS: any failure => immediate sys.exit(1) before the accuracy gate.
PASS requires: no critical failure AND accuracy >= 70.
"""
import argparse
import json
import os
import sys
from datetime import datetime, timedelta, timezone

import psycopg2

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="cowork_gym", user="eigent", password="camel")

PASS_COUNT = 0
FAIL_COUNT = 0
FAILED_CRITICAL = []

# RU+EN tokens for the in-scope product categories (lowercased substrings).
CATEGORY_TOKENS = [
    ["headphones", "наушник"],
    ["audio", "аудио"],
    ["cameras", "камер"],
    ["speakers", "колонк"],
    ["electronics", "электроник"],
    ["tv", "кинотеатр", "телевизор"],
]


def check(name, condition, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS]{' [CRIT]' if critical else ''} {name}")
    else:
        FAIL_COUNT += 1
        if critical:
            FAILED_CRITICAL.append(name)
        detail_str = f": {str(detail)[:200]}" if detail else ""
        print(f"  [FAIL]{' [CRIT]' if critical else ''} {name}{detail_str}")


def _mentions_any_category(text):
    return any(any(tok in text for tok in group) for group in CATEGORY_TOKENS)


def check_ppt(agent_workspace):
    print("\n=== Checking PPT File ===")
    ppt_path = os.path.join(agent_workspace, "Product_Review_Report.pptx")
    check("Product_Review_Report.pptx exists", os.path.isfile(ppt_path),
          f"Expected at {ppt_path}", critical=True)
    if not os.path.isfile(ppt_path):
        return

    try:
        from pptx import Presentation
        prs = Presentation(ppt_path)
    except Exception as e:
        check("PPT file readable", False, str(e), critical=True)
        return

    check("PPT has at least 5 slides", len(prs.slides) >= 5,
          f"Found {len(prs.slides)} slides", critical=True)

    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text.append(shape.text.lower())
    full_text = " ".join(all_text)

    check("PPT contains 'Product Quality Review'",
          "product quality review" in full_text, "Title text not found", critical=True)
    check("PPT contains 'Q1'", "q1" in full_text, "Q1 not found")
    # CRITICAL: broadened RU+EN category grep (store returns RU category names).
    check("PPT mentions at least one in-scope category (RU or EN)",
          _mentions_any_category(full_text),
          "No category names found (RU or EN)", critical=True)
    check("PPT mentions top/best reviewed section",
          "best" in full_text or "top" in full_text or "highest" in full_text
          or "лучш" in full_text or "топ" in full_text,
          "No best-products section found")
    check("PPT mentions recommendations",
          "recommend" in full_text or "action" in full_text or "improve" in full_text
          or "рекоменд" in full_text or "действ" in full_text or "улучш" in full_text,
          "No recommendations found", critical=True)


def check_gcal(launch_time_str=None):
    print("\n=== Checking Google Calendar ===")
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()

    cur.execute("""
        SELECT summary, start_datetime, end_datetime, description
        FROM gcal.events
        WHERE LOWER(summary) LIKE '%product%' AND LOWER(summary) LIKE '%quality%'
    """)
    events = cur.fetchall()
    check("Product Quality Review Meeting event created", len(events) >= 1,
          f"Found {len(events)} matching events", critical=True)

    if events:
        # Non-empty description required by the task.
        has_desc = any((e[3] or "").strip() for e in events)
        check("Review Meeting event has a non-empty description", has_desc,
              "All matching events have empty description")

    if events and launch_time_str:
        try:
            launch_time = datetime.fromisoformat(launch_time_str)
            if launch_time.tzinfo is None:
                launch_time = launch_time.replace(tzinfo=timezone.utc)
            target_date = launch_time + timedelta(days=7)
            for event in events:
                event_start = event[1]
                if event_start.tzinfo is None:
                    event_start = event_start.replace(tzinfo=timezone.utc)
                diff_days = abs((event_start.date() - target_date.date()).days)
                if diff_days <= 2:
                    check("Review Meeting is ~7 days from launch", True)
                    break
            else:
                check("Review Meeting is ~7 days from launch", False,
                      f"Closest event at {events[0][1]}, expected ~{target_date.date()}")
        except Exception as e:
            check("Review Meeting date check", False, str(e))

    cur.close()
    conn.close()


def check_emails():
    print("\n=== Checking Emails ===")
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()

    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
    """)
    all_emails = cur.fetchall()
    conn.close()

    def parse_recipients(to_addr):
        if to_addr is None:
            return []
        if isinstance(to_addr, list):
            return [str(r).strip().lower() for r in to_addr]
        to_str = str(to_addr).strip()
        try:
            parsed = json.loads(to_str)
            if isinstance(parsed, list):
                return [str(r).strip().lower() for r in parsed]
            return [to_str.lower()]
        except (json.JSONDecodeError, TypeError):
            return [to_str.lower()]

    target = "product-team@store.example.com"
    found = None
    for subj, from_addr, to_addr, body in all_emails:
        recipients = parse_recipients(to_addr)
        if target in recipients:
            found = (subj, from_addr, to_addr, body)
            break

    check("Email sent to product-team@store.example.com", found is not None,
          f"No email found for {target}", critical=True)
    if found:
        subj, from_addr, to_addr, body = found
        check("Email from analytics@store.example.com",
              "analytics@store.example.com" in (from_addr or "").lower(),
              f"From: {from_addr}")
        check("Subject contains 'Product Quality Review'",
              "product quality review" in (subj or "").lower(),
              f"Subject: {subj}", critical=True)
        # Body must summarize findings and name a highest/lowest category (RU or EN).
        body_l = (body or "").lower()
        check("Email body is non-empty", len(body_l.strip()) > 0,
              "Empty body")
        check("Email body names at least one in-scope category (RU or EN)",
              _mentions_any_category(body_l),
              "Body does not mention any category (RU or EN)", critical=True)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    print("=" * 70)
    print("WC PRODUCT REVIEW PPT EMAIL GCAL - EVALUATION")
    print("=" * 70)

    check_ppt(args.agent_workspace)
    check_gcal(args.launch_time)
    check_emails()

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100) if total else 0.0

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Accuracy: {accuracy:.1f}%")

    if FAILED_CRITICAL:
        print(f"  Overall: FAIL (critical checks failed: {FAILED_CRITICAL})")
        sys.exit(1)

    overall = accuracy >= 70
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
