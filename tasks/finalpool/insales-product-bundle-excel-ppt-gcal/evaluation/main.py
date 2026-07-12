"""Evaluation for insales-product-bundle-excel-ppt-gcal (InSales / wc.* schema)."""
import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "cowork_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []  # имена провалившихся обязательных проверок

# Обязательные (критичные) проверки — содержательное ядро задачи.
# Любой их провал => итог FAIL независимо от accuracy.
# Имена формируются динамически (по рангу пары), поэтому критичность
# определяется по префиксу через is_critical().
CRITICAL_PREFIXES = (
    "CRITICAL Co-Purchase count",     # точный счётчик совместных покупок (tol=0)
    "CRITICAL Priority #1",           # формула приоритета для топ-1 пары
    "CRITICAL Bundle price rank",     # цена бандла = combined * 0.90 для топ-3
    "CRITICAL Matrix no sub-threshold",  # порог >=2 соблюдён, под-пороговые пары не просочились
    "CRITICAL Calendar date",         # три конкретные даты запуска 03-16/18/20 в 10:00 на 1 час
)


def is_critical(name):
    return any(name.startswith(p) for p in CRITICAL_PREFIXES)


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {str(detail)[:300]}" if detail else ""
        marker = " [CRITICAL]" if is_critical(name) else ""
        print(f"  [FAIL]{marker} {name}{msg}")
        if is_critical(name):
            CRITICAL_FAILS.append(name)


def num_close(a, b, tol=0.5):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def get_expected_data():
    """Compute expected co-purchase data from read-only DB."""
    from collections import defaultdict
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("""
        SELECT o.id as order_id,
               (li->>'product_id')::int as product_id,
               li->>'name' as product_name,
               (li->>'price')::numeric as price
        FROM wc.orders o, jsonb_array_elements(o.line_items) li
        WHERE o.status IN ('completed', 'processing')
    """)
    rows = cur.fetchall()

    # regular_price из каталога: PDF говорит «обычная цена», поэтому агент
    # законно может считать бандл и от line-item price, и от regular_price.
    cur.execute("SELECT id, regular_price FROM wc.products WHERE regular_price IS NOT NULL")
    regular_prices = {pid: float(rp) for pid, rp in cur.fetchall()}

    order_items = defaultdict(list)
    product_info = {}
    for order_id, pid, pname, price in rows:
        order_items[order_id].append((pid, pname, float(price)))
        product_info[pid] = (pname, float(price))

    pair_counts = defaultdict(int)
    pair_revenues = defaultdict(list)
    for order_id, items in order_items.items():
        pids = sorted(set(item[0] for item in items))
        item_map = {item[0]: item for item in items}
        for i in range(len(pids)):
            for j in range(i + 1, len(pids)):
                pa, pb = pids[i], pids[j]
                pair_counts[(pa, pb)] += 1
                combined = item_map[pa][2] + item_map[pb][2]
                pair_revenues[(pa, pb)].append(combined)

    qualified = []
    for (pa, pb), count in pair_counts.items():
        if count >= 2:
            avg_rev = round(sum(pair_revenues[(pa, pb)]) / len(pair_revenues[(pa, pb)]), 2)
            priority = round(count * avg_rev, 2)
            qualified.append({
                'pid_a': pa, 'name_a': product_info[pa][0],
                'price_a': product_info[pa][1],
                'pid_b': pb, 'name_b': product_info[pb][0],
                'price_b': product_info[pb][1],
                'reg_a': regular_prices.get(pa, product_info[pa][1]),
                'reg_b': regular_prices.get(pb, product_info[pb][1]),
                'count': count, 'avg_rev': avg_rev, 'priority': priority
            })

    qualified.sort(key=lambda x: x['priority'], reverse=True)
    cur.close()
    conn.close()
    return qualified


def check_excel(agent_workspace):
    """Check Bundle_Analysis.xlsx."""
    print("\n=== Checking Bundle_Analysis.xlsx ===")
    agent_file = os.path.join(agent_workspace, "Bundle_Analysis.xlsx")
    if not os.path.isfile(agent_file):
        record("Excel file exists", False, f"Not found: {agent_file}")
        return False
    record("Excel file exists", True)

    try:
        wb = openpyxl.load_workbook(agent_file, data_only=True)
    except Exception as e:
        record("Excel readable", False, str(e))
        return False

    expected = get_expected_data()
    all_ok = True

    # --- Check Co-Purchase Matrix sheet ---
    cp_sheet = None
    for name in wb.sheetnames:
        if "co" in name.lower() and "purchase" in name.lower():
            cp_sheet = wb[name]
            break
        if "matrix" in name.lower():
            cp_sheet = wb[name]
            break
    if cp_sheet is None:
        record("Sheet 'Co-Purchase Matrix' exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Sheet 'Co-Purchase Matrix' exists", True)
        rows = list(cp_sheet.iter_rows(min_row=2, values_only=True))
        record("Co-Purchase Matrix has data rows", len(rows) >= len(expected),
               f"Expected {len(expected)}, got {len(rows)}")

        # CRITICAL: ни одна под-пороговая пара (count < 2) не должна просочиться.
        # Любая строка с числовым Co_Purchase_Count < 2 = нарушение порога из PDF.
        sub_threshold = []
        for r in rows:
            if r and len(r) >= 3:
                try:
                    cnt = float(r[2])
                    if cnt < 2:
                        sub_threshold.append((r[0], r[1], cnt))
                except (TypeError, ValueError):
                    pass
        record("CRITICAL Matrix no sub-threshold pairs (count>=2)",
               len(sub_threshold) == 0,
               f"Под-пороговые пары: {sub_threshold[:5]}")
        if sub_threshold:
            all_ok = False

        for rank, ep in enumerate(expected):
            found = False
            for r in rows:
                if r and len(r) >= 5:
                    # Match by product names (partial match). Названия товаров
                    # заморожены на английском (frozen-English), поэтому совпадение
                    # по подстроке остаётся корректным после russification.
                    ra = str(r[0]).lower() if r[0] else ""
                    rb = str(r[1]).lower() if r[1] else ""
                    ea = ep['name_a'].lower()[:30]
                    eb = ep['name_b'].lower()[:30]
                    if (ea[:20] in ra and eb[:20] in rb) or (ea[:20] in rb and eb[:20] in ra):
                        found = True
                        # CRITICAL: точный счётчик совместных покупок (tol=0).
                        ok_count = num_close(r[2], ep['count'], 0)
                        record(f"CRITICAL Co-Purchase count ({ep['name_a'][:25]}...)", ok_count,
                               f"Expected {ep['count']}, got {r[2]}")
                        if not ok_count:
                            all_ok = False
                        # Priority_Score = count * avg_combined_revenue.
                        # Для топ-1 пары — критично (ядро формулы скоринга).
                        prio_name = (f"CRITICAL Priority #1 ({ep['name_a'][:25]}...)"
                                     if rank == 0
                                     else f"Pair priority ({ep['name_a'][:25]}...)")
                        prio_tol = 1.0 if rank == 0 else 5.0
                        ok_priority = num_close(r[4], ep['priority'], prio_tol)
                        record(prio_name, ok_priority,
                               f"Expected {ep['priority']}, got {r[4]}")
                        if not ok_priority:
                            all_ok = False
                        break
            if not found:
                record(f"Pair found: {ep['name_a'][:30]}...", False, "Not found in sheet")
                all_ok = False

    # --- Check Bundle Proposals sheet ---
    bp_sheet = None
    for name in wb.sheetnames:
        if "bundle" in name.lower() and "proposal" in name.lower():
            bp_sheet = wb[name]
            break
    if bp_sheet is None:
        record("Sheet 'Bundle Proposals' exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Sheet 'Bundle Proposals' exists", True)
        rows = list(bp_sheet.iter_rows(min_row=2, values_only=True))
        record("Bundle Proposals has >= 5 rows", len(rows) >= 5,
               f"Got {len(rows)} rows")

        # Check top bundle pricing.
        # Принимаем обе базы расчёта 10%-скидки: line-item цены из заказов
        # ИЛИ regular_price из каталога — оба прочтения PDF легитимны.
        for i, ep in enumerate(expected[:3]):
            bundle_price = round((ep['price_a'] + ep['price_b']) * 0.90, 2)
            bundle_price_reg = round((ep['reg_a'] + ep['reg_b']) * 0.90, 2)
            found = False
            for r in rows:
                if r and len(r) >= 6:
                    rname = str(r[0]).lower() if r[0] else ""
                    # Match both product names to avoid ambiguity
                    name_a_short = ep['name_a'].lower()[:15]
                    name_b_short = ep['name_b'].lower()[:15]
                    if name_a_short in rname and name_b_short in rname:
                        found = True
                        # CRITICAL: правило скидки 10% — Bundle_Price = Combined * 0.90.
                        ok_bp = (num_close(r[4], bundle_price, 0.5)
                                 or num_close(r[4], bundle_price_reg, 0.5))
                        record(f"CRITICAL Bundle price rank {i+1}", ok_bp,
                               f"Expected ~{bundle_price} (line items) "
                               f"or ~{bundle_price_reg} (regular_price), got {r[4]}")
                        if not ok_bp:
                            all_ok = False
                        break
            if not found:
                record(f"CRITICAL Bundle price rank {i+1} found", False,
                       "Not found in proposals")
                all_ok = False

    # --- Check Category Insights sheet ---
    ci_sheet = None
    for name in wb.sheetnames:
        if "category" in name.lower() or "insight" in name.lower():
            ci_sheet = wb[name]
            break
    if ci_sheet is None:
        record("Sheet 'Category Insights' exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Sheet 'Category Insights' exists", True)
        rows = list(ci_sheet.iter_rows(min_row=2, values_only=True))
        record("Category Insights has data rows", len(rows) >= 3,
               f"Got {len(rows)} rows")

        # Check that Electronics appears (most common category).
        # После russification (центральный map в db/zzz_wc_after_init.sql)
        # product_categories.name 'Electronics' -> 'Электроника', поэтому
        # агент читает и пишет 'Электроника'. Принимаем оба варианта.
        has_electronics = any(
            r and r[0] and (
                "electronics" in str(r[0]).lower()
                or "электроник" in str(r[0]).lower()
            )
            for r in rows
        )
        record("Electronics/Электроника category in insights", has_electronics)
        if not has_electronics:
            all_ok = False

    return all_ok


def check_pptx(agent_workspace):
    """Check Bundle_Presentation.pptx."""
    print("\n=== Checking Bundle_Presentation.pptx ===")
    pptx_file = os.path.join(agent_workspace, "Bundle_Presentation.pptx")
    if not os.path.isfile(pptx_file):
        record("PPTX file exists", False, f"Not found: {pptx_file}")
        return False
    record("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_file)
    except Exception as e:
        record("PPTX readable", False, str(e))
        return False

    slide_count = len(prs.slides)
    record("PPTX has >= 5 slides", slide_count >= 5, f"Got {slide_count}")

    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += shape.text.lower() + " "

    record("PPTX mentions 'bundle'", "bundle" in all_text)
    record("PPTX mentions 'revenue' or 'price'",
           "revenue" in all_text or "price" in all_text)
    record("PPTX mentions 'co-purchase' or 'co purchase' or 'purchase'",
           "co-purchase" in all_text or "co purchase" in all_text or "purchase" in all_text)
    record("PPTX mentions 'timeline' or 'implementation' or 'launch'",
           "timeline" in all_text or "implementation" in all_text or "launch" in all_text)

    return slide_count >= 5


def check_calendar():
    """Check Google Calendar for 3 bundle launch meetings."""
    print("\n=== Checking Google Calendar ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        # Агент может озаглавить встречи по-русски, поэтому матчим EN+RU
        # ключевые слова (бандл/запуск/приоритет/категор).
        cur.execute("""
            SELECT summary, start_datetime, end_datetime, description
            FROM gcal.events
            WHERE LOWER(summary) LIKE '%bundle%'
               OR LOWER(summary) LIKE '%launch%'
               OR LOWER(summary) LIKE '%priority%'
               OR LOWER(summary) LIKE '%category%'
               OR LOWER(summary) LIKE '%бандл%'
               OR LOWER(summary) LIKE '%запуск%'
               OR LOWER(summary) LIKE '%приоритет%'
               OR LOWER(summary) LIKE '%категор%'
        """)
        events = cur.fetchall()
        record("At least 3 bundle-related calendar events",
               len(events) >= 3,
               f"Found {len(events)} matching events")

        # CRITICAL: три конкретные даты запуска — Пн 16, Ср 18, Пт 20 марта 2026,
        # каждая начинается в 10:00 и длится ~1 час. Это содержательное ядро
        # требования (а не просто «событие на неделе»).
        target_dates = ("2026-03-16", "2026-03-18", "2026-03-20")
        for td in target_dates:
            ok = False
            for ev in events:
                if not (ev[1] and ev[2]):
                    continue
                start = ev[1]
                start_str = str(start)
                if not start_str.startswith(td):
                    continue
                # начало в 10:00 (час == 10)
                hour_ok = getattr(start, "hour", None) == 10
                duration = (ev[2] - ev[1]).total_seconds()
                dur_ok = 3000 <= duration <= 4500  # 50–75 мин (~1 час)
                if hour_ok and dur_ok:
                    ok = True
                    break
            record(f"CRITICAL Calendar date {td} 10:00 ~1h", ok,
                   "Нет события на эту дату в 10:00 длительностью ~1 час")

        cur.close()
        conn.close()
        return len(events) >= 3
    except Exception as e:
        record("Calendar check", False, str(e))
        return False


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=True)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)
    check_calendar()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\n=== Results: {PASS_COUNT}/{total} passed ({accuracy:.1f}%) ===")
    if CRITICAL_FAILS:
        print(f"Critical fails: {CRITICAL_FAILS}")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
        "critical_fails": CRITICAL_FAILS,
    }
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    # Любой провал критичного чека => FAIL независимо от accuracy.
    if CRITICAL_FAILS:
        print(f"FAIL: критичные чеки провалены ({len(CRITICAL_FAILS)})")
        sys.exit(1)
    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
