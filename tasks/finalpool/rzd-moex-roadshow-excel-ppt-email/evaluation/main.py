"""
Evaluation для задачи roadshow (РЖД + moex-finance).

Источники данных:
- РЖД (схема rzd.*): поезд 752А «Сапсан» Москва-Ленинградская -> Санкт-Петербург-
  Главный, 2026-03-10, отправление 06:50, прибытие 10:50, длительность 04:00,
  бизнес-класс (Бизнес) 14000 руб.
- moex-finance (схема moex.*): тикер MGNT.ME (ПАО Магнит / «Магнит»), сектор
  Consumer Defensive. По последним 5 засеянным дням (2026-05-20..05-26):
  latest close 4182.19, 5_Day_High 4182.19, 5_Day_Low 4025.75, 5_Day_Avg 4082.92.

CRITICAL_CHECKS — семантические проверки сути результата. Любой их провал =>
общий FAIL независимо от accuracy (sys.exit(1) до порога 70%).
Структурные проверки (наличие листов/слайдов, заголовки) — НЕ критические.
"""
import json
import os
import re
import sys
from argparse import ArgumentParser

import psycopg2
import openpyxl


def normalize_ru_numbers(text):
    """RU number normalization for substring/regex checks: collapse digit-group
    separators (space/NBSP/NNBSP/dot/comma before a 3-digit group) and turn
    decimal commas into dots ("31 588" -> "31588", "4 586,91" -> "4586.91")."""
    t = str(text or "")
    t = re.sub(r"(?<=\d)[ \xa0\u202f\u2009.,](?=\d{3}\b)", "", t)
    return re.sub(r"(?<=\d),(?=\d)", ".", t)


# --- Эталонные значения из сидов ---
TRAIN_CODE = "752А"          # кириллица; для сравнения используем normalize_train()
BIZ_PRICE = 14000.0          # руб., бизнес-класс
LATEST_CLOSE = 4182.19
FIVE_DAY_AVG = 4082.92
FIVE_DAY_HIGH = 4182.19
FIVE_DAY_LOW = 4025.75
TO_ADDR = "roadshow@bank.com"


def normalize_train(s):
    """Кириллица → латиница для номеров поездов (752А ≈ 752A)."""
    if s is None:
        return ""
    s = str(s).strip().upper()
    cyr_to_lat = str.maketrans({
        "А": "A", "В": "B", "Е": "E", "К": "K", "М": "M",
        "Н": "H", "О": "O", "Р": "P", "С": "C", "Т": "T",
        "У": "Y", "Х": "X", "Г": "G",
    })
    return s.translate(cyr_to_lat)


def num_close(a, b, rel_tol=0.05, abs_tol=2.0):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


def num_eq(a, b, abs_tol=1.0):
    """Строгое сравнение значения метрики (abs_tol-only, без rel_tol),
    чтобы 5_Day_Avg/High/Low различались по конкретному полю."""
    try:
        return abs(float(a) - float(b)) <= abs_tol
    except (TypeError, ValueError):
        return False


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "cowork_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []

# Критические семантические проверки
CRITICAL_CHECKS = {
    "Travel_Details: поезд 752А",
    "Travel_Details: цена бизнес-класса 14000 руб.",
    "Stock_Summary: тикер MGNT",
    "Stock_Summary: сектор Consumer Defensive",
    "Stock_Summary: последняя цена закрытия ~4182.19",
    "Stock_Summary: 5_Day_Avg ~4082.92",
    "Email: тело содержит 752А, 14000 и последнюю цену закрытия MGNT",
}


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        marker = " [CRITICAL]" if name in CRITICAL_CHECKS else ""
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL]{marker} {name}{msg}")
        if name in CRITICAL_CHECKS:
            CRITICAL_FAILS.append(name)


def check_excel(agent_workspace):
    print("\n=== Проверка 1: Roadshow_Prep.xlsx ===")

    xlsx_path = os.path.join(agent_workspace, "Roadshow_Prep.xlsx")
    if not os.path.exists(xlsx_path):
        record("Roadshow_Prep.xlsx существует", False, f"Не найден: {xlsx_path}")
        return
    record("Roadshow_Prep.xlsx существует", True)

    try:
        # data_only=True: если агент записал формулы (=AVERAGE(...)), читаем
        # вычисленные значения, а не строку формулы.
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    except Exception as e:
        record("Excel читается", False, str(e))
        return
    record("Excel читается", True)

    sheet_names_lower = [s.lower() for s in wb.sheetnames]

    # --- Travel_Details ---
    if "travel_details" not in sheet_names_lower:
        record("Лист Travel_Details существует", False, f"Листы: {wb.sheetnames}")
    else:
        record("Лист Travel_Details существует", True)
        ws = wb[wb.sheetnames[sheet_names_lower.index("travel_details")]]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = [r for r in rows[1:] if any(c for c in r)]
        record("Travel_Details: одна строка данных", len(data_rows) >= 1,
               f"Найдено {len(data_rows)}")

        all_text = " ".join(str(c) for r in rows for c in r if c)
        all_text_norm = normalize_train(all_text)
        # CRITICAL: номер поезда 752А (с учётом кир/лат)
        record("Travel_Details: поезд 752А",
               normalize_train(TRAIN_CODE) in all_text_norm, all_text[:200])
        # Класс обслуживания: Бизнес или Business (некритично)
        low = all_text.lower()
        record("Travel_Details: класс Бизнес/Business",
               "бизнес" in low or "business" in low, all_text[:200])

        numeric_vals = []
        for r in data_rows:
            for c in r:
                try:
                    numeric_vals.append(float(c))
                except (TypeError, ValueError):
                    pass
        # CRITICAL: цена бизнес-класса 14000
        has_price = any(abs(v - BIZ_PRICE) < 1.0 for v in numeric_vals)
        record("Travel_Details: цена бизнес-класса 14000 руб.", has_price,
               f"Числа: {numeric_vals}")

    # --- Stock_Summary ---
    if "stock_summary" not in sheet_names_lower:
        record("Лист Stock_Summary существует", False, f"Листы: {wb.sheetnames}")
    else:
        record("Лист Stock_Summary существует", True)
        ws2 = wb[wb.sheetnames[sheet_names_lower.index("stock_summary")]]
        rows2 = list(ws2.iter_rows(values_only=True))
        all_text2 = " ".join(str(c) for r in rows2 for c in r if c)
        up2 = all_text2.upper()
        # CRITICAL: тикер MGNT
        record("Stock_Summary: тикер MGNT", "MGNT" in up2, all_text2[:200])
        # CRITICAL: сектор Consumer Defensive
        record("Stock_Summary: сектор Consumer Defensive",
               "consumer defensive" in all_text2.lower(), all_text2[:200])

        # Лист Stock_Summary имеет колонки Metric | Value. Делаем построчный
        # lookup Metric -> первое числовое значение в строке, чтобы каждое поле
        # (Latest_Close/High/Low/Avg) проверялось ПО СВОЕЙ ячейке, а не «есть ли
        # где-то такое число». Иначе одно значение 4182.19 при широком допуске
        # удовлетворяло бы сразу всем метрикам.
        def first_num_in_row(r):
            for c in r:
                try:
                    return float(c)
                except (TypeError, ValueError):
                    continue
            return None

        metric_value = {}
        for r in rows2:
            if not r:
                continue
            label = str(r[0]).strip().lower() if r[0] is not None else ""
            if not label:
                continue
            metric_value[label] = first_num_in_row(r[1:])

        # abs_tol-only сравнение конкретной метрики со своим эталоном.
        # CRITICAL: последняя цена закрытия 4182.19 (строка Latest_Close_Price)
        record("Stock_Summary: последняя цена закрытия ~4182.19",
               num_eq(metric_value.get("latest_close_price"), LATEST_CLOSE),
               f"Latest_Close_Price={metric_value.get('latest_close_price')}")
        # CRITICAL: 5_Day_Avg 4082.92 (строка 5_Day_Avg)
        record("Stock_Summary: 5_Day_Avg ~4082.92",
               num_eq(metric_value.get("5_day_avg"), FIVE_DAY_AVG),
               f"5_Day_Avg={metric_value.get('5_day_avg')}")
        # 5-day high/low (некритично) — тоже построчно
        record("Stock_Summary: 5_Day_High ~4182.19",
               num_eq(metric_value.get("5_day_high"), FIVE_DAY_HIGH),
               f"5_Day_High={metric_value.get('5_day_high')}")
        record("Stock_Summary: 5_Day_Low ~4025.75",
               num_eq(metric_value.get("5_day_low"), FIVE_DAY_LOW),
               f"5_Day_Low={metric_value.get('5_day_low')}")


def check_ppt(agent_workspace):
    print("\n=== Проверка 2: Roadshow_Agenda.pptx ===")

    pptx_path = os.path.join(agent_workspace, "Roadshow_Agenda.pptx")
    if not os.path.exists(pptx_path):
        record("Roadshow_Agenda.pptx существует", False, f"Не найден: {pptx_path}")
        return
    record("Roadshow_Agenda.pptx существует", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        record("PPTX читается", False, str(e))
        return
    record("PPTX читается", True)

    record("PPTX: ровно 4 слайда", len(prs.slides) == 4, f"Найдено {len(prs.slides)}")

    all_titles = []
    for slide in prs.slides:
        title_shape = slide.shapes.title
        if title_shape:
            all_titles.append(title_shape.text)

    titles_lower = " ".join(all_titles).lower()
    record("Слайд: заголовок Investor Roadshow",
           "roadshow" in titles_lower, f"Заголовки: {all_titles}")
    record("Слайд: заголовок Journey Details",
           "journey" in titles_lower, f"Заголовки: {all_titles}")
    record("Слайд: заголовок Portfolio Overview",
           "portfolio" in titles_lower, f"Заголовки: {all_titles}")
    record("Слайд: заголовок Meeting Schedule",
           "meeting" in titles_lower or "schedule" in titles_lower,
           f"Заголовки: {all_titles}")

    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    all_text += para.text + " "
    all_text_norm = normalize_train(all_text)
    record("PPTX: упоминается поезд 752А",
           normalize_train(TRAIN_CODE) in all_text_norm, all_text[:200])
    record("PPTX: упоминается MGNT",
           "MGNT" in all_text.upper(), all_text[:200])
    record("PPTX: слайд 4 содержит TBD",
           "tbd" in all_text.lower(), all_text[-200:])


def _to_str(to_addr):
    if to_addr is None:
        return ""
    if isinstance(to_addr, list):
        return " ".join(str(r).lower() for r in to_addr)
    try:
        parsed = json.loads(str(to_addr))
        if isinstance(parsed, list):
            return " ".join(str(r).lower() for r in parsed)
    except Exception:
        pass
    return str(to_addr).lower()


def check_email():
    print("\n=== Проверка 3: письмо на roadshow@bank.com ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT subject, to_addr, body_text FROM email.messages
        WHERE subject ILIKE '%roadshow%' OR subject ILIKE '%travel%prep%'
    """)
    messages = cur.fetchall()
    cur.close()
    conn.close()

    record("Письмо роадшоу отправлено", len(messages) >= 1,
           f"Найдено {len(messages)} писем")

    if not messages:
        record("Письмо отправлено на roadshow@bank.com", False, "нет писем")
        record("Email: тело содержит 752А, 14000 и последнюю цену закрытия MGNT",
               False, "нет писем")
        return

    target = None
    for m in messages:
        if TO_ADDR in _to_str(m[1]):
            target = m
            break
    record("Письмо отправлено на roadshow@bank.com", target is not None,
           f"Адресаты: {[_to_str(m[1])[:60] for m in messages]}")
    if target is None:
        target = messages[0]

    body = str(target[2] or "")
    body_norm = normalize_train(body)
    body_digits = normalize_ru_numbers(body)
    has_train = normalize_train(TRAIN_CODE) in body_norm
    # цена бизнес-класса: 14000 (RU-разряды схлопнуты нормализацией)
    has_price = "14000" in body_digits
    # последняя цена закрытия MGNT ~4182.19 (целая часть 4182 как минимум)
    has_close = "4182" in body_digits
    record("Email: тело содержит 752А, 14000 и последнюю цену закрытия MGNT",
           has_train and has_price and has_close,
           f"752А={has_train}, 14000={has_price}, 4182={has_close}; body={body[:150]}")


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_ppt(args.agent_workspace)
    check_email()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: проверки не выполнены.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nИтого: {PASS_COUNT}/{total} проверок пройдено ({accuracy:.1f}%)")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
        "critical_fails": CRITICAL_FAILS,
    }
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if CRITICAL_FAILS:
        print(f"=== RESULT: FAIL (провалены критические проверки: {CRITICAL_FAILS}) ===")
        sys.exit(1)

    if accuracy >= 70:
        print("=== RESULT: PASS ===")
        sys.exit(0)
    print("=== RESULT: FAIL (accuracy < 70%) ===")
    sys.exit(1)


if __name__ == "__main__":
    main()
