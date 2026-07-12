"""Evaluation for insales-shipping-zone-ppt."""
import argparse
import os
import sys

import openpyxl
import psycopg2

DB = {"host": os.environ.get("PGHOST", "localhost"), "port": 5432, "dbname": "cowork_gym", "user": "eigent", "password": "camel"}

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []

# Критичные проверки: провал любой => FAIL независимо от accuracy.
CRITICAL_CHECKS = {
    "CRITICAL: 'Доставка по РФ' Order_Count + Total_Shipping_Cost",
    "CRITICAL: 'Москва' Order_Count + Total_Shipping_Cost",
    "CRITICAL: 'Международная' Order_Count + Total_Shipping_Cost",
    "CRITICAL: Summary Total_Orders / Total_Shipping_Revenue / Zones_Count",
    "CRITICAL: Google Sheet 'Shipping Performance Dashboard' с данными по зонам",
    "CRITICAL: PPTX содержит заголовок 'Shipping Zone Performance Review' и дату 2026-03-06",
}


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        marker = " [CRITICAL]" if name in CRITICAL_CHECKS else ""
        print(f"  [FAIL]{marker} {name}{msg}")
        if name in CRITICAL_CHECKS:
            CRITICAL_FAILS.append(name)


def num_close(a, b, tol=0.5):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def get_expected_zone_data():
    """Compute expected zone performance from read-only DB."""
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()
    cur.execute("""
        SELECT
          CASE
            WHEN shipping->>'country' = 'Россия' AND shipping->>'city' = 'Москва' THEN 'Москва'
            WHEN shipping->>'country' = 'Россия' THEN 'Доставка по РФ'
            ELSE 'Международная'
          END as zone_name,
          COUNT(*) as order_count,
          ROUND(SUM(shipping_total)::numeric, 2) as total_shipping,
          ROUND(AVG(shipping_total)::numeric, 2) as avg_shipping
        FROM wc.orders
        GROUP BY 1
        ORDER BY COUNT(*) DESC
    """)
    rows = cur.fetchall()
    cur.close()
    conn.close()
    return rows


def check_excel(agent_workspace, groundtruth_workspace):
    """Check Shipping_Zone_Report.xlsx."""
    print("\n=== Checking Shipping_Zone_Report.xlsx ===")

    agent_file = os.path.join(agent_workspace, "Shipping_Zone_Report.xlsx")
    if not os.path.isfile(agent_file):
        record("Excel file exists", False, f"Not found: {agent_file}")
        return False
    record("Excel file exists", True)

    try:
        wb = openpyxl.load_workbook(agent_file, data_only=True)
    except Exception as e:
        record("Excel readable", False, str(e))
        return False

    all_ok = True
    expected_zones = get_expected_zone_data()

    # Check Zone Performance sheet
    zp_sheet = None
    for name in wb.sheetnames:
        if "zone" in name.lower() and "perform" in name.lower():
            zp_sheet = wb[name]
            break
    if zp_sheet is None:
        record("Sheet 'Zone Performance' exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Sheet 'Zone Performance' exists", True)
        rows = list(zp_sheet.iter_rows(min_row=2, values_only=True))
        record("Zone Performance has data rows", len(rows) >= 1, f"Got {len(rows)} rows")

        for ez in expected_zones:
            zone_name, exp_count, exp_total, exp_avg = ez
            found = False
            zone_ok = True
            for r in rows:
                if r and r[0] and zone_name.lower() in str(r[0]).lower():
                    found = True
                    ok_count = num_close(r[1], exp_count, 1)
                    record(f"'{zone_name}' Order_Count", ok_count,
                           f"Expected {exp_count}, got {r[1]}")
                    if not ok_count:
                        all_ok = False
                        zone_ok = False
                    ok_total = num_close(r[2], exp_total, 1.0)
                    record(f"'{zone_name}' Total_Shipping_Cost", ok_total,
                           f"Expected {exp_total}, got {r[2]}")
                    if not ok_total:
                        all_ok = False
                        zone_ok = False
                    ok_avg = num_close(r[3], exp_avg, 0.5)
                    record(f"'{zone_name}' Avg_Shipping_Cost", ok_avg,
                           f"Expected {exp_avg}, got {r[3]}")
                    if not ok_avg:
                        all_ok = False
                    break
            if not found:
                record(f"Zone '{zone_name}' found in sheet", False, "Missing")
                all_ok = False
                zone_ok = False
            # CRITICAL: ключевые числа по каждой зоне совпадают с live-расчётом
            record(f"CRITICAL: '{zone_name}' Order_Count + Total_Shipping_Cost",
                   zone_ok, "Zone numbers mismatch or zone missing")
            if not zone_ok:
                all_ok = False

    # Check Summary sheet
    sum_sheet = None
    for name in wb.sheetnames:
        if "summary" in name.lower():
            sum_sheet = wb[name]
            break
    if sum_sheet is None:
        record("Sheet 'Summary' exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Sheet 'Summary' exists", True)
        summary = {}
        for row in sum_sheet.iter_rows(min_row=2, values_only=True):
            if row and row[0]:
                summary[str(row[0]).strip().lower()] = row[1]

        total_orders = sum(ez[1] for ez in expected_zones)
        total_revenue = sum(ez[2] for ez in expected_zones)
        zones_count = len(expected_zones)

        summary_ok = True
        seen_orders = seen_revenue = seen_zones = False
        for key in summary:
            if "total_order" in key or "total order" in key:
                seen_orders = True
                ok = num_close(summary[key], total_orders, 1)
                record("Summary Total_Orders", ok,
                       f"Expected {total_orders}, got {summary[key]}")
                if not ok:
                    all_ok = False
                    summary_ok = False
            elif "total_shipping" in key or "shipping_revenue" in key or "total shipping" in key:
                seen_revenue = True
                ok = num_close(summary[key], total_revenue, 1.0)
                record("Summary Total_Shipping_Revenue", ok,
                       f"Expected {total_revenue}, got {summary[key]}")
                if not ok:
                    all_ok = False
                    summary_ok = False
            elif "zones_count" in key or ("zone" in key and "count" in key):
                seen_zones = True
                ok = num_close(summary[key], zones_count, 0)
                record("Summary Zones_Count", ok,
                       f"Expected {zones_count}, got {summary[key]}")
                if not ok:
                    all_ok = False
                    summary_ok = False

        if not (seen_orders and seen_revenue and seen_zones):
            summary_ok = False
            all_ok = False
        record("CRITICAL: Summary Total_Orders / Total_Shipping_Revenue / Zones_Count",
               summary_ok,
               f"orders={seen_orders} revenue={seen_revenue} zones={seen_zones}")
        if not summary_ok:
            all_ok = False

    return all_ok


def check_pptx(agent_workspace):
    """Check Shipping_Review.pptx."""
    print("\n=== Checking Shipping_Review.pptx ===")
    from pptx import Presentation

    pptx_file = os.path.join(agent_workspace, "Shipping_Review.pptx")
    if not os.path.isfile(pptx_file):
        record("PPTX file exists", False, f"Not found: {pptx_file}")
        return False
    record("PPTX file exists", True)

    try:
        prs = Presentation(pptx_file)
    except Exception as e:
        record("PPTX readable", False, str(e))
        return False

    all_ok = True
    slide_count = len(prs.slides)
    ok_slides = slide_count >= 4
    record("PPTX has >= 4 slides", ok_slides, f"Got {slide_count}")
    if not ok_slides:
        all_ok = False

    all_text = ""
    raw_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += shape.text.lower() + " "
                raw_text += shape.text + " "

    # Тематические ключи (RU+EN), т.к. слайды могут быть на русском
    ok_topic = ("shipping" in all_text or "доставк" in all_text)
    record("PPTX упоминает доставку (shipping/доставка)", ok_topic,
           "No mention of shipping/доставка")
    if not ok_topic:
        all_ok = False
    ok_perf = (any(k in all_text for k in
                   ("performance", "review", "эффективн", "обзор", "производительн")))
    record("PPTX упоминает эффективность/обзор (performance/review/эффективность/обзор)",
           ok_perf, "No mention of performance/review/эффективность/обзор")
    if not ok_perf:
        all_ok = False

    # CRITICAL: точный заголовок (английский идентификатор) + дата подзаголовка
    has_title = "Shipping Zone Performance Review" in raw_text
    has_date = "2026-03-06" in raw_text
    record("CRITICAL: PPTX содержит заголовок 'Shipping Zone Performance Review' и дату 2026-03-06",
           has_title and has_date,
           f"title={has_title} date={has_date}")
    if not (has_title and has_date):
        all_ok = False

    return all_ok


def check_gsheet():
    """Check Google Sheet with shipping data."""
    print("\n=== Checking Google Sheet ===")
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()

    cur.execute("SELECT id, title FROM gsheet.spreadsheets WHERE title ILIKE '%shipping%'")
    rows = cur.fetchall()
    if not rows:
        record("GSheet with 'shipping' in title", False, "No matching spreadsheet found")
        record("CRITICAL: Google Sheet 'Shipping Performance Dashboard' с данными по зонам",
               False, "No matching spreadsheet found")
        cur.close()
        conn.close()
        return False
    record("GSheet with 'shipping' in title", True)

    ss_id = rows[0][0]
    cur.execute("SELECT id FROM gsheet.sheets WHERE spreadsheet_id = %s", (ss_id,))
    sheets = cur.fetchall()
    has_sheet = len(sheets) >= 1
    record("GSheet has at least one sheet", has_sheet, f"Got {len(sheets)}")

    row_count = 0
    if sheets:
        sheet_id = sheets[0][0]
        cur.execute("SELECT COUNT(DISTINCT row_index) FROM gsheet.cells WHERE spreadsheet_id = %s AND sheet_id = %s",
                    (ss_id, sheet_id))
        row_count = cur.fetchone()[0]
        record("GSheet has data rows", row_count >= 2, f"Got {row_count} rows")

    cur.close()
    conn.close()

    # CRITICAL: дашборд существует с листом и данными по зонам (>= заголовок + 1 строка зоны)
    gs_ok = has_sheet and row_count >= 2
    record("CRITICAL: Google Sheet 'Shipping Performance Dashboard' с данными по зонам",
           gs_ok, f"has_sheet={has_sheet} rows={row_count}")
    return gs_ok


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    excel_ok = check_excel(args.agent_workspace, gt_dir)
    pptx_ok = check_pptx(args.agent_workspace)
    gsheet_ok = check_gsheet()

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100) if total else 0.0

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Accuracy: {accuracy:.1f}%")
    if CRITICAL_FAILS:
        print(f"  Critical fails: {CRITICAL_FAILS}")

    overall = (not CRITICAL_FAILS) and accuracy >= 70 and excel_ok and pptx_ok and gsheet_ok
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
