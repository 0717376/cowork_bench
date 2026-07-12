"""Evaluation for insales-coupon-performance-notion-ppt."""
import argparse
import os
import sys
import psycopg2
import openpyxl


DB = {"host": os.environ.get("PGHOST", "localhost"), "port": 5432, "dbname": "cowork_gym", "user": "eigent", "password": "camel"}


# Семантические проверки сути результата. Любой провал критической проверки => FAIL,
# независимо от итоговой accuracy.
CRITICAL_CHECKS = {
    "Coupon Analysis: топ-промокоды с верными Times_Used и сортировка по убыванию",
    "Coupon Analysis: SAVE20 Usage_Rate_Pct и корректное N/A для пустого лимита",
    "Summary: все четыре метрики корректны (incl. Highest_Discount)",
    "Teamly: страница стратегии по промокодам с упоминанием лидирующего промокода",
    "Email: письмо на marketing@company.com отправлено (тема coupon/campaign)",
}

PASS_COUNT = 0
FAIL_COUNT = 0
FAILED_NAMES = []


def check(name, ok, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if ok:
        PASS_COUNT += 1
        print(f"    PASS: {name}")
    else:
        FAIL_COUNT += 1
        FAILED_NAMES.append(name)
        print(f"    FAIL: {name}" + (f" — {detail}" if detail else ""))


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def get_coupon_data():
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()
    cur.execute("SELECT code, discount_type, amount, usage_count, usage_limit FROM wc.coupons ORDER BY usage_count DESC")
    rows = cur.fetchall()
    cur.close()
    conn.close()
    total_uses = sum(r[3] for r in rows)
    most_used = max(rows, key=lambda x: x[3])
    highest_disc = max(rows, key=lambda x: float(x[2]))
    by_code = {}
    for code, dtype, amount, used, limit in rows:
        rate = "N/A"
        if limit is not None:
            try:
                rate = round(float(used) / float(limit) * 100.0, 2)
            except (TypeError, ZeroDivisionError):
                rate = "N/A"
        by_code[str(code).strip().upper()] = {
            "discount_type": dtype, "amount": amount, "used": used,
            "limit": limit, "rate": rate,
        }
    return {
        "coupons": rows,
        "by_code": by_code,
        "total_coupons": len(rows),
        "most_used": most_used[0],
        "highest_discount": highest_disc[0],
        "total_uses": total_uses,
    }


def check_teamly(top_coupon=None):
    """Return (page_found, mentions_top_coupon).

    The agent-created page lives in teamly.pages (id > 3 are non-seed pages).
    Title may be Russian and/or English; we accept either the English anchor
    'coupon strategy overview 2026' or the Russian 'обзор стратегии по промокодам 2026'.
    """
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("SELECT to_regclass('teamly.pages')")
        if not cur.fetchone()[0]:
            cur.close()
            conn.close()
            return False, False
        cur.execute("SELECT id, title, COALESCE(body, '') FROM teamly.pages WHERE id > 3")
        pages = cur.fetchall()
        cur.close()
        conn.close()
    except Exception:
        return False, False

    title_markers = ["coupon strategy overview 2026", "обзор стратегии по промокодам 2026"]
    found = None
    for pid, title, body in pages:
        tl = (title or "").lower()
        if any(m in tl for m in title_markers):
            found = (title, body)
            break
    if found is None:
        return False, False
    if top_coupon is None:
        return True, True
    combined = (found[0] + " " + found[1]).upper()
    return True, (top_coupon.upper() in combined)


def check_email_sent():
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("SELECT count(*) FROM email.messages WHERE LOWER(to_addr::text) LIKE '%marketing%' AND (LOWER(subject) LIKE '%coupon%' OR LOWER(subject) LIKE '%campaign%')")
        cnt = cur.fetchone()[0]
        cur.close()
        conn.close()
        return cnt >= 1
    except Exception:
        return False


def check_ppt_file(agent_workspace):
    ppt_path = os.path.join(agent_workspace, "Coupon_Analysis_Presentation.pptx")
    if not os.path.exists(ppt_path):
        return False, "Coupon_Analysis_Presentation.pptx not found"
    try:
        from pptx import Presentation
        prs = Presentation(ppt_path)
        if len(prs.slides) < 7:
            return False, f"PPT has only {len(prs.slides)} slides, expected >= 7 (1 title + coupons + 1 summary)"
        return True, ""
    except Exception as e:
        # If pptx not importable, just check file exists
        return True, ""


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    agent_file = os.path.join(args.agent_workspace, "Coupon_Performance.xlsx")
    gt_file = os.path.join(gt_dir, "Coupon_Performance.xlsx")

    if not os.path.exists(agent_file):
        print(f"FAIL: Agent output not found: {agent_file}")
        sys.exit(1)
    if not os.path.exists(gt_file):
        print(f"FAIL: Groundtruth not found: {gt_file}")
        sys.exit(1)

    # Источник истины — живые данные из wc.coupons. Критические проверки зависят
    # от БД, поэтому НЕ маскируем сбой захардкоженным fallback'ом.
    try:
        coupon_data = get_coupon_data()
    except Exception as e:
        print(f"FAIL: Could not query wc.coupons: {e}")
        sys.exit(1)

    agent_wb = openpyxl.load_workbook(agent_file, data_only=True)
    gt_wb = openpyxl.load_workbook(gt_file, data_only=True)

    # ---- Coupon Analysis sheet ----
    print("  Checking Coupon Analysis sheet...")
    a_rows = load_sheet_rows(agent_wb, "Coupon Analysis")
    a_lookup = {}
    data_rows = []
    if a_rows is None:
        check("Coupon Analysis sheet exists", False, "лист не найден")
    else:
        check("Coupon Analysis sheet exists", True)
        data_rows = [r for r in a_rows[1:] if r and any(c is not None for c in r)]
        check("Coupon Analysis has all coupon rows",
              len(data_rows) >= coupon_data["total_coupons"],
              f"{len(data_rows)} rows, expected >= {coupon_data['total_coupons']}")
        for row in data_rows:
            if row and row[0] is not None:
                a_lookup[str(row[0]).strip().upper()] = row

    # CRITICAL: топовые промокоды с верным Times_Used + сортировка по убыванию.
    top_ok = True
    top_detail = []
    for code, expected in (("HOLIDAY30", 50), ("SAVE20", 39), ("VIP20", 41)):
        row = a_lookup.get(code)
        if not row or len(row) < 4 or not num_close(row[3], expected, 0):
            top_ok = False
            got = (row[3] if row and len(row) >= 4 else "<missing>")
            top_detail.append(f"{code}.Times_Used={got} vs {expected}")
    # Сортировка Times_Used по убыванию (4-я колонка, индекс 3).
    used_seq = [r[3] for r in data_rows if r and len(r) >= 4 and isinstance(r[3], (int, float))]
    sorted_desc = all(used_seq[i] >= used_seq[i + 1] for i in range(len(used_seq) - 1)) if used_seq else False
    if not sorted_desc:
        top_ok = False
        top_detail.append(f"не отсортировано по убыванию Times_Used: {used_seq}")
    check("Coupon Analysis: топ-промокоды с верными Times_Used и сортировка по убыванию",
          top_ok, "; ".join(top_detail))

    # CRITICAL: SAVE20 usage rate ≈ 78.00 и N/A для пустого usage_limit.
    rate_ok = True
    rate_detail = []
    save20 = a_lookup.get("SAVE20")
    if not save20 or len(save20) < 6:
        rate_ok = False
        rate_detail.append("SAVE20 строка/колонка Usage_Rate_Pct отсутствует")
    else:
        rv = save20[5]
        if str(rv).strip().upper() == "N/A" or not num_close(rv, 78.00, 0.5):
            rate_ok = False
            rate_detail.append(f"SAVE20.Usage_Rate_Pct={rv} vs ~78.00")
    # N/A для промокодов с пустым usage_limit в источнике.
    na_codes = [c for c, v in coupon_data["by_code"].items() if v["limit"] is None]
    for c in na_codes:
        row = a_lookup.get(c)
        if row and len(row) >= 6 and str(row[5]).strip().upper() != "N/A":
            rate_ok = False
            rate_detail.append(f"{c}.Usage_Rate_Pct={row[5]} ожидалось N/A")
    check("Coupon Analysis: SAVE20 Usage_Rate_Pct и корректное N/A для пустого лимита",
          rate_ok, "; ".join(rate_detail))

    # ---- Summary sheet ----
    print("  Checking Summary sheet...")
    a_rows = load_sheet_rows(agent_wb, "Summary")
    if a_rows is None:
        check("Summary sheet exists", False, "лист не найден")
        check("Summary: все четыре метрики корректны (incl. Highest_Discount)", False, "нет листа Summary")
    else:
        check("Summary sheet exists", True)
        a_data = {str(r[0]).strip().lower(): r[1] for r in a_rows[1:] if r and r[0] is not None}
        sum_ok = True
        sum_detail = []

        tc = a_data.get("total_coupons")
        if tc is None or not num_close(tc, coupon_data["total_coupons"], 0):
            sum_ok = False
            sum_detail.append(f"Total_Coupons={tc} vs {coupon_data['total_coupons']}")

        mu = a_data.get("most_used_coupon")
        if mu is None or str(mu).strip().upper() != coupon_data["most_used"].upper():
            sum_ok = False
            sum_detail.append(f"Most_Used_Coupon={mu} vs {coupon_data['most_used']}")

        hd = a_data.get("highest_discount")
        if hd is None or str(hd).strip().upper() != coupon_data["highest_discount"].upper():
            sum_ok = False
            sum_detail.append(f"Highest_Discount={hd} vs {coupon_data['highest_discount']}")

        tu = a_data.get("total_coupon_uses")
        if tu is None or not num_close(tu, coupon_data["total_uses"], 0):
            sum_ok = False
            sum_detail.append(f"Total_Coupon_Uses={tu} vs {coupon_data['total_uses']}")

        check("Summary: все четыре метрики корректны (incl. Highest_Discount)",
              sum_ok, "; ".join(sum_detail))

    # ---- PPT file ----
    print("  Checking PPT file...")
    ok, detail = check_ppt_file(args.agent_workspace)
    check("PowerPoint presentation present with enough slides", ok, detail)

    # ---- Teamly page (CRITICAL) ----
    print("  Checking Teamly page...")
    page_found, mentions_top = check_teamly(coupon_data["most_used"])
    check("Teamly: страница стратегии по промокодам с упоминанием лидирующего промокода",
          page_found and mentions_top,
          "страница не найдена" if not page_found else "не упомянут лидирующий промокод")

    # ---- Email (CRITICAL) ----
    print("  Checking email to marketing...")
    check("Email: письмо на marketing@company.com отправлено (тема coupon/campaign)",
          check_email_sent(), "письмо не найдено")

    # ---- Итог: критический гейт ПЕРЕД порогом accuracy ----
    total = PASS_COUNT + FAIL_COUNT
    accuracy = PASS_COUNT / total * 100 if total else 0
    critical_failed = [n for n in FAILED_NAMES if n in CRITICAL_CHECKS]

    print(f"\n=== Results: {PASS_COUNT}/{total} passed ({accuracy:.1f}%) ===")
    if critical_failed:
        print(f"CRITICAL FAILURES: {len(critical_failed)}")
        for n in critical_failed:
            print(f"  - {n}")

    if args.res_log_file:
        try:
            import json
            with open(args.res_log_file, "w") as f:
                json.dump({
                    "total_passed": PASS_COUNT, "total_checks": total,
                    "accuracy": accuracy, "critical_failed": critical_failed,
                }, f, indent=2)
        except Exception:
            pass

    if critical_failed or accuracy < 70:
        print("\n=== RESULT: FAIL ===")
        sys.exit(1)
    print("\n=== RESULT: PASS ===")
    sys.exit(0)


if __name__ == "__main__":
    main()
