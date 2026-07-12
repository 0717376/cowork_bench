"""Evaluation script for sf-canvas-training-impact-excel-ppt (ClickHouse, RU)."""
import os
import argparse, json, os, sys
import openpyxl

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent", "password": "camel"
}

PASS_COUNT = 0
FAIL_COUNT = 0
# Semantic checks whose failure forces FAIL regardless of accuracy.
CRITICAL_FAILS = []

def check(name, condition, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        detail_str = str(detail)[:200] if detail else ""
        tag = "FAIL-CRITICAL" if critical else "FAIL"
        print(f"  [{tag}] {name}: {detail_str}")
        if critical:
            CRITICAL_FAILS.append(name)

def safe_float(val, default=None):
    try:
        if val is None: return default
        return float(str(val).replace(",", "").replace("%", "").replace("$", "").strip())
    except (ValueError, TypeError):
        return default

def norm_text(v):
    return str(v).strip().lower() if v is not None else ""

def text_match(a, b):
    """Tolerant text match: exact, substring either way (RU or EN values)."""
    a, b = norm_text(a), norm_text(b)
    if not a or not b:
        return False
    return a == b or a in b or b in a

def get_conn():
    import psycopg2
    return psycopg2.connect(**DB_CONFIG)

def sheet_to_dicts(ws):
    """Return (headers, list-of-row-dicts keyed by lowercased header)."""
    headers = [norm_text(c.value) for c in ws[1]]
    rows = []
    for r in ws.iter_rows(min_row=2, values_only=True):
        if all(c is None for c in r):
            continue
        d = {}
        for i, h in enumerate(headers):
            if h and i < len(r):
                d[h] = r[i]
        rows.append(d)
    return headers, rows

def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT, CRITICAL_FAILS
    PASS_COUNT = 0
    FAIL_COUNT = 0
    CRITICAL_FAILS = []

    # ---- Check Training_Impact_Analysis.xlsx ----
    excel_path = os.path.join(agent_workspace, "Training_Impact_Analysis.xlsx")
    xlsx_exists = os.path.exists(excel_path)
    check("Training_Impact_Analysis.xlsx exists", xlsx_exists, critical=True)

    # Key columns used to match rows by KEY (not positional index): Cyrillic
    # collation reorders departments so index-aligned comparison is invalid.
    KEY_COL = {
        "Department_Performance": "department",
        "Course_Overview": "course_name",
        "Impact_Summary": "metric",
    }

    wb = None
    if xlsx_exists:
        wb = openpyxl.load_workbook(excel_path)
        gt_path = os.path.join(groundtruth_workspace, "Training_Impact_Analysis.xlsx")
        gt_wb = openpyxl.load_workbook(gt_path) if os.path.exists(gt_path) else None

        if gt_wb:
            for sheet_name in gt_wb.sheetnames:
                present = sheet_name in wb.sheetnames
                check(f"{sheet_name} sheet exists", present)
                if not present:
                    continue
                ws = wb[sheet_name]
                gt_ws = gt_wb[sheet_name]
                gt_headers, gt_rows = sheet_to_dicts(gt_ws)
                headers, data_rows = sheet_to_dicts(ws)

                # Headers (non-critical structural)
                for h in gt_headers:
                    if h:
                        check(f"{sheet_name} has {h} column", h in headers, f"headers: {headers[:10]}")

                # Row count (non-critical structural)
                min_rows = max(1, len(gt_rows) - 2)
                check(f"{sheet_name} has >= {min_rows} data rows",
                      len(data_rows) >= min_rows, f"got {len(data_rows)}")

                key = KEY_COL.get(sheet_name)
                if not key:
                    continue
                # Build agent index by key
                agent_by_key = {norm_text(r.get(key)): r for r in data_rows if r.get(key) is not None}

                for gt_row in gt_rows:
                    gk = norm_text(gt_row.get(key))
                    if not gk:
                        continue
                    arow = agent_by_key.get(gk)
                    # Substring-tolerant key lookup fallback (RU/EN, partial)
                    if arow is None:
                        for ak, av in agent_by_key.items():
                            if text_match(gk, ak):
                                arow = av
                                break
                    label = f"{sheet_name}[{key}={gk[:40]}]"
                    if arow is None:
                        check(f"{label} row present", False, "missing key", critical=True)
                        continue
                    # Compare every groundtruth column for this keyed row
                    for h, gv in gt_row.items():
                        if gv is None:
                            continue
                        av = arow.get(h)
                        gf = safe_float(gv)
                        af = safe_float(av)
                        if gf is not None and af is not None:
                            tol = max(0.5, abs(gf) * 0.15)
                            check(f"{label} {h} ~{gf:.2f}",
                                  abs(gf - af) <= tol, f"got {av}")
                        else:
                            check(f"{label} {h} text",
                                  text_match(gv, av),
                                  f"expected {norm_text(gv)[:50]}, got {norm_text(av)[:50]}")

    # ---- Check Training_Impact_Presentation.pptx ----
    pptx_path = os.path.join(agent_workspace, "Training_Impact_Presentation.pptx")
    pptx_exists = os.path.exists(pptx_path)
    check("Training_Impact_Presentation.pptx exists", pptx_exists, critical=True)
    agent_slide_count = 0
    if pptx_exists:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        agent_slide_count = len(prs.slides)
        gt_pptx_path = os.path.join(groundtruth_workspace, "Training_Impact_Presentation.pptx")
        gt_slide_count = 4
        if os.path.exists(gt_pptx_path):
            gt_prs = Presentation(gt_pptx_path)
            gt_slide_count = len(gt_prs.slides)
        check(f"Training_Impact_Presentation.pptx has >= {gt_slide_count} slides",
              agent_slide_count >= gt_slide_count, f"got {agent_slide_count} slides")

    # ---- Check Python script exists (terminal usage) ----
    py_files = []
    if os.path.isdir(agent_workspace):
        py_files = [f for f in os.listdir(agent_workspace) if f.endswith(".py")]
    check("Python analysis script exists", len(py_files) >= 1, f"found: {py_files}")

    # ================= CRITICAL SEMANTIC CHECKS =================
    # Load groundtruth aggregates for semantic correctness.
    gt_path = os.path.join(groundtruth_workspace, "Training_Impact_Analysis.xlsx")
    gt_wb = openpyxl.load_workbook(gt_path) if os.path.exists(gt_path) else None

    if wb is not None and gt_wb is not None:
        # --- Department_Performance: per-department Employee_Count + Avg_Performance_Rating
        if "Department_Performance" in wb.sheetnames and "Department_Performance" in gt_wb.sheetnames:
            _, gt_dp = sheet_to_dicts(gt_wb["Department_Performance"])
            _, ag_dp = sheet_to_dicts(wb["Department_Performance"])
            ag_by_dept = {norm_text(r.get("department")): r for r in ag_dp if r.get("department") is not None}
            for gtr in gt_dp:
                dk = norm_text(gtr.get("department"))
                if not dk:
                    continue
                ar = ag_by_dept.get(dk)
                if ar is None:
                    for ak, av in ag_by_dept.items():
                        if text_match(dk, ak):
                            ar = av; break
                ge = safe_float(gtr.get("employee_count"))
                ae = safe_float(ar.get("employee_count")) if ar else None
                check(f"CRITICAL Department_Performance[{dk[:30]}] Employee_Count={ge}",
                      ae is not None and ge is not None and abs(ge - ae) <= max(0.5, abs(ge) * 0.15),
                      f"got {ae}", critical=True)
                gp = safe_float(gtr.get("avg_performance_rating"))
                ap = safe_float(ar.get("avg_performance_rating")) if ar else None
                check(f"CRITICAL Department_Performance[{dk[:30]}] Avg_Performance_Rating~{gp}",
                      ap is not None and gp is not None and abs(gp - ap) <= max(0.5, abs(gp) * 0.15),
                      f"got {ap}", critical=True)

        # --- Impact_Summary: keyed metrics ---
        if "Impact_Summary" in wb.sheetnames and "Impact_Summary" in gt_wb.sheetnames:
            _, gt_is = sheet_to_dicts(gt_wb["Impact_Summary"])
            _, ag_is = sheet_to_dicts(wb["Impact_Summary"])
            gt_metrics = {norm_text(r.get("metric")): r.get("value") for r in gt_is}
            ag_metrics = {norm_text(r.get("metric")): r.get("value") for r in ag_is}

            # Total_Employees, Total_Courses: exact integer values
            for m in ("total_employees", "total_courses"):
                gv = safe_float(gt_metrics.get(m))
                av = safe_float(ag_metrics.get(m))
                check(f"CRITICAL Impact_Summary {m}={gv}",
                      av is not None and gv is not None and abs(gv - av) <= max(0.5, abs(gv) * 0.05),
                      f"got {av}", critical=True)

            # Highest/Lowest_Performing_Dept: correct RUSSIAN department names
            for m in ("highest_performing_dept", "lowest_performing_dept"):
                gv = gt_metrics.get(m)
                av = ag_metrics.get(m)
                check(f"CRITICAL Impact_Summary {m}='{norm_text(gv)[:30]}'",
                      text_match(gv, av),
                      f"expected {norm_text(gv)[:40]}, got {norm_text(av)[:40]}", critical=True)

            # Most_Popular_Course: English Canvas course name preserved
            gv = gt_metrics.get("most_popular_course")
            av = ag_metrics.get("most_popular_course")
            check(f"CRITICAL Impact_Summary most_popular_course='{norm_text(gv)[:30]}'",
                  text_match(gv, av),
                  f"expected {norm_text(gv)[:40]}, got {norm_text(av)[:40]}", critical=True)

        # --- Course_Overview: Most popular course Enrollment_Count (English course names) ---
        if "Course_Overview" in wb.sheetnames and "Course_Overview" in gt_wb.sheetnames:
            _, gt_co = sheet_to_dicts(gt_wb["Course_Overview"])
            _, ag_co = sheet_to_dicts(wb["Course_Overview"])
            ag_by_course = {norm_text(r.get("course_name")): r for r in ag_co if r.get("course_name") is not None}
            # Highest-enrollment groundtruth course
            top = None
            for r in gt_co:
                e = safe_float(r.get("enrollment_count"))
                if e is not None and (top is None or e > top[0]):
                    top = (e, norm_text(r.get("course_name")))
            if top is not None:
                te, tname = top
                ar = ag_by_course.get(tname)
                if ar is None:
                    for ak, av in ag_by_course.items():
                        if text_match(tname, ak):
                            ar = av; break
                ae = safe_float(ar.get("enrollment_count")) if ar else None
                check(f"CRITICAL Course_Overview top course '{tname[:30]}' Enrollment_Count={te}",
                      ae is not None and abs(te - ae) <= max(0.5, abs(te) * 0.15),
                      f"got {ae}", critical=True)

    # --- All three deliverables structurally present ---
    sheets_ok = wb is not None and all(
        s in wb.sheetnames for s in ("Department_Performance", "Course_Overview", "Impact_Summary"))
    check("CRITICAL Training_Impact_Analysis.xlsx has 3 required sheets",
          sheets_ok, critical=True)
    check("CRITICAL Training_Impact_Presentation.pptx has >= 4 slides",
          pptx_exists and agent_slide_count >= 4, f"got {agent_slide_count}", critical=True)
    check("CRITICAL Python analysis script present in workspace",
          len(py_files) >= 1, f"found: {py_files}", critical=True)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100) if total else 0.0
    print(f"\nAccuracy: {accuracy:.1f}% ({PASS_COUNT}/{total})")

    if CRITICAL_FAILS:
        print(f"CRITICAL FAILURES ({len(CRITICAL_FAILS)}): {CRITICAL_FAILS}")
        sys.exit(1)

    success = accuracy >= 70.0
    return success, f"Passed {PASS_COUNT}/{total} checks (accuracy {accuracy:.1f}%)"

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)
    sys.exit(0 if success else 1)

if __name__ == "__main__":
    main()
