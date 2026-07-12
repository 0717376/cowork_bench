"""Evaluation for canvas-semester-summary-ppt."""
import argparse
import os
import sys
import openpyxl
from pptx import Presentation


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")
    all_errors = []
    # Subset of checks that gate PASS/FAIL: any critical failure => FAIL regardless of accuracy.
    # We track the number of critical checks attempted and how many failed.
    critical_total = 0
    critical_failed = 0
    # Track structural (non-critical) checks for accuracy.
    struct_total = 0
    struct_failed = 0

    def crit(ok, msg):
        nonlocal critical_total, critical_failed
        critical_total += 1
        if not ok:
            critical_failed += 1
            all_errors.append("[CRITICAL] " + msg)
        return ok

    def soft(ok, msg):
        nonlocal struct_total, struct_failed
        struct_total += 1
        if not ok:
            struct_failed += 1
            all_errors.append(msg)
        return ok

    # ---- Check Excel ----
    agent_excel = os.path.join(args.agent_workspace, "Semester_Summary.xlsx")
    gt_excel = os.path.join(gt_dir, "Semester_Summary.xlsx")

    if not os.path.exists(agent_excel):
        print(f"FAIL: Agent output Semester_Summary.xlsx not found")
        sys.exit(1)
    if not os.path.exists(gt_excel):
        print(f"FAIL: Groundtruth Semester_Summary.xlsx not found")
        sys.exit(1)

    agent_wb = openpyxl.load_workbook(agent_excel, data_only=True)
    gt_wb = openpyxl.load_workbook(gt_excel, data_only=True)

    # Capture agent Course Overview data rows so Summary aggregates can be derived as a
    # fallback when the agent stored formulas (data_only=True yields None without a cache).
    agent_overview_rows = []

    # Check Course Overview sheet
    print("  Checking Course Overview...")
    a_rows = load_sheet_rows(agent_wb, "Course Overview")
    g_rows = load_sheet_rows(gt_wb, "Course Overview")
    if a_rows is None:
        crit(False, "Sheet 'Course Overview' not found in agent output")
    elif g_rows is None:
        all_errors.append("Sheet 'Course Overview' not found in groundtruth")
    else:
        a_data = a_rows[1:] if len(a_rows) > 1 else []
        g_data = g_rows[1:] if len(g_rows) > 1 else []
        agent_overview_rows = a_data
        # Structural: row count.
        soft(len(a_data) == len(g_data),
             f"Course Overview row count: agent={len(a_data)}, expected={len(g_data)}")

        a_lookup = {}
        for row in a_data:
            if row and row[0] is not None:
                a_lookup[str(row[0]).strip().lower()] = row
        for g_row in g_data:
            if not g_row or g_row[0] is None:
                continue
            key = str(g_row[0]).strip().lower()
            a_row = a_lookup.get(key)
            if not crit(a_row is not None, f"Missing course: {g_row[0]}"):
                continue
            # Col 1: Course Code (string match) -- CRITICAL semantic value
            if len(a_row) > 1 and len(g_row) > 1:
                crit(str_match(a_row[1], g_row[1]),
                     f"{key}.Course Code: {a_row[1]} vs {g_row[1]}")
            # Col 2: Total Students -- CRITICAL semantic value
            if len(a_row) > 2 and len(g_row) > 2:
                crit(num_close(a_row[2], g_row[2], 5),
                     f"{key}.Total Students: {a_row[2]} vs {g_row[2]}")
            # Col 3: Enrollments -- structural/derived
            if len(a_row) > 3 and len(g_row) > 3:
                soft(num_close(a_row[3], g_row[3], 5),
                     f"{key}.Enrollments: {a_row[3]} vs {g_row[3]}")
            # Col 4: Assignments -- structural/derived
            if len(a_row) > 4 and len(g_row) > 4:
                soft(num_close(a_row[4], g_row[4], 1),
                     f"{key}.Assignments: {a_row[4]} vs {g_row[4]}")
            # Col 5: Avg Points Possible -- CRITICAL semantic value
            if len(a_row) > 5 and len(g_row) > 5:
                crit(num_close(a_row[5], g_row[5], 2.0),
                     f"{key}.Avg Points: {a_row[5]} vs {g_row[5]}")

        if not any("Course Overview" in e or "Missing course" in e for e in all_errors):
            print("    PASS")

    # Check Summary sheet
    print("  Checking Summary...")
    a_rows = load_sheet_rows(agent_wb, "Summary")
    g_rows = load_sheet_rows(gt_wb, "Summary")
    if a_rows is None:
        crit(False, "Sheet 'Summary' not found in agent output")
    elif g_rows is None:
        all_errors.append("Sheet 'Summary' not found in groundtruth")
    else:
        a_data = a_rows[1:] if len(a_rows) > 1 else []
        g_data = g_rows[1:] if len(g_rows) > 1 else []

        # Derive Summary aggregates from the agent's own Course Overview data. Used only as a
        # fallback when a Summary cell reads None because the agent stored a formula whose
        # computed value was not cached (data_only=True). Never let that false-fail a correct
        # answer; the fallback is itself computed from the agent's data, so a non-doing agent
        # with an empty overview cannot benefit (the aggregates would be wrong/zero).
        def _num(v):
            try:
                return float(v)
            except (TypeError, ValueError):
                return None

        derived = {}
        students = [_num(r[2]) for r in agent_overview_rows if len(r) > 2 and _num(r[2]) is not None]
        enrolls = [_num(r[3]) for r in agent_overview_rows if len(r) > 3 and _num(r[3]) is not None]
        assigns = [_num(r[4]) for r in agent_overview_rows if len(r) > 4 and _num(r[4]) is not None]
        n_courses = sum(1 for r in agent_overview_rows if r and r[0] is not None)
        if n_courses:
            derived["total courses"] = n_courses
            derived["total students"] = sum(students)
            derived["total enrollments"] = sum(enrolls)
            derived["total assignments"] = sum(assigns)
            derived["average assignments per course"] = round(sum(assigns) / n_courses, 1)

        a_lookup = {}
        for row in a_data:
            if row and row[0] is not None:
                a_lookup[str(row[0]).strip().lower()] = row
        for g_row in g_data:
            if not g_row or g_row[0] is None:
                continue
            key = str(g_row[0]).strip().lower()
            a_row = a_lookup.get(key)
            if not crit(a_row is not None, f"Missing metric: {g_row[0]}"):
                continue
            # Summary aggregate values are CRITICAL semantic deliverables.
            if len(a_row) > 1 and len(g_row) > 1:
                a_val = a_row[1]
                # Formula-with-no-cache fallback: derive from agent's Course Overview.
                if a_val is None and key in derived:
                    a_val = derived[key]
                crit(num_close(a_val, g_row[1], 2.0),
                     f"{key}.Value: {a_val} vs {g_row[1]}")

        if not any("Summary" in e or "Missing metric" in e for e in all_errors):
            print("    PASS")

    # ---- Check PPT ----
    agent_ppt = os.path.join(args.agent_workspace, "Semester_Summary.pptx")
    if not os.path.exists(agent_ppt):
        crit(False, "Agent output Semester_Summary.pptx not found")
    else:
        print("  Checking Semester_Summary.pptx...")
        prs = Presentation(agent_ppt)
        slides = list(prs.slides)

        # Minimum: title + overview + 7 courses + takeaways = 10 (structural)
        soft(len(slides) >= 10,
             f"PPT has {len(slides)} slides, expected at least 10")

        # Title slide text -- 'Fall 2014' is a preserved English literal (CRITICAL).
        # Checked unconditionally; if the deck is empty the literal is absent and fails.
        title_text = ""
        if slides:
            for shape in slides[0].shapes:
                if shape.has_text_frame:
                    title_text += shape.text_frame.text.lower() + " "
        crit("fall 2014" in title_text,
             f"Title slide missing 'Fall 2014'. Found: {title_text[:100]}")

        # Last slide text (takeaways). Structural takeaway-term check (RU+EN alternatives).
        last_text = ""
        if slides:
            for shape in slides[-1].shapes:
                if shape.has_text_frame:
                    last_text += shape.text_frame.text.lower() + " "
        takeaway_terms = ["takeaway", "key", "summary", "conclusion",
                          "итог", "сводк", "вывод", "заключ"]
        soft(any(t in last_text for t in takeaway_terms),
             "Last slide missing takeaways/summary content")

        # Check that course names appear in presentation (CRITICAL: all 7 present).
        all_ppt_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_ppt_text += shape.text_frame.text.lower() + " "

        expected_courses = ["прикладная аналитика", "биохими", "креативные вычислен",
                            "на основе данных", "экологическая экономик",
                            "основы финансов", "глобальное управлени"]
        for course in expected_courses:
            crit(course in all_ppt_text, f"PPT missing course: {course}")

        # Check largest/smallest mentioned in last slide (CRITICAL semantic).
        crit("креативные вычислен" in last_text or "2498" in last_text,
             "Last slide should mention Creative Computing (largest by students)")
        crit("прикладная аналитика" in last_text or "365" in last_text,
             "Last slide should mention Applied Analytics (smallest by students)")

        if not any("PPT" in e or "ppt" in e.lower() or "slide" in e.lower() for e in all_errors):
            print("    PASS")

    # ---- Gating: critical checks first, then accuracy threshold ----
    total_checks = critical_total + struct_total
    total_failed = critical_failed + struct_failed
    accuracy = 100.0 if total_checks == 0 else 100.0 * (total_checks - total_failed) / total_checks

    if all_errors:
        print("\n--- Errors ---")
        for e in all_errors[:30]:
            print(f"  {e}")

    print(f"\nCritical: {critical_total - critical_failed}/{critical_total} passed")
    print(f"Accuracy: {accuracy:.1f}% ({total_checks - total_failed}/{total_checks})")

    if critical_failed > 0:
        print(f"\n=== RESULT: FAIL (critical check failed) ===")
        sys.exit(1)

    if accuracy < 70:
        print(f"\n=== RESULT: FAIL (accuracy {accuracy:.1f}% < 70%) ===")
        sys.exit(1)

    print("\n=== RESULT: PASS ===")
    sys.exit(0)


if __name__ == "__main__":
    main()
