"""
Скрипт оценки для задачи ppt-canvas-course-summary.

Динамически запрашивает схему Canvas в PostgreSQL для вычисления ожидаемых
значений, затем проверяет выходные файлы агента на корректность. Никакие
изменчивые значения не захардкожены — всё вычисляется из canvas.* во время
оценки.

Canvas русифицирован глобально: НАЗВАНИЕ КУРСА теперь на русском
("Прикладная аналитика и алгоритмы (Осень 2013)"), но task.md также даёт
английский вариант в кавычках — поэтому проверки названия курса принимают
ОБА варианта (RU «аналитик» / EN "applied analytics"). НАЗВАНИЯ ЗАДАНИЙ —
это коды ("TMA 1753", "CMA …", "Final Exam …"), они остаются английскими и
сопоставляются дословно. Проза слайдов/таблицы пишется агентом по-русски
(task.md русифицирован), поэтому проверки PPTX и gsheet терпимы к русскому
тексту (RU+EN).

Критические чеки (CRITICAL_CHECKS): любой их fail => задача FAIL независимо от
общей accuracy. Порог прохождения в остальном: accuracy >= 70% (и нет фейла
критического чека).
"""

from argparse import ArgumentParser
import sys
import os
from pathlib import Path

PASS_COUNT = 0
FAIL_COUNT = 0
FAILED_NAMES = []

# Критические чеки: любой fail => общий FAIL независимо от accuracy.
# Это семантическое ядро задачи: правильные значения, выведенные из Canvas,
# и идентичность основного результата (название курса + ключевые показатели).
CRITICAL_CHECKS = {
    "Excel: Enrollment Stats counts correct",
    "Excel: Assignment Performance correct",
    "Excel: Grade Distribution correct",
    "GSheet: dashboard contains correct course summary values",
    "PPTX: title slide has full course name and 2013",
    "PPTX: highest-average assignment named on findings slide",
}


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        FAILED_NAMES.append(name)
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


def get_expected_data():
    """Запрос PostgreSQL для вычисления ожидаемых данных курса из схемы Canvas."""
    import psycopg2

    conn = psycopg2.connect(
        host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="cowork_gym",
        user="eigent", password="camel"
    )
    cur = conn.cursor()

    # Название курса
    cur.execute("SELECT name FROM canvas.courses WHERE id = 1")
    row = cur.fetchone()
    course_name = row[0] if row else ""

    # Состав участников
    cur.execute("""
        SELECT type, COUNT(*)
        FROM canvas.enrollments
        WHERE course_id = 1
        GROUP BY type ORDER BY type
    """)
    enrollment_raw = dict(cur.fetchall())
    enrollments = {
        "student": enrollment_raw.get("StudentEnrollment", 0),
        "teacher": enrollment_raw.get("TeacherEnrollment", 0),
        "ta": enrollment_raw.get("TaEnrollment", 0),
    }

    # Результаты по заданиям
    cur.execute("""
        SELECT a.name,
               ROUND(AVG(s.score)::numeric, 2)::float as avg_score,
               MAX(s.score)::float as max_score,
               MIN(s.score)::float as min_score,
               COUNT(*) as sub_count
        FROM canvas.assignments a
        JOIN canvas.submissions s ON a.id = s.assignment_id
        WHERE a.course_id = 1 AND s.score IS NOT NULL
        GROUP BY a.name
        ORDER BY a.name
    """)
    assignments = cur.fetchall()

    # Распределение оценок
    cur.execute("""
        SELECT
          CASE
            WHEN avg_score >= 90 THEN 'A (90-100)'
            WHEN avg_score >= 80 THEN 'B (80-89)'
            WHEN avg_score >= 70 THEN 'C (70-79)'
            WHEN avg_score >= 60 THEN 'D (60-69)'
            ELSE 'F (<60)'
          END as grade_range,
          COUNT(*) as student_count
        FROM (
          SELECT s.user_id, AVG(s.score)::float as avg_score
          FROM canvas.submissions s
          JOIN canvas.assignments a ON s.assignment_id = a.id
          WHERE a.course_id = 1 AND s.score IS NOT NULL
          GROUP BY s.user_id
        ) sub
        GROUP BY grade_range
        ORDER BY grade_range
    """)
    grades = cur.fetchall()
    total_graded_students = sum(g[1] for g in grades)

    conn.close()

    # Производные сводные показатели для gsheet/pptx
    overall_class_avg = round(sum(a[1] for a in assignments) / len(assignments), 2) if assignments else 0.0
    n_assignments = len(assignments)
    # Задание с наивысшим средним баллом (для findings-слайда)
    top_assignment = max(assignments, key=lambda a: a[1])[0] if assignments else ""

    summary = {
        "course_name": course_name,
        "total_students": enrollments["student"],
        "n_assignments": n_assignments,
        "overall_class_avg": overall_class_avg,
        "top_assignment": top_assignment,
    }
    return enrollments, assignments, grades, total_graded_students, summary


def check_excel(workspace, enrollments, assignments, grades, total_graded_students):
    """Проверка Course_Summary_AAA_F13.xlsx на корректность."""
    import openpyxl

    xlsx_path = Path(workspace) / "Course_Summary_AAA_F13.xlsx"
    if not xlsx_path.exists():
        record("Excel: file exists", False, f"Course_Summary_AAA_F13.xlsx not found in {workspace}")
        # без файла все критические excel-чеки тоже падают
        record("Excel: Enrollment Stats counts correct", False, "no file")
        record("Excel: Assignment Performance correct", False, "no file")
        record("Excel: Grade Distribution correct", False, "no file")
        return
    record("Excel: file exists", True)

    wb = openpyxl.load_workbook(xlsx_path, data_only=True)

    def get_sheet(name):
        for sn in wb.sheetnames:
            if sn.strip().lower() == name.lower():
                return wb[sn]
        return None

    # --- Имена листов (структурный, не критический) ---
    required_sheets = ["Enrollment Stats", "Assignment Performance", "Grade Distribution"]
    missing = [s for s in required_sheets if get_sheet(s) is None]
    record("Excel: all three sheets present", not missing,
           f"Missing {missing}. Found: {wb.sheetnames}" if missing else "")

    # --- Enrollment Stats (КРИТИЧЕСКИЙ) ---
    ws1 = get_sheet("Enrollment Stats")
    ok = True
    detail = ""
    if ws1 is None:
        ok, detail = False, "sheet missing"
    else:
        rows1 = list(ws1.iter_rows(values_only=True))
        if len(rows1) < 2:
            ok, detail = False, "no data rows"
        else:
            header1 = [str(h).strip().lower() if h else "" for h in rows1[0]]
            need = ["student_count", "teacher_count", "ta_count"]
            if any(c not in header1 for c in need):
                ok, detail = False, f"missing columns. Found: {[str(h) for h in rows1[0]]}"
            else:
                dr = rows1[1]
                vals = {
                    "student": int(dr[header1.index("student_count")]),
                    "teacher": int(dr[header1.index("teacher_count")]),
                    "ta": int(dr[header1.index("ta_count")]),
                }
                for k in ("student", "teacher", "ta"):
                    if vals[k] != enrollments[k]:
                        ok = False
                        detail = f"{k}: expected {enrollments[k]}, got {vals[k]}"
                        break
    record("Excel: Enrollment Stats counts correct", ok, detail)

    # --- Assignment Performance (КРИТИЧЕСКИЙ) ---
    ws2 = get_sheet("Assignment Performance")
    ok = True
    detail = ""
    if ws2 is None:
        ok, detail = False, "sheet missing"
    else:
        rows2 = list(ws2.iter_rows(values_only=True))
        if len(rows2) < 2:
            ok, detail = False, "no data rows"
        else:
            header2 = [str(h).strip().lower() if h else "" for h in rows2[0]]
            need = ["assignment_name", "avg_score", "max_score", "min_score", "submission_count"]
            if any(c not in header2 for c in need):
                ok, detail = False, f"missing columns. Found: {[str(h) for h in rows2[0]]}"
            else:
                idx = {c: header2.index(c) for c in need}
                data_rows2 = rows2[1:]
                if len(data_rows2) != len(assignments):
                    ok, detail = False, f"expected {len(assignments)} rows, got {len(data_rows2)}"
                else:
                    for i, (exp_name, exp_avg, exp_max, exp_min, exp_count) in enumerate(assignments):
                        r = data_rows2[i]
                        name_val = str(r[idx["assignment_name"]]).strip() if r[idx["assignment_name"]] else ""
                        if name_val.lower() != exp_name.lower():
                            ok, detail = False, f"row {i+1}: expected name '{exp_name}', got '{name_val}'"
                            break
                        avg_val = float(r[idx["avg_score"]]) if r[idx["avg_score"]] is not None else None
                        if avg_val is None or abs(avg_val - exp_avg) > 0.5:
                            ok, detail = False, f"'{exp_name}' avg: expected {exp_avg}, got {avg_val}"
                            break
                        max_val = float(r[idx["max_score"]]) if r[idx["max_score"]] is not None else None
                        if max_val is None or abs(max_val - exp_max) > 0.5:
                            ok, detail = False, f"'{exp_name}' max: expected {exp_max}, got {max_val}"
                            break
                        min_val = float(r[idx["min_score"]]) if r[idx["min_score"]] is not None else None
                        if min_val is None or abs(min_val - exp_min) > 0.5:
                            ok, detail = False, f"'{exp_name}' min: expected {exp_min}, got {min_val}"
                            break
                        count_val = int(r[idx["submission_count"]]) if r[idx["submission_count"]] is not None else None
                        if count_val is None or count_val != exp_count:
                            ok, detail = False, f"'{exp_name}' count: expected {exp_count}, got {count_val}"
                            break
    record("Excel: Assignment Performance correct", ok, detail)

    # --- Grade Distribution (КРИТИЧЕСКИЙ) ---
    ws3 = get_sheet("Grade Distribution")
    ok = True
    detail = ""
    if ws3 is None:
        ok, detail = False, "sheet missing"
    else:
        rows3 = list(ws3.iter_rows(values_only=True))
        if len(rows3) < 2:
            ok, detail = False, "no data rows"
        else:
            header3 = [str(h).strip().lower() if h else "" for h in rows3[0]]
            need = ["grade_range", "student_count", "percentage"]
            if any(c not in header3 for c in need):
                ok, detail = False, f"missing columns. Found: {[str(h) for h in rows3[0]]}"
            else:
                idx = {c: header3.index(c) for c in need}
                data_rows3 = rows3[1:]
                if len(data_rows3) != len(grades):
                    ok, detail = False, f"expected {len(grades)} rows, got {len(data_rows3)}"
                else:
                    for i, (exp_range, exp_count) in enumerate(grades):
                        r = data_rows3[i]
                        range_val = str(r[idx["grade_range"]]).strip() if r[idx["grade_range"]] else ""
                        if range_val.lower() != exp_range.lower():
                            ok, detail = False, f"row {i+1}: expected '{exp_range}', got '{range_val}'"
                            break
                        count_val = int(r[idx["student_count"]]) if r[idx["student_count"]] is not None else None
                        if count_val is None or count_val != exp_count:
                            ok, detail = False, f"'{exp_range}' count: expected {exp_count}, got {count_val}"
                            break
                        exp_pct = round(exp_count / total_graded_students * 100, 1)
                        pct_val = float(r[idx["percentage"]]) if r[idx["percentage"]] is not None else None
                        if pct_val is None or abs(pct_val - exp_pct) > 0.2:
                            ok, detail = False, f"'{exp_range}' pct: expected {exp_pct}, got {pct_val}"
                            break
    record("Excel: Grade Distribution correct", ok, detail)

    wb.close()


def check_pptx(workspace, enrollments, assignments, summary):
    """Проверка Course_Summary_AAA_F13.pptx на корректность.

    Греп ключевых слов терпим к русскому тексту: проза слайдов пишется
    агентом по-русски, поэтому проверяются RU+EN варианты. Английскими
    остаются только засеянные идентификаторы Canvas (название курса/заданий).
    """
    from pptx import Presentation

    pptx_path = Path(workspace) / "Course_Summary_AAA_F13.pptx"
    if not pptx_path.exists():
        record("PPTX: file exists", False, f"not found in {workspace}")
        record("PPTX: title slide has full course name and 2013", False, "no file")
        record("PPTX: highest-average assignment named on findings slide", False, "no file")
        return
    record("PPTX: file exists", True)

    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)

    record("PPTX: at least 5 slides", len(slides) >= 5, f"got {len(slides)}")

    all_text = []
    for slide in slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_texts.append(paragraph.text)
        all_text.append("\n".join(slide_texts))

    full_text = "\n".join(all_text).lower()

    # --- Титульный слайд: полное название курса + 2013 (КРИТИЧЕСКИЙ) ---
    # Название курса в Canvas теперь русифицировано ("Прикладная аналитика и
    # алгоритмы (Осень 2013)"), но task.md также даёт английский вариант в
    # кавычках. Принимаем любой из них: дискриминирующая RU-подстрока
    # "аналитик" уникально идентифицирует именно этот курс.
    first_slide_text = all_text[0].lower()
    course_name_ok = (("applied analytics" in first_slide_text and "algorithms" in first_slide_text)
                      or "аналитик" in first_slide_text)
    title_ok = course_name_ok and ("2013" in first_slide_text)
    record("PPTX: title slide has full course name and 2013", title_ok,
           f"title text: {all_text[0][:200]}")

    # --- Количество студентов где-то присутствует (структурный) ---
    student_str = str(enrollments["student"])
    record("PPTX: student count appears", student_str in full_text,
           f"{student_str} not found")

    # --- Все названия заданий присутствуют (английские, засеяны Canvas) ---
    missing_a = [a[0] for a in assignments if a[0].lower() not in full_text]
    record("PPTX: all assignment names appear", not missing_a,
           f"missing: {missing_a}")

    # --- Контент распределения оценок (RU+EN) ---
    grade_keywords = [
        "grade", "distribution", "a (90", "b (80", "c (70", "d (60", "f (<60",
        # русские варианты
        "оцен", "распредел", "диапазон",
    ]
    grade_found = sum(1 for kw in grade_keywords if kw in full_text)
    record("PPTX: grade distribution content present", grade_found >= 3,
           f"found {grade_found} keywords")

    # --- Слайд выводов/рекомендаций (RU+EN) ---
    findings_keywords = [
        "finding", "recommendation", "key", "summary", "overall", "average",
        # русские варианты
        "вывод", "рекоменд", "ключ", "итог", "сводк", "средн",
    ]
    findings_found = sum(1 for kw in findings_keywords if kw in full_text)
    record("PPTX: findings/recommendations content present", findings_found >= 2,
           f"found {findings_found} keywords")

    # --- Задание с наивысшим средним баллом названо в презентации (КРИТИЧЕСКИЙ) ---
    top = (summary.get("top_assignment") or "").lower()
    record("PPTX: highest-average assignment named on findings slide",
           bool(top) and top in full_text,
           f"top assignment '{summary.get('top_assignment')}' not found")


def check_gsheet(summary):
    """Проверка, что 'AAA F13 Course Dashboard' создан в схеме gsheet
    и содержит корректные значения сводки (название курса, число студентов,
    число заданий и средний балл по курсу). Сопоставление значений терпимо
    к русским/английским подписям — проверяются именно числовые значения.
    """
    import psycopg2

    conn = psycopg2.connect(
        host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="cowork_gym",
        user="eigent", password="camel"
    )
    cur = conn.cursor()

    cur.execute("""
        SELECT id, title FROM gsheet.spreadsheets
        WHERE LOWER(title) LIKE '%aaa%' AND LOWER(title) LIKE '%f13%'
    """)
    rows = cur.fetchall()
    if not rows:
        cur.execute("""
            SELECT id, title FROM gsheet.spreadsheets
            WHERE LOWER(title) LIKE '%course%dashboard%'
               OR LOWER(title) LIKE '%aaa%dashboard%'
               OR LOWER(title) LIKE '%analytics%algorithms%'
        """)
        rows = cur.fetchall()

    if not rows:
        conn.close()
        record("GSheet: dashboard spreadsheet exists", False,
               "No 'AAA F13 Course Dashboard' spreadsheet found")
        record("GSheet: dashboard contains correct course summary values", False, "no spreadsheet")
        return

    spreadsheet_id = rows[0][0]
    spreadsheet_title = rows[0][1]
    record("GSheet: dashboard spreadsheet exists", True)
    print(f"    Found spreadsheet: '{spreadsheet_title}' (id={spreadsheet_id})")

    # Первый лист
    cur.execute("""
        SELECT id FROM gsheet.sheets
        WHERE spreadsheet_id = %s ORDER BY id LIMIT 1
    """, (spreadsheet_id,))
    sh = cur.fetchone()

    # Все значения ячеек первого листа (или всей таблицы как fallback)
    if sh:
        cur.execute("""
            SELECT value FROM gsheet.cells
            WHERE spreadsheet_id = %s AND sheet_id = %s
        """, (spreadsheet_id, sh[0]))
    else:
        cur.execute("""
            SELECT value FROM gsheet.cells WHERE spreadsheet_id = %s
        """, (spreadsheet_id,))
    cell_values = [r[0] for r in cur.fetchall() if r[0] is not None]
    conn.close()

    joined = "\n".join(str(v) for v in cell_values).lower()

    # Значение-чек (КРИТИЧЕСКИЙ): название курса + число студентов +
    # число заданий + средний балл по курсу.
    failures = []

    # Название курса — Canvas русифицирован ("Прикладная аналитика и алгоритмы"),
    # но task.md также даёт английский вариант. Принимаем любой; RU-подстрока
    # "аналитик" дискриминирует именно этот курс.
    if not (("applied analytics" in joined and "algorithms" in joined)
            or "аналитик" in joined):
        failures.append("course name not found")

    def numeric_cell_match(target, abs_tol):
        for v in cell_values:
            try:
                if abs(float(str(v).strip()) - float(target)) <= abs_tol:
                    return True
            except (TypeError, ValueError):
                continue
        return False

    if not numeric_cell_match(summary["total_students"], 0):
        failures.append(f"total students {summary['total_students']} not found")
    if not numeric_cell_match(summary["n_assignments"], 0):
        failures.append(f"assignment count {summary['n_assignments']} not found")
    if not numeric_cell_match(summary["overall_class_avg"], 0.2):
        failures.append(f"overall class avg {summary['overall_class_avg']} not found")

    record("GSheet: dashboard contains correct course summary values",
           not failures, "; ".join(failures))


if __name__ == "__main__":
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False, help="Launch time")
    args = parser.parse_args()

    workspace = args.agent_workspace
    if not workspace:
        print("Error: --agent_workspace is required")
        sys.exit(1)

    print("Fetching expected data from database...")
    try:
        enrollments, assignments, grades, total_graded_students, summary = get_expected_data()
        print(f"  Enrollments: S={enrollments['student']}, T={enrollments['teacher']}, TA={enrollments['ta']}")
        print(f"  Assignments: {len(assignments)}")
        print(f"  Grade ranges: {len(grades)}, total graded students: {total_graded_students}")
        print(f"  Summary: {summary}")
    except Exception as e:
        print(f"Error querying database: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

    print("\n--- Check 1: Excel File ---")
    try:
        check_excel(workspace, enrollments, assignments, grades, total_graded_students)
    except Exception as e:
        record("Excel: check ran without error", False, str(e))
        import traceback
        traceback.print_exc()

    print("\n--- Check 2: PowerPoint File ---")
    try:
        check_pptx(workspace, enrollments, assignments, summary)
    except Exception as e:
        record("PPTX: check ran without error", False, str(e))
        import traceback
        traceback.print_exc()

    print("\n--- Check 3: Online Spreadsheet Dashboard ---")
    try:
        check_gsheet(summary)
    except Exception as e:
        record("GSheet: check ran without error", False, str(e))
        import traceback
        traceback.print_exc()

    total = PASS_COUNT + FAIL_COUNT
    accuracy = PASS_COUNT / total * 100 if total > 0 else 0
    print(f"\n=== Results: {PASS_COUNT}/{total} passed ({accuracy:.1f}%) ===")

    critical_failed = [n for n in FAILED_NAMES if n in CRITICAL_CHECKS]
    if critical_failed:
        print(f"CRITICAL FAILURES: {critical_failed}")
        print("FAIL (critical check failed)")
        sys.exit(1)

    if accuracy >= 70:
        print("Pass all tests!")
        sys.exit(0)
    else:
        print(f"FAIL (accuracy {accuracy:.1f}% < 70%)")
        sys.exit(1)
