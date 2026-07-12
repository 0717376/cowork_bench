"""Evaluation for yf-gold-vs-stocks-ppt."""
import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="cowork_gym", user="eigent", password="camel")
PASS_COUNT = 0
FAIL_COUNT = 0
FAILED_NAMES = []

# Critical checks (semantic): any failure => overall FAIL regardless of accuracy.
# These anchor core data extraction and tie the two deliverables to the correct winner.
CRITICAL_CHECKS = {
    "Month 2025-03 gold close",
    "Month 2026-03 gold close",
    "Month 2025-03 IMOEX close",
    "Month 2026-03 IMOEX close",
    "Conclusion names correct winner",
}


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1; print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1; FAILED_NAMES.append(name); print(f"  [FAIL] {name}: {str(detail)[:300]}")


def num_close(a, b, tol=1.0):
    try: return abs(float(a) - float(b)) <= tol
    except: return False


def get_expected():
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()
    monthly = {}
    for symbol in ['GLDRUB_TOM', 'IMOEX.ME']:
        cur.execute("""SELECT date, close FROM moex.stock_prices
            WHERE symbol=%s AND date>='2025-03-06' AND date<='2026-03-05' ORDER BY date""", (symbol,))
        by_month = {}
        for d, c in cur.fetchall():
            mk = d.strftime("%Y-%m")
            by_month[mk] = float(c)
        monthly[symbol] = by_month
    conn.close()
    months = sorted(set(list(monthly['GLDRUB_TOM'].keys()) + list(monthly['IMOEX.ME'].keys())))
    prices = []
    for m in months:
        prices.append({"month": m, "gold": monthly['GLDRUB_TOM'].get(m), "imoex": monthly['IMOEX.ME'].get(m)})
    returns = []
    for i in range(1, len(prices)):
        gr = dr = None
        if prices[i-1]["gold"] and prices[i]["gold"] and prices[i-1]["gold"] != 0:
            gr = round((prices[i]["gold"] - prices[i-1]["gold"]) / prices[i-1]["gold"] * 100, 2)
        if prices[i-1]["imoex"] and prices[i]["imoex"] and prices[i-1]["imoex"] != 0:
            dr = round((prices[i]["imoex"] - prices[i-1]["imoex"]) / prices[i-1]["imoex"] * 100, 2)
        returns.append({"month": prices[i]["month"], "gold_ret": gr, "imoex_ret": dr})
    # Derive the true better-performing asset LIVE from moex data (full-period change first->last).
    gold_prices = [p["gold"] for p in prices if p["gold"]]
    imoex_prices = [p["imoex"] for p in prices if p["imoex"]]
    winner = None
    gold_change = imoex_change = None
    if len(gold_prices) >= 2 and len(imoex_prices) >= 2 and gold_prices[0] and imoex_prices[0]:
        gold_change = (gold_prices[-1] - gold_prices[0]) / gold_prices[0] * 100
        imoex_change = (imoex_prices[-1] - imoex_prices[0]) / imoex_prices[0] * 100
        winner = "gold" if gold_change > imoex_change else "imoex"
    return {"prices": prices, "returns": returns, "months": months,
            "winner": winner, "gold_change": gold_change, "imoex_change": imoex_change}


def sheet_dicts(wb, name):
    for sn in wb.sheetnames:
        if sn.strip().lower() == name.strip().lower():
            ws = wb[sn]
            rows = list(ws.iter_rows(values_only=True))
            if len(rows) < 2: return []
            hdrs = [str(h).strip() if h else "" for h in rows[0]]
            return [{hdrs[i]: row[i] for i in range(len(hdrs))} for row in rows[1:] if not all(v is None for v in row)]
    return None


def check_excel(ws_path, exp):
    print("\n=== Checking Excel ===")
    p = os.path.join(ws_path, "Gold_vs_IMOEX.xlsx")
    if not os.path.isfile(p):
        record("Excel file exists", False, p); return
    record("Excel file exists", True)
    wb = openpyxl.load_workbook(p, data_only=True)

    # Monthly Prices
    d = sheet_dicts(wb, "Monthly Prices")
    if d is None:
        record("Sheet Monthly Prices", False, str(wb.sheetnames))
    else:
        record("Sheet Monthly Prices", True)
        record("Monthly Prices row count", len(d) >= 12, f"Got {len(d)}")
        for ep in exp["prices"][:3] + exp["prices"][-2:]:
            m = next((r for r in d if str(r.get("Month","")).strip() == ep["month"]), None)
            if not m:
                record(f"Month {ep['month']} present", False, "Missing"); continue
            record(f"Month {ep['month']} present", True)
            if ep["gold"]:
                record(f"Month {ep['month']} gold close",
                       num_close(m.get("Gold_Close"), ep["gold"], 20.0),
                       f"{m.get('Gold_Close')} vs {ep['gold']}")
            if ep["imoex"]:
                record(f"Month {ep['month']} IMOEX close",
                       num_close(m.get("IMOEX_Close"), ep["imoex"], 15.0),
                       f"{m.get('IMOEX_Close')} vs {ep['imoex']}")

    # Returns
    d = sheet_dicts(wb, "Returns")
    if d is None:
        record("Sheet Returns", False, str(wb.sheetnames))
    else:
        record("Sheet Returns", True)
        record("Returns row count", len(d) >= 11, f"Got {len(d)}")
        for er in exp["returns"][:2] + exp["returns"][-2:]:
            m = next((r for r in d if str(r.get("Month","")).strip() == er["month"]), None)
            if not m:
                record(f"Return {er['month']} present", False, "Missing"); continue
            if er["gold_ret"] is not None:
                record(f"Return {er['month']} gold",
                       num_close(m.get("Gold_Return_Pct"), er["gold_ret"], 2.0),
                       f"{m.get('Gold_Return_Pct')} vs {er['gold_ret']}")
            if er["imoex_ret"] is not None:
                record(f"Return {er['month']} IMOEX",
                       num_close(m.get("IMOEX_Return_Pct"), er["imoex_ret"], 2.0),
                       f"{m.get('IMOEX_Return_Pct')} vs {er['imoex_ret']}")
    wb.close()


def check_pptx(ws_path, exp):
    print("\n=== Checking PPTX ===")
    p = os.path.join(ws_path, "Gold_vs_Stocks.pptx")
    if not os.path.isfile(p):
        record("PPTX file exists", False, p); return
    record("PPTX file exists", True)
    try:
        from pptx import Presentation
        prs = Presentation(p)
        slides = list(prs.slides)
        record("Slide count >= 3", len(slides) >= 3, f"Got {len(slides)}")
        if len(slides) >= 1:
            title_shape = slides[0].shapes.title
            if title_shape:
                record("Slide 1 has title", True)
                t = title_shape.text.lower()
                record("Slide 1 title mentions gold", "gold" in t, title_shape.text)
            else:
                # Check all shapes for title text
                all_text = " ".join(sh.text for sh in slides[0].shapes if sh.has_text_frame).lower()
                record("Slide 1 mentions gold", "gold" in all_text, all_text[:200])
        if len(slides) >= 3:
            # Use ORIGINAL text lowercased (NOT normalized) for RU keyword checks.
            all_text = " ".join(sh.text for sh in slides[2].shapes if sh.has_text_frame).lower()
            # Conclusion content: RU + EN keywords (agent writes the conclusion in Russian).
            conc_kw = ["conclu", "perform", "better", "outperform",
                       "вывод", "заключ", "лучше", "выше", "опередил", "превзош", "динамик"]
            record("Slide 3 has conclusion content",
                   any(k in all_text for k in conc_kw), all_text[:200])

            # SEMANTIC: conclusion must name the CORRECT better-performing asset
            # (winner derived live from moex data; gold rose ~16% vs IMOEX ~4%).
            winner = exp.get("winner")
            if winner is None:
                record("Conclusion names correct winner", False, "Could not derive winner from moex data")
            else:
                gold_kw = any(k in all_text for k in ["gold", "золот"])
                imoex_kw = any(k in all_text for k in ["imoex", "moex", "мосбирж", "индекс", "акци", "фондов"])
                if winner == "gold":
                    # Gold must be named as the winner; reject if only the loser (IMOEX) is mentioned.
                    record("Conclusion names correct winner", gold_kw,
                           f"winner=gold goldKw={gold_kw} imoexKw={imoex_kw} text={all_text[:200]}")
                else:
                    record("Conclusion names correct winner", imoex_kw,
                           f"winner=imoex goldKw={gold_kw} imoexKw={imoex_kw} text={all_text[:200]}")
    except ImportError:
        record("python-pptx available", False, "Cannot import pptx")
    except Exception as e:
        record("PPTX readable", False, str(e))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", default=".")
    parser.add_argument("--groundtruth_workspace", default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    exp = get_expected()
    check_excel(args.agent_workspace, exp)
    check_pptx(args.agent_workspace, exp)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = PASS_COUNT / total * 100 if total else 0
    critical_failed = [n for n in FAILED_NAMES if n in CRITICAL_CHECKS]

    print(f"\n=== SUMMARY: {PASS_COUNT}/{total} passed ({accuracy:.1f}%) ===")
    if critical_failed:
        print(f"CRITICAL FAILURES: {len(critical_failed)}")
        for n in critical_failed:
            print(f"  - {n}")

    success = (not critical_failed) and accuracy >= 70
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump({"passed": PASS_COUNT, "failed": FAIL_COUNT, "total": total,
                       "accuracy": accuracy, "critical_failed": critical_failed,
                       "success": success}, f)
    sys.exit(0 if success else 1)

if __name__ == "__main__":
    main()
