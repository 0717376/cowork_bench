"""Evaluation for moex-peer-comparison-excel-ppt-email.

Покрытие — 6 эмитентов MOEX: SBER.ME, GAZP.ME, LKOH.ME, TCSG.ME, MGNT.ME, MTSS.ME.

Структура чеков: структурные (NON-critical) + CRITICAL.
Любой провал CRITICAL => FAIL независимо от accuracy. Иначе PASS при accuracy >= 70%.

CRITICAL-чеки (семантические, отражают суть deliverable):
  C1: Три письма отправлены нужным получателям с ТОЧНЫМИ темами
      (portfolio_managers/research_team/compliance @firm.com,
       "Peer Comparison Summary" / "...Detailed Findings" / "...Risk Review").
  C2: Scoring.Overall_Rating совпадает с эталоном для ВСЕХ эмитентов (Strong Buy/Buy/Hold/Sell).
  C3: Scoring.Weighted_Score совпадает с эталоном (tol 0.5) для ВСЕХ эмитентов.
  C4: Company Profiles: Market_Cap и YTD_Return_Pct совпадают с эталоном для ВСЕХ эмитентов.
  C5: Financial Comparison: Revenue и Free_Cash_Flow совпадают с эталоном (в пределах 5%) для ВСЕХ эмитентов.
"""

import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "cowork_gym",
    "user": "eigent",
    "password": "camel",
}

SYMBOLS = ["GAZP.ME", "LKOH.ME", "MGNT.ME", "MTSS.ME", "SBER.ME", "TCSG.ME"]
N = len(SYMBOLS)

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []
CRITICAL_CHECKS = set()


def crit(name):
    CRITICAL_CHECKS.add(name)
    return name


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        marker = " [CRITICAL]" if name in CRITICAL_CHECKS else ""
        detail_str = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL]{marker} {name}{detail_str}")
        if name in CRITICAL_CHECKS:
            CRITICAL_FAILS.append(name)


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def rel_close(a, b, frac=0.05):
    try:
        a = float(a); b = float(b)
        if b == 0:
            return abs(a) < 1e-6
        return abs(a - b) / abs(b) <= frac
    except (TypeError, ValueError):
        return False


def load_gt_sheet(gt_ws, sheet_name):
    gt_path = os.path.join(gt_ws, "Peer_Comparison.xlsx")
    if not os.path.exists(gt_path):
        return None
    wb = openpyxl.load_workbook(gt_path, data_only=True)
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def check_excel(agent_workspace, groundtruth_workspace):
    print("\n=== Checking Excel Output ===")
    excel_path = os.path.join(agent_workspace, "Peer_Comparison.xlsx")
    check("Excel file exists", os.path.isfile(excel_path), f"Expected {excel_path}")
    if not os.path.isfile(excel_path):
        return

    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
    except Exception as e:
        check("Excel file readable", False, str(e))
        return

    # --- Sheet 1: Company Profiles ---
    ws1 = None
    for s in wb.sheetnames:
        if "company" in s.lower() and "profile" in s.lower():
            ws1 = wb[s]
            break
    if ws1 is None:
        for s in wb.sheetnames:
            if "company" in s.lower() or "profile" in s.lower():
                ws1 = wb[s]
                break

    check("Sheet 'Company Profiles' exists", ws1 is not None, f"Sheets: {wb.sheetnames}")
    if ws1 is not None:
        rows = list(ws1.iter_rows(min_row=2, values_only=True))
        data_rows = [r for r in rows if r and r[0] is not None]
        check(f"Company Profiles has {N} rows", len(data_rows) == N, f"Got {len(data_rows)}")

        symbols_found = {str(r[0]).strip().upper() for r in data_rows if r[0]}
        for sym in SYMBOLS:
            check(f"{sym} in Company Profiles", sym in symbols_found, f"Found: {symbols_found}")

        sym_list = [str(r[0]).strip().upper() for r in data_rows if r[0]]
        check("Symbols sorted alphabetically", sym_list == sorted(sym_list), f"Order: {sym_list}")

        gt_rows = load_gt_sheet(groundtruth_workspace, "Company Profiles")
        if gt_rows:
            gt_data = {str(r[0]).strip().upper(): r for r in gt_rows[1:] if r and r[0]}
            agent_data = {str(r[0]).strip().upper(): r for r in data_rows if r and r[0]}

            # CRITICAL C4: Market_Cap & YTD_Return_Pct for ALL symbols
            mcap_ok = True
            ytd_ok = True
            mcap_detail = []
            ytd_detail = []
            for sym in SYMBOLS:
                ar = agent_data.get(sym); gr = gt_data.get(sym)
                if not ar or not gr:
                    mcap_ok = False; ytd_ok = False
                    mcap_detail.append(f"{sym}:missing"); ytd_detail.append(f"{sym}:missing")
                    continue
                # Market Cap col 3, tol 1M
                if len(ar) > 3 and len(gr) > 3:
                    if not num_close(ar[3], gr[3], 1e6):
                        mcap_ok = False
                        mcap_detail.append(f"{sym}:A={ar[3]},GT={gr[3]}")
                else:
                    mcap_ok = False; mcap_detail.append(f"{sym}:short")
                # YTD col 10, tol 0.5
                if len(ar) > 10 and len(gr) > 10:
                    if not num_close(ar[10], gr[10], 0.5):
                        ytd_ok = False
                        ytd_detail.append(f"{sym}:A={ar[10]},GT={gr[10]}")
                else:
                    ytd_ok = False; ytd_detail.append(f"{sym}:short")
            check(crit("Company Profiles: Market_Cap matches GT for ALL symbols"),
                  mcap_ok, "; ".join(mcap_detail))
            check(crit("Company Profiles: YTD_Return_Pct matches GT for ALL symbols"),
                  ytd_ok, "; ".join(ytd_detail))

            # NON-critical spot: Trailing_PE
            for sym in SYMBOLS:
                ar = agent_data.get(sym); gr = gt_data.get(sym)
                if ar and gr and len(ar) > 4 and len(gr) > 4 and ar[4] is not None:
                    check(f"{sym} Trailing_PE", num_close(ar[4], gr[4], 0.5),
                          f"Agent={ar[4]}, GT={gr[4]}")

    # --- Sheet 2: Financial Comparison ---
    ws2 = None
    for s in wb.sheetnames:
        if "financial" in s.lower() and "comparison" in s.lower():
            ws2 = wb[s]
            break
    if ws2 is None:
        for s in wb.sheetnames:
            if "financial" in s.lower():
                ws2 = wb[s]
                break

    check("Sheet 'Financial Comparison' exists", ws2 is not None, f"Sheets: {wb.sheetnames}")
    if ws2 is not None:
        rows2 = list(ws2.iter_rows(min_row=2, values_only=True))
        data_rows2 = [r for r in rows2 if r and r[0] is not None]
        check(f"Financial Comparison has {N} rows", len(data_rows2) == N, f"Got {len(data_rows2)}")

        gt_rows2 = load_gt_sheet(groundtruth_workspace, "Financial Comparison")
        if gt_rows2:
            gt_fin = {str(r[0]).strip().upper(): r for r in gt_rows2[1:] if r and r[0]}
            agent_fin = {str(r[0]).strip().upper(): r for r in data_rows2 if r and r[0]}

            # CRITICAL C5: Revenue & Free_Cash_Flow within 5% for ALL symbols
            rev_ok = True; fcf_ok = True
            rev_detail = []; fcf_detail = []
            for sym in SYMBOLS:
                ar = agent_fin.get(sym); gr = gt_fin.get(sym)
                if not ar or not gr:
                    rev_ok = False; fcf_ok = False
                    rev_detail.append(f"{sym}:missing"); fcf_detail.append(f"{sym}:missing")
                    continue
                # Revenue col 1
                if len(ar) > 1 and len(gr) > 1 and ar[1] is not None and gr[1] is not None:
                    if not rel_close(ar[1], gr[1], 0.05):
                        rev_ok = False; rev_detail.append(f"{sym}:A={ar[1]},GT={gr[1]}")
                else:
                    rev_ok = False; rev_detail.append(f"{sym}:short/none")
                # Free_Cash_Flow col 4
                if len(ar) > 4 and len(gr) > 4 and ar[4] is not None and gr[4] is not None:
                    if not rel_close(ar[4], gr[4], 0.05):
                        fcf_ok = False; fcf_detail.append(f"{sym}:A={ar[4]},GT={gr[4]}")
                else:
                    fcf_ok = False; fcf_detail.append(f"{sym}:short/none")
            check(crit("Financial Comparison: Revenue within 5% for ALL symbols"),
                  rev_ok, "; ".join(rev_detail))
            check(crit("Financial Comparison: Free_Cash_Flow within 5% for ALL symbols"),
                  fcf_ok, "; ".join(fcf_detail))

    # --- Sheet 3: Scoring ---
    ws3 = None
    for s in wb.sheetnames:
        if "scor" in s.lower():
            ws3 = wb[s]
            break

    check("Sheet 'Scoring' exists", ws3 is not None, f"Sheets: {wb.sheetnames}")
    if ws3 is not None:
        rows3 = list(ws3.iter_rows(min_row=2, values_only=True))
        data_rows3 = [r for r in rows3 if r and r[0] is not None]
        check(f"Scoring has {N} rows", len(data_rows3) == N, f"Got {len(data_rows3)}")

        gt_rows3 = load_gt_sheet(groundtruth_workspace, "Scoring")
        if gt_rows3:
            gt_score = {str(r[0]).strip().upper(): r for r in gt_rows3[1:] if r and r[0]}
            agent_score = {str(r[0]).strip().upper(): r for r in data_rows3 if r and r[0]}

            # CRITICAL C2: Overall_Rating exact for ALL symbols
            rating_ok = True; rating_detail = []
            # CRITICAL C3: Weighted_Score tol 0.5 for ALL symbols
            ws_ok = True; ws_detail = []
            for sym in SYMBOLS:
                ar = agent_score.get(sym); gr = gt_score.get(sym)
                if not ar or not gr:
                    rating_ok = False; ws_ok = False
                    rating_detail.append(f"{sym}:missing"); ws_detail.append(f"{sym}:missing")
                    continue
                if len(ar) > 7 and len(gr) > 7 and ar[7] is not None:
                    if str(ar[7]).strip().lower() != str(gr[7]).strip().lower():
                        rating_ok = False
                        rating_detail.append(f"{sym}:A='{ar[7]}',GT='{gr[7]}'")
                else:
                    rating_ok = False; rating_detail.append(f"{sym}:short/none")
                if len(ar) > 6 and len(gr) > 6 and ar[6] is not None:
                    if not num_close(ar[6], gr[6], 0.5):
                        ws_ok = False
                        ws_detail.append(f"{sym}:A={ar[6]},GT={gr[6]}")
                else:
                    ws_ok = False; ws_detail.append(f"{sym}:short/none")
            check(crit("Scoring: Overall_Rating matches GT for ALL symbols"),
                  rating_ok, "; ".join(rating_detail))
            check(crit("Scoring: Weighted_Score within 0.5 for ALL symbols"),
                  ws_ok, "; ".join(ws_detail))


def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint Output ===")
    pptx_path = os.path.join(agent_workspace, "Investor_Presentation.pptx")
    check("PPTX file exists", os.path.isfile(pptx_path), f"Expected {pptx_path}")
    if not os.path.isfile(pptx_path):
        return

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        check("PPTX file readable", False, str(e))
        return

    slide_count = len(prs.slides)
    check("PPTX has >= 6 slides", slide_count >= 6, f"Got {slide_count} slides")

    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text.lower() + " "

    check("PPTX contains 'peer'/'comparison'/'сравнен'",
          any(k in all_text for k in ("peer", "comparison", "сравнен")),
          f"Text sample: {all_text[:200]}")
    check("PPTX mentions a MOEX issuer (sber/сбер/газпром/gazp/...)",
          any(k in all_text for k in ("sber", "сбер", "gazp", "газпром", "lkoh", "лукойл",
                                       "tcsg", "mgnt", "магнит", "mtss", "мтс")),
          f"Text sample: {all_text[:200]}")
    check("PPTX contains scoring/recommendation content",
          any(k in all_text for k in ("score", "rank", "recommendation", "buy", "hold",
                                      "оценк", "рейтинг", "ранг", "рекоменд")),
          f"Text sample: {all_text[:200]}")


SUBJECTS = {
    "portfolio_managers@firm.com": "Peer Comparison Summary",
    "research_team@firm.com": "Peer Comparison Detailed Findings",
    "compliance@firm.com": "Peer Comparison Risk Review",
}


def check_emails():
    print("\n=== Checking Emails ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    all_ok = True
    detail = []
    for addr, subj in SUBJECTS.items():
        cur.execute("""
            SELECT subject FROM email.messages
            WHERE to_addr::text ILIKE %s
            ORDER BY id DESC LIMIT 25
        """, (f"%{addr}%",))
        rows = cur.fetchall()
        subjects = [str(r[0]).strip() for r in rows]
        recipient_ok = len(rows) > 0
        subject_ok = any(s.lower() == subj.lower() for s in subjects)
        check(f"Email to {addr} exists", recipient_ok, "No email found")
        check(f"Email to {addr} has exact subject '{subj}'", subject_ok,
              f"Subjects found: {subjects}")
        if not (recipient_ok and subject_ok):
            all_ok = False
            detail.append(f"{addr}->{subjects}")

    check(crit("All 3 emails to correct recipients with exact subjects"),
          all_ok, "; ".join(detail))

    cur.close()
    conn.close()


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    gt_ws = args.groundtruth_workspace or os.path.join(
        os.path.dirname(__file__), "..", "groundtruth_workspace")

    check_excel(args.agent_workspace, gt_ws)
    check_pptx(args.agent_workspace)
    check_emails()

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100) if total else 0.0

    print(f"\n=== SUMMARY ===")
    print(f"  Checks - Passed: {PASS_COUNT}, Failed: {FAIL_COUNT}  (accuracy {accuracy:.1f}%)")
    if CRITICAL_FAILS:
        print(f"  CRITICAL FAILS ({len(CRITICAL_FAILS)}): {CRITICAL_FAILS}")

    success = (not CRITICAL_FAILS) and accuracy >= 70

    if args.res_log_file:
        result = {
            "passed": PASS_COUNT,
            "failed": FAIL_COUNT,
            "accuracy": accuracy,
            "critical_fails": CRITICAL_FAILS,
            "success": success,
        }
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if CRITICAL_FAILS:
        print("Overall: FAIL (critical check failed)")
        sys.exit(1)
    if accuracy >= 70:
        print("Overall: PASS")
        sys.exit(0)
    print("Overall: FAIL (accuracy below 70%)")
    sys.exit(1)


if __name__ == "__main__":
    main()
