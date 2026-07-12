"""Evaluation for sf-sales-segment-ppt-gform.

ClickHouse Sales DW (sf_data schema) segment values are russified centrally
by db/zzz_clickhouse_after_init.sql:
  Consumer   -> 'Частные клиенты'   (top revenue segment, total 839609.20)
  Enterprise -> 'Корпоративный'
  Government -> 'Государственный'
  SMB        -> 'Малый и средний бизнес'
The agent reads RU segment names from the DWH and writes them into the PPT/form,
so all segment-name checks accept RU first, EN as fallback. Numeric revenue
amounts are NOT remapped (839609.20 stays).
"""
import argparse
import json
import os
import re
import sys

import psycopg2


def normalize_ru_numbers(text):
    """RU number normalization for substring/regex checks: collapse digit-group
    separators (space/NBSP/NNBSP/dot/comma before a 3-digit group) and turn
    decimal commas into dots ("31 588" -> "31588", "4 586,91" -> "4586.91")."""
    t = str(text or "")
    t = re.sub(r"(?<=\d)[ \xa0\u202f\u2009.,](?=\d{3}\b)", "", t)
    return re.sub(r"(?<=\d),(?=\d)", ".", t)

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "cowork_gym",
    "user": "eigent",
    "password": "camel",
}

# Segment names as written by the agent: RU (from russified DWH) or EN fallback.
SEG_CONSUMER = ["Частные клиенты", "Consumer"]
SEG_ENTERPRISE = ["Корпоративный", "Enterprise"]
SEG_GOVERNMENT = ["Государственный", "Government"]
SEG_SMB = ["Малый и средний бизнес", "SMB"]
ALL_SEGMENTS = [SEG_CONSUMER, SEG_ENTERPRISE, SEG_GOVERNMENT, SEG_SMB]

# Consumer / 'Частные клиенты' is the top-revenue segment.
REVENUE_VARIANTS = ["839609.20", "839609", "839,609.20", "839,609"]


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    agent_ws = args.agent_workspace or task_root

    all_errors = []          # non-critical (counted toward accuracy)
    critical_errors = []     # any one => sys.exit(1) before accuracy gate
    passed = 0
    total = 0

    def check(desc, cond, detail="", critical=False):
        nonlocal passed, total
        total += 1
        if cond:
            passed += 1
            print(f"  [OK] {desc}")
        else:
            tag = "CRITICAL" if critical else "FAIL"
            print(f"  [{tag}] {desc} :: {detail}")
            (critical_errors if critical else all_errors).append(f"{desc} :: {detail}")

    def any_in(text, variants):
        return any(v.lower() in text.lower() for v in variants)

    # --- Check 1: PowerPoint ---
    print("Checking PowerPoint presentation...")
    pptx_path = os.path.join(agent_ws, "Segment_Performance.pptx")
    all_text = ""
    first_slide_text = ""
    slide_count = 0
    if not os.path.exists(pptx_path):
        check("Segment_Performance.pptx exists in agent workspace", False,
              f"not found at {pptx_path}", critical=True)
    else:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += " " + shape.text
        if prs.slides:
            for shape in prs.slides[0].shapes:
                if hasattr(shape, "text"):
                    first_slide_text += " " + shape.text

        # Non-critical structural: slide count
        check("PPT has at least 4 slides", slide_count >= 4,
              f"got {slide_count} slides")

        # CRITICAL: first slide title (EN per task wording, RU equivalent accepted)
        title_ok = any_in(first_slide_text, [
            "Customer Segment Performance Report", "Customer Segment",
            "Segment Performance", "результатах клиентских сегментов",
            "Отчёт о результатах", "клиентских сегментов",
        ])
        check("First slide titled 'Customer Segment Performance Report' (RU equiv accepted)",
              title_ok, f"first slide text: {first_slide_text[:120]!r}", critical=True)

        # CRITICAL: top-revenue segment name present (Частные клиенты / Consumer)
        check("PPT contains top-revenue segment name 'Частные клиенты'/'Consumer'",
              any_in(all_text, SEG_CONSUMER), "missing Consumer segment", critical=True)

        # Non-critical: other segments present
        check("PPT contains 'Корпоративный'/'Enterprise' segment",
              any_in(all_text, SEG_ENTERPRISE), "missing Enterprise segment")

        # CRITICAL: Consumer total revenue figure 839609.20
        check("PPT contains Consumer total revenue figure (839609.20)",
              any_in(normalize_ru_numbers(all_text), REVENUE_VARIANTS),
              "missing 839609.20", critical=True)

        # Non-critical: required slide titles for sections 2-4
        check("PPT contains 'Revenue by Segment' slide", "Revenue by Segment".lower() in all_text.lower(),
              "missing 'Revenue by Segment' title")
        check("PPT contains 'Key Insights' slide", "Key Insights".lower() in all_text.lower(),
              "missing 'Key Insights' title")
        check("PPT contains 'Strategy Recommendations' slide",
              "Strategy Recommendations".lower() in all_text.lower(),
              "missing 'Strategy Recommendations' title")

    # --- Check 2: Form ---
    print("Checking feedback form (gform schema)...")
    form_id = None
    questions = []
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("SELECT id, title FROM gform.forms WHERE LOWER(title) LIKE '%sales strategy%'")
        rows = cur.fetchall()
        if rows:
            form_id = rows[0][0]
        check("Form 'Sales Strategy Feedback' exists", form_id is not None,
              "no form matching '%sales strategy%' in gform.forms")

        if form_id is not None:
            cur.execute(
                "SELECT title, question_type, config FROM gform.questions "
                "WHERE form_id = %s ORDER BY position ASC", (form_id,))
            questions = cur.fetchall()
        cur.close()
        conn.close()
    except Exception as e:
        check("Form query succeeds", False, str(e))

    def opt_values(config):
        vals = []
        cfg = config if isinstance(config, dict) else {}
        opts = cfg.get("options")
        if isinstance(opts, list):
            for o in opts:
                if isinstance(o, dict) and "value" in o:
                    vals.append(str(o["value"]))
                elif isinstance(o, dict) and "label" in o:
                    vals.append(str(o["label"]))
                else:
                    vals.append(str(o))
        return vals

    # Non-critical: at least 3 questions
    check("Form has at least 3 questions", len(questions) >= 3,
          f"found {len(questions)} questions")

    # CRITICAL: an investment question lists all four segment names as options.
    seg_question_ok = False
    for q in questions:
        vals = opt_values(q[2])
        if not vals:
            continue
        blob = " | ".join(vals)
        if all(any_in(blob, seg) for seg in ALL_SEGMENTS):
            seg_question_ok = True
            break
    check("Form has a question listing all four segment names as options",
          seg_question_ok,
          "no question lists Частные клиенты/Корпоративный/Государственный/Малый и средний бизнес",
          critical=True)

    # --- Check 3: Email ---
    print("Checking email to sales_director@company.com...")
    target = None
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute(
            "SELECT subject, to_addr, body_text FROM email.messages ORDER BY date DESC")
        for subject, to_addr, body in cur.fetchall():
            to_str = json.dumps(to_addr).lower() if to_addr else ""
            if "sales_director@company.com" in to_str:
                target = (subject or "", body or "")
                break
        cur.close()
        conn.close()
    except Exception as e:
        check("Email query succeeds", False, str(e))

    check("Email sent to sales_director@company.com", target is not None,
          "no message addressed to sales_director@company.com", critical=True)

    if target is not None:
        subject, body = target
        combined = (subject + " " + body)
        # CRITICAL: body mentions top segment AND total revenue figure
        body_seg_ok = any_in(body, SEG_CONSUMER)
        body_rev_ok = any_in(normalize_ru_numbers(body), REVENUE_VARIANTS)
        check("Email body mentions top segment (Частные клиенты/Consumer) and revenue figure",
              body_seg_ok and body_rev_ok,
              f"seg={body_seg_ok} revenue={body_rev_ok}", critical=True)
        # Non-critical: subject references the segment report
        subj_ok = any(w in subject.lower() for w in
                      ["segment", "сегмент", "report", "отчёт", "отчет"])
        check("Email subject references the segment report", subj_ok,
              f"subject: {subject[:100]!r}")

    # --- Critical gate ---
    if critical_errors:
        print(f"\n=== RESULT: FAIL ({len(critical_errors)} CRITICAL errors) ===")
        for e in critical_errors:
            print(f"  [CRITICAL] {e}")
        for e in all_errors:
            print(f"  [fail] {e}")
        sys.exit(1)

    # --- Accuracy gate ---
    accuracy = (passed / total * 100) if total else 0.0
    print(f"\nPassed {passed}/{total} checks (accuracy {accuracy:.1f}%)")
    if accuracy >= 70:
        print("=== RESULT: PASS ===")
        sys.exit(0)
    else:
        print(f"=== RESULT: FAIL (accuracy {accuracy:.1f}% < 70%) ===")
        for e in all_errors:
            print(f"  {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
