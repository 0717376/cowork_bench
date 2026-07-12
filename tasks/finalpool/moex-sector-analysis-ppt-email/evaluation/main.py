"""Evaluation for moex-sector-analysis-ppt-email.

Expected numeric values are recomputed LIVE from the moex.stock_prices schema in
Postgres (the same source the agent reads via the moex-finance MCP), NOT from a
stale groundtruth file. This keeps the eval in sync with the seeded MOEX data.

Baseline-selection rule (must match task.md):
  - Latest_Close          = close on the most recent available date
  - YTD_Return_Pct        = (latest - first_close_in_March_2026) / first_close_in_March_2026 * 100
  - One_Year_Return_Pct   = (latest - first_close_in_April_2026) / first_close_in_April_2026 * 100

Checks split into structural (non-critical) and CRITICAL.
Any CRITICAL fail -> overall FAIL regardless of accuracy %.
Otherwise PASS if accuracy >= 70%.
"""
import argparse
import json
import os
import sys

import openpyxl

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent",
    "password": "camel",
}

TICKERS = ["GAZP.ME", "LKOH.ME", "MGNT.ME", "MTSS.ME", "SBER.ME", "TCSG.ME"]
SECTORS = {
    "SBER.ME": "Financials",
    "TCSG.ME": "Financials",
    "GAZP.ME": "Energy",
    "LKOH.ME": "Energy",
    "MGNT.ME": "Consumer",
    "MTSS.ME": "Telecom",
}

EMAIL_RECIPIENT = "investments@company.com"
EMAIL_SUBJECT = "Sector Analysis Report - 2026-03-06"

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []
CRITICAL_CHECKS = {
    "Stock Performance: 6 тикеров MOEX с верным сектором и Latest_Close",
    "Stock Performance: YTD/One_Year доходности совпадают с пересчётом по базам март/апрель 2026",
    "Sector Summary: Num_Stocks и Avg_YTD_Return по секторам корректны",
    "Email отправлено на investments@company.com с темой 'Sector Analysis Report - 2026-03-06'",
    "PPTX (Key Findings): названы корректные лучший/худший сектор и лучшая/худшая акция по YTD",
}


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    msg = f": {detail[:300]}" if detail else ""
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        marker = " [CRITICAL]" if name in CRITICAL_CHECKS else ""
        print(f"  [FAIL]{marker} {name}{msg}")
        if name in CRITICAL_CHECKS:
            CRITICAL_FAILS.append(name)


def num_close(a, b, tol):
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def str_match(a, b):
    if a is None or b is None:
        return False
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def compute_expected():
    """Recompute expected per-ticker and per-sector values from moex.stock_prices."""
    import psycopg2
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    expected = {}
    for sym in TICKERS:
        cur.execute(
            "SELECT date, close FROM moex.stock_prices WHERE symbol = %s ORDER BY date",
            (sym,),
        )
        rows = [(str(d), float(c)) for d, c in cur.fetchall() if c is not None]
        if not rows:
            expected[sym] = None
            continue
        latest_close = rows[-1][1]
        mar = next((c for d, c in rows if d.startswith("2026-03")), None)
        apr = next((c for d, c in rows if d.startswith("2026-04")), None)
        ytd = round((latest_close - mar) / mar * 100, 2) if mar else None
        oneyr = round((latest_close - apr) / apr * 100, 2) if apr else None
        expected[sym] = {
            "sector": SECTORS[sym],
            "latest_close": round(latest_close, 2),
            "ytd": ytd,
            "oneyr": oneyr,
        }
    cur.close()
    conn.close()

    # Sector aggregation
    sector_agg = {}
    for sym, v in expected.items():
        if v is None:
            continue
        sec = v["sector"]
        sector_agg.setdefault(sec, {"ytds": [], "oneyrs": []})
        if v["ytd"] is not None:
            sector_agg[sec]["ytds"].append(v["ytd"])
        if v["oneyr"] is not None:
            sector_agg[sec]["oneyrs"].append(v["oneyr"])
    sectors = {}
    for sec, agg in sector_agg.items():
        sectors[sec] = {
            "num": len(agg["ytds"]),
            "avg_ytd": round(sum(agg["ytds"]) / len(agg["ytds"]), 2) if agg["ytds"] else None,
            "avg_oneyr": round(sum(agg["oneyrs"]) / len(agg["oneyrs"]), 2) if agg["oneyrs"] else None,
        }

    # Best/worst leaders by YTD (for Key Findings substance check)
    valid = {s: v for s, v in expected.items() if v and v["ytd"] is not None}
    best_stock = max(valid, key=lambda s: valid[s]["ytd"]) if valid else None
    worst_stock = min(valid, key=lambda s: valid[s]["ytd"]) if valid else None
    sec_valid = {s: v["avg_ytd"] for s, v in sectors.items() if v["avg_ytd"] is not None}
    best_sector = max(sec_valid, key=lambda s: sec_valid[s]) if sec_valid else None
    worst_sector = min(sec_valid, key=lambda s: sec_valid[s]) if sec_valid else None

    return {
        "stocks": expected,
        "sectors": sectors,
        "leaders": {
            "best_stock": best_stock,
            "worst_stock": worst_stock,
            "best_sector": best_sector,
            "worst_sector": worst_sector,
        },
    }


def check_excel(agent_workspace, exp):
    print("\n=== Check 1: Excel Sector_Analysis.xlsx ===")
    path = os.path.join(agent_workspace, "Sector_Analysis.xlsx")
    if not os.path.exists(path):
        record("Sector_Analysis.xlsx exists", False, f"Not found at {path}")
        record("Stock Performance: 6 тикеров MOEX с верным сектором и Latest_Close", False, "xlsx missing")
        record("Stock Performance: YTD/One_Year доходности совпадают с пересчётом по базам март/апрель 2026", False, "xlsx missing")
        record("Sector Summary: Num_Stocks и Avg_YTD_Return по секторам корректны", False, "xlsx missing")
        return
    record("Sector_Analysis.xlsx exists", True)

    try:
        wb = openpyxl.load_workbook(path, data_only=True)
    except Exception as e:
        record("Excel readable", False, str(e))
        return
    record("Excel readable", True)

    # ----- Stock Performance -----
    a_rows = load_sheet_rows(wb, "Stock Performance")
    record("Лист 'Stock Performance' существует", a_rows is not None)
    sector_ok = True
    close_ok = True
    returns_ok = True
    if a_rows is None:
        sector_ok = close_ok = returns_ok = False
    else:
        a_data = a_rows[1:] if len(a_rows) > 1 else []
        lookup = {}
        for row in a_data:
            if row and row[0] is not None:
                lookup[str(row[0]).strip().upper()] = row
        for sym in TICKERS:
            e = exp["stocks"][sym]
            row = lookup.get(sym.upper())
            if e is None:
                continue
            if row is None:
                sector_ok = close_ok = returns_ok = False
                record(f"  {sym}: строка отсутствует", False)
                continue
            # Sector (col 1)
            if not (len(row) > 1 and str_match(row[1], e["sector"])):
                sector_ok = False
                print(f"      {sym}.Sector: {row[1] if len(row) > 1 else None} != {e['sector']}")
            # Latest_Close (col 2, tol 5.0)
            if not (len(row) > 2 and num_close(row[2], e["latest_close"], 5.0)):
                close_ok = False
                print(f"      {sym}.Latest_Close: {row[2] if len(row) > 2 else None} != {e['latest_close']} (tol=5.0)")
            # YTD (col 3, tol 3.0)
            if not (len(row) > 3 and num_close(row[3], e["ytd"], 3.0)):
                returns_ok = False
                print(f"      {sym}.YTD: {row[3] if len(row) > 3 else None} != {e['ytd']} (tol=3.0)")
            # One_Year (col 4, tol 3.0)
            if not (len(row) > 4 and num_close(row[4], e["oneyr"], 3.0)):
                returns_ok = False
                print(f"      {sym}.One_Year: {row[4] if len(row) > 4 else None} != {e['oneyr']} (tol=3.0)")

    record("Stock Performance: 6 тикеров MOEX с верным сектором и Latest_Close", sector_ok and close_ok)
    record("Stock Performance: YTD/One_Year доходности совпадают с пересчётом по базам март/апрель 2026", returns_ok)

    # ----- Sector Summary -----
    s_rows = load_sheet_rows(wb, "Sector Summary")
    record("Лист 'Sector Summary' существует", s_rows is not None)
    summary_ok = True
    if s_rows is None:
        summary_ok = False
    else:
        s_data = s_rows[1:] if len(s_rows) > 1 else []
        lookup = {}
        for row in s_data:
            if row and row[0] is not None:
                lookup[str(row[0]).strip().lower()] = row
        for sec, e in exp["sectors"].items():
            row = lookup.get(sec.lower())
            if row is None:
                summary_ok = False
                print(f"      Сектор отсутствует: {sec}")
                continue
            if not (len(row) > 1 and num_close(row[1], e["num"], 0)):
                summary_ok = False
                print(f"      {sec}.Num_Stocks: {row[1] if len(row) > 1 else None} != {e['num']}")
            if not (len(row) > 2 and num_close(row[2], e["avg_ytd"], 3.0)):
                summary_ok = False
                print(f"      {sec}.Avg_YTD: {row[2] if len(row) > 2 else None} != {e['avg_ytd']} (tol=3.0)")

    record("Sector Summary: Num_Stocks и Avg_YTD_Return по секторам корректны", summary_ok)


def check_pptx(agent_workspace, exp):
    print("\n=== Check 2: PPTX Sector_Analysis.pptx ===")
    path = os.path.join(agent_workspace, "Sector_Analysis.pptx")
    if not os.path.exists(path):
        record("Sector_Analysis.pptx exists", False, f"Not found at {path}")
        record("PPTX (Key Findings): названы корректные лучший/худший сектор и лучшая/худшая акция по YTD", False, "pptx missing")
        return
    record("Sector_Analysis.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(path)
    except Exception as e:
        record("PPTX readable", False, str(e))
        record("PPTX (Key Findings): названы корректные лучший/худший сектор и лучшая/худшая акция по YTD", False, str(e))
        return
    record("PPTX readable", True)

    slides = list(prs.slides)
    record("PPTX: >= 4 слайдов", len(slides) >= 4, f"Found {len(slides)}")

    def slide_text(slide):
        chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for r in shape.table.rows:
                    for c in r.cells:
                        chunks.append(c.text or "")
        return " ".join(chunks)

    all_text = " ".join(slide_text(s) for s in slides)
    all_lower = all_text.lower()

    # Title slide mentions 'sector'
    title_lower = slide_text(slides[0]).lower() if slides else ""
    record("PPTX: титульный слайд содержит 'sector'", "sector" in title_lower)

    # All tickers present
    for sym in TICKERS:
        base = sym.split(".")[0].lower()
        record(f"PPTX: упомянут тикер {sym}", base in all_lower)

    # ---- CRITICAL: Key Findings names correct leaders ----
    lead = exp["leaders"]
    # Locate the Key Findings slide (last slide with 'key'/'findings'/'вывод'); fallback to all text.
    kf_text = all_lower
    for s in slides:
        t = slide_text(s).lower()
        if "key" in t or "finding" in t or "вывод" in t or "итог" in t:
            kf_text = t
            break

    best_stock_base = lead["best_stock"].split(".")[0].lower() if lead["best_stock"] else ""
    worst_stock_base = lead["worst_stock"].split(".")[0].lower() if lead["worst_stock"] else ""
    best_sec = (lead["best_sector"] or "").lower()
    worst_sec = (lead["worst_sector"] or "").lower()

    findings_ok = (
        best_stock_base in kf_text
        and worst_stock_base in kf_text
        and best_sec in kf_text
        and worst_sec in kf_text
    )
    detail = (
        f"best_stock={lead['best_stock']}({best_stock_base in kf_text}), "
        f"worst_stock={lead['worst_stock']}({worst_stock_base in kf_text}), "
        f"best_sector={lead['best_sector']}({best_sec in kf_text}), "
        f"worst_sector={lead['worst_sector']}({worst_sec in kf_text})"
    )
    record("PPTX (Key Findings): названы корректные лучший/худший сектор и лучшая/худшая акция по YTD",
           findings_ok, detail)


def check_email(exp):
    print("\n=== Check 3: Email sent ===")
    try:
        import psycopg2
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute(
            "SELECT m.subject, m.to_addr, m.body_text FROM email.messages m "
            "JOIN email.sent_log s ON s.message_id = m.id"
        )
        rows = cur.fetchall()
        cur.close()
        conn.close()
    except Exception as e:
        record("Email DB readable", False, str(e))
        record("Email отправлено на investments@company.com с темой 'Sector Analysis Report - 2026-03-06'", False, str(e))
        return
    record("Email DB readable", True)

    def to_addresses(to_addr):
        if isinstance(to_addr, list):
            return " ".join(str(r).lower() for r in to_addr)
        if to_addr:
            try:
                parsed = json.loads(str(to_addr))
                if isinstance(parsed, list):
                    return " ".join(str(r).lower() for r in parsed)
            except Exception:
                pass
            return str(to_addr).lower()
        return ""

    matches = [
        r for r in rows
        if EMAIL_RECIPIENT in to_addresses(r[1])
        and str_match(r[0], EMAIL_SUBJECT)
    ]
    record(
        "Email отправлено на investments@company.com с темой 'Sector Analysis Report - 2026-03-06'",
        len(matches) >= 1,
        f"sent total={len(rows)}; subjects={[r[0] for r in rows][:5]}",
    )

    # Soft: body mentions best/worst stock by YTD
    if matches:
        body = (matches[0][2] or "").lower()
        lead = exp["leaders"]
        bs = lead["best_stock"].split(".")[0].lower() if lead["best_stock"] else ""
        ws = lead["worst_stock"].split(".")[0].lower() if lead["worst_stock"] else ""
        record("Email body упоминает лучшую и худшую акцию по YTD",
               bool(bs) and bool(ws) and bs in body and ws in body,
               f"best={bs}, worst={ws}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    try:
        exp = compute_expected()
    except Exception as e:
        print(f"FATAL: could not recompute expected values from moex.stock_prices: {e}")
        sys.exit(1)

    check_excel(args.agent_workspace, exp)
    check_pptx(args.agent_workspace, exp)
    check_email(exp)

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
        "critical_fails": CRITICAL_FAILS,
    }
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if CRITICAL_FAILS:
        print(f"FAIL: критичные чеки провалены ({len(CRITICAL_FAILS)}): {CRITICAL_FAILS}")
        sys.exit(1)
    if accuracy >= 70:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)
    print("\n=== RESULT: FAIL ===")
    sys.exit(1)


if __name__ == "__main__":
    main()
