"""Evaluation script for pw-scholarly-trend-analysis-excel-ppt.

Internal-consistency evaluation: metrics are recomputed from the agent's OWN
Data_Analysis sheet (no frozen groundtruth answer key). Citations/Relevance are
derived by the agent from the :30324 dashboard, so we do NOT hardcode their
values; we only enforce that the summary metrics are self-consistent with the
data the agent produced.

CRITICAL_CHECKS (semantic): any failure => overall FAIL regardless of accuracy.
Otherwise PASS requires accuracy >= 70%.
"""
import os
import argparse, json, os, sys, glob as globmod
import openpyxl

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent", "password": "camel"
}

PASS_COUNT = 0
FAIL_COUNT = 0
FAILED_NAMES = []

# Critical semantic checks: any failure => overall FAIL regardless of accuracy.
CRITICAL_CHECKS = {
    "Data_Analysis has >= 4 rows with all required columns populated",
    "Metrics Total_Papers equals Data_Analysis row count",
    "Metrics Avg_Citations equals rounded mean of Citations column",
    "Metrics Top_Area equals Area with highest Relevance_Score",
    "Recommendations: >= 2 actions, all Areas appear in Data_Analysis",
    "Trend_Analysis_Presentation.pptx exists with >= 4 slides and content tie-in",
    "scholarly_trend_results.json exists and is valid JSON",
}


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        FAILED_NAMES.append(name)
        detail_str = str(detail)[:200] if detail else ""
        print(f"  [FAIL] {name}: {detail_str}")


def safe_float(val, default=None):
    try:
        if val is None:
            return default
        return float(str(val).replace(',', '').replace('%', '').replace('$', '').strip())
    except (ValueError, TypeError):
        return default


def col_index(headers, name):
    name = name.lower()
    for i, h in enumerate(headers):
        if h == name:
            return i
    return None


def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT, FAILED_NAMES
    PASS_COUNT = 0
    FAIL_COUNT = 0
    FAILED_NAMES = []

    excel_path = os.path.join(agent_workspace, "Trend_Analysis_Report.xlsx")
    check("Trend_Analysis_Report.xlsx exists", os.path.exists(excel_path))

    da_areas = set()      # Areas present in Data_Analysis
    da_rows = []          # parsed rows: dict with paper_id/title/area/citations/relevance
    if os.path.exists(excel_path):
        wb = openpyxl.load_workbook(excel_path)

        # ---- Data_Analysis ----
        check("Data_Analysis sheet exists", "Data_Analysis" in wb.sheetnames)
        if "Data_Analysis" in wb.sheetnames:
            ws = wb["Data_Analysis"]
            headers = [str(c.value).strip().lower() if c.value else "" for c in ws[1]]
            for expected_col in ['Paper_ID', 'Title', 'Area', 'Citations', 'Relevance_Score']:
                check(f"Data_Analysis has {expected_col} column",
                      expected_col.lower() in headers, f"headers: {headers[:8]}")

            ci = {k: col_index(headers, k) for k in
                  ['paper_id', 'title', 'area', 'citations', 'relevance_score']}
            for r in ws.iter_rows(min_row=2, values_only=True):
                if all(c is None for c in r):
                    continue
                row = {}
                for k, idx in ci.items():
                    row[k] = r[idx] if (idx is not None and idx < len(r)) else None
                da_rows.append(row)
                if row.get('area') is not None and str(row['area']).strip():
                    da_areas.add(str(row['area']).strip().lower())

            check("Data_Analysis has >= 4 rows", len(da_rows) >= 4, f"got {len(da_rows)}")

            # CRITICAL: >= 4 fully-populated rows
            populated = [
                r for r in da_rows
                if r.get('paper_id') and str(r.get('title') or '').strip()
                and str(r.get('area') or '').strip()
                and safe_float(r.get('citations')) is not None
                and safe_float(r.get('relevance_score')) is not None
            ]
            check("Data_Analysis has >= 4 rows with all required columns populated",
                  len(populated) >= 4, f"populated rows: {len(populated)}")

            # Non-critical: sorted alphabetically by Area
            areas_seq = [str(r.get('area') or '').strip().lower() for r in da_rows if r.get('area')]
            check("Data_Analysis sorted alphabetically by Area",
                  areas_seq == sorted(areas_seq), f"{areas_seq}")

        # ---- Metrics (recomputed from agent's own Data_Analysis) ----
        metrics = {}
        check("Metrics sheet exists", "Metrics" in wb.sheetnames)
        if "Metrics" in wb.sheetnames:
            ws = wb["Metrics"]
            headers = [str(c.value).strip().lower() if c.value else "" for c in ws[1]]
            for expected_col in ['Metric', 'Value']:
                check(f"Metrics has {expected_col} column",
                      expected_col.lower() in headers, f"headers: {headers[:8]}")
            mi = col_index(headers, 'metric')
            vi = col_index(headers, 'value')
            for r in ws.iter_rows(min_row=2, values_only=True):
                if mi is None or vi is None or mi >= len(r):
                    continue
                if r[mi] is None:
                    continue
                metrics[str(r[mi]).strip().lower()] = r[vi] if vi < len(r) else None
            check("Metrics has >= 4 rows", len(metrics) >= 4, f"got {len(metrics)}")

        # Recompute expected values from agent's own data
        citations = [safe_float(r.get('citations')) for r in da_rows
                     if safe_float(r.get('citations')) is not None]
        relevances = [safe_float(r.get('relevance_score')) for r in da_rows
                      if safe_float(r.get('relevance_score')) is not None]

        # CRITICAL: Total_Papers == row count
        tp = safe_float(metrics.get('total_papers'))
        check("Metrics Total_Papers equals Data_Analysis row count",
              tp is not None and int(tp) == len(da_rows) and len(da_rows) >= 4,
              f"Total_Papers={metrics.get('total_papers')} rows={len(da_rows)}")

        # CRITICAL: Avg_Citations == rounded mean
        ac = safe_float(metrics.get('avg_citations'))
        exp_ac = round(sum(citations) / len(citations)) if citations else None
        check("Metrics Avg_Citations equals rounded mean of Citations column",
              ac is not None and exp_ac is not None and abs(ac - exp_ac) <= 1,
              f"Avg_Citations={metrics.get('avg_citations')} expected~{exp_ac}")

        # CRITICAL: Top_Area == area with highest Relevance_Score
        top_area_expected = None
        if da_rows:
            best = max(
                (r for r in da_rows if safe_float(r.get('relevance_score')) is not None),
                key=lambda r: safe_float(r.get('relevance_score')), default=None)
            if best is not None:
                top_area_expected = str(best.get('area') or '').strip().lower()
        ta = str(metrics.get('top_area') or '').strip().lower()
        check("Metrics Top_Area equals Area with highest Relevance_Score",
              top_area_expected is not None and ta == top_area_expected,
              f"Top_Area={metrics.get('top_area')} expected={top_area_expected}")

        # Non-critical: Avg_Relevance present & consistent
        avr = safe_float(metrics.get('avg_relevance'))
        exp_avr = round(sum(relevances) / len(relevances)) if relevances else None
        check("Metrics has Avg_Relevance consistent with Relevance_Score column",
              avr is not None and exp_avr is not None and abs(avr - exp_avr) <= 1,
              f"Avg_Relevance={metrics.get('avg_relevance')} expected~{exp_avr}")

        # ---- Recommendations ----
        check("Recommendations sheet exists", "Recommendations" in wb.sheetnames)
        rec_count = 0
        rec_areas_ok = True
        if "Recommendations" in wb.sheetnames:
            ws = wb["Recommendations"]
            headers = [str(c.value).strip().lower() if c.value else "" for c in ws[1]]
            for expected_col in ['Priority', 'Action', 'Area']:
                check(f"Recommendations has {expected_col} column",
                      expected_col.lower() in headers, f"headers: {headers[:8]}")
            ai = col_index(headers, 'area')
            for r in ws.iter_rows(min_row=2, values_only=True):
                if all(c is None for c in r):
                    continue
                rec_count += 1
                if ai is not None and ai < len(r) and r[ai] is not None:
                    if str(r[ai]).strip().lower() not in da_areas:
                        rec_areas_ok = False
            check("Recommendations has >= 2 rows", rec_count >= 2, f"got {rec_count}")

        # CRITICAL: >= 2 actions and all referenced Areas exist in Data_Analysis
        check("Recommendations: >= 2 actions, all Areas appear in Data_Analysis",
              rec_count >= 2 and rec_areas_ok and len(da_areas) > 0,
              f"rec_count={rec_count} areas_ok={rec_areas_ok}")

    # ---- PowerPoint: exact name + >= 4 slides + content tie-in (CRITICAL) ----
    pptx_path = os.path.join(agent_workspace, "Trend_Analysis_Presentation.pptx")
    pptx_ok = False
    pptx_detail = ""
    if os.path.exists(pptx_path):
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            n_slides = len(prs.slides)
            slide_text = []
            for s in prs.slides:
                for sh in s.shapes:
                    if sh.has_text_frame:
                        slide_text.append(sh.text)
            blob = " ".join(slide_text).lower()
            # content tie-in: at least one Area from Data_Analysis appears
            tie_in = any(a and a in blob for a in da_areas) if da_areas else False
            pptx_ok = n_slides >= 4 and tie_in
            pptx_detail = f"slides={n_slides} tie_in={tie_in}"
        except Exception as e:
            pptx_detail = f"error: {e}"
    else:
        # fallback to any pptx for the detail, but exact name is required
        others = globmod.glob(os.path.join(agent_workspace, "*.pptx"))
        pptx_detail = f"exact file missing; other pptx: {len(others)}"
    check("Trend_Analysis_Presentation.pptx exists with >= 4 slides and content tie-in",
          pptx_ok, pptx_detail)

    # ---- Processor script ----
    check("scholarly_trend_processor.py exists",
          os.path.exists(os.path.join(agent_workspace, "scholarly_trend_processor.py")))

    # ---- Results JSON (CRITICAL: exists + valid JSON) ----
    json_path = os.path.join(agent_workspace, "scholarly_trend_results.json")
    json_ok = False
    json_detail = ""
    if os.path.exists(json_path):
        try:
            with open(json_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            json_ok = isinstance(data, (dict, list)) and len(data) > 0
            json_detail = f"type={type(data).__name__}"
        except Exception as e:
            json_detail = f"error: {e}"
    else:
        json_detail = "file missing"
    check("scholarly_trend_results.json exists and is valid JSON", json_ok, json_detail)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = PASS_COUNT / total * 100 if total else 0
    critical_failed = [n for n in FAILED_NAMES if n in CRITICAL_CHECKS]

    print(f"\n=== Results: {PASS_COUNT}/{total} passed ({accuracy:.1f}%) ===")
    if critical_failed:
        print(f"CRITICAL FAILURES: {len(critical_failed)}")
        for n in critical_failed:
            print(f"  - {n}")

    if res_log_file:
        try:
            with open(res_log_file, "w") as f:
                json.dump({
                    "total_passed": PASS_COUNT, "total_checks": total,
                    "accuracy": accuracy, "critical_failed": critical_failed,
                }, f, indent=2)
        except Exception:
            pass

    success = (not critical_failed) and accuracy >= 70
    return success, f"Passed {PASS_COUNT}/{total} checks ({accuracy:.1f}%)"


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
