"""
Evaluation script for sf-department-budget-ppt task (ClickHouse, RU).

Checks:
1. Excel file (budget_vs_actuals.xlsx) - both sheets
2. PowerPoint file (budget_vs_actuals.pptx) - structure and content

Department names in the ClickHouse warehouse are russified centrally, so the
agent's deliverables will contain Russian department labels. All department
matching is done via DEPT_ALIASES -> canonical English key, so the eval is
language-agnostic (accepts RU or EN). EXPECTED_* dicts stay English-keyed.
"""

import argparse
import json
import os
import sys

import openpyxl

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []

DEPARTMENTS = ["Engineering", "Finance", "HR", "Operations", "R&D", "Sales", "Support"]

# Department labels may be written in Russian (warehouse russified centrally)
# or English. Alias map -> canonical English key.
DEPT_ALIASES = {
    "engineering": "Engineering", "инженерия": "Engineering",
    "finance": "Finance", "финансы": "Finance",
    "hr": "HR", "кадры": "HR", "human resources": "HR",
    "operations": "Operations", "операции": "Operations",
    "r&d": "R&D", "rd": "R&D", "ниокр": "R&D", "r and d": "R&D",
    "sales": "Sales", "продажи": "Sales",
    "support": "Support", "поддержка": "Support",
}

# Russian display names per canonical English dept (for PPTX substring checks).
DEPT_RU = {
    "Engineering": "Инженерия", "Finance": "Финансы", "HR": "Кадры",
    "Operations": "Операции", "R&D": "НИОКР", "Sales": "Продажи",
    "Support": "Поддержка",
}

# Expected actual expenditures from the warehouse (numeric, language-neutral).
EXPECTED_ACTUALS = {
    "Engineering": 418604451.00,
    "Finance": 413713333.00,
    "HR": 416980025.00,
    "Operations": 411598247.00,
    "R&D": 410147671.00,
    "Sales": 425710161.00,
    "Support": 423053067.00,
}

EXPECTED_BUDGETS = {
    "Engineering": 633555641.27,
    "Finance": 616980984.11,
    "HR": 697911374.65,
    "Operations": 572668189.63,
    "R&D": 715914847.42,
    "Sales": 642575220.92,
    "Support": 684741486.31,
}


def canon_dept(val):
    if val is None:
        return None
    return DEPT_ALIASES.get(str(val).strip().lower(), str(val).strip())


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def critical(name, passed, detail=""):
    """A semantic check whose failure forces overall FAIL."""
    global PASS_COUNT, FAIL_COUNT, CRITICAL_FAILS
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS][CRITICAL] {name}")
    else:
        FAIL_COUNT += 1
        CRITICAL_FAILS.append(name)
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL][CRITICAL] {name}{msg}")


def str_match(a, b):
    """Language-agnostic department equality via canonicalization."""
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    return canon_dept(a).lower() == canon_dept(b).lower()


def num_close(a, b, rel_tol=0.02):
    try:
        a, b = float(a), float(b)
        if b == 0:
            return abs(a) < 1000
        return abs(a - b) / abs(b) <= rel_tol
    except (TypeError, ValueError):
        return False


# ============================================================================
# Check 1: Excel file
# ============================================================================

def check_excel(agent_workspace, groundtruth_workspace):
    print("\n=== Checking Excel Output ===")

    agent_file = os.path.join(agent_workspace, "budget_vs_actuals.xlsx")

    if not os.path.isfile(agent_file):
        record("Excel file exists", False, f"Not found: {agent_file}")
        return False
    record("Excel file exists", True)

    wb = openpyxl.load_workbook(agent_file, data_only=True)

    def get_sheet(wb, target):
        for name in wb.sheetnames:
            if name.strip().lower() == target.strip().lower():
                return wb[name]
        return None

    # Sheet 1: Department Details
    ws1 = get_sheet(wb, "Department Details")
    if ws1 is None:
        record("Sheet 'Department Details' exists", False, f"Sheets: {wb.sheetnames}")
        return False
    record("Sheet 'Department Details' exists", True)

    headers = [str(c.value).strip() if c.value else "" for c in ws1[1]]
    headers_lc = [h.lower() for h in headers]
    expected = ["Department", "Approved_Budget", "Actual_Expenditure", "Variance",
                "Variance_Pct", "Headcount", "Avg_Salary"]
    for col in expected:
        record(f"Department Details has {col} column", col.lower() in headers_lc,
               f"Got: {headers}")

    def idx(col):
        return headers_lc.index(col.lower()) if col.lower() in headers_lc else None

    i_dep = idx("Department")
    i_bud = idx("Approved_Budget")
    i_act = idx("Actual_Expenditure")
    i_var = idx("Variance")
    i_pct = idx("Variance_Pct")

    rows = list(ws1.iter_rows(min_row=2, values_only=True))
    record("Department Details has 7 rows", len(rows) == 7, f"Got {len(rows)}")

    # Build canonical-keyed map of agent rows
    agent_by_dept = {}
    written_depts = []
    for r in rows:
        if not r or i_dep is None or r[i_dep] is None or str(r[i_dep]).strip() == "":
            continue
        agent_by_dept[canon_dept(r[i_dep])] = r
        written_depts.append(str(r[i_dep]).strip())

    # CRITICAL: all 7 departments present (RU or EN)
    critical("Department Details covers all 7 departments",
             set(DEPARTMENTS).issubset(set(agent_by_dept.keys())),
             f"present={sorted(agent_by_dept.keys())}")

    var_errors = []
    for dept in DEPARTMENTS:
        agent_row = agent_by_dept.get(dept)
        if not agent_row:
            record(f"Dept {dept} present", False, "Missing")
            continue
        record(f"Dept {dept} present", True)

        if i_act is not None:
            record(f"Dept {dept}: Actual Expenditure",
                   num_close(agent_row[i_act], EXPECTED_ACTUALS[dept]),
                   f"Expected ~{EXPECTED_ACTUALS[dept]:.0f}, got {agent_row[i_act]}")
        if i_bud is not None:
            record(f"Dept {dept}: Approved Budget",
                   num_close(agent_row[i_bud], EXPECTED_BUDGETS[dept]),
                   f"Expected ~{EXPECTED_BUDGETS[dept]:.0f}, got {agent_row[i_bud]}")

        # Variance = budget - actuals (per-row arithmetic, used for critical check)
        expected_var = EXPECTED_BUDGETS[dept] - EXPECTED_ACTUALS[dept]
        var_ok = i_var is not None and num_close(agent_row[i_var], expected_var)
        record(f"Dept {dept}: Variance", var_ok,
               f"Expected ~{expected_var:.0f}, got {agent_row[i_var] if i_var is not None else 'NA'}")
        if not var_ok:
            var_errors.append((dept, agent_row[i_var] if i_var is not None else None))

    # CRITICAL: per-department Variance == Budget - Actuals for all 7 rows
    critical("Per-department Variance == Approved_Budget - Actual_Expenditure (all 7 rows)",
             len(var_errors) == 0 and len(agent_by_dept) >= 7,
             f"errors={var_errors[:5]}")

    # CRITICAL: Variance_Pct correct for extreme depts (R&D ~42.71, Operations ~28.13)
    pct_ok = True
    pct_detail = []
    if i_pct is not None:
        for dept, exp_pct in (("R&D", 42.71), ("Operations", 28.13)):
            row = agent_by_dept.get(dept)
            got = row[i_pct] if row else None
            try:
                if got is None or abs(float(got) - exp_pct) > 0.5:
                    pct_ok = False
                    pct_detail.append((dept, got, exp_pct))
            except (TypeError, ValueError):
                pct_ok = False
                pct_detail.append((dept, got, exp_pct))
    else:
        pct_ok = False
        pct_detail.append("Variance_Pct column missing")
    critical("Variance_Pct correct for R&D (~42.71) and Operations (~28.13)",
             pct_ok, f"{pct_detail}")

    # NON-critical: alphabetical sort (RU and EN sort orders both acceptable)
    record("Department Details sorted alphabetically by Department",
           written_depts == sorted(written_depts),
           f"order={written_depts}")

    # Sheet 2: Totals
    ws2 = get_sheet(wb, "Totals")
    if ws2 is None:
        record("Sheet 'Totals' exists", False, f"Sheets: {wb.sheetnames}")
        critical("Totals sheet present", False, "missing")
        return True
    record("Sheet 'Totals' exists", True)

    summary = {}
    for row in ws2.iter_rows(min_row=1, values_only=True):
        if row and row[0]:
            summary[str(row[0]).strip().lower()] = row[1]

    record("Totals: Total_Departments = 7",
           str(summary.get("total_departments", "")).strip() == "7",
           f"Got {summary.get('total_departments')}")

    total_budget = sum(EXPECTED_BUDGETS.values())
    total_actuals = sum(EXPECTED_ACTUALS.values())
    total_variance = total_budget - total_actuals  # ~1,644,540,789.31

    record("Totals: Total_Budget",
           num_close(summary.get("total_budget", 0), total_budget),
           f"Expected ~{total_budget:.0f}, got {summary.get('total_budget')}")

    record("Totals: Total_Actuals",
           num_close(summary.get("total_actuals", 0), total_actuals),
           f"Expected ~{total_actuals:.0f}, got {summary.get('total_actuals')}")

    # CRITICAL: Total_Variance == Total_Budget - Total_Actuals (only if present)
    if "total_variance" in summary:
        critical("Totals: Total_Variance == Total_Budget - Total_Actuals (~1,644,540,789.31)",
                 num_close(summary.get("total_variance"), total_variance),
                 f"Expected ~{total_variance:.2f}, got {summary.get('total_variance')}")
    else:
        record("Totals: Total_Variance present", False, "row missing (non-critical)")

    # CRITICAL: correct winner/loser ranking (language-agnostic)
    critical("Totals: Most_Under_Budget == R&D",
             str_match(summary.get("most_under_budget", ""), "R&D"),
             f"Got {summary.get('most_under_budget')}")
    critical("Totals: Most_Over_Budget == Operations",
             str_match(summary.get("most_over_budget", ""), "Operations"),
             f"Got {summary.get('most_over_budget')}")

    return True


# ============================================================================
# Check 2: PowerPoint file
# ============================================================================

def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint ===")

    pptx_path = os.path.join(agent_workspace, "budget_vs_actuals.pptx")

    if not os.path.isfile(pptx_path):
        record("PPTX file exists", False, f"Not found: {pptx_path}")
        critical("PPTX file exists", False, f"Not found: {pptx_path}")
        return False
    record("PPTX file exists", True)

    try:
        from pptx import Presentation
    except ImportError:
        record("python-pptx installed", False, "Cannot import pptx")
        return False

    prs = Presentation(pptx_path)
    slides = list(prs.slides)

    # Should have: 1 title + 1 overview + 7 departments + 1 summary = 10 slides
    record("PPTX has at least 10 slides", len(slides) >= 10,
           f"Found {len(slides)} slides")

    # Collect all text
    all_text = ""
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text + "\n"
    all_lower = all_text.lower()

    record("PPTX: has title 'budget vs actuals'",
           "budget vs actuals" in all_lower or "fy2026" in all_lower,
           "Title not found")

    record("PPTX: mentions March 2026",
           "march 2026" in all_lower,
           "March 2026 not found")

    # Each department must appear (RU display name OR English alias).
    for dept in DEPARTMENTS:
        ru = DEPT_RU[dept].lower()
        en = dept.lower()
        record(f"PPTX: mentions {dept}",
               ru in all_lower or en in all_lower,
               f"{dept}/{DEPT_RU[dept]} not found")

    record("PPTX: has Summary slide",
           "summary" in all_lower or "сводка" in all_lower or "итог" in all_lower,
           "Summary not found")

    record("PPTX: mentions 'variance' or 'budget'",
           any(t in all_lower for t in ("variance", "budget", "отклонение", "бюджет")),
           "Budget/variance terms not found")

    # CRITICAL: at least 10 slides AND Summary slide names both R&D and Operations.
    rd_ok = DEPT_RU["R&D"].lower() in all_lower or "r&d" in all_lower
    ops_ok = DEPT_RU["Operations"].lower() in all_lower or "operations" in all_lower
    critical("PPTX has >=10 slides and Summary names R&D and Operations",
             len(slides) >= 10 and rd_ok and ops_ok,
             f"slides={len(slides)} rd={rd_ok} ops={ops_ok}")

    return True


# ============================================================================
# Main
# ============================================================================

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    excel_ok = check_excel(args.agent_workspace, args.groundtruth_workspace)
    pptx_ok = check_pptx(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = (PASS_COUNT / total * 100.0) if total else 0.0

    structural_ok = excel_ok and pptx_ok

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Accuracy: {accuracy:.1f}%")

    # CRITICAL gate: any critical failure => hard FAIL regardless of accuracy.
    if CRITICAL_FAILS:
        print(f"  CRITICAL FAILURES: {CRITICAL_FAILS}")

    all_passed = structural_ok and not CRITICAL_FAILS and accuracy >= 70.0
    print(f"  Overall: {'PASS' if all_passed else 'FAIL'}")

    if args.res_log_file:
        result = {
            "passed": PASS_COUNT,
            "failed": FAIL_COUNT,
            "accuracy": accuracy,
            "critical_fails": CRITICAL_FAILS,
            "success": all_passed,
        }
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(0 if all_passed else 1)


if __name__ == "__main__":
    main()
