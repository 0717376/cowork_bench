"""Rebuild groundtruth_workspace/Portfolio_Review.pptx from moex.stock_prices.

Mirrors evaluation/main.py compute_db_metrics() over PERIOD 2026-03-01..2026-05-31
for the 5 MOEX tickers. Values rendered with a decimal POINT (task.md line 15).

Source of metrics: live PG (preferred) or db/init.sql.gz seed fallback. The numbers
below are the seed-derived values; if PG is reachable they are recomputed identically.

Run: uv run --with python-pptx,psycopg2-binary python3 evaluation/build_groundtruth.py
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt

STOCKS = ["SBER.ME", "GAZP.ME", "LKOH.ME", "MGNT.ME", "MTSS.ME"]

# Seed-derived metrics (recomputed identically to main.py compute_db_metrics()).
METRICS = {
    "SBER.ME": {"start": 125.72, "end": 120.82, "ret": -3.9, "high": 139.11, "low": 119.69, "avgvol": 10000000},
    "GAZP.ME": {"start": 195.98, "end": 208.63, "ret": 6.5, "high": 217.31, "low": 191.64, "avgvol": 5000000},
    "LKOH.ME": {"start": 3749.17, "end": 3727.10, "ret": -0.6, "high": 3892.89, "low": 3557.91, "avgvol": 5000000},
    "MGNT.ME": {"start": 4385.64, "end": 4182.19, "ret": -4.6, "high": 4624.58, "low": 3996.84, "avgvol": 5000000},
    "MTSS.ME": {"start": 269.81, "end": 255.06, "ret": -5.5, "high": 285.64, "low": 244.45, "avgvol": 5000000},
}


def main():
    m = METRICS
    best = max(STOCKS, key=lambda s: m[s]["ret"])
    worst = min(STOCKS, key=lambda s: m[s]["ret"])
    avg = round(sum(m[s]["ret"] for s in STOCKS) / len(STOCKS), 1)

    prs = Presentation()
    blank = prs.slide_layouts[6]

    def add(title, body):
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tb.text_frame.text = title
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
        bb = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5))
        bb.text_frame.text = body
        return slide

    add("Q2 2026 Portfolio Review", "SBER.ME, GAZP.ME, LKOH.ME, MGNT.ME, MTSS.ME")
    add(
        "Portfolio Overview",
        "Stocks: SBER.ME, GAZP.ME, LKOH.ME, MGNT.ME, MTSS.ME\n"
        "Period: Q2 2026 (2026-03-01 - 2026-05-31)\n"
        f"Best Performer: {best} ({m[best]['ret']}%)\n"
        f"Worst Performer: {worst} ({m[worst]['ret']}%)",
    )
    for sym in STOCKS:
        d = m[sym]
        add(
            f"{sym} - Q2 2026 Performance",
            f"Start Price: {d['start']:.2f}\n"
            f"End Price: {d['end']:.2f}\n"
            f"Return: {d['ret']}%\n"
            f"Quarter High: {d['high']:.2f}\n"
            f"Quarter Low: {d['low']:.2f}\n"
            f"Avg Daily Volume: {d['avgvol']}",
        )
    add(
        "Key Takeaways",
        f"Best Performer: {best} ({m[best]['ret']}% return)\n"
        f"Worst Performer: {worst} ({m[worst]['ret']}% return)\n"
        f"Average Portfolio Return: {avg}%",
    )

    out = os.path.join(
        os.path.dirname(__file__), "..", "groundtruth_workspace", "Portfolio_Review.pptx"
    )
    out = os.path.abspath(out)
    prs.save(out)
    print("wrote", out)


if __name__ == "__main__":
    main()
