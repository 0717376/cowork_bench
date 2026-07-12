"""
Evaluation for yf-portfolio-ppt-gcal (Russified: MOEX portfolio review).

Структура: structural checks + CRITICAL semantic checks.
CRITICAL (любой провал => общий FAIL независимо от accuracy):
  C1: квартальная доходность каждого из 5 тикеров (пересчитана из moex.stock_prices)
      присутствует в тексте PPTX (с допуском форматирования +/-0.5).
  C2: лучшая и худшая по доходности акции (из БД) названы в PPTX.
  C3: средняя доходность портфеля (round(mean,1)) присутствует на слайде выводов.
  C4: по каждому тикеру начальная цена, конечная цена, максимум (MAX(high)) и
      минимум (MIN(low)) квартала присутствуют в тексте PPTX (с допуском).
  C5: событие совещания в gcal.events существует и корректно: дата 2026-06-15,
      окно 14:00-15:00, локация "Переговорная комната А".

Иначе PASS, если accuracy >= 70%.
Объёмы торгов проверяются мягко (структурно, не critical).
"""
import argparse
import json
import os
import sys

import psycopg2
from pptx import Presentation

DB = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "cowork_gym"),
    "user": "eigent",
    "password": "camel",
}

# Период обзора (данные moex.stock_prices покрывают 2026-02-25..2026-05-26).
PERIOD_LO = "2026-03-01"
PERIOD_HI = "2026-05-31"
STOCKS = ["SBER.ME", "GAZP.ME", "LKOH.ME", "MGNT.ME", "MTSS.ME"]

PASS_COUNT = 0
FAIL_COUNT = 0
CRITICAL_FAILS = []
CRITICAL_CHECKS = set()  # заполняется по ходу (имена critical-чеков)


def record(name, passed, detail="", critical=False):
    global PASS_COUNT, FAIL_COUNT
    if critical:
        CRITICAL_CHECKS.add(name)
    msg = f": {detail[:300]}" if detail else ""
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        marker = " [CRITICAL]" if critical else ""
        print(f"  [FAIL]{marker} {name}{msg}")
        if critical:
            CRITICAL_FAILS.append(name)


def value_in_text(value, text, tol=0.5):
    """Число присутствует в тексте PPTX с допуском форматирования +/-tol.

    Перебираем форматы с 0..2 знаками после точки и небольшую дельту,
    чтобы покрыть округления при отображении на слайдах.
    """
    candidates = set()
    for delta_steps in range(-int(tol * 10), int(tol * 10) + 1):
        v = value + delta_steps / 10.0
        for prec in (0, 1, 2):
            s = f"{v:.{prec}f}"
            candidates.add(s)
            # вариант без ведущего знака минус (на случай "доходность -3.9" vs "3.9%")
            if s.startswith("-"):
                candidates.add(s[1:])
    return any(c in text for c in candidates)


def compute_db_metrics():
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()
    metrics = {}
    for sym in STOCKS:
        cur.execute(
            "SELECT close FROM moex.stock_prices WHERE symbol=%s AND date="
            "(SELECT MIN(date) FROM moex.stock_prices WHERE symbol=%s AND date>=%s)",
            (sym, sym, PERIOD_LO),
        )
        sp = cur.fetchone()
        cur.execute(
            "SELECT close FROM moex.stock_prices WHERE symbol=%s AND date="
            "(SELECT MAX(date) FROM moex.stock_prices WHERE symbol=%s AND date<=%s)",
            (sym, sym, PERIOD_HI),
        )
        ep = cur.fetchone()
        cur.execute(
            "SELECT MAX(high), MIN(low), AVG(volume) FROM moex.stock_prices "
            "WHERE symbol=%s AND date>=%s AND date<=%s",
            (sym, PERIOD_LO, PERIOD_HI),
        )
        hi, lo, av = cur.fetchone()
        start = float(sp[0])
        end = float(ep[0])
        ret = round((end - start) / start * 100, 1)
        metrics[sym] = {
            "start": start,
            "end": end,
            "ret": ret,
            "high": float(hi),
            "low": float(lo),
            "avgvol": float(av),
        }
    conn.close()
    return metrics


def check_ppt(agent_ppt):
    print("\n=== Check PPTX Portfolio_Review.pptx ===")
    if not os.path.exists(agent_ppt):
        record("Portfolio_Review.pptx существует", False, f"не найден: {agent_ppt}")
        # critical-чеки должны быть оценены явно
        for sym in STOCKS:
            record(f"PPTX: доходность {sym}", False, "PPTX отсутствует", critical=True)
            record(f"PPTX: цены/максимум/минимум {sym}", False, "PPTX отсутствует", critical=True)
        record("PPTX: названы лучшая и худшая акции", False, "PPTX отсутствует", critical=True)
        record("PPTX: средняя доходность портфеля на слайде выводов", False, "PPTX отсутствует", critical=True)
        return

    record("Portfolio_Review.pptx существует", True)
    prs = Presentation(agent_ppt)
    slides = list(prs.slides)

    record("PPTX: >= 8 слайдов (title+overview+5+takeaways)", len(slides) >= 8,
           f"найдено {len(slides)}")

    # Собрать весь текст
    all_text = ""
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text + " "
    all_lower = all_text.lower()

    # Титульный слайд
    title_text = ""
    if slides:
        for shape in slides[0].shapes:
            if shape.has_text_frame:
                title_text += shape.text_frame.text.lower() + " "
    record("PPTX: титульный слайд содержит 'Q2 2025/2026' и 'portfolio'",
           ("q2 2026" in title_text or "q2 2025" in title_text) and "portfolio" in title_text,
           f"title: {title_text[:120]}")

    # Все тикеры присутствуют
    for sym in STOCKS:
        base = sym.split(".")[0].lower()
        record(f"PPTX: упомянут тикер {sym}", base in all_lower)

    # Слайд выводов
    last_text = ""
    if len(slides) >= 2:
        for shape in slides[-1].shapes:
            if shape.has_text_frame:
                last_text += shape.text_frame.text.lower() + " "
    takeaway_kw = ["takeaway", "key", "summary", "conclusion",
                   "вывод", "выводы", "итог", "итоги", "ключевые", "резюме"]
    record("PPTX: финальный слайд с выводами",
           any(k in last_text for k in takeaway_kw),
           f"last: {last_text[:120]}")

    # ---- Данные из БД ----
    metrics = compute_db_metrics()

    # CRITICAL C1: доходность каждого тикера
    for sym in STOCKS:
        ret = metrics[sym]["ret"]
        record(f"PPTX: доходность {sym} (~{ret}%)",
               value_in_text(ret, all_text, tol=0.5),
               f"ожидалось ~{ret}", critical=True)

    # CRITICAL C4: цены/максимум/минимум каждого тикера
    for sym in STOCKS:
        m = metrics[sym]
        ok_start = value_in_text(m["start"], all_text, tol=0.5)
        ok_end = value_in_text(m["end"], all_text, tol=0.5)
        ok_high = value_in_text(m["high"], all_text, tol=0.5)
        ok_low = value_in_text(m["low"], all_text, tol=0.5)
        record(
            f"PPTX: цены/максимум/минимум {sym}",
            ok_start and ok_end and ok_high and ok_low,
            f"start={m['start']}({ok_start}) end={m['end']}({ok_end}) "
            f"high={m['high']}({ok_high}) low={m['low']}({ok_low})",
            critical=True,
        )

    # CRITICAL C2: лучшая/худшая акции
    best_sym = max(STOCKS, key=lambda s: metrics[s]["ret"])
    worst_sym = min(STOCKS, key=lambda s: metrics[s]["ret"])
    best_base = best_sym.split(".")[0].lower()
    worst_base = worst_sym.split(".")[0].lower()
    record(
        "PPTX: названы лучшая и худшая акции",
        best_base in all_lower and worst_base in all_lower,
        f"best={best_sym}({best_base in all_lower}) worst={worst_sym}({worst_base in all_lower})",
        critical=True,
    )

    # CRITICAL C3: средняя доходность портфеля
    avg_ret = round(sum(metrics[s]["ret"] for s in STOCKS) / len(STOCKS), 1)
    record(
        "PPTX: средняя доходность портфеля на слайде выводов",
        value_in_text(avg_ret, last_text, tol=0.3) or value_in_text(avg_ret, all_text, tol=0.3),
        f"ожидалось ~{avg_ret}",
        critical=True,
    )

    # NON-critical: средний объём (мягко, структурно)
    for sym in STOCKS:
        av = int(round(metrics[sym]["avgvol"]))
        # объёмы — крупные целые; ищем без разделителей и грубое присутствие
        av_str = str(av)
        record(f"PPTX: средний объём {sym} (мягко)",
               av_str in all_text.replace(",", "").replace(" ", "").replace(" ", ""),
               f"avgvol={av}")


def check_gcal():
    print("\n=== Check Google Calendar event (CRITICAL) ===")
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute(
            """
            SELECT summary, start_datetime, end_datetime, location
              FROM gcal.events
             WHERE start_datetime >= '2026-06-15' AND start_datetime < '2026-06-16'
            """
        )
        events = cur.fetchall()
        conn.close()
    except Exception as e:
        record("Календарь: событие совещания 2026-06-15", False, f"ошибка БД: {e}", critical=True)
        return

    if not events:
        record("Календарь: событие совещания 2026-06-15", False,
               "событие на 2026-06-15 не найдено", critical=True)
        return

    def ev_ok(ev):
        summary, start_dt, end_dt, location = ev
        s = (summary or "").lower()
        loc = (location or "").lower()
        summary_ok = ("portfolio" in s and "review" in s) or "обзор портфел" in s
        time_ok = start_dt is not None and start_dt.hour == 14 and (
            end_dt is None or end_dt.hour == 15
        )
        loc_ok = "переговорн" in loc or "комната а" in loc or "conference room a" in loc
        return summary_ok and time_ok and loc_ok, (summary_ok, time_ok, loc_ok)

    matches = []
    detail = ""
    for ev in events:
        ok, parts = ev_ok(ev)
        detail = f"summary='{ev[0]}' start={ev[1]} loc='{ev[3]}' parts(summary,time,loc)={parts}"
        if ok:
            matches.append(ev)
    record(
        "Календарь: событие совещания 2026-06-15 14:00-15:00, локация 'Переговорная комната А'",
        len(matches) >= 1,
        detail or f"events: {events}",
        critical=True,
    )


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    agent_ppt = os.path.join(args.agent_workspace, "Portfolio_Review.pptx")
    check_ppt(agent_ppt)
    check_gcal()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: чеки не выполнены.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
        "critical_fails": CRITICAL_FAILS,
    }
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2, ensure_ascii=False)

    if CRITICAL_FAILS:
        print(f"FAIL: критичные чеки провалены ({len(CRITICAL_FAILS)}): {CRITICAL_FAILS}")
        sys.exit(1)
    if accuracy >= 70:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)
    else:
        print(f"\n=== RESULT: FAIL (accuracy {accuracy:.1f}% < 70%) ===")
        sys.exit(1)


if __name__ == "__main__":
    main()
